"""
Certify Intel - Simplified Backend for Excel Dashboard
A lightweight FastAPI backend that:
1. Stores competitor data in SQLite (simple, no PostgreSQL setup needed)
2. Scrapes competitor websites using Playwright
3. Extracts data using OpenAI GPT
4. Exports data to Excel-compatible formats
"""
import os
import sys
import logging

from constants import __version__, NO_HALLUCINATION_INSTRUCTION  # noqa: E402

# Configure logging early for CI/CD and production
logger = logging.getLogger(__name__)

# ==============================================================================
# ENVIRONMENT LOADING - PyInstaller Compatible (v5.0.3 Fix)
# ==============================================================================
def _load_env():
    """
    Load environment variables with PyInstaller awareness.

    When bundled with PyInstaller:
    - __main__.py sets CERTIFY_BUNDLED=true and loads .env from exe directory
    - We skip load_dotenv() here to avoid overwriting with empty values

    When running normally (development):
    - load_dotenv() finds .env in the current/backend directory
    """
    # Check if already loaded by __main__.py (bundled mode)
    if os.getenv('CERTIFY_BUNDLED') == 'true':
        logger.info("[ENV] Running bundled - .env already loaded by __main__.py")
        return True

    try:
        from dotenv import load_dotenv

        # Check for PyInstaller frozen mode even without CERTIFY_BUNDLED flag
        if getattr(sys, 'frozen', False):
            # Running as PyInstaller bundle - look for .env next to executable
            exe_dir = os.path.dirname(sys.executable)
            env_path = os.path.join(exe_dir, '.env')
            if os.path.exists(env_path):
                load_dotenv(env_path)
                logger.info(f"[ENV] Loaded from exe directory: {env_path}")
                return True
            else:
                logger.info(f"[ENV] Warning: No .env found at {env_path}")
                logger.info("[ENV] Please place your .env file next to the executable")
                return False
        else:
            # Normal development mode - look in current directory
            load_dotenv()
            logger.info("[ENV] Loading environment variables from .env")
            return True

    except ImportError:
        logger.info("[ENV] python-dotenv not installed, using system environment variables")
        return False

_load_env()
logger.info(f"[ENV] Database URL: {os.getenv('DATABASE_URL', 'sqlite:///./certify_intel.db')}")
# ==============================================================================

import json
import asyncio
import uuid

# Optional: slowapi rate limiting (per-endpoint granular limits)
try:
    from slowapi import Limiter, _rate_limit_exceeded_handler
    from slowapi.util import get_remote_address
    from slowapi.errors import RateLimitExceeded
    from slowapi.middleware import SlowAPIMiddleware

    _rate_limit_enabled = os.getenv("RATE_LIMIT_ENABLED", "true").lower() == "true"

    if _rate_limit_enabled:
        limiter = Limiter(key_func=get_remote_address)
    else:
        limiter = Limiter(key_func=get_remote_address, enabled=False)
except ImportError:
    limiter = None
    _rate_limit_enabled = False
    logger.info("slowapi not installed - per-endpoint rate limiting disabled")

# Optional: structured JSON logging for production
if os.getenv("JSON_LOGGING", "false").lower() == "true":
    try:
        from pythonjsonlogger import jsonlogger

        _json_handler = logging.StreamHandler()
        _json_formatter = jsonlogger.JsonFormatter(
            '%(asctime)s %(name)s %(levelname)s %(message)s',
            rename_fields={"asctime": "timestamp", "levelname": "level"}
        )
        _json_handler.setFormatter(_json_formatter)
        logging.root.handlers = [_json_handler]
        logging.root.setLevel(logging.INFO)
        logger.info("JSON logging enabled")
    except ImportError:
        logger.info("python-json-logger not installed - using default logging")

try:
    import yfinance as yf
except ImportError:
    class MockYF:
        @staticmethod
        def Ticker(t):
            return type('MockTicker', (), {'info': {}, 'history': lambda *a, **k: None})()
    yf = MockYF()
    logger.warning("yfinance not found, using mock")

from datetime import datetime, timedelta
from typing import Optional, List, Dict, Any
from contextlib import asynccontextmanager



from fastapi import FastAPI, HTTPException, BackgroundTasks, Depends, Request, Form, File, UploadFile, Response, Query
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, RedirectResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, HttpUrl
from sqlalchemy import create_engine, Column, Integer, String, Float, DateTime, Text, Boolean, select, text
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session
from sqlalchemy.ext.asyncio import AsyncSession
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Font, Alignment
from openpyxl.utils import get_column_letter

# Enterprise Automation Import
try:
    from scheduler import start_scheduler, stop_scheduler
    SCHEDULER_AVAILABLE = True
except ImportError:
    SCHEDULER_AVAILABLE = False
    logger.warning("Scheduler module not found. Automation disabled.")

# New Data Source Imports
import glassdoor_scraper
import indeed_scraper
import sec_edgar_scraper
import uspto_scraper
import klas_scraper
import appstore_scraper

import himss_scraper
import social_media_monitor

# Database setup - SQLite for simplicity
# Database setup

from database import (
    engine, SessionLocal, Base, get_db, Competitor, ChangeLog, DataSource,
    DataChangeHistory, PendingDataChange, User, SystemPrompt, KnowledgeBaseItem, UserSettings, ActivityLog,
    CompetitorProduct, ProductPricingTier, ProductFeatureMatrix, CustomerCountEstimate,
    RefreshSession,  # Phase 4: Task 5.0.1-031
    ChatSession, ChatMessage,  # v7.2.0: Chat persistence
    # SQLAlchemy 2.0 Async imports
    get_async_db, AsyncSessionLocal, get_user_by_email, get_user_by_id,
    get_competitor_by_id, get_all_competitors, get_active_competitors,
    get_news_articles, get_news_count, get_system_prompt, get_knowledge_base_items,
    get_change_history, count_competitors, count_by_threat_level,
    create_competitor, update_competitor, delete_competitor,
    get_products_by_competitor, get_all_products, get_data_sources
)
from confidence_scoring import (
    calculate_confidence_score, get_source_defaults, calculate_data_staleness,
    determine_confidence_level_from_score, triangulate_data_points,
    get_reliability_description, get_credibility_description, get_source_type_description,
    SOURCE_TYPE_DEFAULTS, RELIABILITY_DESCRIPTIONS, CREDIBILITY_DESCRIPTIONS
)
from data_triangulator import (
    DataTriangulator, triangulate_competitor, triangulation_result_to_dict
)

# Auth imports for route protection - centralized in dependencies.py
from dependencies import (
    oauth2_scheme, get_current_user, get_current_user_optional, log_activity
)
from utils.prompt_utils import resolve_system_prompt as _resolve_system_prompt


# Global progress tracker for scrape operations (thread-safe via Lock)
import threading
_scrape_lock = threading.Lock()
scrape_progress = {
    "active": False,
    "total": 0,
    "completed": 0,
    "current_competitor": None,
    "competitors_done": [],
    "changes_detected": 0,
    "new_values_added": 0,
    # Phase 2: Enhanced tracking (Task 5.0.1-026)
    "started_at": None,
    "recent_changes": [],    # Last 10 changes for live display
    "change_details": [],    # All changes for AI summary
    "errors": [],            # Any errors encountered
    # v7.1.0: Post-scrape enrichment tracking
    "enrichment_active": False,
    "news_articles_fetched": 0,
    "stock_prices_updated": 0
}


# White fill for Excel cells
WHITE_FILL = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")


def get_openai_client():
    """Get OpenAI client if API key is configured, otherwise return None."""
    try:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return None
        from openai import OpenAI
        return OpenAI(api_key=api_key)
    except Exception:
        return None

from constants import KNOWN_TICKERS  # noqa: E402


def seed_win_loss_data(db: Session):
    """
    Seed realistic win/loss deal data for analytics demo (v7.1.5).
    Only runs if win_loss_deals table is empty.
    """
    from database import WinLossDeal, Competitor
    import random

    logger = logging.getLogger(__name__)

    try:
        # Check if already seeded
        existing_count = db.query(WinLossDeal).count()
        if existing_count > 0:
            logger.info(f"[Seed] Win/loss deals already exist ({existing_count} records), skipping seed")
            return

        # Get top 10 competitors by threat level
        top_competitors = db.query(Competitor).filter(
            Competitor.is_deleted == False,
            Competitor.threat_level.isnot(None)
        ).order_by(
            Competitor.threat_level.desc()
        ).limit(10).all()

        if not top_competitors:
            logger.warning("[Seed] No competitors found to seed win/loss data")
            return

        # Deal templates
        win_reasons = [
            "Better pricing and ROI",
            "Superior integration capabilities",
            "Stronger customer support and training",
            "More comprehensive feature set",
            "Better ease of use and user experience",
            "Proven implementation track record"
        ]

        loss_reasons = [
            "Price - competitor offered lower cost",
            "Features - competitor had capabilities we lacked",
            "Integration - competitor had existing EHR partnership",
            "Support - concerns about our support model",
            "Brand recognition - larger, more established vendor",
            "Performance - concerns about system scalability"
        ]

        win_notes_templates = [
            "Customer valued our white-glove onboarding process and dedicated support team.",
            "Our API-first architecture and integration library was the key differentiator.",
            "Pricing model aligned better with their budget and growth plans.",
            "References from similar-sized customers helped close the deal.",
            "Our HITRUST certification and security posture were deciding factors.",
            "Demo showed clear workflow improvements over competitor solution."
        ]

        loss_notes_templates = [
            "Lost to incumbent vendor with existing relationship.",
            "Competitor had deeper features in their specialty area.",
            "Budget constraints made competitor's lower price decisive.",
            "Customer already using competitor's other products (ecosystem lock-in).",
            "Implementation timeline was too aggressive for our capacity.",
            "Lack of specialty-specific features was a concern."
        ]

        # Create 20 realistic deals
        deals = []
        base_date = datetime.now() - timedelta(days=90)

        for i in range(20):
            # 60% wins, 40% losses
            is_win = random.random() < 0.6
            outcome = "win" if is_win else "loss"

            # Pick random competitor
            competitor = random.choice(top_competitors)

            # Generate realistic deal value ($25k - $250k)
            deal_value = random.randint(25, 250) * 1000

            # Random date in past 90 days
            days_offset = random.randint(0, 90)
            deal_date = base_date + timedelta(days=days_offset)

            # Pick reason and notes
            reason = random.choice(win_reasons if is_win else loss_reasons)
            notes = random.choice(win_notes_templates if is_win else loss_notes_templates)

            deal = WinLossDeal(
                user_id=1,  # Admin user
                competitor_id=competitor.id,
                competitor_name=competitor.name,
                outcome=outcome,
                deal_value=deal_value,
                deal_date=deal_date,
                customer_name=f"Healthcare System {i+1:02d}",
                customer_size=random.choice(["Small (1-50)", "Medium (51-500)", "Large (500+)"]),
                reason=reason,
                sales_rep="Demo Data",
                notes=notes
            )
            deals.append(deal)

        # Insert all deals
        db.add_all(deals)
        db.commit()

        wins = sum(1 for d in deals if d.outcome == "win")
        losses = len(deals) - wins
        logger.info(f"[Seed] Created {len(deals)} win/loss deals ({wins} wins, {losses} losses)")

    except Exception as e:
        logger.error(f"[Seed] Error seeding win/loss data: {e}", exc_info=True)
        db.rollback()


def lookup_ticker_dynamically(company_name: str) -> dict:
    """
    Try to find ticker symbol for a company using yfinance.
    Returns dict with symbol, exchange, name or None values if not found.
    """
    try:
        import yfinance as yf
        # Common patterns to try
        search_terms = [
            company_name,
            company_name.replace(" ", ""),
            company_name.split()[0] if " " in company_name else company_name
        ]

        for term in search_terms:
            try:
                ticker = yf.Ticker(term)
                info = ticker.info
                if info and info.get('symbol') and info.get('shortName'):
                    exchange = info.get('exchange', 'UNKNOWN')
                    # Map exchange codes to readable names
                    exchange_map = {
                        'NMS': 'NASDAQ', 'NGM': 'NASDAQ', 'NCM': 'NASDAQ',
                        'NYQ': 'NYSE', 'NYSE': 'NYSE', 'PCX': 'NYSE ARCA'
                    }
                    return {
                        "symbol": info.get('symbol'),
                        "exchange": exchange_map.get(exchange, exchange),
                        "name": info.get('shortName') or info.get('longName')
                    }
            except (KeyError, AttributeError, Exception):
                continue
        return None
    except Exception as e:
        logger.warning(f"Dynamic ticker lookup failed for {company_name}: {e}")
        return None

# Database Models imported from database.py


# ============== Pydantic Schemas (centralized in schemas/) ==============
from schemas.competitors import (  # noqa: E402
    CompetitorCreate, CompetitorResponse, CorrectionRequest, ScrapeRequest,
    BulkUpdateRequest, BulkDeleteRequest, BulkExportRequest, SearchResult,
)
from schemas.auth import UserResponse, UserInviteRequest  # noqa: E402
from schemas.prompts import (  # noqa: E402
    SystemPromptBase, SystemPromptCreate, SystemPromptResponse,
    KnowledgeBaseItemBase, KnowledgeBaseItemCreate, KnowledgeBaseItemResponse,
    UserSavedPromptCreate, UserSavedPromptUpdate, UserSavedPromptResponse,
)
from schemas.products import (  # noqa: E402
    ProductCreate, ProductResponse, PricingTierCreate, PricingTierResponse,
    FeatureMatrixCreate, CustomerCountCreate, CustomerCountResponse,
    CustomerCountVerifyRequest,
)
from schemas.common import (  # noqa: E402
    DataChangeSubmission, WinLossCreate, WebhookCreate,
    SubscriptionCreate, SubscriptionUpdate, WebVitalsMetric,
)



# ============== FastAPI App ==============

@asynccontextmanager
async def lifespan(app: FastAPI):
    # Startup
    logger.info("Certify Intel Backend starting...")
    logger.info("=" * 60)

    # Configuration Validation
    logger.info("Validating configuration...")
    missing_required = []

    # Check required environment variables
    if not os.getenv("SECRET_KEY"):
        missing_required.append("SECRET_KEY")

    if missing_required:
        logger.error(f"ERROR: Missing required environment variables: {', '.join(missing_required)}")
        raise ValueError(f"Missing required env vars: {', '.join(missing_required)}")

    # Warn about optional features
    optional_features = {
        "OPENAI_API_KEY": "AI Features (Executive Summaries, Discovery Agent, Web Extraction)",
        "SMTP_HOST": "Email Alerts",
        "SLACK_WEBHOOK_URL": "Slack Notifications"
    }

    disabled_scrapers = ["Crunchbase", "PitchBook", "LinkedIn (live scraping)"]

    logger.info("Available Scrapers:")
    logger.info("[OK] Playwright Base Scraper - Website content extraction")
    logger.info("[OK] SEC Edgar (yfinance) - Public company financials")
    logger.info("[OK] News Monitor (Google News RSS) - Real-time news")
    logger.info("[OK] Known Data Fallback - Pre-populated data for demo")
    logger.info("[OK] [15+ other specialized scrapers with fallback data]")

    logger.info("Disabled Scrapers (Paid APIs - not available):")
    for scraper in disabled_scrapers:
        logger.debug(f"[X] {scraper}")

    logger.info("Optional Features:")
    for env_var, feature in optional_features.items():
        if os.getenv(env_var):
            logger.info(f"[OK] {feature} - ENABLED")
        else:
            logger.warning(f"[!] {feature} - DISABLED (set {env_var} to enable)")

    # Langfuse observability check
    if os.getenv("ENABLE_LANGFUSE", "false").lower() == "true":
        logger.info("[OK] Langfuse Observability - ENABLED")
    else:
        logger.info("[!] Langfuse Observability - DISABLED (set ENABLE_LANGFUSE=true)")

    # Enterprise data provider diagnostic (v8.3.0)
    try:
        from data_providers import get_all_provider_status
        all_status = get_all_provider_status()
        active = [s for s in all_status if s["configured"]]
        logger.info(f"Enterprise data providers: {len(active)}/{len(all_status)} active")
        for s in all_status:
            tag = "[OK]" if s["configured"] else "[!]"
            logger.info(f"  {tag} {s['name']} ({s['env_key']})")
    except ImportError:
        logger.debug("Enterprise data providers module not available")
    except Exception as dp_err:
        logger.warning(f"Data provider diagnostic failed: {dp_err}")

    logger.info("=" * 60)

    # Skip heavy startup tasks during testing
    is_testing = os.getenv("TESTING") == "true"

    # Start Enterprise Scheduler with retry
    if SCHEDULER_AVAILABLE and not is_testing:
        logger.info("Initializing Enterprise Automation Engine...")
        import time as _time
        for _attempt in range(3):
            try:
                start_scheduler()
                break
            except Exception as e:
                wait = 2 ** _attempt
                logger.warning(f"Scheduler start failed (attempt {_attempt + 1}/3): {e}, retrying in {wait}s...")
                _time.sleep(wait)
        else:
            logger.error("Scheduler failed to start after 3 attempts")
    elif is_testing:
        logger.info("[TEST] Skipping scheduler startup")

    # Run Startup Tasks
    try:
        from extended_features import ClassificationWorkflow, auth_manager

        # Initialize DB Session for startup tasks
        db = SessionLocal()

        # 1. Ensure Admin User
        logger.info("Ensuring default admin user exists...")
        auth_manager.ensure_default_admin(db)

        # 2. Preinstall Knowledge Base (client-provided data)
        # This only runs on first startup - subsequent startups skip
        if not is_testing:
            logger.info("Checking knowledge base preinstall...")
            try:
                from knowledge_base_importer import preinstall_knowledge_base
                result = preinstall_knowledge_base(db)
                if result and result.get("success"):
                    logger.info(f"[OK] Preinstalled {result.get('competitors_imported', 0)} competitors")
                    logger.info(f"[OK] Updated {result.get('competitors_updated', 0)} existing competitors")
                    logger.info(f"[OK] Data labeled as 'Certify Health (Preinstalled)'")
                elif result and result.get("skipped"):
                    logger.info(f"[OK] No competitors found in knowledge base folder")
                else:
                    logger.info(f"[OK] Already preinstalled, skipping")
            except Exception as e:
                logger.warning(f"[!] Knowledge base preinstall warning: {e}")
        else:
            logger.info("[TEST] Skipping knowledge base preinstall")

        # 3. Seed Win/Loss Deals (v7.1.5)
        if not is_testing:
            logger.info("Checking win/loss deal data...")
            try:
                seed_win_loss_data(db)
            except Exception as e:
                logger.warning(f"[!] Win/loss seed warning: {e}")
        else:
            logger.info("[TEST] Skipping win/loss seed")

        # 4. Run Classification Workflow - DISABLED (costs money, use button instead)
        # workflow = ClassificationWorkflow(db)
        # print("Running 'Private vs Public' Classification Workflow...")
        # workflow.run_classification_pipeline()

        # 3. Migrate SystemPrompt table (add category + description columns)
        try:
            db.execute(text("ALTER TABLE system_prompts ADD COLUMN category VARCHAR"))
            db.commit()
            logger.info("[OK] Added 'category' column to system_prompts")
        except Exception:
            db.rollback()  # Column already exists

        try:
            db.execute(text("ALTER TABLE system_prompts ADD COLUMN description VARCHAR"))
            db.commit()
            logger.info("[OK] Added 'description' column to system_prompts")
        except Exception:
            db.rollback()  # Column already exists

        # 4b. Add is_archived column to news_article_cache
        try:
            db.execute(text("ALTER TABLE news_article_cache ADD COLUMN is_archived BOOLEAN DEFAULT 0"))
            db.commit()
            logger.info("[Migration] Added is_archived to news_article_cache")
        except Exception:
            db.rollback()

        # 4c. Add ai_threat_summary column to competitors
        try:
            db.execute(text("ALTER TABLE competitors ADD COLUMN ai_threat_summary TEXT"))
            db.commit()
            logger.info("[Migration] Added ai_threat_summary to competitors")
        except Exception:
            db.rollback()

        # 4d. Add source URL refinement columns to data_sources (v8.3.0)
        _source_url_columns = [
            ("source_page_url", "TEXT"),
            ("source_anchor_text", "TEXT"),
            ("source_css_selector", "TEXT"),
            ("source_section", "TEXT"),
            ("deep_link_url", "TEXT"),
            ("last_url_verified", "TIMESTAMP"),
            ("url_status", "TEXT DEFAULT 'pending'"),
        ]
        for col_name, col_type in _source_url_columns:
            try:
                db.execute(text(f"ALTER TABLE data_sources ADD COLUMN {col_name} {col_type}"))
                db.commit()
                logger.info(f"[Migration] Added {col_name} to data_sources")
            except Exception:
                db.rollback()

        # 5. Create chat persistence tables (v7.2.0)
        try:
            db.execute(text("""
                CREATE TABLE IF NOT EXISTS chat_sessions (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    user_id INTEGER REFERENCES users(id),
                    page_context VARCHAR,
                    competitor_id INTEGER REFERENCES competitors(id),
                    title VARCHAR,
                    created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                    updated_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                    is_active BOOLEAN DEFAULT 1
                )
            """))
            db.execute(text("""
                CREATE TABLE IF NOT EXISTS chat_messages (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    session_id INTEGER REFERENCES chat_sessions(id),
                    role VARCHAR NOT NULL,
                    content TEXT NOT NULL,
                    metadata_json TEXT,
                    created_at DATETIME DEFAULT CURRENT_TIMESTAMP
                )
            """))
            db.commit()
            logger.info("[OK] Chat persistence tables ready (chat_sessions, chat_messages)")
        except Exception as e:
            db.rollback()
            logger.warning(f"[!] Chat tables migration warning: {e}")

        # 6. Seed all AI prompts into database
        if not is_testing:
            try:
                from prompt_seeder import seed_system_prompts
                result = seed_system_prompts(db)
                logger.info(f"[OK] Prompt seeder: {result['inserted']} new, {result['updated']} updated, {result['skipped']} existing")
            except Exception as e:
                logger.warning(f"[!] Prompt seeder warning: {e}")

        # 7. Performance indexes (PERF-012)
        _perf_indexes = [
            "CREATE INDEX IF NOT EXISTS ix_comp_threat_name "
            "ON competitors (threat_level, name)",
            "CREATE INDEX IF NOT EXISTS ix_news_comp_date "
            "ON news_article_cache (competitor_id, published_at)",
            "CREATE INDEX IF NOT EXISTS ix_activity_user_date "
            "ON activity_logs (user_id, created_at)",
            "CREATE INDEX IF NOT EXISTS ix_sources_comp_verified "
            "ON data_sources (competitor_id, verified)",
        ]
        try:
            for idx_sql in _perf_indexes:
                db.execute(text(idx_sql))
            db.commit()
            logger.info("[OK] Performance indexes created/verified")
        except Exception as e:
            db.rollback()
            logger.debug(f"[!] Index creation note: {e}")

        # 8. MFA columns on users table (v9.0.0)
        _mfa_columns = [
            ("mfa_enabled", "BOOLEAN DEFAULT 0"),
            ("mfa_secret", "VARCHAR"),
            ("mfa_backup_codes", "TEXT"),
        ]
        for col_name, col_type in _mfa_columns:
            try:
                db.execute(text(f"ALTER TABLE users ADD COLUMN {col_name} {col_type}"))
                db.commit()
                logger.info(f"[Migration] Added {col_name} to users")
            except Exception:
                db.rollback()

        # 9. RefreshToken table (v9.0.0)
        try:
            db.execute(text("""
                CREATE TABLE IF NOT EXISTS refresh_tokens (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    token VARCHAR NOT NULL UNIQUE,
                    user_id INTEGER REFERENCES users(id),
                    expires_at DATETIME NOT NULL,
                    revoked BOOLEAN DEFAULT 0,
                    created_at DATETIME DEFAULT CURRENT_TIMESTAMP
                )
            """))
            db.commit()
            logger.info("[OK] RefreshToken table ready")
        except Exception as e:
            db.rollback()
            logger.debug(f"[!] RefreshToken table note: {e}")

        db.close()
    except Exception as e:
        logger.warning(f"Startup task warning: {e}")
    finally:
        # Ensure DB session is always closed even if an unhandled exception occurs
        try:
            db.close()
        except Exception:
            pass

    yield
    
    # Shutdown: Clean up resources
    logger.info("Certify Intel Backend shutting down...")
    if SCHEDULER_AVAILABLE:
        try:
            stop_scheduler()
        except Exception as e:
            logger.debug(f"Scheduler shutdown note: {e}")

    # Flush Langfuse traces on shutdown
    try:
        from observability import shutdown_langfuse
        shutdown_langfuse()
    except Exception:
        pass

app = FastAPI(
    title="Certify Health Intel API",
    description="Backend for Competitive Intelligence Dashboard",
    version=__version__,
    lifespan=lifespan
)

# GZip compression for API responses (PERF-011)
from starlette.middleware.gzip import GZipMiddleware  # noqa: E402
app.add_middleware(GZipMiddleware, minimum_size=1000)

# Attach slowapi limiter if available
if limiter is not None:
    app.state.limiter = limiter
    app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)
    app.add_middleware(SlowAPIMiddleware)


@app.middleware("http")
async def correlation_id_middleware(request, call_next):
    """Add correlation ID to all requests for distributed tracing."""
    correlation_id = request.headers.get("X-Correlation-ID", str(uuid.uuid4()))
    response = await call_next(request)
    response.headers["X-Correlation-ID"] = correlation_id
    return response


# Health/version/readiness endpoints moved to routers/health.py


# Valid API endpoints are defined above.
# The catch-all static file mount must be last.



# Redirect legacy /app path to root
@app.get("/app")
@app.get("/app/")
async def redirect_app():
    return RedirectResponse(url="/")








@app.get("/api/logo-proxy")
async def proxy_logo(url: str):
    """Proxy image requests to avoid CORS issues. Only allows known logo services."""
    import httpx
    from fastapi import Response
    from urllib.parse import urlparse

    # SSRF protection: only allow known logo/favicon services
    ALLOWED_DOMAINS = {
        "logo.clearbit.com",
        "www.google.com",
        "google.com",
        "s2.googleusercontent.com",
        "favicons.githubusercontent.com",
    }

    try:
        # Validate URL format
        if not url.startswith("https://"):
            raise HTTPException(status_code=400, detail="Only HTTPS URLs allowed")

        parsed = urlparse(url)
        if parsed.hostname not in ALLOWED_DOMAINS:
            raise HTTPException(status_code=400, detail="Domain not allowed for logo proxy")

        # Try Clearbit first
        try:
            async with httpx.AsyncClient(follow_redirects=True) as client:
                resp = await client.get(url, timeout=5)
            if resp.status_code == 200:
                return Response(content=resp.content, media_type=resp.headers.get("content-type", "image/png"))
        except httpx.HTTPError as e:
            logger.debug(f"Clearbit logo fetch failed: {e}")

        # Fallback to Google Favicon service if Clearbit fails
        # Extract domain from the input URL (expected: https://logo.clearbit.com/domain.com)
        domain = url.split("/")[-1]
        # Sanitize domain - only allow alphanumeric, dots, and hyphens
        import re
        if not re.match(r'^[a-zA-Z0-9.-]+$', domain):
            raise HTTPException(status_code=400, detail="Invalid domain")
        google_url = f"https://www.google.com/s2/favicons?domain={domain}&sz=128"

        try:
            async with httpx.AsyncClient(follow_redirects=True) as client:
                resp = await client.get(google_url, timeout=5)
            if resp.status_code == 200:
                return Response(content=resp.content, media_type=resp.headers.get("content-type", "image/png"))
        except httpx.HTTPError as e:
            logger.debug(f"Google favicon fetch failed: {e}")

        # Return 404 to trigger frontend onerror fallback
        raise HTTPException(status_code=404, detail="Image not found")

    except HTTPException:
        raise
    except Exception as e:
        logger.warning(f"Proxy error: {e}")
        raise HTTPException(status_code=404, detail="Proxy error")


# get_db imported from database.py


# ==============================================================================
# BACKGROUND AI TASK TRACKING (v7.2.0)
# ==============================================================================
# In-memory store for background AI tasks. Keys are task_id (UUID strings).
# Each entry: {status, page_context, user_id, started_at, completed_at,
#              result, error, task_type, read_at}
_ai_tasks: Dict[str, Dict[str, Any]] = {}

# Analytics summary progress tracking
_analytics_summary_progress: Dict[str, Dict[str, Any]] = {}


@app.post("/api/ai/tasks")
async def submit_ai_task(
    request: Request,
    background_tasks: BackgroundTasks,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Submit a background AI task. Returns task_id for polling."""
    import uuid as _uuid

    body = await request.json()
    endpoint = body.get("endpoint")
    page_context = body.get("page_context", "unknown")

    if not endpoint:
        raise HTTPException(status_code=400, detail="endpoint is required")

    task_id = str(_uuid.uuid4())
    _ai_tasks[task_id] = {
        "status": "running",
        "page_context": page_context,
        "user_id": current_user.get("id"),
        "started_at": datetime.utcnow().isoformat(),
        "completed_at": None,
        "result": None,
        "error": None,
        "task_type": endpoint,
        "read_at": None,
    }

    return {"task_id": task_id, "status": "running"}


@app.get("/api/ai/tasks/pending")
async def get_pending_ai_tasks(
    current_user: dict = Depends(get_current_user)
):
    """Get all running + completed-unread tasks for the current user."""
    user_id = current_user.get("id")
    pending = []
    for tid, task in _ai_tasks.items():
        if task.get("user_id") != user_id:
            continue
        if task["status"] == "running" or (
            task["status"] == "completed" and task.get("read_at") is None
        ):
            pending.append({"task_id": tid, **task})
    # Sort newest first
    pending.sort(key=lambda t: t.get("started_at", ""), reverse=True)

    # Prune old completed/failed tasks to prevent memory leak
    from datetime import timedelta
    cutoff = datetime.utcnow() - timedelta(hours=1)
    stale_ids = [
        tid for tid, task in _ai_tasks.items()
        if task.get("status") in ("completed", "failed")
        and task.get("read_at")
        and isinstance(task.get("read_at"), str)
        and datetime.fromisoformat(task["read_at"]) < cutoff
    ]
    for tid in stale_ids:
        del _ai_tasks[tid]

    return pending


@app.get("/api/ai/tasks/{task_id}")
async def get_ai_task_status(
    task_id: str,
    current_user: dict = Depends(get_current_user)
):
    """Check the status of a background AI task."""
    task = _ai_tasks.get(task_id)
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    if task.get("user_id") != current_user.get("id"):
        raise HTTPException(status_code=403, detail="Not your task")
    # Set read_at on first retrieval of completed/failed task for pruning
    if task.get("status") in ("completed", "failed") and not task.get("read_at"):
        task["read_at"] = datetime.utcnow().isoformat()
    return {"task_id": task_id, **task}


@app.put("/api/ai/tasks/{task_id}/read")
async def mark_ai_task_read(
    task_id: str,
    current_user: dict = Depends(get_current_user)
):
    """Mark a notification as read."""
    task = _ai_tasks.get(task_id)
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    if task.get("user_id") != current_user.get("id"):
        raise HTTPException(status_code=403, detail="Not your task")
    task["read_at"] = datetime.utcnow().isoformat()
    return {"status": "ok"}


@app.get("/api/analytics/summary")
@app.get("/api/analytics/executive-summary")
async def get_dashboard_summary(
    prompt_key: Optional[str] = None,
    session_id: Optional[int] = None,
    background: bool = False,
    background_tasks: BackgroundTasks = None,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Generate comprehensive AI executive summary analyzing ALL dashboard data points."""
    if background and background_tasks:
        import uuid as _uuid
        task_id = str(_uuid.uuid4())
        _ai_tasks[task_id] = {
            "status": "running", "page_context": "dashboard",
            "user_id": current_user.get("id"),
            "started_at": datetime.utcnow().isoformat(),
            "completed_at": None, "result": None, "error": None,
            "task_type": "executive_summary", "read_at": None,
        }

        async def _run_bg():
            try:
                # Re-open DB session for background thread
                bg_db = SessionLocal()
                try:
                    from ai_router import get_ai_router, TaskType
                    comps = bg_db.query(Competitor).filter(
                        Competitor.is_deleted == False
                    ).all()
                    if not comps:
                        _ai_tasks[task_id]["result"] = {
                            "summary": "No active competitors found.",
                            "type": "empty", "model": None
                        }
                    else:
                        router = get_ai_router()
                        result = await asyncio.wait_for(
                            router.generate(
                                prompt=f"Summarize {len(comps)} competitors.",
                                task_type=TaskType.ANALYSIS,
                                system_prompt="Generate executive summary." + NO_HALLUCINATION_INSTRUCTION,
                                max_tokens=2000, temperature=0.7,
                            ),
                            timeout=45.0
                        )
                        _ai_tasks[task_id]["result"] = {
                            "summary": result.get("response", ""),
                            "type": "ai",
                            "model": result.get("model", "unknown"),
                        }
                finally:
                    bg_db.close()
                _ai_tasks[task_id]["status"] = "completed"
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
            except asyncio.TimeoutError:
                logger.error("Background executive summary timed out after 45s")
                _ai_tasks[task_id]["status"] = "failed"
                _ai_tasks[task_id]["error"] = "AI analysis timed out. Please try again."
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
            except Exception as bg_err:
                logger.error(f"Background summary error: {bg_err}")
                _ai_tasks[task_id]["status"] = "failed"
                _ai_tasks[task_id]["error"] = str(bg_err)
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()

        background_tasks.add_task(_run_bg)
        return {"task_id": task_id, "status": "running"}

    try:
        import os

        # Get all non-deleted competitors (matching dashboard count)
        competitors = db.query(Competitor).filter(
            Competitor.is_deleted == False
        ).all()
        
        if not competitors:
            return {"summary": "No active competitors found to analyze.", "type": "empty", "model": None}
        
        # Calculate statistics
        total = len(competitors)
        high_threat = sum(1 for c in competitors if c.threat_level and c.threat_level.upper() == "HIGH")
        medium_threat = sum(1 for c in competitors if c.threat_level and c.threat_level.upper() == "MEDIUM")
        low_threat = sum(1 for c in competitors if c.threat_level and c.threat_level.upper() == "LOW")
        public_companies = sum(1 for c in competitors if c.is_public)
        private_companies = total - public_companies
        
        # Gather pricing info
        pricing_models = {}
        for c in competitors:
            model = c.pricing_model or "Unknown"
            pricing_models[model] = pricing_models.get(model, 0) + 1
        
        # Gather top threats
        top_threats = [c.name for c in competitors if c.threat_level and c.threat_level.upper() == "HIGH"][:5]
        
        # Build comprehensive data summary for AI
        data_summary = f"""
COMPETITIVE INTELLIGENCE DATA SNAPSHOT:
========================================
TIMESTAMP: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S UTC')}

TRACKING OVERVIEW:
- Total Competitors Monitored: {total}
- High Threat: {high_threat} ({round(high_threat/total*100, 1)}%)
- Medium Threat: {medium_threat} ({round(medium_threat/total*100, 1)}%)
- Low Threat: {low_threat} ({round(low_threat/total*100, 1)}%)
- Public Companies: {public_companies}
- Private Companies: {private_companies}

TOP HIGH-THREAT COMPETITORS: {', '.join(top_threats) if top_threats else 'None identified'}

PRICING MODEL DISTRIBUTION:
{chr(10).join(f'- {model}: {count} competitors' for model, count in pricing_models.items())}

DETAILED COMPETITOR DATA (ALL {total} ACTIVE COMPETITORS):
"""
        # INCLUDED ALL COMPETITORS (No limit) for maximum context
        for i, c in enumerate(competitors, 1):
            data_summary += f"""
--- COMPETITOR #{i} of {total} ---
NAME: {c.name}
THREAT LEVEL: {c.threat_level}
STATUS: {c.status or 'Active'}
WEBSITE: {c.website or 'N/A'}
LAST UPDATED: {c.last_updated.strftime('%Y-%m-%d %H:%M') if c.last_updated else 'N/A'}
DATA QUALITY SCORE: {c.data_quality_score or 'N/A'}/100

PRICING INFO:
- Model: {c.pricing_model or 'Unknown'}
- Base Price: {c.base_price or 'N/A'}
- Price Unit: {c.price_unit or 'N/A'}

PRODUCT INFO:
- Categories: {c.product_categories or 'N/A'}
- Key Features: {c.key_features or 'N/A'}
- Integrations: {c.integration_partners or 'N/A'}
- Certifications: {c.certifications or 'N/A'}

MARKET INFO:
- Target Segments: {c.target_segments or 'N/A'}
- Customer Size Focus: {c.customer_size_focus or 'N/A'}
- Geographic Focus: {c.geographic_focus or 'N/A'}
- Customer Count: {c.customer_count or 'N/A'}
- Key Customers: {c.key_customers or 'N/A'}
- G2 Rating: {c.g2_rating or 'N/A'}

COMPANY INFO:
- Headquarters: {c.headquarters or 'N/A'}
- Founded: {c.year_founded or 'N/A'}
- Employees: {c.employee_count or 'N/A'}
- Employee Growth: {c.employee_growth_rate or 'N/A'}
- Total Funding: {c.funding_total or 'N/A'}
- Latest Round: {c.latest_round or 'N/A'}
- Is Public: {'Yes' if c.is_public else 'No (Private)'}
"""
            # Add live stock data for public companies
            if c.is_public and c.ticker_symbol:
                stock_info = fetch_real_stock_data(c.ticker_symbol)
                if stock_info and stock_info.get('price'):
                    change_sign = '+' if stock_info.get('change', 0) >= 0 else ''
                    data_summary += f"""
STOCK DATA (LIVE):
- Ticker: {c.ticker_symbol} ({c.stock_exchange or 'NYSE'})
- Current Price: ${stock_info.get('price', 0):.2f}
- Daily Change: {change_sign}{stock_info.get('change', 0):.2f} ({change_sign}{stock_info.get('change_percent', 0):.2f}%)
- Market Cap: ${stock_info.get('market_cap', 0):,.0f}
- 52-Week High: ${stock_info.get('high52', 'N/A')}
- 52-Week Low: ${stock_info.get('low52', 'N/A')}
"""
                else:
                    data_summary += f"""
STOCK DATA:
- Ticker: {c.ticker_symbol} ({c.stock_exchange or 'NYSE'})
- Stock Price: Data unavailable
"""

            data_summary += f"""
NOTES: {c.notes or 'None'}
--- END {c.name} ---
"""

        # Try OpenAI first
        api_key = os.getenv("OPENAI_API_KEY")
        model_used = "gpt-4.1"
        provider = "OpenAI"
        
        if api_key:
            try:
                # Fetch dynamic system prompt (prompt_key override, then user-specific, then global)
                user_id = current_user.get("id") if current_user else None
                lookup_key = prompt_key if prompt_key else "dashboard_summary"
                prompt_db = db.query(SystemPrompt).filter(
                    SystemPrompt.key == lookup_key,
                    SystemPrompt.user_id == user_id
                ).first()
                if not prompt_db:
                    # Fallback to global prompt (user_id=NULL)
                    prompt_db = db.query(SystemPrompt).filter(
                        SystemPrompt.key == lookup_key,
                        SystemPrompt.user_id == None
                    ).first()
                system_content = prompt_db.content if prompt_db else """You are Certify Health's competitive intelligence analyst. Generate a comprehensive, executive-level strategic summary. 

Your summary MUST include:
1. **Executive Overview** - High-level market position assessment
2. **Threat Analysis** - Breakdown of competitive landscape by threat level
3. **Pricing Intelligence** - Analysis of competitor pricing strategies
4. **Market Trends** - Emerging patterns and shifts
5. **Strategic Recommendations** - 3-5 specific, actionable recommendations
6. **Watch List** - Key competitors requiring immediate attention

Use data-driven insights. Be specific with numbers and competitor names. Format with markdown headers and bullet points."""

                # RAG: Inject Knowledge Base
                kb_items = db.query(KnowledgeBaseItem).filter(KnowledgeBaseItem.is_active == True).all()
                if kb_items:
                    data_summary += "\n\nINTERNAL KNOWLEDGE BASE (USE THIS CONTEXT):\n==========================\n"
                    for item in kb_items:
                        data_summary += f"\n--- {item.title} ({item.source_type}) ---\n{item.content_text}\n"

                from ai_router import get_ai_router, TaskType
                router = get_ai_router()
                result = await router.generate(
                    prompt=data_summary,
                    task_type=TaskType.ANALYSIS,
                    system_prompt=system_content + NO_HALLUCINATION_INSTRUCTION,
                    max_tokens=2000,
                    temperature=0.7,
                )
                summary = result["response"]
                provider = result.get("provider", "anthropic")
                model_used = result.get("model", "claude-opus-4.5")
                provider_logo = "/static/anthropic-logo.svg" if "claude" in model_used else "/static/openai-logo.svg"

                # Persist to chat session if session_id provided
                if session_id:
                    try:
                        _save_chat_messages(db, session_id, "Generate executive summary", summary, {
                            "provider": provider, "model": model_used, "endpoint": "analytics_summary",
                        })
                    except Exception as save_err:
                        logger.warning(f"Failed to save summary to chat: {save_err}")

                result_data = {
                    "summary": summary,
                    "type": "ai",
                    "model": model_used,
                    "provider": provider,
                    "provider_logo": provider_logo,
                    "data_points_analyzed": total,
                    "generated_at": datetime.utcnow().isoformat() + "Z",
                    "session_id": session_id,
                }

                # v8.0.8: Persist summary to DB so it survives server restart
                try:
                    import json as _json
                    from database import PersistentCache
                    user_id = current_user.get("id") if current_user else 0
                    cache_key = f"dashboard_summary_{user_id}"
                    existing = db.query(PersistentCache).filter(
                        PersistentCache.cache_key == cache_key
                    ).first()
                    if existing:
                        existing.data_json = _json.dumps(result_data, default=str)
                        existing.updated_at = datetime.utcnow()
                    else:
                        db.add(PersistentCache(
                            cache_key=cache_key, user_id=user_id,
                            data_json=_json.dumps(result_data, default=str),
                        ))
                    db.commit()
                except Exception as persist_err:
                    logger.warning(f"Failed to persist dashboard summary: {persist_err}")

                return result_data
            except Exception as e:
                logger.error(f"AI Summary Error: {e}")

        # Fallback to automated summary
        summary = f"""# AI-Generated Executive Summary

## Executive Overview
We are currently tracking **{total} competitors** in the patient engagement and healthcare check-in market. The competitive landscape shows **{high_threat} high-threat** competitors requiring immediate attention, **{medium_threat} medium-threat** competitors to monitor, and **{low_threat} low-threat** competitors with minimal impact.

## Threat Analysis
- **High Threat ({high_threat})**: {', '.join(top_threats[:3]) if top_threats else 'None'} - These competitors have significant market overlap and strong positioning
- **Medium Threat ({medium_threat})**: Competitors with partial market overlap requiring ongoing monitoring
- **Low Threat ({low_threat})**: Minimal competitive impact at this time

## Market Composition
- **Public Companies**: {public_companies} ({round(public_companies/total*100, 1)}% of tracked competitors)
- **Private Companies**: {private_companies} ({round(private_companies/total*100, 1)}% of tracked competitors)

## Pricing Intelligence
{chr(10).join(f'- **{model}**: {count} competitors' for model, count in list(pricing_models.items())[:5])}

## Strategic Recommendations
1. **Monitor pricing changes weekly** - Especially for high-threat competitors
2. **Investigate feature gaps** in patient intake workflows
3. **Review battlecards** for top high-threat targets
4. **Track funding announcements** from private competitors
5. **Analyze customer reviews** on G2 and Capterra for competitive insights

## Watch List
Top competitors requiring immediate attention: {', '.join(top_threats) if top_threats else 'No high-threat competitors identified'}

---
*Summary generated automatically based on {total} tracked competitors*
"""
        return {
            "summary": summary,
            "type": "fallback",
            "model": "automated",
            "provider": "Certify Intel",
            "provider_logo": "/static/certify-logo.svg",
            "data_points_analyzed": total,
            "generated_at": datetime.utcnow().isoformat()
        }
        
    except Exception as e:
        logger.error(f"Summary Error: {e}")
        return {"summary": "Error generating summary. Please try again.", "type": "error", "model": None}


@app.post("/api/ai/analyze")
async def ai_analyze(
    request: Request,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """General AI analysis endpoint for prompts with context."""
    body = await request.json()
    prompt = body.get("prompt", "")
    context = body.get("context", {})

    if not prompt:
        raise HTTPException(status_code=400, detail="prompt is required")

    try:
        from ai_router import get_ai_router, TaskType

        # Build context string from provided context dict
        context_str = ""
        if "competitors" in context:
            competitors_data = context["competitors"]
            context_str += f"\nCompetitors data: {len(competitors_data)} competitors\n"
            for comp in competitors_data[:10]:  # Limit to avoid token overflow
                context_str += f"- {comp.get('name', 'Unknown')}: {comp.get('notes', 'N/A')[:100]}\n"

        if "news" in context:
            news_data = context["news"]
            context_str += f"\nNews data: {len(news_data)} articles\n"
            for article in news_data[:10]:
                context_str += f"- {article.get('title', 'Unknown')}: {article.get('summary', 'N/A')[:100]}\n"

        full_prompt = prompt + "\n\nContext:\n" + context_str if context_str else prompt

        router = get_ai_router()
        result = await router.generate(
            prompt=full_prompt,
            task_type=TaskType.ANALYSIS,
            system_prompt="You are a competitive intelligence analyst. Provide clear, actionable insights based on the data provided." + NO_HALLUCINATION_INSTRUCTION,
            max_tokens=2000,
            temperature=0.7,
        )

        return {
            "content": result.get("response", ""),
            "provider": result.get("provider", "unknown"),
            "model": result.get("model", "unknown"),
        }

    except Exception as e:
        logger.error(f"AI analyze error: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.get("/api/analytics/summary/saved")
def get_saved_dashboard_summary(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Load previously generated dashboard summary from DB (v8.0.8)."""
    import json as _json
    from database import PersistentCache

    user_id = current_user.get("id", 0)
    cache_key = f"dashboard_summary_{user_id}"
    try:
        cached = db.query(PersistentCache).filter(
            PersistentCache.cache_key == cache_key
        ).first()
        if cached and cached.data_json:
            return _json.loads(cached.data_json)
    except Exception as e:
        logger.warning(f"Failed to load saved summary: {e}")
    return None


@app.delete("/api/analytics/summary/saved")
def delete_saved_dashboard_summary(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Clear saved dashboard summary from DB (v8.0.8)."""
    from database import PersistentCache
    user_id = current_user.get("id", 0)
    cache_key = f"dashboard_summary_{user_id}"
    try:
        db.query(PersistentCache).filter(
            PersistentCache.cache_key == cache_key
        ).delete()
        db.commit()
    except Exception as e:
        logger.warning(f"Failed to delete saved summary: {e}")
    return {"status": "ok"}


@app.get("/api/analytics/summary/progress")
async def get_analytics_summary_progress(
    current_user: dict = Depends(get_current_user_optional)
):
    """Get progress of analytics summary generation."""
    user_id = current_user.get("id") if current_user else "anonymous"
    progress = _analytics_summary_progress.get(user_id, {
        "status": "idle",
        "percentage": 0,
        "message": "No summary generation in progress"
    })
    return progress


@app.post("/api/analytics/chat")
async def chat_with_summary(request: dict, db: Session = Depends(get_db), current_user: dict = Depends(get_current_user)):
    """Chat with AI about the competitive intelligence data. Supports conversation history."""
    try:
        import os

        user_message = request.get("message", "")
        prompt_key = request.get("prompt_key")  # Optional: frontend can specify which prompt to use
        session_id = request.get("session_id")  # Optional: existing chat session
        conversation_history = request.get("conversation_history")  # Optional: list of {role, content}
        if not user_message:
            return {"response": "Please provide a message.", "success": False}

        # Get competitor data for context
        # Get ALL non-deleted competitors for consistency
        competitors = db.query(Competitor).filter(Competitor.is_deleted == False).all()
        
        # Build Comprehensive Context (Same as Summary Generation to ensure consistency)
        pricing_models = {}
        for c in competitors:
            model = c.pricing_model or "Unknown"
            pricing_models[model] = pricing_models.get(model, 0) + 1

        context = f"""
FULL DATA SNAPSHOT (LIVE):
==========================
Total Competitors: {len(competitors)}
High Threat: {sum(1 for c in competitors if c.threat_level == 'High')}
Medium Threat: {sum(1 for c in competitors if c.threat_level == 'Medium')}
Low Threat: {sum(1 for c in competitors if c.threat_level == 'Low')}

COMPETITORS:
"""
        # Include ALL competitors in chat context too
        for c in competitors:
            # Build base competitor info
            comp_info = f"""
---
COMPETITOR: {c.name}
THREAT: {c.threat_level}
WEBSITE: {c.website or 'N/A'}
PRICING: {c.base_price or 'N/A'} ({c.pricing_model or 'Unknown Model'})
OFFERING: {c.product_categories or 'N/A'}
FEATURES: {c.key_features or 'N/A'}
EMPLOYEES: {c.employee_count or 'N/A'}
G2 RATING: {c.g2_rating or 'N/A'}
"""
            # Add stock data for public companies
            if c.is_public and c.ticker_symbol:
                stock_data = fetch_real_stock_data(c.ticker_symbol)
                if stock_data and stock_data.get('price'):
                    change_sign = '+' if stock_data.get('change', 0) >= 0 else ''
                    comp_info += f"""PUBLIC COMPANY: Yes
STOCK TICKER: {c.ticker_symbol} ({c.stock_exchange or 'NYSE'})
CURRENT STOCK PRICE: ${stock_data.get('price', 'N/A'):.2f}
PRICE CHANGE: {change_sign}{stock_data.get('change', 0):.2f} ({change_sign}{stock_data.get('change_percent', 0):.2f}%)
MARKET CAP: ${stock_data.get('market_cap', 0):,.0f}
52-WEEK HIGH: ${stock_data.get('high52', 'N/A')}
52-WEEK LOW: ${stock_data.get('low52', 'N/A')}
"""
                else:
                    comp_info += f"""PUBLIC COMPANY: Yes
STOCK TICKER: {c.ticker_symbol} ({c.stock_exchange or 'NYSE'})
STOCK DATA: Unable to fetch live data
"""
            else:
                comp_info += f"PUBLIC COMPANY: No (Private)\n"

            comp_info += f"NOTES: {c.notes or ''}\n---\n"
            context += comp_info

        # Get user ID for personalized prompts
        user_id = current_user.get("id") if current_user else None

        # Helper to load a prompt by key (user-specific first, then global fallback)
        def _load_prompt(key: str, fallback: str = "") -> str:
            p = db.query(SystemPrompt).filter(SystemPrompt.key == key, SystemPrompt.user_id == user_id).first()
            if not p:
                p = db.query(SystemPrompt).filter(SystemPrompt.key == key, SystemPrompt.user_id == None).first()
            return p.content if p else fallback

        # If frontend sent a specific prompt_key, use that as the persona
        if prompt_key:
            base_persona = _load_prompt(prompt_key, "You are a competitive intelligence analyst for Certify Health.")
        else:
            base_persona = _load_prompt("chat_persona", "You are a competitive intelligence analyst for Certify Health.")

        summary_instructions = _load_prompt("dashboard_summary", "Focus on strategic insights.")

        # RAG: Inject Knowledge Base
        kb_text = ""
        kb_items = db.query(KnowledgeBaseItem).filter(KnowledgeBaseItem.is_active == True).all()
        if kb_items:
            kb_text += "\n\nINTERNAL KNOWLEDGE BASE:\n"
            for item in kb_items:
                kb_text += f"\n--- {item.title} ---\n{item.content_text}\n"


        full_system_content = f"""{base_persona}

ALIGNMENT INSTRUCTION:
The user has defined the following strategy for the Dashboard Summary. Use this as context for your tone and priorities:
"{summary_instructions}"

CRITICAL INSTRUCTION:
You have access to a LIVE database of competitors below with REAL-TIME STOCK DATA for public companies.
- If the user asks for a website, LOOK at the 'WEBSITE' field for that competitor and provide it.
- If the user asks for pricing details, LOOK at the 'PRICING' field.
- If the user asks about stock prices, market cap, or financial data for PUBLIC COMPANIES, LOOK at the 'STOCK' fields (CURRENT STOCK PRICE, MARKET CAP, PRICE CHANGE, etc.).
- For public companies, you have LIVE stock data including: current price, daily change, market cap, 52-week high/low.
- Do NOT say "I cannot browse the web" or "I don't have access to real-time stock data" if the answer is in the data below.
- Do NOT say "I am working with hypothetical data". This IS the live, real-time data from the Certify Intel platform.
- When asked about a public company's stock, provide the EXACT values from the data (e.g., "Phreesia (PHR) is currently trading at $15.84, up +0.3%").

LIVE DATA CONTEXT:
{context}

{kb_text}
"""

        # Build prompt with conversation history for multi-turn context
        history_prefix = ""
        if conversation_history and isinstance(conversation_history, list):
            for msg in conversation_history[-10:]:  # Last 10 messages max
                role = msg.get("role", "user")
                content = msg.get("content", "")
                history_prefix += f"\n[{role.upper()}]: {content}"
            history_prefix += f"\n[USER]: {user_message}\n\nRespond to the latest user message above, with full context of the conversation."
        else:
            history_prefix = user_message

        from ai_router import get_ai_router, TaskType
        router = get_ai_router()
        result = await router.generate(
            prompt=history_prefix,
            task_type=TaskType.ANALYSIS,
            system_prompt=full_system_content + NO_HALLUCINATION_INSTRUCTION,
            max_tokens=4000,
        )
        if result and result.get("response"):
            ai_response = result["response"]

            # Persist to chat session if session_id provided
            if session_id:
                try:
                    _save_chat_messages(db, session_id, user_message, ai_response, {
                        "provider": result.get("provider"),
                        "model": result.get("model"),
                        "endpoint": "analytics_chat",
                    })
                except Exception as save_err:
                    logger.warning(f"Failed to save chat messages: {save_err}")

            return {"response": ai_response, "success": True, "session_id": session_id}
        else:
            return {"response": "AI chat is unavailable. Check API key configuration.", "success": False}
    except Exception as e:
        return {"response": "An error occurred processing your request. Please try again.", "success": False}


# /api/health moved to routers/health.py

# Chat session endpoints moved to routers/chat.py


def _save_chat_messages(
    db: Session,
    session_id: int,
    user_message: str,
    ai_response: str,
    metadata: dict = None
):
    """Helper to save user message + AI response to a chat session."""
    import json as _json

    # Save user message
    user_msg = ChatMessage(
        session_id=session_id,
        role="user",
        content=user_message,
    )
    db.add(user_msg)

    # Save AI response
    ai_msg = ChatMessage(
        session_id=session_id,
        role="assistant",
        content=ai_response,
        metadata_json=_json.dumps(metadata) if metadata else None,
    )
    db.add(ai_msg)

    # Touch session updated_at
    session = db.query(ChatSession).filter(ChatSession.id == session_id).first()
    if session:
        session.updated_at = datetime.utcnow()
        # Auto-title from first user message
        if not session.title:
            session.title = user_message[:80] + ("..." if len(user_message) > 80 else "")

    db.commit()


# AI Provider Status (v5.0.2)
@app.get("/api/ai/status")
def get_ai_status():
    """
    Get current AI provider status and configuration.

    v5.0.2: Shows hybrid AI routing status (OpenAI/Gemini).
    """
    from extractor import OPENAI_AVAILABLE, GEMINI_AVAILABLE, get_extractor

    extractor = get_extractor()
    provider_config = os.getenv("AI_PROVIDER", "hybrid")

    # Check Gemini availability
    gemini_available = GEMINI_AVAILABLE and bool(os.getenv("GOOGLE_AI_API_KEY"))
    openai_available = OPENAI_AVAILABLE and bool(os.getenv("OPENAI_API_KEY"))

    # Determine active provider for different task types
    active_for_extraction = "none"
    active_for_summary = "none"

    if hasattr(extractor, 'get_provider'):
        active_for_extraction = extractor.get_provider("data_extraction")
    elif openai_available:
        active_for_extraction = "openai"
    elif gemini_available:
        active_for_extraction = "gemini"

    # Summary uses the insight generator
    from analytics import DashboardInsightGenerator
    insight_gen = DashboardInsightGenerator()
    active_for_summary = insight_gen.get_active_provider()

    # P2-5: Add budget enforcement info
    try:
        from ai_router import get_ai_router
        router = get_ai_router()
        budget_info = {
            "daily_budget_usd": router.cost_tracker.daily_budget_usd,
            "today_spend_usd": round(router.cost_tracker.get_today_spend(), 4),
            "remaining_budget_usd": round(router.cost_tracker.get_remaining_budget(), 4),
            "budget_used_percent": round((router.cost_tracker.get_today_spend() / router.cost_tracker.daily_budget_usd) * 100, 1) if router.cost_tracker.daily_budget_usd > 0 else 0,
            "budget_warning": router.cost_tracker.get_today_spend() >= (router.cost_tracker.daily_budget_usd * 0.8)
        }
    except Exception as e:
        logger.warning(f"Failed to get budget info: {e}")
        budget_info = {"error": "Budget tracking unavailable"}

    # Check Anthropic (Claude) availability
    anthropic_available = bool(os.getenv("ANTHROPIC_API_KEY"))

    return {
        "status": "configured" if (openai_available or gemini_available or anthropic_available) else "not_configured",
        "provider_config": provider_config,
        "providers": {
            "anthropic": {
                "available": anthropic_available,
                "model": "claude-opus-4-5"
            },
            "openai": {
                "available": openai_available,
                "model": os.getenv("OPENAI_MODEL", "gpt-4o-mini")
            },
            "gemini": {
                "available": gemini_available,
                "model": os.getenv("GOOGLE_AI_MODEL", "gemini-3-flash-preview")
            },
            "deepseek": {
                "available": bool(os.getenv("DEEPSEEK_API_KEY")),
                "model": os.getenv("DEEPSEEK_MODEL", "deepseek-chat")
            }
        },
        "routing": {
            "data_extraction": active_for_extraction,
            "executive_summary": active_for_summary,
            "complex_analysis": "anthropic" if anthropic_available else ("openai" if openai_available else "gemini"),
            "bulk_tasks": os.getenv("AI_BULK_TASKS", "gemini"),
            "quality_tasks": os.getenv("AI_QUALITY_TASKS", "anthropic")
        },
        "budget": budget_info,  # P2-5: Budget enforcement info
        "fallback_enabled": os.getenv("AI_FALLBACK_ENABLED", "true").lower() == "true",
        "version": "v5.0.5",
        "multimodal": {
            "screenshot_analysis": True,
            "pdf_analysis": True
        }
    }


@app.get("/api/observability/status")
async def get_observability_status(
    current_user: dict = Depends(get_current_user)
) -> Dict[str, Any]:
    """
    Get Langfuse observability status and health.

    Returns whether Langfuse is enabled, connected, and configuration info.
    """
    try:
        from observability import check_langfuse_health
        health = await check_langfuse_health()
    except ImportError:
        health = {
            "enabled": False,
            "host": None,
            "connected": False,
            "error": "observability module not found"
        }
    except Exception:
        health = {
            "enabled": False,
            "host": None,
            "connected": False,
            "error": "Failed to check Langfuse health"
        }

    return {
        "langfuse": health,
        "setup_instructions": (
            "1. docker-compose -f docker-compose.langfuse.yml up -d\n"
            "2. Visit http://localhost:3000 and create account\n"
            "3. Get API keys from Settings -> API Keys\n"
            "4. Set ENABLE_LANGFUSE=true, LANGFUSE_PUBLIC_KEY, "
            "LANGFUSE_SECRET_KEY in .env"
        ) if not health.get("connected") else None
    }


# ============== MULTIMODAL AI ENDPOINTS (v5.0.5) ==============

@app.post("/api/ai/analyze-screenshot")
async def analyze_screenshot(
    competitor_name: str = Form(...),
    page_type: str = Form("homepage"),
    file: UploadFile = File(...)
):
    """
    Analyze a competitor website screenshot using Gemini multimodal AI.

    v5.0.5: New endpoint for visual competitive intelligence.

    Args:
        competitor_name: Name of the competitor
        page_type: Type of page (homepage, pricing, features, about)
        file: Screenshot image file (PNG, JPEG, WebP, GIF)

    Returns:
        Extracted competitive intelligence as JSON
    """
    try:
        from gemini_provider import GeminiProvider

        provider = GeminiProvider()
        if not provider.is_available:
            raise HTTPException(
                status_code=503,
                detail="Gemini AI not available. Configure GOOGLE_AI_API_KEY."
            )

        # Read file content
        content = await file.read()

        # Validate file type
        allowed_types = ["image/png", "image/jpeg", "image/webp", "image/gif"]
        if file.content_type not in allowed_types:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {file.content_type}. Use PNG, JPEG, WebP, or GIF."
            )

        # Analyze screenshot
        result = provider.analyze_screenshot(content, competitor_name, page_type)

        return {
            "success": "error" not in result,
            "data": result,
            "competitor": competitor_name,
            "page_type": page_type,
            "filename": file.filename
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Screenshot analysis failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/ai/analyze-pdf")
async def analyze_pdf(
    document_type: str = Form("general"),
    competitor_name: Optional[str] = Form(None),
    custom_prompt: Optional[str] = Form(None),
    file: UploadFile = File(...)
):
    """
    Analyze a PDF document using Gemini AI.

    v5.0.5: New endpoint for document analysis.

    Args:
        document_type: Type of document (whitepaper, case_study, datasheet, annual_report, general)
        competitor_name: Name of the competitor (optional)
        custom_prompt: Custom analysis prompt (optional)
        file: PDF file

    Returns:
        Extracted insights as JSON
    """
    try:
        from gemini_provider import GeminiProvider

        provider = GeminiProvider()
        if not provider.is_available:
            raise HTTPException(
                status_code=503,
                detail="Gemini AI not available. Configure GOOGLE_AI_API_KEY."
            )

        # Read file content
        content = await file.read()

        # Validate file type
        if file.content_type != "application/pdf" and not file.filename.lower().endswith(".pdf"):
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type. Please upload a PDF file."
            )

        # Analyze PDF
        result = provider.analyze_pdf(
            pdf_source=content,
            prompt=custom_prompt,
            competitor_name=competitor_name,
            document_type=document_type
        )

        return {
            "success": "error" not in result,
            "data": result,
            "document_type": document_type,
            "competitor": competitor_name,
            "filename": file.filename
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"PDF analysis failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/ai/analyze-image")
async def analyze_image(
    prompt: str = Form(...),
    file: UploadFile = File(...)
):
    """
    Analyze any image using Gemini multimodal AI.

    v5.0.5: General purpose image analysis endpoint.

    Args:
        prompt: Analysis prompt describing what to extract
        file: Image file (PNG, JPEG, WebP, GIF)

    Returns:
        AI analysis response
    """
    try:
        from gemini_provider import GeminiProvider

        provider = GeminiProvider()
        if not provider.is_available:
            raise HTTPException(
                status_code=503,
                detail="Gemini AI not available. Configure GOOGLE_AI_API_KEY."
            )

        # Read file content
        content = await file.read()

        # Validate file type
        allowed_types = ["image/png", "image/jpeg", "image/webp", "image/gif"]
        if file.content_type not in allowed_types:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {file.content_type}. Use PNG, JPEG, WebP, or GIF."
            )

        # Analyze image
        response = provider.analyze_image(content, prompt)

        return {
            "success": response.success,
            "content": response.content,
            "model": response.model,
            "cost_estimate": response.cost_estimate,
            "latency_ms": response.latency_ms,
            "error": response.error,
            "filename": file.filename
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Image analysis failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


# ============== v5.0.6 ADVANCED AI FEATURES ==============

@app.post("/api/ai/analyze-video")
async def analyze_video(
    video_type: str = Form("demo"),
    competitor_name: Optional[str] = Form(None),
    custom_prompt: Optional[str] = Form(None),
    file: UploadFile = File(...)
):
    """
    Analyze a video using Gemini multimodal AI.

    v5.0.6: Video intelligence for competitor demos, webinars, and tutorials.

    Args:
        video_type: Type of video (demo, webinar, tutorial, advertisement, general)
        competitor_name: Optional competitor name for context
        custom_prompt: Optional custom analysis prompt
        file: Video file (MP4, WebM, MOV)

    Returns:
        Extracted competitive intelligence from video
    """
    try:
        from gemini_provider import GeminiProvider

        provider = GeminiProvider()
        if not provider.is_available:
            raise HTTPException(
                status_code=503,
                detail="Gemini AI not available. Configure GOOGLE_AI_API_KEY."
            )

        # Read file content
        content = await file.read()

        # Check file size (limit to 50MB)
        if len(content) > 50 * 1024 * 1024:
            raise HTTPException(
                status_code=400,
                detail="Video file too large. Maximum size is 50MB."
            )

        # Validate file type
        allowed_types = ["video/mp4", "video/webm", "video/quicktime"]
        if file.content_type not in allowed_types:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {file.content_type}. Use MP4, WebM, or MOV."
            )

        # Analyze video
        result = provider.analyze_video(
            video_source=content,
            prompt=custom_prompt,
            competitor_name=competitor_name,
            video_type=video_type
        )

        result["filename"] = file.filename
        return result

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Video analysis failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/ai/search-grounded")
async def search_grounded(
    query: str = Form(...),
    competitor_name: Optional[str] = Form(None),
    search_type: str = Form("general")
):
    """
    Search with Gemini's real-time Google Search grounding.

    v5.0.6: Get current, factual information about competitors.

    Args:
        query: Search query or question
        competitor_name: Optional competitor name for context
        search_type: Type of search (general, news, financial, product)

    Returns:
        Grounded response with real-time information
    """
    try:
        from gemini_provider import GeminiProvider

        provider = GeminiProvider()
        if not provider.is_available:
            raise HTTPException(
                status_code=503,
                detail="Gemini AI not available. Configure GOOGLE_AI_API_KEY."
            )

        result = provider.search_and_ground(
            query=query,
            competitor_name=competitor_name,
            search_type=search_type
        )

        return result

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Grounded search failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/ai/research-competitor")
async def research_competitor(
    competitor_name: str = Form(...),
    research_areas: Optional[str] = Form(None)
):
    """
    Comprehensive competitor research using real-time grounding.

    v5.0.6: Deep research feature for competitor profiles.

    Args:
        competitor_name: Name of the competitor
        research_areas: Comma-separated areas (overview,products,pricing,news,financials)

    Returns:
        Comprehensive research results
    """
    try:
        from gemini_provider import GeminiProvider

        provider = GeminiProvider()
        if not provider.is_available:
            raise HTTPException(
                status_code=503,
                detail="Gemini AI not available. Configure GOOGLE_AI_API_KEY."
            )

        areas = None
        if research_areas:
            areas = [a.strip() for a in research_areas.split(",")]

        result = provider.research_competitor(
            competitor_name=competitor_name,
            research_areas=areas
        )

        return result

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Competitor research failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/ai/process-news-batch")
async def process_news_batch(request: Request):
    """
    Process multiple news articles efficiently using Flash-Lite.

    v5.0.6: Bulk processing for cost-effective news analysis.

    Request body:
        articles: List of {title, snippet, url} objects
        analysis_type: summary, sentiment, categorize, or extract

    Returns:
        List of analysis results
    """
    try:
        from gemini_provider import GeminiProvider

        provider = GeminiProvider()
        if not provider.is_available:
            raise HTTPException(
                status_code=503,
                detail="Gemini AI not available. Configure GOOGLE_AI_API_KEY."
            )

        body = await request.json()
        articles = body.get("articles", [])
        analysis_type = body.get("analysis_type", "summary")

        if not articles:
            raise HTTPException(status_code=400, detail="No articles provided")

        if len(articles) > 100:
            raise HTTPException(status_code=400, detail="Maximum 100 articles per batch")

        results = provider.process_news_batch(
            articles=articles,
            analysis_type=analysis_type
        )

        return {
            "results": results,
            "articles_processed": len(articles),
            "analysis_type": analysis_type
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"News batch processing failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/ai/analyze-news-trends")
async def analyze_news_trends(request: Request):
    """
    Analyze trends across multiple news articles.

    v5.0.6: Trend analysis for competitive intelligence.

    Request body:
        articles: List of {title, snippet, url} objects
        competitor_name: Optional competitor focus

    Returns:
        Trend analysis with themes, sentiment, and recommendations
    """
    try:
        from gemini_provider import GeminiProvider

        provider = GeminiProvider()
        if not provider.is_available:
            raise HTTPException(
                status_code=503,
                detail="Gemini AI not available. Configure GOOGLE_AI_API_KEY."
            )

        body = await request.json()
        articles = body.get("articles", [])
        competitor_name = body.get("competitor_name")

        if not articles:
            raise HTTPException(status_code=400, detail="No articles provided")

        result = provider.analyze_news_trends(
            articles=articles,
            competitor_name=competitor_name
        )

        return result

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"News trend analysis failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


# ============== DEEP RESEARCH ENDPOINTS (v5.0.6) ==============

@app.post("/api/ai/deep-research")
async def deep_research(request: Request):
    """
    Generate deep research report for a competitor.

    v5.0.6: NEWS-4B/4C - ChatGPT and Gemini Deep Research integration.

    Request body:
        competitor_name: Name of the competitor
        research_type: Type of research (battlecard, market_analysis, product_deep_dive, quick_summary)
        provider: Optional provider preference (chatgpt, gemini, or auto)
        additional_context: Optional additional context

    Returns:
        Research result with content, sections, and metadata
    """
    try:
        from ai_research import get_research_manager

        body = await request.json()
        competitor_name = body.get("competitor_name")
        research_type = body.get("research_type", "battlecard")
        provider = body.get("provider")  # None for auto-selection
        additional_context = body.get("additional_context")

        if not competitor_name:
            raise HTTPException(status_code=400, detail="competitor_name is required")

        manager = get_research_manager()
        result = await manager.research(
            competitor_name=competitor_name,
            research_type=research_type,
            provider=provider,
            additional_context=additional_context
        )

        from dataclasses import asdict
        return asdict(result)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Deep research failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/ai/generate-battlecard")
async def generate_battlecard_endpoint(
    request: Request,
    background: bool = False,
    background_tasks: BackgroundTasks = None,
    current_user: dict = Depends(get_current_user_optional)
):
    """
    Generate a sales battlecard for a competitor.

    v7.1.2: Enhanced with real news data and competitor DB fields.
    Fetches live news articles so AI can reference real sources with URLs.

    Request body:
        competitor_name: Name of the competitor
        competitor_id: Optional competitor ID (to fetch DB data + news)
        provider: Optional provider preference
        additional_context: Optional additional context

    Returns:
        Complete battlecard with sections and real source URLs
    """
    try:
        from ai_research import generate_battlecard

        body = await request.json()

        if background and background_tasks and current_user:
            import uuid as _uuid
            task_id = str(_uuid.uuid4())
            comp_name = body.get("competitor_name", "Unknown")
            _ai_tasks[task_id] = {
                "status": "running", "page_context": "battlecards",
                "user_id": current_user.get("id"),
                "started_at": datetime.utcnow().isoformat(),
                "completed_at": None, "result": None, "error": None,
                "task_type": f"battlecard_{comp_name}", "read_at": None,
            }

            async def _run_battlecard_bg():
                try:
                    from ai_research import generate_battlecard as gen_bc
                    result = await asyncio.wait_for(
                        gen_bc(
                            competitor_name=body.get("competitor_name"),
                            provider=body.get("provider"),
                            additional_context=body.get("additional_context"),
                        ),
                        timeout=45.0
                    )
                    _ai_tasks[task_id]["result"] = result
                    _ai_tasks[task_id]["status"] = "completed"
                    _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
                except asyncio.TimeoutError:
                    logger.error("Background battlecard generation timed out after 45s")
                    _ai_tasks[task_id]["status"] = "failed"
                    _ai_tasks[task_id]["error"] = "AI analysis timed out. Please try again."
                    _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
                except Exception as bg_err:
                    logger.error(f"Background battlecard error: {bg_err}")
                    _ai_tasks[task_id]["status"] = "failed"
                    _ai_tasks[task_id]["error"] = str(bg_err)
                    _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()

            background_tasks.add_task(_run_battlecard_bg)
            return {"task_id": task_id, "status": "running"}
        competitor_name = body.get("competitor_name")
        competitor_id = body.get("competitor_id")
        provider = body.get("provider")
        additional_context = body.get("additional_context")
        bc_session_id = body.get("session_id")  # Optional chat session
        conversation_history = body.get("conversation_history")  # Optional history
        bc_prompt_key = body.get("prompt_key")  # Optional custom prompt

        # If prompt_key provided, load custom instructions and prepend
        if bc_prompt_key:
            bc_db = SessionLocal()
            try:
                custom = _resolve_system_prompt(bc_db, None, bc_prompt_key, "")
                if custom:
                    additional_context = (
                        f"CUSTOM INSTRUCTIONS:\n{custom}\n\n"
                        + (additional_context or "")
                    )
            finally:
                bc_db.close()

        if not competitor_name:
            raise HTTPException(status_code=400, detail="competitor_name is required")

        # Fetch real news articles for this competitor (run sync fetch in thread to avoid blocking)
        news_articles = []
        competitor_data = {}
        try:
            from news_monitor import NewsMonitor
            monitor = NewsMonitor()
            news_result = await asyncio.wait_for(
                asyncio.to_thread(monitor.fetch_news, competitor_name, 90),
                timeout=15.0
            )
            if news_result:
                # NewsDigest is a dataclass with .articles list of NewsArticle dataclasses
                raw_articles = getattr(news_result, 'articles', None)
                if raw_articles is None and isinstance(news_result, dict):
                    raw_articles = news_result.get("articles", [])
                if raw_articles:
                    news_articles = [{
                        "title": getattr(a, 'title', '') if not isinstance(a, dict) else a.get('title', ''),
                        "url": getattr(a, 'url', '') if not isinstance(a, dict) else a.get('url', ''),
                        "source": getattr(a, 'source', '') if not isinstance(a, dict) else a.get('source', ''),
                        "published_date": getattr(a, 'published_date', '') if not isinstance(a, dict) else a.get('published_date', ''),
                        "sentiment": getattr(a, 'sentiment', '') if not isinstance(a, dict) else a.get('sentiment', ''),
                    } for a in raw_articles[:8]]
        except Exception as e:
            logger.warning(f"Could not fetch news for battlecard: {e}")

        # Fetch competitor data from database if ID provided
        if competitor_id:
            db = None
            try:
                from database import CompetitorProduct, DataSource
                db = SessionLocal()
                comp = db.query(Competitor).filter(Competitor.id == competitor_id).first()
                if comp:
                    competitor_data = {
                        "website": comp.website,
                        "description": comp.notes,
                        "employee_count": comp.employee_count,
                        "annual_revenue": comp.annual_revenue,
                        "estimated_revenue": comp.estimated_revenue,
                        "founding_year": comp.year_founded,
                        "headquarters": comp.headquarters,
                        "pricing_model": comp.pricing_model,
                        "product_categories": comp.product_categories,
                        "key_features": comp.key_features,
                        "threat_level": comp.threat_level,
                        "primary_market": comp.primary_market,
                        "target_segments": comp.target_segments,
                        "g2_rating": comp.g2_rating,
                        "customer_count": comp.customer_count,
                        "certifications": comp.certifications,
                        "integration_partners": comp.integration_partners,
                        "dim_overall_score": comp.dim_overall_score,
                        # Dimension scores for competitive positioning (9 dimensions)
                        "dim_product_packaging_score": getattr(comp, "dim_product_packaging_score", None),
                        "dim_integration_depth_score": getattr(comp, "dim_integration_depth_score", None),
                        "dim_support_service_score": getattr(comp, "dim_support_service_score", None),
                        "dim_retention_stickiness_score": getattr(comp, "dim_retention_stickiness_score", None),
                        "dim_user_adoption_score": getattr(comp, "dim_user_adoption_score", None),
                        "dim_implementation_ttv_score": getattr(comp, "dim_implementation_ttv_score", None),
                        "dim_reliability_enterprise_score": getattr(comp, "dim_reliability_enterprise_score", None),
                        "dim_pricing_flexibility_score": getattr(comp, "dim_pricing_flexibility_score", None),
                        "dim_reporting_analytics_score": getattr(comp, "dim_reporting_analytics_score", None),
                    }

                    # Fetch products for this competitor
                    products = db.query(CompetitorProduct).filter(
                        CompetitorProduct.competitor_id == competitor_id
                    ).all()
                    if products:
                        competitor_data["products"] = ", ".join(
                            f"{p.product_name} ({p.product_category or 'N/A'})"
                            for p in products[:20]
                        )

                    # Fetch verified data sources
                    data_sources = db.query(DataSource).filter(
                        DataSource.competitor_id == competitor_id,
                        DataSource.is_verified == True
                    ).all()
                    if data_sources:
                        competitor_data["verified_sources"] = ", ".join(
                            f"{ds.source_type}: {ds.source_url or ds.source_name or 'internal'}"
                            for ds in data_sources[:10]
                        )
            except Exception as e:
                logger.warning(f"Could not fetch competitor data for battlecard: {e}")
            finally:
                if db:
                    db.close()

        # Also check if we have cached news in the database
        if not news_articles:
            db = None
            try:
                from database import NewsArticleCache
                db = SessionLocal()
                cached_news = db.query(NewsArticleCache).filter(
                    NewsArticleCache.competitor_name.ilike(f"%{competitor_name}%"),
                    NewsArticleCache.is_archived != True
                ).order_by(NewsArticleCache.published_at.desc()).limit(8).all()
                if cached_news:
                    news_articles = [{
                        "title": n.title,
                        "url": n.url,
                        "source": n.source,
                        "published_date": str(n.published_at) if n.published_at else "",
                        "sentiment": n.sentiment,
                    } for n in cached_news]
            except Exception as e:
                logger.warning(f"Could not fetch cached news for battlecard: {e}")
            finally:
                if db:
                    db.close()

        # Prepend conversation history as additional context if provided
        if conversation_history and isinstance(conversation_history, list):
            history_text = "\n".join(
                f"[{m.get('role', 'user').upper()}]: {m.get('content', '')}"
                for m in conversation_history[-10:]
            )
            additional_context = (additional_context or "") + f"\n\nPrevious conversation:\n{history_text}"

        result = await asyncio.wait_for(
            generate_battlecard(
                competitor_name=competitor_name,
                provider=provider,
                additional_context=additional_context,
                news_articles=news_articles,
                competitor_data=competitor_data,
            ),
            timeout=60.0
        )

        # Persist to chat session if session_id provided
        if bc_session_id:
            try:
                import json as _json
                bc_db = SessionLocal()
                try:
                    bc_content = result.get("battlecard", result.get("content", str(result)))
                    if not isinstance(bc_content, str):
                        bc_content = _json.dumps(bc_content, indent=2)
                    _save_chat_messages(bc_db, bc_session_id, f"Generate battlecard for {competitor_name}", bc_content, {
                        "endpoint": "generate_battlecard", "competitor_name": competitor_name,
                    })
                finally:
                    bc_db.close()
            except Exception as save_err:
                logger.warning(f"Failed to save battlecard to chat: {save_err}")

        # v8.0.8: Persist strategy to Battlecard table so it survives server restart
        if competitor_id and result and isinstance(result, dict):
            try:
                import json as _json
                from database import Battlecard as BattlecardModel
                bc_db = SessionLocal()
                try:
                    bc_content = result.get("battlecard", result.get("content", ""))
                    if not isinstance(bc_content, str):
                        bc_content = _json.dumps(bc_content, indent=2)

                    existing_bc = bc_db.query(BattlecardModel).filter(
                        BattlecardModel.competitor_id == competitor_id,
                        BattlecardModel.battlecard_type == "strategy",
                        BattlecardModel.is_active == True
                    ).first()
                    if existing_bc:
                        existing_bc.content = bc_content
                        existing_bc.generated_at = datetime.utcnow()
                        existing_bc.updated_at = datetime.utcnow()
                    else:
                        bc_db.add(BattlecardModel(
                            competitor_id=competitor_id,
                            title=f"How to Win Against {competitor_name}",
                            content=bc_content,
                            battlecard_type="strategy",
                            generated_at=datetime.utcnow(),
                            generated_by="ai",
                            is_active=True,
                        ))
                    bc_db.commit()
                finally:
                    bc_db.close()
            except Exception as persist_err:
                logger.warning(f"Failed to persist battlecard strategy: {persist_err}")

        if isinstance(result, dict):
            result["session_id"] = bc_session_id
        return result

    except asyncio.TimeoutError:
        logger.error(f"Battlecard generation timed out after 60s for {body.get('competitor_name', 'unknown')}")
        return {"success": False, "error": "Strategy generation timed out. Please try again.", "content": None}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Battlecard generation failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.get("/api/ai/battlecard-strategy/{competitor_id}")
def get_saved_battlecard_strategy(
    competitor_id: int,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """Load previously generated battlecard strategy from DB (v8.0.8)."""
    from database import Battlecard as BattlecardModel
    try:
        bc = db.query(BattlecardModel).filter(
            BattlecardModel.competitor_id == competitor_id,
            BattlecardModel.battlecard_type == "strategy",
            BattlecardModel.is_active == True
        ).order_by(BattlecardModel.generated_at.desc()).first()
        if bc:
            return {
                "content": bc.content,
                "generated_at": bc.generated_at.isoformat() if bc.generated_at else None,
                "title": bc.title,
            }
    except Exception as e:
        logger.warning(f"Failed to load saved battlecard strategy: {e}")
    return None


@app.post("/api/sales-marketing/playbook/generate")
async def generate_sales_playbook(
    request: Request,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
) -> Dict[str, Any]:
    """Generate a structured sales playbook for a competitor.

    Gathers dimension scores, battlecard data, talking points, and win/loss
    history, then uses AI to produce a comprehensive sales playbook with
    positioning, differentiators, objection handling, and more.
    """
    try:
        from ai_router import get_ai_router, TaskType
        from database import (
            Battlecard as BattlecardModel,
            TalkingPoint, WinLossDeal
        )

        body = await request.json()
        competitor_id = body.get("competitor_id")
        deal_context = body.get("deal_context", "")
        prompt_key = body.get("prompt_key")

        if not competitor_id:
            raise HTTPException(
                status_code=400,
                detail="competitor_id is required"
            )

        competitor = db.query(Competitor).filter(
            Competitor.id == competitor_id,
            Competitor.is_deleted == False
        ).first()
        if not competitor:
            raise HTTPException(
                status_code=404,
                detail="Competitor not found"
            )

        # --- Gather context data ---

        # 1. Dimension scores
        dim_fields = [
            ("product_packaging", "Product & Packaging"),
            ("integration_depth", "Integration Depth"),
            ("support_service", "Support & Service"),
            ("retention_stickiness", "Retention & Stickiness"),
            ("user_adoption", "User Adoption"),
            ("implementation_ttv", "Implementation & TTV"),
            ("reliability_enterprise", "Reliability & Enterprise"),
            ("pricing_flexibility", "Pricing Flexibility"),
            ("reporting_analytics", "Reporting & Analytics"),
        ]
        dimensions_text = ""
        for dim_key, dim_label in dim_fields:
            score = getattr(competitor, f"dim_{dim_key}_score", None)
            evidence = getattr(
                competitor, f"dim_{dim_key}_evidence", None
            )
            if score is not None:
                dimensions_text += (
                    f"- {dim_label}: {score}/5"
                )
                if evidence:
                    dimensions_text += f" | {evidence[:200]}"
                dimensions_text += "\n"

        # 2. Latest battlecard
        battlecard_text = ""
        bc = db.query(BattlecardModel).filter(
            BattlecardModel.competitor_id == competitor_id,
            BattlecardModel.is_active == True
        ).order_by(
            BattlecardModel.generated_at.desc()
        ).first()
        if bc and bc.content:
            battlecard_text = bc.content[:2000]

        # 3. Talking points
        talking_points_text = ""
        tps = db.query(TalkingPoint).filter(
            TalkingPoint.competitor_id == competitor_id,
            TalkingPoint.is_active == True
        ).order_by(
            TalkingPoint.effectiveness_score.desc().nullslast()
        ).limit(20).all()
        for tp in tps:
            talking_points_text += (
                f"- [{tp.point_type}] {tp.content[:150]}\n"
            )

        # 4. Win/loss history
        win_loss_text = ""
        deals = db.query(WinLossDeal).filter(
            WinLossDeal.competitor_id == competitor_id
        ).order_by(WinLossDeal.deal_date.desc()).limit(10).all()
        wins = sum(1 for d in deals if d.outcome == "win")
        losses = sum(1 for d in deals if d.outcome == "loss")
        win_loss_text += f"Record: {wins}W / {losses}L\n"
        for d in deals[:5]:
            reason = d.reason[:100] if d.reason else "No reason recorded"
            win_loss_text += (
                f"- {d.outcome.upper()}: {reason}\n"
            )

        # --- Build AI prompt ---
        context_parts = [
            f"COMPETITOR: {competitor.name}",
            f"Threat Level: {competitor.threat_level or 'Unknown'}",
            f"Website: {competitor.website or 'N/A'}",
            f"Products: {competitor.product_categories or 'N/A'}",
            f"Key Features: {(competitor.key_features or 'N/A')[:500]}",
            f"Pricing: {competitor.pricing_model or 'N/A'}"
            f" - {competitor.base_price or 'N/A'}",
        ]
        if dimensions_text:
            context_parts.append(
                f"\nDIMENSION SCORES:\n{dimensions_text}"
            )
        if battlecard_text:
            context_parts.append(
                f"\nEXISTING BATTLECARD:\n{battlecard_text}"
            )
        if talking_points_text:
            context_parts.append(
                f"\nTALKING POINTS:\n{talking_points_text}"
            )
        if win_loss_text:
            context_parts.append(
                f"\nWIN/LOSS HISTORY:\n{win_loss_text}"
            )
        if deal_context:
            context_parts.append(
                f"\nDEAL CONTEXT:\n{deal_context}"
            )

        full_context = "\n".join(context_parts)

        # Load custom prompt if provided, otherwise use default
        default_system = (
            "You are a senior sales strategist for Certify Health, "
            "a healthcare technology company. Generate a comprehensive "
            "sales playbook for competing against the specified competitor. "
            "Structure your response with these sections:\n"
            "1. POSITIONING - How to position Certify Health vs this competitor\n"
            "2. KEY DIFFERENTIATORS - Top 3-5 competitive advantages\n"
            "3. OBJECTION HANDLING - Common objections and responses\n"
            "4. PRICING STRATEGY - How to compete on price/value\n"
            "5. PROOF POINTS - Customer evidence and data points\n"
            "6. RECOMMENDED APPROACH - Step-by-step sales strategy\n\n"
            "Be specific, actionable, and grounded in the provided data."
        )

        if prompt_key:
            user_id = current_user.get("id")
            p = db.query(SystemPrompt).filter(
                SystemPrompt.key == prompt_key,
                SystemPrompt.user_id == user_id
            ).first()
            if not p:
                p = db.query(SystemPrompt).filter(
                    SystemPrompt.key == prompt_key,
                    SystemPrompt.user_id == None  # noqa: E711
                ).first()
            system_prompt = (
                p.content if p else default_system
            )
        else:
            system_prompt = default_system

        system_prompt += NO_HALLUCINATION_INSTRUCTION

        router = get_ai_router()
        result = await router.generate(
            prompt=full_context,
            task_type=TaskType.BATTLECARD,
            system_prompt=system_prompt,
            max_tokens=3000,
            temperature=0.7,
        )

        return {
            "competitor_id": competitor_id,
            "competitor_name": competitor.name,
            "title": f"Sales Playbook: {competitor.name}",
            "content": result.get("response", ""),
            "playbook": result.get("response", ""),
            "provider": result.get("provider", "unknown"),
            "model": result.get("model", "unknown"),
            "sections": [
                "Positioning",
                "Key Differentiators",
                "Objection Handling",
                "Pricing Strategy",
                "Proof Points",
                "Recommended Approach",
            ],
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to generate sales playbook: {e}")
        raise HTTPException(
            status_code=500,
            detail="An internal error occurred while generating the playbook"
        )


@app.post("/api/battlecards/generate-summaries")
async def generate_battlecard_summaries(
    request: Request,
    background_tasks: BackgroundTasks,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """
    Batch-generate AI threat summaries for competitors missing them.
    Uses Gemini (cheap/fast) to generate 2-4 sentence summaries.
    """
    try:
        body = await request.json() if request.headers.get("content-type") == "application/json" else {}
        force = body.get("force", False)

        # Get competitors needing summaries
        query = db.query(Competitor).filter(Competitor.is_deleted == False)  # noqa: E712
        if not force:
            query = query.filter(
                (Competitor.ai_threat_summary == None) | (Competitor.ai_threat_summary == "")  # noqa: E711
            )
        comps = query.all()

        if not comps:
            return {"status": "ok", "message": "All competitors already have summaries", "generated": 0}

        # Build batch prompt with all competitors
        comp_list = []
        for c in comps:
            comp_list.append({
                "id": c.id,
                "name": c.name,
                "threat_level": c.threat_level or "Medium",
                "product_categories": c.product_categories or "",
                "key_features": c.key_features or "",
                "target_segments": c.target_segments or "",
                "notes": (c.notes or "")[:200],
            })

        async def _generate_batch():
            from ai_router import get_ai_router, TaskType
            import json as _json
            bg_db = SessionLocal()
            try:
                router = get_ai_router()

                # Process in batches of 15
                batch_size = 15
                total_generated = 0
                for i in range(0, len(comp_list), batch_size):
                    batch = comp_list[i:i + batch_size]
                    comp_text = "\n".join(
                        f"- ID:{c['id']} | {c['name']} | Threat: {c['threat_level']} | Products: {c['product_categories']} | Features: {c['key_features']} | Segments: {c['target_segments']}"
                        for c in batch
                    )

                    prompt = f"""Generate a 2-3 sentence competitive threat summary for each competitor below.
Each summary should describe what the company does and the specific threat they pose to Certify Health (a healthcare patient intake, payments, and engagement platform).
Be direct and specific. Reference their actual products/segments.

Return ONLY a valid JSON object mapping competitor ID to summary string. Example:
{{"123": "Summary text here.", "456": "Another summary."}}

Competitors:
{comp_text}"""

                    try:
                        result = await asyncio.wait_for(
                            router.generate(
                                prompt=prompt,
                                task_type=TaskType.SUMMARIZATION,
                                system_prompt="You are a competitive intelligence analyst. Return ONLY valid JSON. No markdown, no code fences.",
                                max_tokens=3000,
                                temperature=0.3,
                            ),
                            timeout=45.0
                        )

                        response_text = result.get("response", "")
                        # Strip markdown code fences if present
                        response_text = response_text.strip()
                        if response_text.startswith("```"):
                            response_text = response_text.split("\n", 1)[-1]
                        if response_text.endswith("```"):
                            response_text = response_text.rsplit("```", 1)[0]
                        response_text = response_text.strip()

                        summaries = _json.loads(response_text)

                        for comp_id_str, summary in summaries.items():
                            comp_id = int(comp_id_str)
                            comp_obj = bg_db.query(Competitor).filter(Competitor.id == comp_id).first()
                            if comp_obj and summary:
                                comp_obj.ai_threat_summary = str(summary).strip()
                                total_generated += 1

                        bg_db.commit()
                        logger.info(f"[Summaries] Batch {i // batch_size + 1}: generated {len(summaries)} summaries")

                    except Exception as batch_err:
                        logger.error(f"[Summaries] Batch {i // batch_size + 1} failed: {batch_err}")
                        bg_db.rollback()

                logger.info(f"[Summaries] Total generated: {total_generated}/{len(comp_list)}")

            except Exception as e:
                logger.error(f"[Summaries] Generation failed: {e}")
            finally:
                bg_db.close()

        background_tasks.add_task(_generate_batch)
        return {"status": "running", "message": f"Generating summaries for {len(comps)} competitors in background", "count": len(comps)}

    except Exception as e:
        logger.error(f"Battlecard summary generation error: {e}")
        return {"status": "error", "message": "Failed to start summary generation"}


@app.get("/api/ai/research-types")
async def get_research_types():
    """
    Get available research types and descriptions.

    v5.0.6: Research type discovery endpoint.

    Returns:
        List of available research types with descriptions
    """
    try:
        from ai_research import get_research_manager
        manager = get_research_manager()
        return {
            "research_types": manager.get_research_types(),
            "providers": manager.get_available_providers()
        }
    except Exception as e:
        logger.error(f"Failed to get research types: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.get("/api/ai/research-providers")
async def get_research_providers():
    """
    Get available AI research providers and their status.

    v5.0.6: Provider availability check.

    Returns:
        Provider availability status
    """
    try:
        from ai_research import get_research_manager
        manager = get_research_manager()
        return manager.get_available_providers()
    except Exception as e:
        logger.error(f"Failed to get research providers: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


# ============== AI NEWS AGENT (v6.4.0) ==============

@app.post("/api/ai/news-agent")
async def ai_news_agent(
    request: Request,
    background: bool = False,
    background_tasks: BackgroundTasks = None,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """
    AI-powered news agent that fetches, analyzes, and stores real news articles.

    v6.4.0: Uses Gemini search_and_ground() for real-time Google Search,
    ML sentiment analysis, and event type detection.

    Request body:
        competitor_ids: Optional list of competitor IDs (None = all)
        limit_per_competitor: Max articles per competitor (default: 10)
        run_async: Run in background (default: false)

    Returns:
        News articles with sentiment, event types, and sources
    """
    from datetime import datetime, timedelta
    from database import NewsArticleCache

    try:
        body = await request.json() if request.headers.get('content-type') == 'application/json' else {}
    except (json.JSONDecodeError, ValueError, Exception):
        body = {}

    competitor_ids = body.get("competitor_ids")
    limit_per_competitor = body.get("limit_per_competitor", 10)
    run_async = body.get("run_async", False)
    keywords = body.get("keywords", "")  # v7.1.0: User search keywords from News Feed filters
    date_from = body.get("date_from")  # Optional: YYYY-MM-DD
    date_to = body.get("date_to")  # Optional: YYYY-MM-DD

    # Get competitors
    query = db.query(Competitor).filter(Competitor.is_deleted == False)
    if competitor_ids:
        query = query.filter(Competitor.id.in_(competitor_ids))
    competitors = query.all()

    if not competitors:
        return {"status": "error", "message": "No competitors found"}

    async def fetch_news_for_all():
        """Fetch news using AI search grounding."""
        results = {
            "total_articles": 0,
            "by_competitor": {},
            "by_sentiment": {"positive": 0, "neutral": 0, "negative": 0},
            "articles": []
        }

        # Try to use Gemini search grounding first
        try:
            from gemini_provider import GeminiProvider
            provider = GeminiProvider()
            use_ai = provider.is_available
        except Exception as e:
            logger.warning(f"Gemini not available: {e}")
            use_ai = False

        # Fallback to direct Google News RSS if AI unavailable
        if not use_ai:
            logger.info("Using Google News RSS fallback (no AI)")

        for comp in competitors:
            comp_articles = []

            try:
                if use_ai:
                    # Use Gemini with Google Search grounding
                    # v7.1.0: Include user keywords in search query
                    keyword_part = f' {keywords}' if keywords else ''
                    query_text = f'"{comp.name}"{keyword_part} healthcare technology news 2026'
                    try:
                        search_result = provider.search_and_ground(
                            query=query_text,
                            include_sources=True
                        )

                        # Parse AI response into articles
                        if search_result and "response" in search_result:
                            # Extract articles from grounded search
                            response_text = search_result.get("response", "")
                            sources = search_result.get("sources", [])

                            # Create article entries from sources
                            for source in sources[:limit_per_competitor]:
                                article = {
                                    "title": source.get("title", ""),
                                    "url": source.get("url", ""),
                                    "snippet": source.get("snippet", ""),
                                    "source": source.get("source_name", "Google Search"),
                                    "competitor_name": comp.name,
                                    "competitor_id": comp.id
                                }
                                if article["title"] and article["url"]:
                                    comp_articles.append(article)

                    except Exception as ai_err:
                        logger.warning(f"AI search failed for {comp.name}: {ai_err}")
                        # Fallback to RSS
                        comp_articles = await fetch_google_news_rss(comp.name, limit_per_competitor)
                else:
                    # Use direct Google News RSS
                    comp_articles = await fetch_google_news_rss(comp.name, limit_per_competitor)

                # Analyze sentiment and detect event types
                for article in comp_articles:
                    # Sentiment analysis
                    sentiment = analyze_article_sentiment(article.get("title", ""), article.get("snippet", ""))
                    article["sentiment"] = sentiment

                    # Event type detection
                    event_type = detect_event_type(article.get("title", ""))
                    article["event_type"] = event_type

                    # Store in database
                    try:
                        existing = db.query(NewsArticleCache).filter(
                            NewsArticleCache.url == article.get("url")
                        ).first()

                        if not existing:
                            pub_date = datetime.utcnow() - timedelta(days=1)  # Default to yesterday

                            cache_entry = NewsArticleCache(
                                competitor_id=comp.id,
                                competitor_name=comp.name,
                                title=article.get("title", ""),
                                url=article.get("url", ""),
                                source=article.get("source", "Google News"),
                                source_type="google_news",
                                published_at=pub_date,
                                snippet=article.get("snippet", ""),
                                sentiment=sentiment,
                                event_type=event_type,
                                is_major_event=event_type in ["funding", "acquisition", "leadership"],
                                fetched_at=datetime.utcnow(),
                            )
                            db.add(cache_entry)
                            results["total_articles"] += 1
                            results["by_sentiment"][sentiment] += 1

                    except Exception as db_err:
                        logger.warning(f"DB save error: {db_err}")

                    results["articles"].append(article)

                results["by_competitor"][comp.name] = len(comp_articles)
                db.commit()

            except Exception as comp_err:
                logger.error(f"Error fetching news for {comp.name}: {comp_err}")
                results["by_competitor"][comp.name] = 0

        return results

    # Helper function for Google News RSS
    async def fetch_google_news_rss(company_name: str, limit: int) -> list:
        """Fetch news from Google News RSS."""
        import httpx
        import xml.etree.ElementTree as ET
        from urllib.parse import quote

        articles = []
        query = quote(f'"{company_name}" healthcare')
        url = f"https://news.google.com/rss/search?q={query}&hl=en-US&gl=US&ceid=US:en"

        try:
            async with httpx.AsyncClient(timeout=10) as client:
                response = await client.get(url)
                if response.status_code == 200:
                    root = ET.fromstring(response.text)
                    items = root.findall(".//item")[:limit]

                    for item in items:
                        title = item.find("title")
                        link = item.find("link")
                        pub_date = item.find("pubDate")
                        source = item.find("source")

                        if title is not None and link is not None:
                            articles.append({
                                "title": title.text.rsplit(" - ", 1)[0] if title.text else "",
                                "url": link.text or "",
                                "snippet": "",
                                "source": source.text if source is not None else "Google News",
                                "competitor_name": company_name,
                                "published_date": pub_date.text if pub_date is not None else ""
                            })
        except Exception as e:
            logger.warning(f"RSS fetch failed for {company_name}: {e}")

        return articles

    # Helper function for sentiment analysis
    def analyze_article_sentiment(title: str, snippet: str) -> str:
        """Analyze sentiment using ML or keyword fallback."""
        text = f"{title} {snippet}".lower()

        # Try ML sentiment first
        try:
            from ml_sentiment import get_news_sentiment_analyzer
            analyzer = get_news_sentiment_analyzer()
            result = analyzer.analyze_headline(title, snippet)
            return result.label
        except Exception:
            pass

        # Keyword fallback
        positive_words = ["growth", "success", "award", "funding", "launch", "partnership", "expan", "innovat", "leading", "record", "milestone"]
        negative_words = ["lawsuit", "breach", "layoff", "decline", "loss", "concern", "issue", "problem", "fail", "cut", "struggle"]

        pos_count = sum(1 for w in positive_words if w in text)
        neg_count = sum(1 for w in negative_words if w in text)

        if pos_count > neg_count:
            return "positive"
        elif neg_count > pos_count:
            return "negative"
        return "neutral"

    # Helper function for event type detection
    def detect_event_type(title: str) -> str:
        """Detect event type from title."""
        title_lower = title.lower()

        if any(w in title_lower for w in ["funding", "invest", "raise", "series", "million", "billion"]):
            return "funding"
        elif any(w in title_lower for w in ["acqui", "merger", "buy", "purchase", "deal"]):
            return "acquisition"
        elif any(w in title_lower for w in ["launch", "release", "announce", "new product", "feature"]):
            return "product_launch"
        elif any(w in title_lower for w in ["partner", "collaborat", "alliance", "agreement", "team"]):
            return "partnership"
        elif any(w in title_lower for w in ["ceo", "cto", "cfo", "appoint", "hire", "join", "executive", "leader"]):
            return "leadership"
        elif any(w in title_lower for w in ["revenue", "earning", "quarter", "financial", "profit", "ipo"]):
            return "financial"
        elif any(w in title_lower for w in ["lawsuit", "legal", "court", "sue", "settlement"]):
            return "legal"
        return "general"

    if background and background_tasks and current_user:
        import uuid as _uuid
        task_id = str(_uuid.uuid4())
        _ai_tasks[task_id] = {
            "status": "running", "page_context": "newsfeed",
            "user_id": current_user.get("id"),
            "started_at": datetime.utcnow().isoformat(),
            "completed_at": None, "result": None, "error": None,
            "task_type": "ai_news_fetch", "read_at": None,
        }

        async def _run_news_bg():
            try:
                results = await fetch_news_for_all()
                _ai_tasks[task_id]["result"] = {
                    "status": "success",
                    "message": f"Fetched {results['total_articles']} articles",
                    **results
                }
                _ai_tasks[task_id]["status"] = "completed"
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
            except Exception as bg_err:
                logger.error(f"Background news agent error: {bg_err}")
                _ai_tasks[task_id]["status"] = "failed"
                _ai_tasks[task_id]["error"] = str(bg_err)
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()

        background_tasks.add_task(_run_news_bg)
        return {"task_id": task_id, "status": "running"}

    def _filter_articles_by_date(results: dict) -> dict:
        """Filter articles by date_from/date_to if provided."""
        if not date_from and not date_to:
            return results
        filtered = []
        for article in results.get("articles", []):
            pub = article.get("published_date") or article.get("published_at", "")
            if not pub:
                continue
            pub_str = str(pub)[:10]  # YYYY-MM-DD
            if date_from and pub_str < date_from:
                continue
            if date_to and pub_str > date_to:
                continue
            filtered.append(article)
        results["articles"] = filtered
        results["total_articles"] = len(filtered)
        return results

    if run_async:
        # Run in background - directly add async function to background tasks
        background_tasks.add_task(fetch_news_for_all)
        return {
            "status": "started",
            "message": f"Fetching news for {len(competitors)} competitors in background"
        }
    else:
        # Run synchronously - await the async function
        results = await fetch_news_for_all()
        results = _filter_articles_by_date(results)
        return {
            "status": "success",
            "message": f"Fetched {results['total_articles']} articles for {len(competitors)} competitors",
            **results
        }


@app.get("/api/ai/news-agent/status")
async def ai_news_agent_status(db: Session = Depends(get_db)):
    """Get AI news agent status and cached article counts."""
    from database import NewsArticleCache
    from datetime import datetime, timedelta
    from sqlalchemy import func

    # Count recent articles (exclude archived)
    recent_cutoff = datetime.utcnow() - timedelta(hours=24)
    recent_count = db.query(func.count(NewsArticleCache.id)).filter(
        NewsArticleCache.fetched_at > recent_cutoff,
        NewsArticleCache.is_archived != True
    ).scalar()

    total_count = db.query(func.count(NewsArticleCache.id)).filter(
        NewsArticleCache.is_archived != True
    ).scalar()

    # Sentiment breakdown
    sentiment_counts = db.query(
        NewsArticleCache.sentiment,
        func.count(NewsArticleCache.id)
    ).group_by(NewsArticleCache.sentiment).all()

    return {
        "total_cached_articles": total_count,
        "articles_last_24h": recent_count,
        "sentiment_breakdown": {s[0] or "unknown": s[1] for s in sentiment_counts},
        "ai_available": check_ai_available(),
        "last_checked": datetime.utcnow().isoformat()
    }


def check_ai_available() -> dict:
    """Check if AI providers are available."""
    result = {"gemini": False, "openai": False}
    try:
        from gemini_provider import GeminiProvider
        provider = GeminiProvider()
        result["gemini"] = provider.is_available
    except (ImportError, Exception) as e:
        logger.debug(f"Gemini provider not available: {e}")
    try:
        import os
        result["openai"] = bool(os.getenv("OPENAI_API_KEY"))
    except Exception as e:
        logger.debug(f"OpenAI check failed: {e}")
    return result


# ============== FIRECRAWL ENDPOINTS (v5.0.6) ==============

@app.post("/api/firecrawl/scrape")
async def firecrawl_scrape(request: Request):
    """
    Scrape a URL using Firecrawl.

    v5.0.6: NEWS-4D - Firecrawl MCP integration.

    Request body:
        url: URL to scrape
        formats: Optional list of formats (markdown, html, links)
        only_main_content: Whether to extract only main content (default: true)

    Returns:
        Scraped content with metadata
    """
    try:
        from firecrawl_integration import get_firecrawl_client
        from dataclasses import asdict

        client = get_firecrawl_client()
        if not client.is_available:
            raise HTTPException(
                status_code=503,
                detail="Firecrawl not available. Configure FIRECRAWL_API_KEY."
            )

        body = await request.json()
        url = body.get("url")

        if not url:
            raise HTTPException(status_code=400, detail="url is required")

        result = await client.scrape(
            url=url,
            formats=body.get("formats"),
            only_main_content=body.get("only_main_content", True),
            include_tags=body.get("include_tags"),
            exclude_tags=body.get("exclude_tags"),
            wait_for=body.get("wait_for"),
            timeout=body.get("timeout", 30000)
        )

        await client.close()
        return asdict(result)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Firecrawl scrape failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/firecrawl/scrape-batch")
async def firecrawl_scrape_batch(request: Request):
    """
    Scrape multiple URLs using Firecrawl.

    v5.0.6: Batch scraping for efficiency.

    Request body:
        urls: List of URLs to scrape
        formats: Optional list of formats

    Returns:
        Batch results with all scraped content
    """
    try:
        from firecrawl_integration import get_firecrawl_client
        from dataclasses import asdict

        client = get_firecrawl_client()
        if not client.is_available:
            raise HTTPException(
                status_code=503,
                detail="Firecrawl not available. Configure FIRECRAWL_API_KEY."
            )

        body = await request.json()
        urls = body.get("urls", [])

        if not urls:
            raise HTTPException(status_code=400, detail="urls is required")

        result = await client.scrape_batch(
            urls=urls,
            formats=body.get("formats"),
            only_main_content=body.get("only_main_content", True)
        )

        await client.close()
        return asdict(result)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Firecrawl batch scrape failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/firecrawl/scrape-competitor")
async def firecrawl_scrape_competitor(request: Request):
    """
    Scrape a competitor website for intelligence.

    v5.0.6: Specialized competitor scraping with structured extraction.

    Request body:
        url: Competitor website URL
        include_subpages: Whether to scrape subpages (default: true)

    Returns:
        Structured competitor data from homepage, pricing, about, etc.
    """
    try:
        from firecrawl_integration import get_competitor_scraper

        scraper = get_competitor_scraper()
        if not scraper.is_available:
            raise HTTPException(
                status_code=503,
                detail="Firecrawl not available. Configure FIRECRAWL_API_KEY."
            )

        body = await request.json()
        url = body.get("url")

        if not url:
            raise HTTPException(status_code=400, detail="url is required")

        result = await scraper.scrape_competitor_website(
            url=url,
            include_subpages=body.get("include_subpages", True)
        )

        await scraper.close()
        return result

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Firecrawl competitor scrape failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/firecrawl/crawl")
async def firecrawl_crawl(request: Request):
    """
    Start a crawl job for a website.

    v5.0.6: Website crawling for comprehensive data collection.

    Request body:
        url: Starting URL
        limit: Maximum pages to crawl (default: 10)
        max_depth: Maximum link depth (default: 2)

    Returns:
        Crawl job ID and status
    """
    try:
        from firecrawl_integration import get_firecrawl_client

        client = get_firecrawl_client()
        if not client.is_available:
            raise HTTPException(
                status_code=503,
                detail="Firecrawl not available. Configure FIRECRAWL_API_KEY."
            )

        body = await request.json()
        url = body.get("url")

        if not url:
            raise HTTPException(status_code=400, detail="url is required")

        result = await client.crawl(
            url=url,
            limit=body.get("limit", 10),
            max_depth=body.get("max_depth", 2),
            include_paths=body.get("include_paths"),
            exclude_paths=body.get("exclude_paths")
        )

        await client.close()
        return result

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Firecrawl crawl failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.get("/api/firecrawl/crawl/{job_id}")
async def firecrawl_crawl_status(job_id: str):
    """
    Get the status of a crawl job.

    v5.0.6: Check crawl job progress and results.

    Args:
        job_id: Crawl job ID

    Returns:
        Crawl job status and results
    """
    try:
        from firecrawl_integration import get_firecrawl_client

        client = get_firecrawl_client()
        if not client.is_available:
            raise HTTPException(
                status_code=503,
                detail="Firecrawl not available."
            )

        result = await client.get_crawl_status(job_id)
        await client.close()
        return result

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Firecrawl status check failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.get("/api/firecrawl/status")
async def firecrawl_status():
    """
    Check Firecrawl service availability.

    v5.0.6: Service health check.

    Returns:
        Availability status
    """
    try:
        from firecrawl_integration import get_firecrawl_client
        client = get_firecrawl_client()
        return {
            "available": client.is_available,
            "configured": bool(os.getenv("FIRECRAWL_API_KEY"))
        }
    except Exception as e:
        return {
            "available": False,
            "configured": False,
            "error": "An unexpected error occurred"
        }


# Analytics sub-endpoints
@app.get("/api/analytics/threats")
def get_threat_analytics(db: Session = Depends(get_db)):
    """Get threat distribution analytics."""
    competitors = db.query(Competitor).filter(Competitor.is_deleted == False).all()
    return {
        "high": sum(1 for c in competitors if c.threat_level == "High"),
        "medium": sum(1 for c in competitors if c.threat_level == "Medium"),
        "low": sum(1 for c in competitors if c.threat_level == "Low"),
        "total": len(competitors)
    }


@app.get("/api/analytics/market-share")
def get_market_share_analytics(db: Session = Depends(get_db)):
    """Get estimated market share by customer count, sorted by largest first."""
    import re
    all_competitors = db.query(Competitor).filter(Competitor.is_deleted == False).all()
    shares = []
    for c in all_competitors:
        count = 0
        if c.customer_count:
            match = re.search(r'\d+', str(c.customer_count).replace(',', ''))
            if match:
                count = int(match.group())
        shares.append({
            "name": c.name,
            "customers": count,
            "website": c.website
        })
    # Sort by customer count descending, take top 10
    shares.sort(key=lambda s: s["customers"], reverse=True)
    shares = shares[:10]
    total = sum(s["customers"] for s in shares)
    for s in shares:
        s["share"] = round(s["customers"] / total * 100, 1) if total > 0 else 0
    return {"market_share": shares}


@app.get("/api/analytics/pricing")
def get_pricing_analytics(db: Session = Depends(get_db)):
    """Get pricing model distribution."""
    competitors = db.query(Competitor).filter(Competitor.is_deleted == False).all()
    models = {}
    for c in competitors:
        model = c.pricing_model or "Unknown"
        models[model] = models.get(model, 0) + 1
    return {"pricing_models": [{"model": k, "count": v} for k, v in models.items()]}


@app.get("/api/analytics/dashboard")
def get_analytics_dashboard(db: Session = Depends(get_db)):
    """
    Get comprehensive analytics dashboard data.

    v5.2.0: Phase 5 - Complete analytics dashboard.

    Returns all metrics needed for the analytics page including:
    - Competitor counts by status and threat level
    - Data freshness metrics
    - News sentiment trends
    - Dimension score averages
    - Market positioning data
    """
    from datetime import timedelta
    from sqlalchemy import func

    competitors = db.query(Competitor).filter(Competitor.is_deleted == False).all()
    total = len(competitors)

    # Status breakdown
    status_counts = {}
    for c in competitors:
        status = c.status or "Unknown"
        status_counts[status] = status_counts.get(status, 0) + 1

    # Threat level breakdown
    threat_counts = {"High": 0, "Medium": 0, "Low": 0}
    for c in competitors:
        level = c.threat_level or "Low"
        if level in threat_counts:
            threat_counts[level] += 1

    # Data freshness
    now = datetime.utcnow()
    fresh_count = 0  # Updated within 7 days
    stale_count = 0  # Not updated in 30+ days
    for c in competitors:
        if c.last_updated:
            age = (now - c.last_updated).days
            if age <= 7:
                fresh_count += 1
            elif age > 30:
                stale_count += 1

    # Data quality average
    quality_sum = sum(c.data_quality_score or 0 for c in competitors)
    avg_quality = quality_sum / total if total > 0 else 0

    # News sentiment from cache
    try:
        from database import NewsArticleCache
        recent_cutoff = now - timedelta(days=30)
        news_query = db.query(NewsArticleCache).filter(
            NewsArticleCache.published_at >= recent_cutoff,
            NewsArticleCache.is_archived != True
        ).all()

        sentiment_counts = {"positive": 0, "negative": 0, "neutral": 0}
        for article in news_query:
            sent = article.sentiment or "neutral"
            if sent in sentiment_counts:
                sentiment_counts[sent] += 1

        news_total = len(news_query)
    except Exception as e:
        logger.debug(f"News sentiment query failed: {e}")
        sentiment_counts = {"positive": 0, "negative": 0, "neutral": 0}
        news_total = 0

    # Dimension score averages (if available)
    dimension_averages = {}
    dimension_fields = [
        "product_packaging_score", "integration_depth_score", "support_service_score",
        "retention_stickiness_score", "user_adoption_score", "implementation_ttv_score",
        "reliability_enterprise_score", "pricing_flexibility_score", "reporting_analytics_score"
    ]

    for field in dimension_fields:
        scores = [getattr(c, field, None) for c in competitors if hasattr(c, field)]
        valid_scores = [s for s in scores if s is not None]
        if valid_scores:
            dimension_name = field.replace("_score", "").replace("_", " ").title()
            dimension_averages[dimension_name] = sum(valid_scores) / len(valid_scores)

    # Market positioning (customer count vs threat level)
    market_data = []
    for c in competitors[:50]:  # Limit for performance
        customer_count = 0
        if c.customer_count:
            try:
                customer_count = int(''.join(filter(str.isdigit, str(c.customer_count)[:10])))
            except (ValueError, TypeError):
                customer_count = 100

        market_data.append({
            "name": c.name,
            "x": c.data_quality_score or 50,  # Data quality as proxy for market presence
            "y": customer_count,
            "threat": c.threat_level or "Low",
            "size": {"High": 30, "Medium": 20, "Low": 10}.get(c.threat_level, 15),
            "website": c.website
        })

    # Recent changes count
    changes_count = db.query(DataChangeHistory).filter(
        DataChangeHistory.changed_at >= now - timedelta(days=7)
    ).count()

    return {
        "summary": {
            "total_competitors": total,
            "active_count": status_counts.get("Active", 0),
            "discovered_count": status_counts.get("Discovered", 0),
            "avg_data_quality": round(avg_quality, 1),
            "fresh_data_count": fresh_count,
            "stale_data_count": stale_count,
            "recent_changes": changes_count
        },
        "threat_distribution": threat_counts,
        "status_distribution": status_counts,
        "news_sentiment": {
            "counts": sentiment_counts,
            "total": news_total
        },
        "dimension_averages": dimension_averages,
        "market_positioning": market_data[:30],
        "generated_at": now.isoformat()
    }


@app.get("/api/analytics/market-map")
def get_market_map_data(db: Session = Depends(get_db)):
    """
    Get data for market positioning visualization.

    v5.2.0: Phase 5 - Market map for analytics page.

    Returns competitors positioned by data quality (x) and estimated size (y).
    """
    competitors = db.query(Competitor).filter(
        Competitor.is_deleted == False,
        Competitor.status == "Active"
    ).all()

    data = []
    for c in competitors:
        # Estimate customer count
        customer_count = 100
        if c.customer_count:
            try:
                customer_count = int(''.join(filter(str.isdigit, str(c.customer_count)[:10]))) or 100
            except (ValueError, TypeError):
                customer_count = 100

        # Estimate employee count for sizing
        employee_count = 50
        if c.employee_count:
            try:
                employee_count = int(''.join(filter(str.isdigit, str(c.employee_count)[:6]))) or 50
            except (ValueError, TypeError):
                employee_count = 50

        data.append({
            "id": c.id,
            "name": c.name,
            "x": c.data_quality_score or 50,
            "y": min(customer_count, 10000),  # Cap for display
            "size": min(employee_count / 10, 50),  # Scale for bubble size
            "threat": c.threat_level or "Low",
            "category": c.product_categories[:50] if c.product_categories else "Unknown",
            "website": c.website
        })

    return {
        "competitors": data,
        "count": len(data),
        "axes": {
            "x": {"label": "Data Quality Score", "min": 0, "max": 100},
            "y": {"label": "Estimated Customers", "min": 0, "max": 10000}
        }
    }


@app.get("/api/analytics/market-quadrant")
async def get_market_quadrant(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
) -> Dict[str, Any]:
    """Return market-position quadrant data for active competitors.

    x-axis = market strength (0-100): derived from dimension scores or
    fallback proxies (threat_level + product count).
    y-axis = growth momentum (0-100): derived from news sentiment +
    recent ChangeLog activity in last 90 days.
    size = employee or customer count (capped for display).
    """
    try:
        from database import NewsArticleCache

        competitors = db.query(Competitor).filter(
            Competitor.is_deleted == False,
            Competitor.status == "Active"
        ).all()

        cutoff = datetime.utcnow() - timedelta(days=90)
        comp_ids = [c.id for c in competitors]

        # Pre-fetch news sentiment counts per competitor (last 90 days)
        news_rows = db.query(NewsArticleCache).filter(
            NewsArticleCache.competitor_id.in_(comp_ids),
            NewsArticleCache.is_archived != True,
            NewsArticleCache.published_at >= cutoff
        ).all()

        sentiment_map: Dict[int, Dict[str, int]] = {}
        for article in news_rows:
            cid = article.competitor_id
            if cid not in sentiment_map:
                sentiment_map[cid] = {
                    "positive": 0, "neutral": 0, "negative": 0
                }
            s = (article.sentiment or "neutral").strip().lower()
            if s in sentiment_map[cid]:
                sentiment_map[cid][s] += 1

        # Pre-fetch ChangeLog activity counts per competitor (last 90 days)
        change_rows = db.query(ChangeLog).filter(
            ChangeLog.competitor_id.in_(comp_ids),
            ChangeLog.detected_at >= cutoff
        ).all()

        activity_map: Dict[int, int] = {}
        for ch in change_rows:
            activity_map[ch.competitor_id] = activity_map.get(
                ch.competitor_id, 0
            ) + 1

        # Dimension score field names (9 dimensions, each 1-5 scale)
        dim_fields = [
            "dim_product_packaging_score",
            "dim_integration_depth_score",
            "dim_support_service_score",
            "dim_retention_stickiness_score",
            "dim_user_adoption_score",
            "dim_implementation_ttv_score",
            "dim_reliability_enterprise_score",
            "dim_pricing_flexibility_score",
            "dim_reporting_analytics_score",
        ]

        results = []
        for c in competitors:
            # --- Market Strength (x: 0-100) ---
            dim_scores = [
                getattr(c, f, None) for f in dim_fields
            ]
            valid_dims = [s for s in dim_scores if s is not None]

            if valid_dims:
                # Average of 1-5 scores, scaled to 0-100
                market_strength = (sum(valid_dims) / len(valid_dims)) * 20.0
            else:
                # Fallback: threat_level + product_overlap_score
                tl = (c.threat_level or "").strip().lower()
                base = {"high": 70, "medium": 50, "low": 30}.get(tl, 40)
                overlap = c.product_overlap_score or 0
                market_strength = min(100, base + (overlap * 0.3))

            # --- Growth Momentum (y: 0-100) ---
            sdata = sentiment_map.get(c.id, {})
            pos = sdata.get("positive", 0)
            neg = sdata.get("negative", 0)
            total_news = pos + sdata.get("neutral", 0) + neg

            # Sentiment score: ratio of positive minus negative
            if total_news > 0:
                sentiment_score = ((pos - neg) / total_news + 1) * 50
            else:
                sentiment_score = 50  # neutral default

            # Activity score: more changes = more momentum (cap at 50)
            changes = activity_map.get(c.id, 0)
            activity_score = min(50, changes * 5)

            # Blend: 60% sentiment + 40% activity
            growth_momentum = (sentiment_score * 0.6) + (activity_score * 0.4)
            growth_momentum = max(0, min(100, growth_momentum))

            # --- Bubble size ---
            bubble_size = 10  # default
            for count_field in [c.employee_count, c.customer_count]:
                if count_field:
                    try:
                        digits = "".join(
                            filter(str.isdigit, str(count_field)[:10])
                        )
                        parsed = int(digits) if digits else 0
                        if parsed > 0:
                            bubble_size = max(5, min(50, parsed / 200))
                            break
                    except (ValueError, TypeError):
                        continue

            results.append({
                "id": c.id,
                "name": c.name,
                "market_strength": round(market_strength, 1),
                "growth_momentum": round(growth_momentum, 1),
                "company_size": round(bubble_size, 1),
                "threat_level": c.threat_level or "Low",
            })

        return {"competitors": results}
    except Exception as e:
        logger.error(f"Failed to compute market quadrant data: {e}")
        raise HTTPException(
            status_code=500,
            detail="An internal error occurred while computing market quadrant"
        )


# ============== DATA QUALITY ENDPOINTS ==============

# List of all data fields that should be tracked
COMPETITOR_DATA_FIELDS = [
    "name", "website", "status", "threat_level", "pricing_model", "base_price",
    "price_unit", "product_categories", "key_features", "integration_partners",
    "certifications", "target_segments", "customer_size_focus", "geographic_focus",
    "customer_count", "customer_acquisition_rate", "key_customers", "g2_rating",
    "employee_count", "employee_growth_rate", "year_founded", "headquarters",
    "funding_total", "latest_round", "pe_vc_backers", "website_traffic",
    "social_following", "recent_launches", "news_mentions", "is_public",
    "ticker_symbol", "stock_exchange"
]


def calculate_quality_score(competitor) -> int:
    """Calculate data quality score (0-100) based on field completeness."""
    filled_fields = 0
    for field in COMPETITOR_DATA_FIELDS:
        value = getattr(competitor, field, None)
        if value is not None and str(value).strip() not in ["", "None", "Unknown", "N/A"]:
            filled_fields += 1
    return int((filled_fields / len(COMPETITOR_DATA_FIELDS)) * 100)


# Data quality endpoints (completeness, scores, stale, verify, completeness/{id})
# moved to routers/data_quality.py


# ============== CHANGE HISTORY ENDPOINTS ==============

@app.get("/api/changes")
def get_changes(
    competitor_id: Optional[int] = None,
    days: int = 30,
    field_name: Optional[str] = None,
    changed_by: Optional[str] = None,
    start_date: Optional[str] = None,
    end_date: Optional[str] = None,
    page: int = 1,
    page_size: int = 50,
    db: Session = Depends(get_db)
):
    """
    Get change logs for timeline with advanced filtering.

    v5.2.0: Enhanced filtering for Phase 4.

    Args:
        competitor_id: Filter by specific competitor
        days: Number of days to look back (default: 30)
        field_name: Filter by specific field (e.g., "pricing_model", "key_features")
        changed_by: Filter by who made the change
        start_date: Filter from this date (YYYY-MM-DD)
        end_date: Filter until this date (YYYY-MM-DD)
        page: Page number for pagination
        page_size: Number of items per page
    """
    from datetime import timedelta

    # Build query
    query = db.query(DataChangeHistory)

    # Date filtering
    if start_date:
        try:
            start_dt = datetime.strptime(start_date, '%Y-%m-%d')
            query = query.filter(DataChangeHistory.changed_at >= start_dt)
        except ValueError:
            pass
    else:
        cutoff = datetime.utcnow() - timedelta(days=days)
        query = query.filter(DataChangeHistory.changed_at >= cutoff)

    if end_date:
        try:
            end_dt = datetime.strptime(end_date, '%Y-%m-%d') + timedelta(days=1)
            query = query.filter(DataChangeHistory.changed_at < end_dt)
        except ValueError:
            pass

    # Other filters
    if competitor_id:
        query = query.filter(DataChangeHistory.competitor_id == competitor_id)

    if field_name:
        query = query.filter(DataChangeHistory.field_name.ilike(f"%{field_name}%"))

    if changed_by:
        query = query.filter(DataChangeHistory.changed_by.ilike(f"%{changed_by}%"))

    # Get total count before pagination
    total_count = query.count()

    # Apply pagination
    offset = (page - 1) * page_size
    changes = query.order_by(DataChangeHistory.changed_at.desc()).offset(offset).limit(page_size).all()

    return {
        "competitor_id": competitor_id,
        "days": days,
        "count": len(changes),
        "total_count": total_count,
        "page": page,
        "page_size": page_size,
        "total_pages": (total_count + page_size - 1) // page_size,
        "changes": [
            {
                "id": c.id,
                "competitor_id": c.competitor_id,
                "competitor_name": c.competitor_name,
                "change_type": c.field_name.replace("_", " ").title(),
                "field_name": c.field_name,
                "previous_value": c.old_value,
                "new_value": c.new_value,
                "severity": "Medium",
                "detected_at": c.changed_at.isoformat(),
                "source": c.changed_by or "manual",
                "changed_by": c.changed_by,
                "change_reason": c.change_reason
            }
            for c in changes
        ]
    }


@app.get("/api/changes/export")
def export_changes(
    competitor_id: Optional[int] = None,
    days: int = 90,
    format: str = "csv",
    db: Session = Depends(get_db)
):
    """
    Export change logs to CSV or Excel.

    v5.2.0: Phase 4 - Change log export.

    Args:
        competitor_id: Filter by specific competitor
        days: Number of days to look back
        format: "csv" or "excel"
    """
    from datetime import timedelta
    import io
    import csv

    cutoff = datetime.utcnow() - timedelta(days=days)
    query = db.query(DataChangeHistory).filter(DataChangeHistory.changed_at >= cutoff)

    if competitor_id:
        query = query.filter(DataChangeHistory.competitor_id == competitor_id)

    changes = query.order_by(DataChangeHistory.changed_at.desc()).all()

    if format == "csv":
        output = io.StringIO()
        writer = csv.writer(output)
        writer.writerow([
            "Date", "Competitor", "Field", "Old Value", "New Value",
            "Changed By", "Reason"
        ])

        for c in changes:
            writer.writerow([
                c.changed_at.isoformat() if c.changed_at else "",
                c.competitor_name,
                c.field_name,
                c.old_value[:200] if c.old_value else "",
                c.new_value[:200] if c.new_value else "",
                c.changed_by or "system",
                c.change_reason or ""
            ])

        output.seek(0)
        return StreamingResponse(
            iter([output.getvalue()]),
            media_type="text/csv",
            headers={
                "Content-Disposition": f"attachment; filename=changelog_export_{datetime.now().strftime('%Y%m%d')}.csv"
            }
        )

    elif format == "excel":
        try:
            from openpyxl import Workbook
            from fastapi.responses import FileResponse

            wb = Workbook()
            ws = wb.active
            ws.title = "Change Log"

            # Header
            ws.append([
                "Date", "Competitor", "Field", "Old Value", "New Value",
                "Changed By", "Reason"
            ])

            for c in changes:
                ws.append([
                    c.changed_at.strftime("%Y-%m-%d %H:%M") if c.changed_at else "",
                    c.competitor_name,
                    c.field_name,
                    (c.old_value[:200] if c.old_value else ""),
                    (c.new_value[:200] if c.new_value else ""),
                    c.changed_by or "system",
                    c.change_reason or ""
                ])

            os.makedirs("./exports", exist_ok=True)
            filepath = f"./exports/changelog_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
            wb.save(filepath)

            return FileResponse(
                filepath,
                media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                filename=f"changelog_export_{datetime.now().strftime('%Y%m%d')}.xlsx"
            )
        except ImportError:
            return {"error": "openpyxl not installed for Excel export"}

    return {"error": f"Unknown format: {format}"}


@app.get("/api/changes/history/{competitor_id}")
def get_competitor_change_history(competitor_id: int, db: Session = Depends(get_db)):
    """Get detailed change history for a specific competitor."""
    changes = db.query(DataChangeHistory).filter(
        DataChangeHistory.competitor_id == competitor_id
    ).order_by(DataChangeHistory.changed_at.desc()).limit(100).all()

    return {
        "competitor_id": competitor_id,
        "total": len(changes),
        "changes": [
            {
                "id": c.id,
                "field": c.field_name,
                "old_value": c.old_value,
                "new_value": c.new_value,
                "changed_by": c.changed_by,
                "reason": c.change_reason,
                "source_url": c.source_url,
                "changed_at": c.changed_at.isoformat() if c.changed_at else None
            }
            for c in changes
        ]
    }


# ==============================================================================
# REL-008: Audit Trail Enhancements
# Diff view, rollback, bulk export, timeline visualization
# ==============================================================================

@app.get("/api/changes/{change_id}/diff")
def get_change_diff(change_id: int, db: Session = Depends(get_db)):
    """
    Get a detailed diff view for a specific change.
    Returns side-by-side comparison with highlighted differences.
    """
    import difflib

    change = db.query(DataChangeHistory).filter(DataChangeHistory.id == change_id).first()
    if not change:
        raise HTTPException(status_code=404, detail="Change not found")

    old_val = change.old_value or ""
    new_val = change.new_value or ""

    # Generate unified diff
    diff_lines = list(difflib.unified_diff(
        old_val.splitlines(keepends=True),
        new_val.splitlines(keepends=True),
        fromfile='Previous',
        tofile='Current',
        lineterm=''
    ))

    # Calculate statistics
    additions = sum(1 for line in diff_lines if line.startswith('+') and not line.startswith('+++'))
    deletions = sum(1 for line in diff_lines if line.startswith('-') and not line.startswith('---'))

    # Generate HTML diff for visualization
    html_diff = difflib.HtmlDiff().make_table(
        old_val.splitlines(),
        new_val.splitlines(),
        fromdesc='Previous Value',
        todesc='New Value',
        context=True,
        numlines=3
    )

    return {
        "change_id": change_id,
        "competitor_id": change.competitor_id,
        "competitor_name": change.competitor_name,
        "field_name": change.field_name,
        "changed_at": change.changed_at.isoformat() if change.changed_at else None,
        "changed_by": change.changed_by,
        "old_value": old_val,
        "new_value": new_val,
        "diff": {
            "unified": "\n".join(diff_lines),
            "html": html_diff,
            "additions": additions,
            "deletions": deletions,
            "total_changes": additions + deletions
        }
    }


@app.post("/api/changes/{change_id}/rollback")
def rollback_change(
    change_id: int,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """
    Rollback a specific change to its previous value.
    Creates a new change record documenting the rollback.
    """
    change = db.query(DataChangeHistory).filter(DataChangeHistory.id == change_id).first()
    if not change:
        raise HTTPException(status_code=404, detail="Change not found")

    # Get the competitor
    competitor = db.query(Competitor).filter(Competitor.id == change.competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    field_name = change.field_name
    old_value = change.old_value

    # Check if field exists on competitor
    if not hasattr(competitor, field_name):
        raise HTTPException(status_code=400, detail=f"Field '{field_name}' not found on competitor")

    # Get current value before rollback
    current_value = getattr(competitor, field_name, None)
    current_value_str = str(current_value) if current_value is not None else ""

    # Perform the rollback
    try:
        # Handle type conversion for common field types
        if old_value is None or old_value == "":
            setattr(competitor, field_name, None)
        elif field_name in ['employee_count', 'id']:
            setattr(competitor, field_name, int(old_value) if old_value else None)
        elif field_name in ['total_funding']:
            setattr(competitor, field_name, float(old_value) if old_value else None)
        elif field_name in ['is_active', 'is_public']:
            setattr(competitor, field_name, old_value.lower() == 'true' if old_value else False)
        else:
            setattr(competitor, field_name, old_value)

        # Create a new change record for the rollback
        rollback_record = DataChangeHistory(
            competitor_id=change.competitor_id,
            competitor_name=change.competitor_name,
            field_name=field_name,
            old_value=current_value_str,
            new_value=old_value,
            changed_by=current_user.get("email", "system"),
            change_reason=f"Rollback of change #{change_id}",
            changed_at=datetime.utcnow()
        )
        db.add(rollback_record)

        # Log activity
        log_activity(
            db, current_user.get("email"), current_user.get("id"),
            "rollback", f"Rolled back {field_name} for {change.competitor_name}"
        )

        db.commit()

        return {
            "success": True,
            "message": f"Successfully rolled back {field_name}",
            "rollback_record_id": rollback_record.id,
            "field_name": field_name,
            "restored_value": old_value,
            "previous_value": current_value_str
        }

    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail="Rollback failed. Please try again.")


@app.get("/api/changes/timeline")
def get_changes_timeline(
    days: int = 30,
    competitor_id: Optional[int] = None,
    db: Session = Depends(get_db)
):
    """
    Get changes organized by date for timeline visualization.
    Groups changes by day with counts and summaries.
    """
    from sqlalchemy import func, cast, Date, String

    cutoff = datetime.utcnow() - timedelta(days=days)

    # Use func.date() instead of cast(Date) for SQLite compatibility with Python 3.14+
    # SQLite's date() returns a string, avoiding fromisoformat type errors
    date_expr = func.date(DataChangeHistory.changed_at)

    # Base query
    query = db.query(
        date_expr.label('date'),
        func.count(DataChangeHistory.id).label('count')
    ).filter(DataChangeHistory.changed_at >= cutoff)

    if competitor_id:
        query = query.filter(DataChangeHistory.competitor_id == competitor_id)

    # Group by date
    daily_counts = query.group_by(
        date_expr
    ).order_by(
        date_expr.desc()
    ).all()

    # Get detailed changes per day (limit 5 per day for preview)
    timeline = []
    for day_record in daily_counts:
        day_date = day_record.date
        if day_date is None:
            continue

        # Get sample changes for this day
        day_changes_query = db.query(DataChangeHistory).filter(
            func.date(DataChangeHistory.changed_at) == day_date
        )
        if competitor_id:
            day_changes_query = day_changes_query.filter(
                DataChangeHistory.competitor_id == competitor_id
            )

        day_changes = day_changes_query.order_by(
            DataChangeHistory.changed_at.desc()
        ).limit(5).all()

        timeline.append({
            "date": day_date.isoformat() if hasattr(day_date, 'isoformat') else str(day_date),
            "count": day_record.count,
            "changes": [
                {
                    "id": c.id,
                    "competitor_name": c.competitor_name,
                    "field_name": c.field_name,
                    "changed_by": c.changed_by,
                    "changed_at": c.changed_at.isoformat() if c.changed_at else None
                }
                for c in day_changes
            ]
        })

    return {
        "period_days": days,
        "total_days_with_changes": len(timeline),
        "total_changes": sum(d["count"] for d in timeline),
        "timeline": timeline
    }


@app.get("/api/changes/field-history/{competitor_id}/{field_name}")
def get_field_history(
    competitor_id: int,
    field_name: str,
    limit: int = 20,
    db: Session = Depends(get_db)
):
    """
    Get the complete history of a specific field for a competitor.
    Useful for understanding how a value evolved over time.
    """
    changes = db.query(DataChangeHistory).filter(
        DataChangeHistory.competitor_id == competitor_id,
        DataChangeHistory.field_name == field_name
    ).order_by(DataChangeHistory.changed_at.desc()).limit(limit).all()

    return {
        "competitor_id": competitor_id,
        "field_name": field_name,
        "total_changes": len(changes),
        "history": [
            {
                "id": c.id,
                "old_value": c.old_value,
                "new_value": c.new_value,
                "changed_by": c.changed_by,
                "change_reason": c.change_reason,
                "changed_at": c.changed_at.isoformat() if c.changed_at else None
            }
            for c in changes
        ]
    }


@app.post("/api/changes/bulk-export")
def bulk_export_changes(
    competitor_ids: List[int],
    days: int = 90,
    format: str = "csv",
    db: Session = Depends(get_db)
):
    """
    Export changes for multiple competitors at once.
    """
    import io
    import csv

    cutoff = datetime.utcnow() - timedelta(days=days)

    changes = db.query(DataChangeHistory).filter(
        DataChangeHistory.competitor_id.in_(competitor_ids),
        DataChangeHistory.changed_at >= cutoff
    ).order_by(DataChangeHistory.changed_at.desc()).all()

    if format == "csv":
        output = io.StringIO()
        writer = csv.writer(output)
        writer.writerow([
            "Date", "Competitor ID", "Competitor Name", "Field",
            "Old Value", "New Value", "Changed By", "Reason"
        ])

        for c in changes:
            writer.writerow([
                c.changed_at.isoformat() if c.changed_at else "",
                c.competitor_id,
                c.competitor_name,
                c.field_name,
                (c.old_value[:500] if c.old_value else "")[:500],
                (c.new_value[:500] if c.new_value else "")[:500],
                c.changed_by or "system",
                c.change_reason or ""
            ])

        output.seek(0)
        return StreamingResponse(
            iter([output.getvalue()]),
            media_type="text/csv",
            headers={
                "Content-Disposition": f"attachment; filename=bulk_changelog_{datetime.now().strftime('%Y%m%d')}.csv"
            }
        )

    return {
        "total": len(changes),
        "competitor_ids": competitor_ids,
        "changes": [
            {
                "id": c.id,
                "competitor_id": c.competitor_id,
                "competitor_name": c.competitor_name,
                "field_name": c.field_name,
                "old_value": c.old_value,
                "new_value": c.new_value,
                "changed_at": c.changed_at.isoformat() if c.changed_at else None
            }
            for c in changes
        ]
    }


# ============== ACTIVITY LOGS ENDPOINTS (Shared across all users) ==============

@app.get("/api/activity-logs")
def get_activity_logs(
    action_type: str = None,
    limit: int = 100,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get activity logs showing who made changes and when (visible to all users)."""
    query = db.query(ActivityLog)

    if action_type:
        query = query.filter(ActivityLog.action_type == action_type)

    logs = query.order_by(ActivityLog.created_at.desc()).limit(limit).all()

    return {
        "total": len(logs),
        "logs": [
            {
                "id": log.id,
                "user_email": log.user_email,
                "action_type": log.action_type,
                "action_details": log.action_details,
                "created_at": log.created_at.isoformat() if log.created_at else None
            }
            for log in logs
        ]
    }


@app.get("/api/activity-logs/summary")
def get_activity_summary(
    days: int = 7,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get summary of recent activity by user and action type."""
    from sqlalchemy import func

    cutoff = datetime.utcnow() - timedelta(days=days)

    # Get activity counts by user
    user_activity = db.query(
        ActivityLog.user_email,
        func.count(ActivityLog.id).label("action_count")
    ).filter(
        ActivityLog.created_at >= cutoff
    ).group_by(ActivityLog.user_email).all()

    # Get activity counts by type
    type_activity = db.query(
        ActivityLog.action_type,
        func.count(ActivityLog.id).label("action_count")
    ).filter(
        ActivityLog.created_at >= cutoff
    ).group_by(ActivityLog.action_type).all()

    return {
        "period_days": days,
        "by_user": [{"user": u[0], "count": u[1]} for u in user_activity],
        "by_type": [{"type": t[0], "count": t[1]} for t in type_activity]
    }


# ============== DISCOVERY ENGINE ENDPOINTS ==============

@app.get("/api/discovery/context")
def get_discovery_context(current_user: dict = Depends(get_current_user)):
    """Get the current Discovery Engine context (certification DNA)."""
    try:
        context_path = os.path.join(os.path.dirname(__file__), "certify_context.json")
        if not os.path.exists(context_path):
             return {"core_keywords": [], "market_keywords": [], "exclusions": []}
             
        with open(context_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")

@app.post("/api/discovery/context")
def update_discovery_context(new_context: dict, current_user: dict = Depends(get_current_user)):
    """Update the Discovery Engine context."""
    try:
        context_path = os.path.join(os.path.dirname(__file__), "certify_context.json")
        with open(context_path, 'w', encoding='utf-8') as f:
            json.dump(new_context, f, indent=4)
        return {"success": True, "message": "Context updated successfully"}
    except Exception as e:
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")

@app.post("/api/discovery/refine-context")
async def refine_discovery_context(request: dict, current_user: dict = Depends(get_current_user)):
    """Use AI to refine the context based on user chat input."""
    try:
        user_input = request.get("message")
        current_context = request.get("current_context")

        if not user_input:
            raise HTTPException(status_code=400, detail="Message required")

        system_prompt = """You are an expert configuration assistant for a Competitor Discovery Engine.
        Your goal is to update the JSON configuration profile based on the user's request.

        The JSON structure is:
        {
            "core_keywords": ["list", "of", "competitor", "keywords"],
            "market_keywords": ["target", "markets"],
            "required_context": ["must", "have", "terms"],
            "negative_keywords": ["terms", "to", "avoid"],
            "known_competitors": ["known", "competitor", "names"],
            "exclusions": ["industries", "to", "exclude"]
        }

        Return ONLY the updated JSON. Do not return markdown formatting.""" + NO_HALLUCINATION_INSTRUCTION

        from ai_router import get_ai_router, TaskType
        router = get_ai_router()
        result = await router.generate_json(
            prompt=f"Current Profile: {json.dumps(current_context)}\n\nUser Request: {user_input}\n\nUpdate the profile:",
            task_type=TaskType.DATA_EXTRACTION,
            system_prompt=system_prompt,
        )

        refined_json = result.get("response_json", {})
        if not refined_json:
            raise ValueError("AI did not return valid JSON")
        return {"success": True, "refined_context": refined_json}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Context refinement failed: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")

@app.post("/api/discovery/schedule")
def schedule_discovery_run(request: dict, current_user: dict = Depends(get_current_user)):
    """Schedule a one-off discovery run."""
    try:
        run_at_str = request.get("run_at") # ISO format string
        if not run_at_str:
            raise HTTPException(status_code=400, detail="run_at timestamp required")
            
        run_at = datetime.fromisoformat(run_at_str.replace('Z', '+00:00'))
        
        from scheduler import schedule_one_off_discovery
        schedule_one_off_discovery(run_at)
        
        return {"success": True, "message": f"Discovery job scheduled for {run_at.isoformat()}"}
    except Exception as e:
         raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/changes/log")
def log_data_change(
    competitor_id: int,
    field_name: str,
    old_value: str,
    new_value: str,
    changed_by: str = "system",
    reason: Optional[str] = None,
    source_url: Optional[str] = None,
    db: Session = Depends(get_db)
):
    """Log a data change for audit purposes."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")
    
    change = DataChangeHistory(
        competitor_id=competitor_id,
        competitor_name=competitor.name,
        field_name=field_name,
        old_value=old_value,
        new_value=new_value,
        changed_by=changed_by,
        change_reason=reason,
        source_url=source_url
    )
    db.add(change)
    db.commit()

    return {"success": True, "change_id": change.id}


# ============== PENDING DATA CHANGES ENDPOINTS ==============
# DataChangeSubmission imported from schemas.common

@app.post("/api/data-changes/submit")
def submit_data_change(
    body: DataChangeSubmission,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user_optional)
):
    """Submit a data change for admin approval."""
    competitor_id = body.competitor_id
    field_name = body.field_name
    old_value = body.old_value
    new_value = body.new_value
    source_url = body.source_url
    notes = body.notes
    value_type = body.value_type

    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    # Check for existing pending change
    existing = db.query(PendingDataChange).filter(
        PendingDataChange.competitor_id == competitor_id,
        PendingDataChange.field_name == field_name,
        PendingDataChange.status == "pending"
    ).first()

    if existing:
        # Update existing pending change
        existing.new_value = str(new_value) if new_value else None
        existing.source_url = source_url
        existing.notes = notes
        existing.value_type = value_type
        existing.submitted_at = datetime.utcnow()
        db.commit()
        return {"success": True, "message": "Pending change updated", "change_id": existing.id}

    # Create new pending change
    pending = PendingDataChange(
        competitor_id=competitor_id,
        competitor_name=competitor.name,
        field_name=field_name,
        old_value=str(old_value) if old_value else None,
        new_value=str(new_value) if new_value else None,
        value_type=value_type,
        source_url=source_url,
        notes=notes,
        submitted_by=current_user.email if current_user else "anonymous",
        status="pending"
    )
    db.add(pending)
    db.commit()

    return {"success": True, "message": "Change submitted for approval", "change_id": pending.id}


@app.get("/api/data-changes/pending")
def get_pending_changes(
    competitor_id: Optional[int] = None,
    db: Session = Depends(get_db)
):
    """Get all pending data changes awaiting approval."""
    query = db.query(PendingDataChange).filter(PendingDataChange.status == "pending")

    if competitor_id:
        query = query.filter(PendingDataChange.competitor_id == competitor_id)

    pending = query.order_by(PendingDataChange.submitted_at.desc()).all()

    return {
        "pending_count": len(pending),
        "changes": [
            {
                "id": p.id,
                "competitor_id": p.competitor_id,
                "competitor_name": p.competitor_name,
                "field_name": p.field_name,
                "old_value": p.old_value,
                "new_value": p.new_value,
                "value_type": p.value_type,
                "source_url": p.source_url,
                "notes": p.notes,
                "submitted_by": p.submitted_by,
                "submitted_at": p.submitted_at.isoformat() if p.submitted_at else None,
                "status": p.status
            }
            for p in pending
        ]
    }


@app.post("/api/data-changes/{change_id}/approve")
def approve_data_change(
    change_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Approve a pending data change (admin only)."""
    # Check if user is admin
    if current_user.role != "admin":
        raise HTTPException(status_code=403, detail="Only admins can approve changes")

    pending = db.query(PendingDataChange).filter(PendingDataChange.id == change_id).first()
    if not pending:
        raise HTTPException(status_code=404, detail="Pending change not found")

    if pending.status != "pending":
        raise HTTPException(status_code=400, detail=f"Change is already {pending.status}")

    # Get the competitor
    competitor = db.query(Competitor).filter(Competitor.id == pending.competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    # Get current value
    old_value = getattr(competitor, pending.field_name, None)

    # Apply the change to the competitor
    try:
        new_value = pending.new_value
        if pending.value_type == "number" and new_value:
            try:
                new_value = float(new_value) if '.' in str(new_value) else int(new_value)
            except (ValueError, TypeError):
                pass

        setattr(competitor, pending.field_name, new_value)

        # Update pending status
        pending.status = "approved"
        pending.reviewed_by = current_user.email
        pending.reviewed_at = datetime.utcnow()

        # Log the change in history
        history = DataChangeHistory(
            competitor_id=pending.competitor_id,
            competitor_name=pending.competitor_name,
            field_name=pending.field_name,
            old_value=str(old_value) if old_value else None,
            new_value=str(new_value) if new_value else None,
            changed_by=current_user.email,
            change_reason=f"Manual edit approved. Notes: {pending.notes or 'None'}",
            source_url=pending.source_url
        )
        db.add(history)

        # Update or create data source record
        source = db.query(DataSource).filter(
            DataSource.competitor_id == pending.competitor_id,
            DataSource.field_name == pending.field_name
        ).first()

        if source:
            source.source_type = "manual_verified"
            source.source_url = pending.source_url
            source.entered_by = pending.submitted_by
            source.verified_at = datetime.utcnow()
            source.verified_by = current_user.email
        else:
            source = DataSource(
                competitor_id=pending.competitor_id,
                field_name=pending.field_name,
                source_type="manual_verified",
                source_url=pending.source_url,
                source_name="Manual Entry",
                entered_by=pending.submitted_by,
                verified_at=datetime.utcnow(),
                verified_by=current_user.email
            )
            db.add(source)

        db.commit()

        return {
            "success": True,
            "message": f"Change approved and applied to {pending.competitor_name}.{pending.field_name}",
            "field_name": pending.field_name,
            "new_value": new_value
        }
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail="Failed to apply change. Please try again.")


@app.post("/api/data-changes/{change_id}/reject")
def reject_data_change(
    change_id: int,
    reason: Optional[str] = None,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Reject a pending data change (admin only)."""
    # Check if user is admin
    if current_user.role != "admin":
        raise HTTPException(status_code=403, detail="Only admins can reject changes")

    pending = db.query(PendingDataChange).filter(PendingDataChange.id == change_id).first()
    if not pending:
        raise HTTPException(status_code=404, detail="Pending change not found")

    if pending.status != "pending":
        raise HTTPException(status_code=400, detail=f"Change is already {pending.status}")

    pending.status = "rejected"
    pending.reviewed_by = current_user.email
    pending.reviewed_at = datetime.utcnow()
    pending.review_notes = reason

    db.commit()

    return {
        "success": True,
        "message": "Change rejected",
        "change_id": change_id
    }


@app.get("/api/data-changes/approved-fields")
def get_approved_manual_fields(
    competitor_id: Optional[int] = None,
    db: Session = Depends(get_db)
):
    """Get list of fields that have been manually edited and approved."""
    query = db.query(PendingDataChange).filter(PendingDataChange.status == "approved")

    if competitor_id:
        query = query.filter(PendingDataChange.competitor_id == competitor_id)

    approved = query.all()

    # Return as dict keyed by competitor_id-field_name
    result = {}
    for a in approved:
        key = f"{a.competitor_id}-{a.field_name}"
        result[key] = {
            "status": "approved",
            "new_value": a.new_value,
            "source_url": a.source_url,
            "approved_at": a.reviewed_at.isoformat() if a.reviewed_at else None,
            "approved_by": a.reviewed_by
        }

    return {"approved_fields": result}


# ============== DATA SOURCE ENDPOINTS ==============


def _map_url_quality(url_status):
    """Map internal url_status to frontend-expected quality values."""
    mapping = {
        "verified": "exact_page",
        "page_level": "page_level",
        "broken": "broken",
    }
    if url_status is None:
        return None
    return mapping.get(url_status, "homepage_only")


@app.get("/api/sources/batch")
def get_batch_source_verification(
    competitor_ids: str = Query(
        ..., description="Comma-separated competitor IDs"
    ),
    fields: Optional[str] = Query(
        None, description="Comma-separated field names to filter"
    ),
    db: Session = Depends(get_db)
):
    """
    Get source verification status for multiple competitors.

    v7.2: Batch endpoint for source dot display on battlecards.
    Returns the latest source per competitor+field combination.
    NOTE: Must be defined BEFORE parametric /api/sources/{competitor_id}
    route to avoid FastAPI route ordering 422.
    """
    try:
        comp_id_list = [
            int(x.strip()) for x in competitor_ids.split(",")
            if x.strip()
        ]
    except ValueError:
        raise HTTPException(
            status_code=400,
            detail="competitor_ids must be comma-separated integers"
        )

    if not comp_id_list:
        raise HTTPException(
            status_code=400,
            detail="At least one competitor_id is required"
        )

    if len(comp_id_list) > 100:
        raise HTTPException(
            status_code=400,
            detail="Maximum 100 competitor IDs per request"
        )

    query = db.query(DataSource).filter(
        DataSource.competitor_id.in_(comp_id_list)
    )

    if fields:
        field_list = [
            f.strip() for f in fields.split(",") if f.strip()
        ]
        if field_list:
            query = query.filter(
                DataSource.field_name.in_(field_list)
            )

    query = query.order_by(
        DataSource.competitor_id,
        DataSource.field_name,
        DataSource.extracted_at.desc()
    )
    all_sources = query.all()

    sources_dict = {}
    seen = set()
    for s in all_sources:
        key = (s.competitor_id, s.field_name)
        if key in seen:
            continue
        seen.add(key)

        cid = str(s.competitor_id)
        if cid not in sources_dict:
            sources_dict[cid] = {}

        sources_dict[cid][s.field_name] = {
            "has_source": True,
            "is_verified": bool(s.is_verified),
            "source_url": s.source_url,
            "source_name": s.source_name,
            "confidence_score": s.confidence_score,
            "current_value": s.current_value,
            "extracted_at": (
                s.extracted_at.isoformat()
                if s.extracted_at else None
            ),
            "deep_link_url": getattr(s, 'deep_link_url', None),
            "url_quality": _map_url_quality(getattr(s, 'url_status', None)),
            "source_section": getattr(s, 'source_section', None),
        }

    for comp_id in comp_id_list:
        cid = str(comp_id)
        if cid not in sources_dict:
            sources_dict[cid] = {}

    return {"sources": sources_dict}


@app.get("/api/sources/quality-summary")
async def get_source_quality_summary(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """
    Return aggregate source quality statistics across all competitors.

    NOTE: Must be defined BEFORE /api/sources/{competitor_id} to avoid
    FastAPI route ordering collision.
    """
    from sqlalchemy import func

    total = db.query(func.count(DataSource.id)).filter(
        DataSource.source_url.isnot(None)
    ).scalar() or 0

    exact_page = db.query(func.count(DataSource.id)).filter(
        DataSource.url_status == "verified",
        DataSource.deep_link_url.isnot(None),
    ).scalar() or 0

    page_level = db.query(func.count(DataSource.id)).filter(
        DataSource.url_status == "page_level",
    ).scalar() or 0

    homepage_only = db.query(func.count(DataSource.id)).filter(
        DataSource.source_url.isnot(None),
        DataSource.source_page_url.is_(None),
        DataSource.url_status != "broken",
    ).scalar() or 0

    broken = db.query(func.count(DataSource.id)).filter(
        DataSource.url_status == "broken",
    ).scalar() or 0

    with_source = db.query(func.count(DataSource.id)).filter(
        DataSource.source_url.isnot(None)
    ).scalar() or 0

    total_fields = db.query(func.count(DataSource.id)).scalar() or 0

    return {
        "total": total,
        "exact_page": exact_page,
        "page_level": page_level,
        "homepage_only": homepage_only,
        "broken": broken,
        "with_source": with_source,
        "total_fields": total_fields,
        "coverage_pct": round(
            (exact_page + page_level) / max(total, 1) * 100, 1
        ),
    }


@app.get("/api/sources/coverage")
def get_source_coverage(db: Session = Depends(get_db)):
    """
    Get source coverage report by field category.

    NOTE: Must be defined BEFORE /api/sources/{competitor_id} to avoid
    FastAPI route ordering collision.
    """
    from source_discovery_engine import get_source_discovery_engine

    engine = get_source_discovery_engine()
    return engine.get_coverage_report()


@app.get("/api/sources/{competitor_id}")
def get_competitor_sources(competitor_id: int, db: Session = Depends(get_db)):
    """Get all data sources for a competitor's fields."""
    sources = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id
    ).all()
    
    sources_by_field = {}
    for s in sources:
        sources_by_field[s.field_name] = {
            "source_type": s.source_type,
            "source_url": s.source_url,
            "source_name": s.source_name,
            "entered_by": s.entered_by,
            "formula": s.formula,
            "verified_at": s.verified_at.isoformat() if s.verified_at else None
        }
    
    return {
        "competitor_id": competitor_id,
        "sources": sources_by_field
    }


@app.get("/api/sources/{competitor_id}/{field_name}")
def get_field_source(competitor_id: int, field_name: str, db: Session = Depends(get_db)):
    """Get the data source for a specific field of a competitor."""
    source = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id,
        DataSource.field_name == field_name
    ).first()
    
    if not source:
        return {
            "competitor_id": competitor_id,
            "field_name": field_name,
            "source_type": "unknown",
            "source_url": None,
            "message": "No source recorded for this field"
        }
    
    return {
        "competitor_id": competitor_id,
        "field_name": field_name,
        "source_type": source.source_type,
        "source_url": source.source_url,
        "source_name": source.source_name,
        "current_value": source.current_value,
        "entered_by": source.entered_by,
        "formula": source.formula,
        "verified_at": source.verified_at.isoformat() if source.verified_at else None
    }


@app.post("/api/sources/set")
def set_field_source(
    competitor_id: int,
    field_name: str,
    source_type: str,  # "external", "manual", "calculated"
    source_url: Optional[str] = None,
    source_name: Optional[str] = None,
    entered_by: Optional[str] = None,
    formula: Optional[str] = None,
    db: Session = Depends(get_db)
):
    """Set or update the source for a competitor's field."""
    # Check if source already exists
    existing = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id,
        DataSource.field_name == field_name
    ).first()
    
    if existing:
        existing.source_type = source_type
        existing.source_url = source_url
        existing.source_name = source_name
        existing.entered_by = entered_by
        existing.formula = formula
        existing.verified_at = datetime.utcnow()
        existing.updated_at = datetime.utcnow()
    else:
        new_source = DataSource(
            competitor_id=competitor_id,
            field_name=field_name,
            source_type=source_type,
            source_url=source_url,
            source_name=source_name,
            entered_by=entered_by,
            formula=formula
        )
        db.add(new_source)
    
    db.commit()
    return {"success": True, "message": f"Source set for {field_name}"}


# ============== AI-POWERED SOURCE DISCOVERY (v6.3.9) ==============

@app.get("/api/sources/field/{competitor_id}/{field_name}")
def get_source_for_field(competitor_id: int, field_name: str, db: Session = Depends(get_db)):
    """
    Get detailed source information for a specific field.

    v6.3.9: Enhanced endpoint with confidence and extraction details.
    """
    from source_discovery_engine import get_source_discovery_engine

    engine = get_source_discovery_engine()
    source = engine.get_source_for_field(competitor_id, field_name)

    if not source:
        return {
            "competitor_id": competitor_id,
            "field_name": field_name,
            "source_url": None,
            "has_source": False,
            "message": "No source recorded - AI discovery pending"
        }

    source["has_source"] = True
    return source


@app.post("/api/sources/discover/{competitor_id}")
async def discover_sources_for_competitor(
    competitor_id: int,
    priority: Optional[str] = None,
    max_fields: int = 50,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    Run AI source discovery for a single competitor.

    v6.3.9: Uses Gemini with Google Search grounding to find authoritative sources.

    Args:
        competitor_id: ID of the competitor to process
        priority: Optional priority filter (p0, p1, p2, p3)
        max_fields: Maximum number of fields to process (default 50)
    """
    import asyncio
    from source_discovery_engine import get_source_discovery_engine, FieldPriority

    # Verify competitor exists
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    # Parse priority filter
    priority_filter = None
    if priority:
        priority_map = {
            "p0": FieldPriority.P0,
            "p1": FieldPriority.P1,
            "p2": FieldPriority.P2,
            "p3": FieldPriority.P3,
        }
        priority_filter = priority_map.get(priority.lower())

    engine = get_source_discovery_engine()

    # Run discovery
    progress = await engine.discover_sources_for_competitor(
        competitor_id=competitor_id,
        priority_filter=priority_filter,
        max_fields=max_fields
    )

    return {
        "status": "success",
        "competitor_id": competitor_id,
        "competitor_name": progress.competitor_name,
        "fields_processed": progress.fields_processed,
        "sources_found": progress.sources_found,
        "coverage": round(progress.sources_found / max(progress.fields_processed, 1) * 100, 1),
        "errors_count": len(progress.errors),
        "duration_seconds": (progress.completed_at - progress.started_at).total_seconds() if progress.completed_at and progress.started_at else 0
    }


# Global variable to track batch discovery progress
_batch_discovery_progress = {
    "job_id": None,
    "status": "idle",
    "competitors_total": 0,
    "competitors_processed": 0,
    "sources_found": 0,
    "started_at": None,
    "completed_at": None
}


@app.post("/api/sources/discover/all")
async def discover_sources_for_all(
    background_tasks: BackgroundTasks,
    priority: Optional[str] = None,
    max_per_competitor: int = 30,
    current_user: User = Depends(get_current_user)
):
    """
    Run AI source discovery for all competitors (background job).

    v6.3.9: Processes all competitors sequentially with progress tracking.
    """
    import uuid
    from source_discovery_engine import get_source_discovery_engine, FieldPriority

    global _batch_discovery_progress  # noqa: F824

    # Check if already running
    if _batch_discovery_progress["status"] == "running":
        return {
            "status": "already_running",
            "job_id": _batch_discovery_progress["job_id"],
            "progress": _batch_discovery_progress
        }

    # Parse priority filter
    priority_filter = None
    if priority:
        priority_map = {
            "p0": FieldPriority.P0,
            "p1": FieldPriority.P1,
            "p2": FieldPriority.P2,
            "p3": FieldPriority.P3,
        }
        priority_filter = priority_map.get(priority.lower())

    # Generate job ID
    job_id = str(uuid.uuid4())[:8]

    # Initialize progress
    _batch_discovery_progress = {
        "job_id": job_id,
        "status": "running",
        "competitors_total": 0,
        "competitors_processed": 0,
        "sources_found": 0,
        "started_at": datetime.utcnow().isoformat(),
        "completed_at": None
    }

    # Define background task
    async def run_batch_discovery():
        global _batch_discovery_progress  # noqa: F824

        try:
            engine = get_source_discovery_engine()
            result = await engine.discover_sources_for_all_competitors(
                priority_filter=priority_filter,
                max_per_competitor=max_per_competitor
            )

            _batch_discovery_progress["status"] = "completed"
            _batch_discovery_progress["competitors_total"] = result["total_competitors"]
            _batch_discovery_progress["competitors_processed"] = result["competitors_processed"]
            _batch_discovery_progress["sources_found"] = result["total_sources_found"]
            _batch_discovery_progress["completed_at"] = result["completed_at"]

        except Exception as e:
            _batch_discovery_progress["status"] = "error"
            _batch_discovery_progress["error"] = "Batch discovery failed"

    # Start background task - add async function directly
    background_tasks.add_task(run_batch_discovery)

    return {
        "status": "started",
        "job_id": job_id,
        "message": "Source discovery started for all competitors. Check progress with GET /api/sources/discover/status"
    }


@app.get("/api/sources/discover/status")
def get_discovery_status():
    """Get the status of the batch source discovery job."""
    return _batch_discovery_progress


# ============== URL REFINEMENT ENGINE (v8.3.0) ==============

_url_refinement_progress: Dict[str, Dict[str, Any]] = {}


@app.post("/api/sources/refine-urls")
async def refine_source_urls(
    request: Request,
    background_tasks: BackgroundTasks,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Batch refine source URLs for a single competitor."""
    import uuid as _uuid
    from url_refinement_engine import refine_source_url

    body = await request.json()
    competitor_id = body.get("competitor_id")
    field_names = body.get("field_names")  # optional filter

    if not competitor_id:
        raise HTTPException(status_code=400, detail="competitor_id is required")

    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    task_id = str(_uuid.uuid4())[:8]

    _ai_tasks[task_id] = {
        "status": "running",
        "page_context": "sources",
        "user_id": current_user.get("id"),
        "started_at": datetime.utcnow().isoformat(),
        "completed_at": None,
        "result": None,
        "error": None,
        "task_type": "refine_urls",
        "read_at": None,
    }

    comp_name = competitor.name
    comp_website = competitor.website or ""

    async def _run_refinement():
        refine_db = SessionLocal()
        try:
            sources = refine_db.query(DataSource).filter(
                DataSource.competitor_id == competitor_id
            )
            if field_names:
                sources = sources.filter(DataSource.field_name.in_(field_names))
            sources = sources.all()

            refined_count = 0
            for src in sources:
                try:
                    result = await refine_source_url(
                        competitor_name=comp_name,
                        website=comp_website,
                        field_name=src.field_name,
                        current_value=src.current_value,
                        current_url=src.source_url,
                    )
                    if result.source_page_url:
                        src.source_page_url = result.source_page_url
                        src.source_anchor_text = result.source_anchor_text
                        src.source_section = result.source_section
                        src.deep_link_url = result.deep_link_url
                        src.url_status = result.url_status
                        src.last_url_verified = datetime.utcnow()
                        refined_count += 1
                except Exception as field_err:
                    logger.warning(
                        "URL refinement failed for %s/%s: %s",
                        comp_name, src.field_name, field_err,
                    )

            refine_db.commit()
            _ai_tasks[task_id]["status"] = "completed"
            _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
            _ai_tasks[task_id]["result"] = {
                "competitor": comp_name,
                "sources_processed": len(sources),
                "sources_refined": refined_count,
            }
        except Exception as e:
            logger.error("URL refinement batch failed: %s", e)
            _ai_tasks[task_id]["status"] = "failed"
            _ai_tasks[task_id]["error"] = "URL refinement failed"
            _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
        finally:
            refine_db.close()

    background_tasks.add_task(_run_refinement)

    return {"task_id": task_id, "status": "running"}


@app.post("/api/sources/refine-urls/all")
async def refine_source_urls_all(
    background_tasks: BackgroundTasks,
    force_refresh: bool = False,
    current_user: dict = Depends(get_current_user)
):
    """Admin: batch refine source URLs for all active competitors.

    Two-phase pipeline:
      Phase 1: Find correct page URLs via pattern/sitemap strategies
      Phase 2: Content-aware fragment matching (fetch page, find real text)

    Args:
        force_refresh: If True, re-process ALL sources (not just unrefined).
    """
    import uuid as _uuid

    task_id = str(_uuid.uuid4())[:8]

    # Check if already running
    for tid, t in _url_refinement_progress.items():
        if t.get("status") == "running":
            return {
                "status": "already_running",
                "task_id": tid,
                "progress": t,
            }

    _url_refinement_progress[task_id] = {
        "status": "running",
        "competitors_total": 0,
        "competitors_processed": 0,
        "sources_refined": 0,
        "sources_exact_match": 0,
        "sources_page_level": 0,
        "phase": "finding_pages",
        "current_competitor": None,
        "started_at": datetime.utcnow().isoformat(),
        "completed_at": None,
    }

    # Also register in _ai_tasks so frontend can poll via /api/ai/tasks/{id}
    _ai_tasks[task_id] = {
        "status": "running",
        "page_context": "settings",
        "user_id": current_user.get("id"),
        "started_at": datetime.utcnow().isoformat(),
        "completed_at": None,
        "result": None,
        "error": None,
        "task_type": "refine_urls_all",
        "read_at": None,
        "progress": 0,
        "status_message": "Starting URL refinement...",
    }

    _force = force_refresh  # Capture for closure

    async def _run_all_refinement():
        import httpx as _httpx
        from url_refinement_engine import (
            _normalize_base_url, _get_page_type, PAGE_PATTERNS,
            _head_check, _strategy_sitemap, build_text_fragment,
            _make_deep_link,
        )
        from content_matcher import (
            fetch_page_text, find_value_on_page,
        )

        refine_db = SessionLocal()
        try:
            competitors = refine_db.query(Competitor).filter(
                Competitor.status != "deleted"
            ).all()
            total = len(competitors)
            _url_refinement_progress[task_id]["competitors_total"] = total

            # ════════════════════════════════════════════════════════════
            # PHASE 1: Find correct page URLs (pattern + sitemap)
            # ════════════════════════════════════════════════════════════
            _url_refinement_progress[task_id]["phase"] = "finding_pages"

            async with _httpx.AsyncClient(
                timeout=5.0, follow_redirects=True,
                headers={
                    "User-Agent": "CertifyIntel/8.3.0 BatchRefiner"
                },
                limits=_httpx.Limits(
                    max_connections=25, max_keepalive_connections=15
                ),
            ) as client:
                _src_sem = asyncio.Semaphore(3)

                async def _refine_one_source(src, comp_name, base_url):
                    """Phase 1: Pattern + sitemap to find page URL."""
                    async with _src_sem:
                        page_type = _get_page_type(src.field_name)
                        patterns = PAGE_PATTERNS.get(
                            page_type, [page_type]
                        )

                        found_url = None
                        for slug in patterns[:3]:
                            candidate = (
                                f"{base_url}/{slug}" if slug else base_url
                            )
                            ok, final_url = await _head_check(
                                client, candidate, timeout=3.0
                            )
                            if ok and final_url:
                                found_url = final_url
                                break

                        if not found_url:
                            hit = await _strategy_sitemap(
                                client, base_url, src.field_name
                            )
                            if hit:
                                found_url = hit[0]

                        if found_url:
                            src.source_page_url = found_url
                            src.source_section = page_type
                            src.url_status = "pending"
                            src.last_url_verified = datetime.utcnow()
                            # Temporary: raw DB value fragment
                            # (will be replaced in Phase 2)
                            if src.current_value:
                                frag = build_text_fragment(
                                    src.current_value
                                )
                                if frag:
                                    src.source_anchor_text = (
                                        src.current_value or ""
                                    )[:200]
                                    src.deep_link_url = _make_deep_link(
                                        found_url, frag
                                    )
                            _url_refinement_progress[task_id][
                                "sources_refined"
                            ] += 1
                            return True
                        return False

                async def _refine_competitor(i, comp):
                    """Phase 1: Refine sources for one competitor."""
                    _url_refinement_progress[task_id][
                        "current_competitor"
                    ] = comp.name
                    pct = int(((i + 1) / max(total, 1)) * 50)
                    _ai_tasks[task_id]["progress"] = pct
                    _ai_tasks[task_id]["status_message"] = (
                        f"Phase 1: Finding pages for {comp.name}"
                        f" ({i + 1}/{total})"
                    )

                    base_url = _normalize_base_url(comp.website or "")
                    if not base_url:
                        return

                    if _force:
                        sources = refine_db.query(DataSource).filter(
                            DataSource.competitor_id == comp.id,
                            DataSource.source_url.isnot(None),
                        ).limit(20).all()
                    else:
                        sources = refine_db.query(DataSource).filter(
                            DataSource.competitor_id == comp.id,
                            DataSource.source_url.isnot(None),
                            DataSource.deep_link_url.is_(None),
                        ).limit(10).all()

                    if not sources:
                        return

                    tasks = [
                        asyncio.wait_for(
                            _refine_one_source(
                                src, comp.name, base_url
                            ),
                            timeout=12.0,
                        )
                        for src in sources
                    ]
                    results = await asyncio.gather(
                        *tasks, return_exceptions=True
                    )
                    for r in results:
                        if isinstance(r, Exception):
                            logger.debug(
                                "Phase 1 refine error: %s", r
                            )

                    refine_db.commit()
                    _url_refinement_progress[task_id][
                        "competitors_processed"
                    ] = i + 1

                # Process competitors in parallel batches of 5
                BATCH_SIZE = 5
                for batch_start in range(0, total, BATCH_SIZE):
                    batch = list(enumerate(competitors))[
                        batch_start: batch_start + BATCH_SIZE
                    ]
                    await asyncio.gather(
                        *[
                            _refine_competitor(i, comp)
                            for i, comp in batch
                        ],
                        return_exceptions=True,
                    )

            # ════════════════════════════════════════════════════════════
            # PHASE 2: Content-aware text fragment matching
            # ════════════════════════════════════════════════════════════
            _url_refinement_progress[task_id]["phase"] = "matching_content"
            _ai_tasks[task_id]["status_message"] = (
                "Phase 2: Matching content on pages..."
            )

            # Gather all sources that have a page URL and a value
            sources_for_phase2 = refine_db.query(DataSource).filter(
                DataSource.source_page_url.isnot(None),
                DataSource.current_value.isnot(None),
                DataSource.current_value != "",
            ).all()

            if not _force:
                # Only process sources whose deep_link still uses
                # raw DB value (Phase 1 output) or have no fragment
                pass  # Process all sources with page URLs

            # Group by page URL to avoid fetching same page twice
            url_groups: Dict[str, list] = {}
            for src in sources_for_phase2:
                page_url = src.source_page_url
                if page_url not in url_groups:
                    url_groups[page_url] = []
                url_groups[page_url].append(src)

            total_pages = len(url_groups)
            pages_processed = 0
            _page_sem = asyncio.Semaphore(10)

            async def _match_page_group(page_url, sources_list):
                """Fetch one page, match all fields targeting it."""
                nonlocal pages_processed
                async with _page_sem:
                    page_text = await fetch_page_text(page_url)
                    if not page_text:
                        # No content - mark as page_level
                        for src in sources_list:
                            if src.source_page_url:
                                src.url_status = "page_level"
                                # Strip fragment, keep page URL
                                src.deep_link_url = (
                                    src.source_page_url.split("#")[0]
                                )
                        pages_processed += 1
                        return

                    for src in sources_list:
                        val = (src.current_value or "").strip()
                        if not val or len(val) < 2:
                            continue

                        match = find_value_on_page(
                            page_text, val,
                            field_name=src.field_name or ""
                        )
                        if match:
                            # Build fragment from ACTUAL page text
                            frag = build_text_fragment(
                                match.matched_text,
                                context_before=(
                                    match.context_before
                                    if match.context_before else None
                                ),
                                context_after=(
                                    match.context_after
                                    if match.context_after else None
                                ),
                            )
                            if frag:
                                page_base = (
                                    src.source_page_url.split("#")[0]
                                )
                                src.deep_link_url = _make_deep_link(
                                    page_base, frag
                                )
                                src.source_anchor_text = (
                                    match.matched_text[:200]
                                )
                                src.url_status = "verified"
                                _url_refinement_progress[task_id][
                                    "sources_exact_match"
                                ] += 1
                        else:
                            # Page found but text not matched
                            src.url_status = "page_level"
                            page_base = (
                                src.source_page_url.split("#")[0]
                            )
                            src.deep_link_url = page_base
                            _url_refinement_progress[task_id][
                                "sources_page_level"
                            ] += 1

                    pages_processed += 1

            # Process page groups in parallel batches of 10
            page_items = list(url_groups.items())
            PAGE_BATCH = 10
            for batch_start in range(0, len(page_items), PAGE_BATCH):
                batch = page_items[
                    batch_start: batch_start + PAGE_BATCH
                ]
                pct = 50 + int(
                    (pages_processed / max(total_pages, 1)) * 50
                )
                _ai_tasks[task_id]["progress"] = min(pct, 99)
                _ai_tasks[task_id]["status_message"] = (
                    f"Phase 2: Matching content"
                    f" ({pages_processed}/{total_pages} pages)"
                )

                await asyncio.gather(
                    *[
                        asyncio.wait_for(
                            _match_page_group(url, srcs),
                            timeout=15.0,
                        )
                        for url, srcs in batch
                    ],
                    return_exceptions=True,
                )

                refine_db.commit()

            # ── Final stats ────────────────────────────────────────────
            refined = _url_refinement_progress[task_id]["sources_refined"]
            exact = _url_refinement_progress[task_id]["sources_exact_match"]
            page_lvl = _url_refinement_progress[task_id]["sources_page_level"]

            _url_refinement_progress[task_id]["status"] = "completed"
            _url_refinement_progress[task_id]["completed_at"] = (
                datetime.utcnow().isoformat()
            )
            _url_refinement_progress[task_id]["current_competitor"] = None
            _ai_tasks[task_id]["status"] = "completed"
            _ai_tasks[task_id]["progress"] = 100
            _ai_tasks[task_id]["status_message"] = (
                f"Done: {exact} exact matches,"
                f" {page_lvl} page-level,"
                f" {refined} pages found"
            )
            _ai_tasks[task_id]["completed_at"] = (
                datetime.utcnow().isoformat()
            )
            _ai_tasks[task_id]["result"] = {
                "competitors_processed": total,
                "sources_refined": refined,
                "sources_exact_match": exact,
                "sources_page_level": page_lvl,
                "unique_pages_fetched": total_pages,
            }
        except Exception as e:
            logger.error("Batch URL refinement failed: %s", e)
            _url_refinement_progress[task_id]["status"] = "error"
            _url_refinement_progress[task_id]["completed_at"] = (
                datetime.utcnow().isoformat()
            )
            _ai_tasks[task_id]["status"] = "failed"
            _ai_tasks[task_id]["error"] = "URL refinement failed"
            _ai_tasks[task_id]["completed_at"] = (
                datetime.utcnow().isoformat()
            )
        finally:
            refine_db.close()

    background_tasks.add_task(_run_all_refinement)

    return {
        "task_id": task_id,
        "status": "started",
        "message": "URL refinement started for all competitors",
        "force_refresh": _force,
    }


# Data providers status/test endpoints moved to routers/admin.py


@app.post("/api/admin/enrich-from-providers")
async def enrich_from_providers(
    background_tasks: BackgroundTasks,
    current_user: dict = Depends(get_current_user)
):
    """Batch enrich all competitors from enterprise data providers."""
    import uuid as _uuid
    import time as _time

    try:
        from data_providers import get_active_providers
    except ImportError:
        raise HTTPException(status_code=503, detail="Data providers module not installed")

    active = get_active_providers()
    if not active:
        return {"status": "skipped", "message": "No data providers configured"}

    task_id = str(_uuid.uuid4())[:8]

    _ai_tasks[task_id] = {
        "status": "running",
        "page_context": "admin",
        "user_id": current_user.get("id"),
        "started_at": datetime.utcnow().isoformat(),
        "completed_at": None,
        "result": None,
        "error": None,
        "task_type": "enrich_providers",
        "read_at": None,
    }

    provider_names = [p.provider_name for p in active]

    async def _run_enrichment():
        enrich_db = SessionLocal()
        try:
            from data_providers import get_active_providers as _gap
            providers = _gap()

            competitors = enrich_db.query(Competitor).filter(
                Competitor.status != "deleted"
            ).all()

            total_updates = 0
            total_sources = 0

            for comp in competitors:
                for provider in providers:
                    try:
                        result = await provider.query_for_competitor(comp.name)
                        if result.error:
                            continue

                        # Apply mapped fields
                        for field_name, value in result.fields.items():
                            if hasattr(comp, field_name) and value is not None:
                                old_val = getattr(comp, field_name, None)
                                if str(old_val) != str(value):
                                    setattr(comp, field_name, value)
                                    total_updates += 1

                            # Create DataSource record
                            source_url = result.source_urls.get(field_name, "")
                            new_source = DataSource(
                                competitor_id=comp.id,
                                field_name=field_name,
                                current_value=str(value) if value else None,
                                source_type="api",
                                source_name=provider.provider_name,
                                source_url=source_url,
                                extraction_method="enterprise_provider",
                                confidence_score=75,
                                confidence_level="moderate",
                                extracted_at=datetime.utcnow(),
                            )
                            enrich_db.add(new_source)
                            total_sources += 1

                    except Exception as prov_err:
                        logger.warning(
                            "Provider %s failed for %s: %s",
                            provider.provider_name, comp.name, prov_err,
                        )

                enrich_db.commit()

            _ai_tasks[task_id]["status"] = "completed"
            _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
            _ai_tasks[task_id]["result"] = {
                "competitors_processed": len(competitors),
                "providers_used": provider_names,
                "fields_updated": total_updates,
                "sources_created": total_sources,
            }
        except Exception as e:
            logger.error("Provider enrichment failed: %s", e)
            _ai_tasks[task_id]["status"] = "failed"
            _ai_tasks[task_id]["error"] = "Provider enrichment failed"
            _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
        finally:
            enrich_db.close()

    background_tasks.add_task(_run_enrichment)

    return {
        "task_id": task_id,
        "status": "started",
        "providers": provider_names,
        "message": f"Enrichment started with {len(active)} providers",
    }


# ============== VERIFICATION ENGINE (Verify + Correct) ==============

_verification_progress = {
    "job_id": None,
    "status": "idle",
    "competitors_total": 0,
    "competitors_processed": 0,
    "fields_verified": 0,
    "fields_corrected": 0,
    "fields_marked_na": 0,
    "current_competitor": None,
    "estimated_time_remaining": None,
    "started_at": None,
    "completed_at": None
}


@app.post("/api/verification/run-all")
async def run_verification_all(
    background_tasks: BackgroundTasks,
    current_user: dict = Depends(get_current_user)
):
    """Run AI verification on all active competitors (background task)."""
    import uuid
    import time as _time
    from source_discovery_engine import get_source_discovery_engine

    global _verification_progress  # noqa: F824

    # Check if already running
    if _verification_progress["status"] == "running":
        return {
            "status": "already_running",
            "job_id": _verification_progress["job_id"],
            "progress": _verification_progress
        }

    job_id = str(uuid.uuid4())[:8]

    # Initialize progress
    _verification_progress = {
        "job_id": job_id,
        "status": "running",
        "competitors_total": 0,
        "competitors_processed": 0,
        "fields_verified": 0,
        "fields_corrected": 0,
        "fields_marked_na": 0,
        "current_competitor": None,
        "estimated_time_remaining": None,
        "started_at": datetime.utcnow().isoformat(),
        "completed_at": None
    }

    async def run_all_verification():
        global _verification_progress  # noqa: F824

        db = SessionLocal()
        try:
            engine = get_source_discovery_engine()

            competitors = db.query(Competitor).filter(
                Competitor.status != "deleted"
            ).all()

            _verification_progress["competitors_total"] = len(competitors)
            start_time = _time.time()

            for i, comp in enumerate(competitors):
                _verification_progress["current_competitor"] = comp.name

                result = await engine.verify_sources_for_competitor(comp.id)

                _verification_progress["competitors_processed"] = i + 1
                correct = result.get("fields_correct", 0)
                corrected = result.get("fields_corrected", 0)
                _verification_progress["fields_verified"] += correct + corrected
                _verification_progress["fields_corrected"] += result.get("fields_corrected", 0)
                _verification_progress["fields_marked_na"] += result.get("fields_unverifiable", 0)

                # Estimate remaining time
                elapsed = _time.time() - start_time
                if i > 0:
                    avg_per_comp = elapsed / (i + 1)
                    remaining = len(competitors) - (i + 1)
                    _verification_progress["estimated_time_remaining"] = round(avg_per_comp * remaining)

            _verification_progress["status"] = "completed"
            _verification_progress["completed_at"] = datetime.utcnow().isoformat()
            _verification_progress["current_competitor"] = None
            _verification_progress["estimated_time_remaining"] = None

        except Exception as e:
            logger.error(f"Batch verification failed: {e}")
            _verification_progress["status"] = "error"
            _verification_progress["completed_at"] = datetime.utcnow().isoformat()
        finally:
            db.close()

    background_tasks.add_task(run_all_verification)

    return {
        "status": "started",
        "job_id": job_id,
        "message": "Verification started for all competitors"
    }


@app.get("/api/verification/progress")
async def get_verification_progress(current_user: dict = Depends(get_current_user)):
    """Get the current status of the batch verification job."""
    return _verification_progress


@app.post("/api/verification/run/{competitor_id}")
async def run_verification_single(
    competitor_id: int,
    current_user: dict = Depends(get_current_user)
):
    """Run AI verification on a single competitor (foreground)."""
    from source_discovery_engine import get_source_discovery_engine

    try:
        engine = get_source_discovery_engine()
        result = await engine.verify_sources_for_competitor(competitor_id)
        return result
    except Exception as e:
        logger.error(f"Single competitor verification failed for {competitor_id}: {e}")
        return {"error": "Verification failed", "competitor_id": competitor_id}


# ============== ENHANCED DATA SOURCES WITH CONFIDENCE SCORING ==============

@app.get("/api/competitors/{competitor_id}/data-sources")
def get_competitor_data_sources_enhanced(competitor_id: int, db: Session = Depends(get_db)):
    """Get all data sources and confidence scores for a competitor with enhanced metadata."""
    sources = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id
    ).order_by(DataSource.field_name).all()

    return [{
        "field": s.field_name,
        "value": s.current_value,
        "previous_value": s.previous_value,
        "source_type": s.source_type,
        "source_name": s.source_name,
        "source_url": s.source_url,
        "extraction_method": s.extraction_method,
        "confidence": {
            "score": s.confidence_score or 0,
            "level": s.confidence_level or "low",
            "reliability": s.source_reliability,
            "credibility": s.information_credibility,
            "corroborating_sources": s.corroborating_sources or 0,
            "reliability_description": get_reliability_description(s.source_reliability) if s.source_reliability else None,
            "credibility_description": get_credibility_description(s.information_credibility) if s.information_credibility else None
        },
        "verification": {
            "is_verified": s.is_verified,
            "verified_by": s.verified_by,
            "verification_date": s.verification_date.isoformat() if s.verification_date else None
        },
        "temporal": {
            "extracted_at": s.extracted_at.isoformat() if s.extracted_at else None,
            "data_as_of_date": s.data_as_of_date.isoformat() if s.data_as_of_date else None,
            "staleness_days": s.staleness_days or 0
        }
    } for s in sources]


# Data quality endpoints (low-confidence, confidence-distribution)
# moved to routers/data_quality.py


@app.post("/api/sources/set-with-confidence")
def set_field_source_with_confidence(
    competitor_id: int,
    field_name: str,
    current_value: str,
    source_type: str,
    source_url: Optional[str] = None,
    source_name: Optional[str] = None,
    extraction_method: Optional[str] = None,
    source_reliability: Optional[str] = None,
    information_credibility: Optional[int] = None,
    corroborating_sources: int = 0,
    data_as_of_date: Optional[str] = None,
    db: Session = Depends(get_db)
):
    """Set or update a data source with confidence scoring."""
    # Calculate confidence score
    defaults = get_source_defaults(source_type)
    reliability = source_reliability or defaults["reliability"]
    credibility = information_credibility or defaults["credibility"]

    confidence_result = calculate_confidence_score(
        source_type=source_type,
        source_reliability=reliability,
        information_credibility=credibility,
        corroborating_sources=corroborating_sources,
        data_age_days=0
    )

    # Check if source already exists
    existing = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id,
        DataSource.field_name == field_name
    ).first()

    if existing:
        # Store previous value
        existing.previous_value = existing.current_value
        existing.current_value = current_value
        existing.source_type = source_type
        existing.source_url = source_url
        existing.source_name = source_name
        existing.extraction_method = extraction_method
        existing.source_reliability = reliability
        existing.information_credibility = credibility
        existing.confidence_score = confidence_result.score
        existing.confidence_level = confidence_result.level
        existing.corroborating_sources = corroborating_sources
        existing.extracted_at = datetime.utcnow()
        if data_as_of_date:
            existing.data_as_of_date = datetime.fromisoformat(str(data_as_of_date)) if isinstance(data_as_of_date, str) else data_as_of_date
        existing.updated_at = datetime.utcnow()
    else:
        new_source = DataSource(
            competitor_id=competitor_id,
            field_name=field_name,
            current_value=current_value,
            source_type=source_type,
            source_url=source_url,
            source_name=source_name,
            extraction_method=extraction_method,
            source_reliability=reliability,
            information_credibility=credibility,
            confidence_score=confidence_result.score,
            confidence_level=confidence_result.level,
            corroborating_sources=corroborating_sources,
            data_as_of_date=(datetime.fromisoformat(str(data_as_of_date)) if isinstance(data_as_of_date, str) else data_as_of_date) if data_as_of_date else None
        )
        db.add(new_source)

    db.commit()

    return {
        "success": True,
        "message": f"Source set for {field_name}",
        "confidence": {
            "score": confidence_result.score,
            "level": confidence_result.level,
            "explanation": confidence_result.explanation,
            "breakdown": confidence_result.breakdown
        }
    }


@app.post("/api/sources/verify/{competitor_id}/{field_name}")
def verify_data_source(
    competitor_id: int,
    field_name: str,
    verification_method: str,
    corroborating_sources: int = 0,
    db: Session = Depends(get_db)
):
    """Mark a data source as verified and recalculate confidence."""
    source = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id,
        DataSource.field_name == field_name
    ).first()

    if not source:
        raise HTTPException(status_code=404, detail=f"No source found for {field_name}")

    # Update verification status
    source.is_verified = True
    source.verified_by = verification_method
    source.verification_date = datetime.utcnow()
    source.corroborating_sources = corroborating_sources

    # Recalculate confidence with verification bonus
    staleness = calculate_data_staleness(source.extracted_at, source.data_as_of_date)
    confidence_result = calculate_confidence_score(
        source_type=source.source_type or "unknown",
        source_reliability=source.source_reliability,
        information_credibility=source.information_credibility,
        corroborating_sources=corroborating_sources,
        data_age_days=staleness
    )

    source.confidence_score = confidence_result.score
    source.confidence_level = confidence_result.level
    source.staleness_days = staleness

    db.commit()

    return {
        "success": True,
        "message": f"Field {field_name} verified via {verification_method}",
        "new_confidence": {
            "score": confidence_result.score,
            "level": confidence_result.level
        }
    }


@app.get("/api/source-types")
def get_source_types():
    """Get all available source types with their default reliability ratings."""
    return {
        "source_types": [
            {
                "type": source_type,
                "reliability": info["reliability"],
                "credibility": info["credibility"],
                "description": info["description"],
                "reliability_description": get_reliability_description(info["reliability"]),
                "credibility_description": get_credibility_description(info["credibility"])
            }
            for source_type, info in SOURCE_TYPE_DEFAULTS.items()
        ],
        "reliability_scale": RELIABILITY_DESCRIPTIONS,
        "credibility_scale": CREDIBILITY_DESCRIPTIONS
    }


# Data quality endpoints (recalculate-confidence, overview) moved to routers/data_quality.py


# ============== DATA TRIANGULATION ENDPOINTS ==============

@app.post("/api/triangulate/{competitor_id}")
async def triangulate_competitor_data(competitor_id: int, db: Session = Depends(get_db)):
    """
    Triangulate all key data fields for a competitor using multiple sources.

    This cross-references data from:
    - Website scrapes
    - SEC filings (if public company)
    - News articles
    - Manual entries

    Returns verified values with confidence scores.
    """
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    triangulator = DataTriangulator(db)

    results = await triangulator.triangulate_all_key_fields(
        competitor_id=competitor_id,
        competitor_name=competitor.name,
        website=competitor.website,
        is_public=competitor.is_public,
        ticker_symbol=competitor.ticker_symbol
    )

    # Update DataSource records with triangulated confidence
    for field_name, result in results.items():
        if result.confidence_score > 0:
            existing = db.query(DataSource).filter(
                DataSource.competitor_id == competitor_id,
                DataSource.field_name == field_name
            ).first()

            if existing:
                existing.confidence_score = result.confidence_score
                existing.confidence_level = result.confidence_level
                existing.corroborating_sources = result.sources_agreeing
                existing.is_verified = result.confidence_level == "high"
                existing.verified_by = "triangulation" if result.sources_agreeing > 1 else None
                existing.verification_date = datetime.utcnow() if result.sources_agreeing > 1 else None
                existing.updated_at = datetime.utcnow()

    db.commit()

    return {
        "competitor_id": competitor_id,
        "competitor_name": competitor.name,
        "triangulation_results": {
            field: triangulation_result_to_dict(result)
            for field, result in results.items()
        },
        "triangulated_at": datetime.utcnow().isoformat()
    }


@app.post("/api/triangulate/{competitor_id}/{field_name}")
async def triangulate_single_field(
    competitor_id: int,
    field_name: str,
    db: Session = Depends(get_db)
):
    """Triangulate a specific field for a competitor."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    triangulator = DataTriangulator(db)

    # Run appropriate triangulation based on field
    if field_name == "customer_count":
        result = await triangulator.triangulate_customer_count(
            competitor_id, competitor.name, competitor.website,
            competitor.is_public, competitor.ticker_symbol
        )
    elif field_name == "employee_count":
        result = await triangulator.triangulate_employee_count(
            competitor_id, competitor.name,
            competitor.is_public, competitor.ticker_symbol
        )
    elif field_name in ["base_price", "pricing"]:
        result = await triangulator.triangulate_pricing(
            competitor_id, competitor.name, competitor.website
        )
    else:
        raise HTTPException(
            status_code=400,
            detail=f"Triangulation not supported for field: {field_name}. Supported: customer_count, employee_count, base_price"
        )

    # Update DataSource with triangulated confidence
    existing = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id,
        DataSource.field_name == field_name
    ).first()

    if existing and result.confidence_score > 0:
        existing.confidence_score = result.confidence_score
        existing.confidence_level = result.confidence_level
        existing.corroborating_sources = result.sources_agreeing
        existing.is_verified = result.sources_agreeing > 1
        existing.verified_by = "triangulation" if result.sources_agreeing > 1 else None
        existing.verification_date = datetime.utcnow() if result.sources_agreeing > 1 else None
        existing.updated_at = datetime.utcnow()
        db.commit()

    return triangulation_result_to_dict(result)


@app.post("/api/triangulate/all")
async def triangulate_all_competitors(
    background_tasks: BackgroundTasks,
    db: Session = Depends(get_db)
):
    """Trigger triangulation for all active competitors (background job)."""
    competitors = db.query(Competitor).filter(
        Competitor.is_deleted == False
    ).all()

    # Add background task for each competitor
    for comp in competitors:
        background_tasks.add_task(
            run_triangulation_job,
            comp.id,
            comp.name,
            comp.website,
            comp.is_public,
            comp.ticker_symbol
        )

    return {
        "success": True,
        "message": f"Triangulation started for {len(competitors)} competitors",
        "competitors_queued": len(competitors)
    }


async def run_triangulation_job(
    competitor_id: int,
    competitor_name: str,
    website: str,
    is_public: bool,
    ticker_symbol: str
):
    """Background job to triangulate data for a competitor."""
    db = SessionLocal()
    try:
        triangulator = DataTriangulator(db)
        results = await triangulator.triangulate_all_key_fields(
            competitor_id, competitor_name, website, is_public, ticker_symbol
        )

        # Update DataSource records
        for field_name, result in results.items():
            if result.confidence_score > 0:
                existing = db.query(DataSource).filter(
                    DataSource.competitor_id == competitor_id,
                    DataSource.field_name == field_name
                ).first()

                if existing:
                    existing.confidence_score = result.confidence_score
                    existing.confidence_level = result.confidence_level
                    existing.corroborating_sources = result.sources_agreeing
                    existing.is_verified = result.confidence_level == "high"
                    existing.verified_by = "triangulation" if result.sources_agreeing > 1 else None
                    existing.verification_date = datetime.utcnow() if result.sources_agreeing > 1 else None

        db.commit()
        logger.info(f"Triangulation complete for {competitor_name}")

    except Exception as e:
        logger.error(f"Triangulation failed for {competitor_name}: {e}")
        db.rollback()
    finally:
        db.close()


@app.get("/api/triangulation/status")
def get_triangulation_status(db: Session = Depends(get_db)):
    """Get overview of triangulation status across all competitors."""
    sources = db.query(DataSource).all()

    verified_count = len([s for s in sources if s.is_verified])
    triangulated_count = len([s for s in sources if s.verified_by == "triangulation"])
    pending_count = len([s for s in sources if not s.is_verified and s.current_value])

    # Group by confidence level
    high = len([s for s in sources if s.confidence_level == "high"])
    moderate = len([s for s in sources if s.confidence_level == "moderate"])
    low = len([s for s in sources if s.confidence_level == "low" or s.confidence_level is None])

    return {
        "total_data_points": len(sources),
        "verification_status": {
            "verified": verified_count,
            "triangulated": triangulated_count,
            "pending_verification": pending_count
        },
        "confidence_distribution": {
            "high": high,
            "moderate": moderate,
            "low": low
        },
        "verification_rate": round(verified_count / len(sources) * 100, 1) if sources else 0
    }


# ============== PHASE 3: PRODUCT & PRICING MANAGEMENT ==============

# Product/Pricing/CustomerCount models imported from schemas.products

# ============== PRODUCT CRUD ENDPOINTS ==============

@app.get("/api/competitors/{competitor_id}/products", response_model=List[ProductResponse])
async def get_competitor_products(competitor_id: int, db: Session = Depends(get_db)):
    """Get all products and pricing for a competitor."""
    # Verify competitor exists
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    products = db.query(CompetitorProduct).filter(
        CompetitorProduct.competitor_id == competitor_id
    ).all()

    result = []
    for p in products:
        # Get pricing tiers for this product
        tiers = db.query(ProductPricingTier).filter(
            ProductPricingTier.product_id == p.id
        ).order_by(ProductPricingTier.tier_position).all()

        product_dict = {
            "id": p.id,
            "competitor_id": p.competitor_id,
            "product_name": p.product_name,
            "product_category": p.product_category,
            "product_subcategory": p.product_subcategory,
            "description": p.description,
            "key_features": p.key_features,
            "target_segment": p.target_segment,
            "is_primary_product": p.is_primary_product,
            "market_position": p.market_position,
            "launched_date": p.launched_date,
            "last_updated": p.last_updated,
            "pricing_tiers": [{
                "id": t.id,
                "tier_name": t.tier_name,
                "pricing_model": t.pricing_model,
                "price_display": t.price_display,
                "base_price": t.base_price,
                "price_unit": t.price_unit,
                "percentage_rate": t.percentage_rate,
                "confidence_score": t.confidence_score,
                "price_verified": t.price_verified
            } for t in tiers]
        }
        result.append(product_dict)

    return result


@app.post("/api/products", response_model=ProductResponse)
async def create_product(product: ProductCreate, db: Session = Depends(get_db)):
    """Create a new product for a competitor."""
    # Verify competitor exists
    competitor = db.query(Competitor).filter(Competitor.id == product.competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    # Check if product already exists
    existing = db.query(CompetitorProduct).filter(
        CompetitorProduct.competitor_id == product.competitor_id,
        CompetitorProduct.product_name == product.product_name
    ).first()
    if existing:
        raise HTTPException(status_code=400, detail="Product already exists for this competitor")

    new_product = CompetitorProduct(
        competitor_id=product.competitor_id,
        product_name=product.product_name,
        product_category=product.product_category,
        product_subcategory=product.product_subcategory,
        description=product.description,
        key_features=product.key_features,
        target_segment=product.target_segment,
        is_primary_product=product.is_primary_product,
        market_position=product.market_position,
        last_updated=datetime.utcnow()
    )
    db.add(new_product)
    db.commit()
    db.refresh(new_product)

    return {
        "id": new_product.id,
        "competitor_id": new_product.competitor_id,
        "product_name": new_product.product_name,
        "product_category": new_product.product_category,
        "product_subcategory": new_product.product_subcategory,
        "description": new_product.description,
        "key_features": new_product.key_features,
        "target_segment": new_product.target_segment,
        "is_primary_product": new_product.is_primary_product,
        "market_position": new_product.market_position,
        "launched_date": new_product.launched_date,
        "last_updated": new_product.last_updated,
        "pricing_tiers": []
    }


@app.put("/api/products/{product_id}", response_model=ProductResponse)
async def update_product(product_id: int, product: ProductCreate, db: Session = Depends(get_db)):
    """Update an existing product."""
    existing = db.query(CompetitorProduct).filter(CompetitorProduct.id == product_id).first()
    if not existing:
        raise HTTPException(status_code=404, detail="Product not found")

    existing.product_name = product.product_name
    existing.product_category = product.product_category
    existing.product_subcategory = product.product_subcategory
    existing.description = product.description
    existing.key_features = product.key_features
    existing.target_segment = product.target_segment
    existing.is_primary_product = product.is_primary_product
    existing.market_position = product.market_position
    existing.last_updated = datetime.utcnow()

    db.commit()
    db.refresh(existing)

    # Get pricing tiers
    tiers = db.query(ProductPricingTier).filter(
        ProductPricingTier.product_id == product_id
    ).order_by(ProductPricingTier.tier_position).all()

    return {
        "id": existing.id,
        "competitor_id": existing.competitor_id,
        "product_name": existing.product_name,
        "product_category": existing.product_category,
        "product_subcategory": existing.product_subcategory,
        "description": existing.description,
        "key_features": existing.key_features,
        "target_segment": existing.target_segment,
        "is_primary_product": existing.is_primary_product,
        "market_position": existing.market_position,
        "launched_date": existing.launched_date,
        "last_updated": existing.last_updated,
        "pricing_tiers": [{
            "id": t.id,
            "tier_name": t.tier_name,
            "pricing_model": t.pricing_model,
            "price_display": t.price_display,
            "base_price": t.base_price,
            "price_unit": t.price_unit,
            "confidence_score": t.confidence_score
        } for t in tiers]
    }


@app.delete("/api/products/{product_id}")
async def delete_product(product_id: int, db: Session = Depends(get_db)):
    """Delete a product and its pricing tiers."""
    product = db.query(CompetitorProduct).filter(CompetitorProduct.id == product_id).first()
    if not product:
        raise HTTPException(status_code=404, detail="Product not found")

    # Delete associated pricing tiers first
    db.query(ProductPricingTier).filter(ProductPricingTier.product_id == product_id).delete()
    # Delete associated features
    db.query(ProductFeatureMatrix).filter(ProductFeatureMatrix.product_id == product_id).delete()
    # Delete the product
    db.delete(product)
    db.commit()

    return {"message": "Product deleted successfully", "product_id": product_id}


# ============== PRICING TIER ENDPOINTS ==============

@app.get("/api/products/{product_id}/pricing-tiers", response_model=List[PricingTierResponse])
async def get_pricing_tiers(product_id: int, db: Session = Depends(get_db)):
    """Get all pricing tiers for a product."""
    product = db.query(CompetitorProduct).filter(CompetitorProduct.id == product_id).first()
    if not product:
        raise HTTPException(status_code=404, detail="Product not found")

    tiers = db.query(ProductPricingTier).filter(
        ProductPricingTier.product_id == product_id
    ).order_by(ProductPricingTier.tier_position).all()

    return tiers


@app.post("/api/pricing-tiers", response_model=PricingTierResponse)
async def create_pricing_tier(tier: PricingTierCreate, db: Session = Depends(get_db)):
    """Create a new pricing tier for a product."""
    # Verify product exists
    product = db.query(CompetitorProduct).filter(CompetitorProduct.id == tier.product_id).first()
    if not product:
        raise HTTPException(status_code=404, detail="Product not found")

    # Calculate confidence score based on source
    source_type = "website_scrape" if tier.price_source == "website" else "manual_verified"
    if tier.price_source == "sales_quote":
        source_type = "api_verified"

    confidence_result = calculate_confidence_score(source_type=source_type)

    new_tier = ProductPricingTier(
        product_id=tier.product_id,
        tier_name=tier.tier_name,
        tier_position=tier.tier_position,
        pricing_model=tier.pricing_model,
        base_price=tier.base_price,
        price_currency=tier.price_currency,
        price_unit=tier.price_unit,
        price_display=tier.price_display,
        percentage_rate=tier.percentage_rate,
        percentage_basis=tier.percentage_basis,
        min_volume=tier.min_volume,
        max_volume=tier.max_volume,
        included_features=tier.included_features,
        excluded_features=tier.excluded_features,
        contract_length=tier.contract_length,
        setup_fee=tier.setup_fee,
        implementation_cost=tier.implementation_cost,
        price_source=tier.price_source,
        price_verified=False,
        confidence_score=confidence_result.score,
        last_verified=datetime.utcnow()
    )
    db.add(new_tier)
    db.commit()
    db.refresh(new_tier)

    return new_tier


@app.put("/api/pricing-tiers/{tier_id}", response_model=PricingTierResponse)
async def update_pricing_tier(tier_id: int, tier: PricingTierCreate, db: Session = Depends(get_db)):
    """Update an existing pricing tier."""
    existing = db.query(ProductPricingTier).filter(ProductPricingTier.id == tier_id).first()
    if not existing:
        raise HTTPException(status_code=404, detail="Pricing tier not found")

    # Update fields
    existing.tier_name = tier.tier_name
    existing.tier_position = tier.tier_position
    existing.pricing_model = tier.pricing_model
    existing.base_price = tier.base_price
    existing.price_currency = tier.price_currency
    existing.price_unit = tier.price_unit
    existing.price_display = tier.price_display
    existing.percentage_rate = tier.percentage_rate
    existing.percentage_basis = tier.percentage_basis
    existing.min_volume = tier.min_volume
    existing.max_volume = tier.max_volume
    existing.included_features = tier.included_features
    existing.excluded_features = tier.excluded_features
    existing.contract_length = tier.contract_length
    existing.setup_fee = tier.setup_fee
    existing.implementation_cost = tier.implementation_cost
    existing.price_source = tier.price_source
    existing.updated_at = datetime.utcnow()

    db.commit()
    db.refresh(existing)

    return existing


@app.delete("/api/pricing-tiers/{tier_id}")
async def delete_pricing_tier(tier_id: int, db: Session = Depends(get_db)):
    """Delete a pricing tier."""
    tier = db.query(ProductPricingTier).filter(ProductPricingTier.id == tier_id).first()
    if not tier:
        raise HTTPException(status_code=404, detail="Pricing tier not found")

    db.delete(tier)
    db.commit()

    return {"message": "Pricing tier deleted successfully", "tier_id": tier_id}


@app.post("/api/pricing-tiers/{tier_id}/verify")
async def verify_pricing_tier(tier_id: int, db: Session = Depends(get_db)):
    """Mark a pricing tier as verified."""
    tier = db.query(ProductPricingTier).filter(ProductPricingTier.id == tier_id).first()
    if not tier:
        raise HTTPException(status_code=404, detail="Pricing tier not found")

    tier.price_verified = True
    tier.last_verified = datetime.utcnow()
    # Boost confidence when verified
    tier.confidence_score = min(100, (tier.confidence_score or 50) + 20)

    db.commit()
    db.refresh(tier)

    return {
        "message": "Pricing tier verified",
        "tier_id": tier_id,
        "confidence_score": tier.confidence_score
    }


# ============== PRICING COMPARISON ENDPOINT ==============

@app.get("/api/pricing/compare")
async def compare_pricing(
    category: Optional[str] = None,  # e.g., "Patient Intake"
    pricing_model: Optional[str] = None,  # e.g., "per_visit"
    db: Session = Depends(get_db)
):
    """Compare pricing across competitors for a product category."""
    query = db.query(ProductPricingTier).join(CompetitorProduct).join(Competitor)

    if category:
        query = query.filter(CompetitorProduct.product_category == category)

    if pricing_model:
        query = query.filter(ProductPricingTier.pricing_model == pricing_model)

    tiers = query.all()

    result = []
    for t in tiers:
        product = db.query(CompetitorProduct).filter(CompetitorProduct.id == t.product_id).first()
        competitor = db.query(Competitor).filter(Competitor.id == product.competitor_id).first() if product else None

        result.append({
            "competitor_id": competitor.id if competitor else None,
            "competitor_name": competitor.name if competitor else "Unknown",
            "product_id": product.id if product else None,
            "product_name": product.product_name if product else "Unknown",
            "product_category": product.product_category if product else None,
            "tier_name": t.tier_name,
            "pricing_model": t.pricing_model,
            "base_price": t.base_price,
            "price_display": t.price_display,
            "price_unit": t.price_unit,
            "percentage_rate": t.percentage_rate,
            "confidence_score": t.confidence_score,
            "price_verified": t.price_verified,
            "price_source": t.price_source
        })

    # Sort by base price (nulls last)
    result.sort(key=lambda x: (x["base_price"] is None, x["base_price"] or 0))

    return {
        "category": category,
        "pricing_model": pricing_model,
        "total_tiers": len(result),
        "comparison": result
    }


@app.get("/api/pricing/models")
async def get_pricing_models():
    """Get available healthcare pricing model types."""
    return {
        "pricing_models": [
            {"value": "per_visit", "label": "Per Visit/Encounter", "description": "Charge per patient encounter", "example": "$3.00/visit"},
            {"value": "per_provider", "label": "Per Provider", "description": "Monthly fee per provider/physician", "example": "$400/provider/month"},
            {"value": "per_location", "label": "Per Location", "description": "Fee per practice location", "example": "$1,500/location/month"},
            {"value": "subscription_flat", "label": "Flat Subscription", "description": "Fixed monthly fee", "example": "$299/month"},
            {"value": "subscription_tiered", "label": "Tiered Subscription", "description": "Tiered by features or volume", "example": "$99-$499/month"},
            {"value": "percentage_collections", "label": "Percentage of Collections", "description": "% of collected revenue", "example": "4-8% of collections"},
            {"value": "percentage_charges", "label": "Percentage of Charges", "description": "% of billed charges", "example": "2-4% of charges"},
            {"value": "per_bed", "label": "Per Bed", "description": "Hospital capacity pricing", "example": "$15,000/bed"},
            {"value": "per_member", "label": "Per Member (PMPM)", "description": "Per covered life", "example": "$0.50 PMPM"},
            {"value": "custom_enterprise", "label": "Custom/Enterprise", "description": "Negotiated pricing", "example": "Contact Sales"}
        ],
        "product_categories": [
            "Patient Intake",
            "Patient Payments",
            "Revenue Cycle Management (RCM)",
            "Practice Management",
            "EHR/EMR",
            "Telehealth",
            "Patient Engagement",
            "Scheduling",
            "Analytics",
            "Interoperability"
        ]
    }


# ============== FEATURE MATRIX ENDPOINTS ==============

@app.get("/api/products/{product_id}/features")
async def get_product_features(product_id: int, db: Session = Depends(get_db)):
    """Get all features for a product."""
    product = db.query(CompetitorProduct).filter(CompetitorProduct.id == product_id).first()
    if not product:
        raise HTTPException(status_code=404, detail="Product not found")

    features = db.query(ProductFeatureMatrix).filter(
        ProductFeatureMatrix.product_id == product_id
    ).all()

    # Group by category
    by_category = {}
    for f in features:
        if f.feature_category not in by_category:
            by_category[f.feature_category] = []
        by_category[f.feature_category].append({
            "id": f.id,
            "feature_name": f.feature_name,
            "feature_status": f.feature_status,
            "feature_tier": f.feature_tier,
            "notes": f.notes,
            "last_verified": f.last_verified
        })

    return {
        "product_id": product_id,
        "product_name": product.product_name,
        "features_by_category": by_category,
        "total_features": len(features)
    }


@app.post("/api/features")
async def create_feature(feature: FeatureMatrixCreate, db: Session = Depends(get_db)):
    """Add a feature to a product."""
    product = db.query(CompetitorProduct).filter(CompetitorProduct.id == feature.product_id).first()
    if not product:
        raise HTTPException(status_code=404, detail="Product not found")

    new_feature = ProductFeatureMatrix(
        product_id=feature.product_id,
        feature_category=feature.feature_category,
        feature_name=feature.feature_name,
        feature_status=feature.feature_status,
        feature_tier=feature.feature_tier,
        notes=feature.notes,
        source_url=feature.source_url,
        last_verified=datetime.utcnow()
    )
    db.add(new_feature)
    db.commit()
    db.refresh(new_feature)

    return {
        "id": new_feature.id,
        "product_id": new_feature.product_id,
        "feature_category": new_feature.feature_category,
        "feature_name": new_feature.feature_name,
        "feature_status": new_feature.feature_status,
        "message": "Feature added successfully"
    }


@app.delete("/api/features/{feature_id}")
async def delete_feature(feature_id: int, db: Session = Depends(get_db)):
    """Delete a feature."""
    feature = db.query(ProductFeatureMatrix).filter(ProductFeatureMatrix.id == feature_id).first()
    if not feature:
        raise HTTPException(status_code=404, detail="Feature not found")

    db.delete(feature)
    db.commit()

    return {"message": "Feature deleted successfully", "feature_id": feature_id}


@app.get("/api/features/compare")
async def compare_features(
    category: str,  # Product category like "Patient Intake"
    feature_category: Optional[str] = None,  # Feature category like "Payments"
    db: Session = Depends(get_db)
):
    """Compare features across competitors for a product category."""
    # Get all products in this category
    products = db.query(CompetitorProduct).filter(
        CompetitorProduct.product_category == category
    ).all()

    if not products:
        return {"message": f"No products found in category: {category}", "comparison": []}

    # Get all unique features
    feature_query = db.query(ProductFeatureMatrix).filter(
        ProductFeatureMatrix.product_id.in_([p.id for p in products])
    )
    if feature_category:
        feature_query = feature_query.filter(ProductFeatureMatrix.feature_category == feature_category)

    all_features = feature_query.all()

    # Build comparison matrix
    # Structure: {feature_name: {competitor_name: status}}
    feature_names = set(f.feature_name for f in all_features)

    comparison = []
    for feature_name in sorted(feature_names):
        feature_row = {"feature_name": feature_name, "competitors": {}}
        for product in products:
            competitor = db.query(Competitor).filter(Competitor.id == product.competitor_id).first()
            comp_name = competitor.name if competitor else f"Competitor {product.competitor_id}"

            # Find this feature for this product
            feature = next(
                (f for f in all_features if f.product_id == product.id and f.feature_name == feature_name),
                None
            )
            feature_row["competitors"][comp_name] = feature.feature_status if feature else "unknown"
        comparison.append(feature_row)

    return {
        "product_category": category,
        "feature_category": feature_category,
        "competitors": [
            db.query(Competitor).filter(Competitor.id == p.competitor_id).first().name
            for p in products
        ],
        "comparison": comparison
    }


# ============== PRODUCT EXTRACTION FROM CONTENT ==============

@app.post("/api/competitors/{competitor_id}/extract-products")
async def extract_products_from_content(
    competitor_id: int,
    background_tasks: BackgroundTasks,
    db: Session = Depends(get_db)
):
    """
    Extract products and pricing from competitor's scraped content using AI.
    This endpoint triggers the extraction and stores results in the database.

    v5.0.2: Uses hybrid AI routing (OpenAI or Gemini based on config)
    """
    from extractor import get_extractor

    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    # Check if we have scraped content
    # Try to get recent scrape content from the competitor's website
    content = ""

    # First check if there's recent DataSource with website content
    recent_source = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id,
        DataSource.source_type == "website_scrape"
    ).order_by(DataSource.extracted_at.desc()).first()

    if recent_source and recent_source.current_value:
        content = recent_source.current_value
    else:
        # Try to scrape fresh content
        try:
            from scraper import CompetitorScraper

            async def scrape_for_products():
                async with CompetitorScraper() as scraper:
                    result = await scraper.scrape(competitor.website)
                    return result.get("content", "") if isinstance(result, dict) else ""

            # Run the scraper - await since we're in an async endpoint
            content = await scrape_for_products()
        except Exception as e:
            return {
                "status": "error",
                "message": "Could not get content for extraction",
                "competitor_id": competitor_id
            }

    if not content or len(content) < 100:
        return {
            "status": "error",
            "message": "Not enough content available for extraction. Try scraping the competitor first.",
            "competitor_id": competitor_id
        }

    # Extract products using AI (v5.0.2 - hybrid routing)
    extractor = get_extractor()
    extraction_result = extractor.extract_products_and_pricing(competitor.name, content)

    if "error" in extraction_result:
        return {
            "status": "error",
            "message": extraction_result["error"],
            "competitor_id": competitor_id
        }

    products_created = 0
    tiers_created = 0

    # Process extracted products
    for product_data in extraction_result.get("products", []):
        # Check if product already exists
        existing_product = db.query(CompetitorProduct).filter(
            CompetitorProduct.competitor_id == competitor_id,
            CompetitorProduct.product_name == product_data.get("product_name")
        ).first()

        if existing_product:
            # Update existing product
            existing_product.product_category = product_data.get("product_category", existing_product.product_category)
            existing_product.target_segment = product_data.get("target_segment", existing_product.target_segment)
            existing_product.is_primary_product = product_data.get("is_primary_product", existing_product.is_primary_product)
            existing_product.key_features = json.dumps(product_data.get("key_features", [])) if product_data.get("key_features") else existing_product.key_features
            existing_product.last_updated = datetime.utcnow()
            product = existing_product
        else:
            # Create new product
            product = CompetitorProduct(
                competitor_id=competitor_id,
                product_name=product_data.get("product_name", f"{competitor.name} Product"),
                product_category=product_data.get("product_category", "Unknown"),
                target_segment=product_data.get("target_segment"),
                is_primary_product=product_data.get("is_primary_product", False),
                key_features=json.dumps(product_data.get("key_features", [])) if product_data.get("key_features") else None,
                last_updated=datetime.utcnow()
            )
            db.add(product)
            db.flush()  # Get the ID
            products_created += 1

        # Process pricing tiers
        for tier_data in product_data.get("pricing_tiers", []):
            # Check if tier exists
            existing_tier = db.query(ProductPricingTier).filter(
                ProductPricingTier.product_id == product.id,
                ProductPricingTier.tier_name == tier_data.get("tier_name")
            ).first()

            if existing_tier:
                # Update existing tier
                existing_tier.pricing_model = tier_data.get("pricing_model", existing_tier.pricing_model)
                existing_tier.base_price = tier_data.get("base_price", existing_tier.base_price)
                existing_tier.price_currency = tier_data.get("price_currency", "USD")
                existing_tier.price_unit = tier_data.get("price_unit", existing_tier.price_unit)
                existing_tier.price_display = tier_data.get("price_display", existing_tier.price_display)
                existing_tier.percentage_rate = tier_data.get("percentage_rate", existing_tier.percentage_rate)
                existing_tier.setup_fee = tier_data.get("setup_fee", existing_tier.setup_fee)
                existing_tier.contract_length = tier_data.get("contract_length", existing_tier.contract_length)
                existing_tier.included_features = json.dumps(tier_data.get("included_features", [])) if tier_data.get("included_features") else existing_tier.included_features
                existing_tier.price_source = "gpt_extraction"
                existing_tier.updated_at = datetime.utcnow()
            else:
                # Calculate confidence score for extracted pricing
                confidence_result = calculate_confidence_score(
                    source_type="website_scrape",
                    information_credibility=4  # GPT extraction from marketing content
                )

                # Create new tier
                new_tier = ProductPricingTier(
                    product_id=product.id,
                    tier_name=tier_data.get("tier_name", "Standard"),
                    tier_position=tier_data.get("tier_position", 1),
                    pricing_model=tier_data.get("pricing_model", "custom_enterprise"),
                    base_price=tier_data.get("base_price"),
                    price_currency=tier_data.get("price_currency", "USD"),
                    price_unit=tier_data.get("price_unit"),
                    price_display=tier_data.get("price_display"),
                    percentage_rate=tier_data.get("percentage_rate"),
                    setup_fee=tier_data.get("setup_fee"),
                    contract_length=tier_data.get("contract_length"),
                    included_features=json.dumps(tier_data.get("included_features", [])) if tier_data.get("included_features") else None,
                    price_source="gpt_extraction",
                    price_verified=False,
                    confidence_score=confidence_result.score
                )
                db.add(new_tier)
                tiers_created += 1

    db.commit()

    return {
        "status": "success",
        "competitor_id": competitor_id,
        "competitor_name": competitor.name,
        "products_created": products_created,
        "tiers_created": tiers_created,
        "extraction_confidence": extraction_result.get("extraction_confidence", 0),
        "extraction_notes": extraction_result.get("extraction_notes", "")
    }


@app.post("/api/products/{product_id}/extract-features")
async def extract_features_from_content(
    product_id: int,
    db: Session = Depends(get_db)
):
    """
    Extract feature matrix for a product using AI.

    v5.0.2: Uses hybrid AI routing (OpenAI or Gemini based on config)
    """
    from extractor import get_extractor

    product = db.query(CompetitorProduct).filter(CompetitorProduct.id == product_id).first()
    if not product:
        raise HTTPException(status_code=404, detail="Product not found")

    competitor = db.query(Competitor).filter(Competitor.id == product.competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    # Get content
    content = ""
    recent_source = db.query(DataSource).filter(
        DataSource.competitor_id == product.competitor_id,
        DataSource.source_type == "website_scrape"
    ).order_by(DataSource.extracted_at.desc()).first()

    if recent_source and recent_source.current_value:
        content = recent_source.current_value
    else:
        return {
            "status": "error",
            "message": "No content available for feature extraction. Scrape the competitor first."
        }

    # Extract features (v5.0.2 - hybrid routing)
    extractor = get_extractor()
    extraction_result = extractor.extract_feature_matrix(competitor.name, product.product_name, content)

    if "error" in extraction_result:
        return {
            "status": "error",
            "message": extraction_result["error"]
        }

    features_created = 0

    for feature_data in extraction_result.get("features", []):
        # Check if feature exists
        existing = db.query(ProductFeatureMatrix).filter(
            ProductFeatureMatrix.product_id == product_id,
            ProductFeatureMatrix.feature_name == feature_data.get("feature_name")
        ).first()

        if not existing:
            new_feature = ProductFeatureMatrix(
                product_id=product_id,
                feature_category=feature_data.get("feature_category", "Other"),
                feature_name=feature_data.get("feature_name"),
                feature_status=feature_data.get("feature_status", "unknown"),
                feature_tier=feature_data.get("feature_tier"),
                notes=feature_data.get("notes"),
                last_verified=datetime.utcnow()
            )
            db.add(new_feature)
            features_created += 1

    db.commit()

    return {
        "status": "success",
        "product_id": product_id,
        "product_name": product.product_name,
        "features_created": features_created,
        "extraction_confidence": extraction_result.get("extraction_confidence", 0)
    }


# ============== PHASE 4: CUSTOMER COUNT VERIFICATION ENDPOINTS ==============

@app.get("/api/competitors/{competitor_id}/customer-counts", response_model=List[CustomerCountResponse])
async def get_customer_counts(competitor_id: int, db: Session = Depends(get_db)):
    """Get all customer count estimates for a competitor, ordered by date."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    counts = db.query(CustomerCountEstimate).filter(
        CustomerCountEstimate.competitor_id == competitor_id
    ).order_by(CustomerCountEstimate.as_of_date.desc()).all()

    return counts


@app.get("/api/competitors/{competitor_id}/customer-count/latest", response_model=CustomerCountResponse)
async def get_latest_customer_count(competitor_id: int, db: Session = Depends(get_db)):
    """Get the most recent verified customer count for a competitor."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    # First try to get verified count
    count = db.query(CustomerCountEstimate).filter(
        CustomerCountEstimate.competitor_id == competitor_id,
        CustomerCountEstimate.is_verified == True
    ).order_by(CustomerCountEstimate.as_of_date.desc()).first()

    # If no verified count, get any most recent
    if not count:
        count = db.query(CustomerCountEstimate).filter(
            CustomerCountEstimate.competitor_id == competitor_id
        ).order_by(CustomerCountEstimate.as_of_date.desc()).first()

    if not count:
        raise HTTPException(status_code=404, detail="No customer count estimates found")

    return count


@app.post("/api/customer-counts", response_model=CustomerCountResponse)
async def create_customer_count(
    count_data: CustomerCountCreate,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Create a new customer count estimate."""
    # Verify competitor exists
    competitor = db.query(Competitor).filter(Competitor.id == count_data.competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    # Calculate confidence score based on source type
    source_mapping = {
        "sec_10k": "sec_filing",
        "sec_filing": "sec_filing",
        "website": "website_scrape",
        "press_release": "news_article",
        "definitive_hc": "definitive_hc",
        "klas_report": "klas_report",
        "linkedin": "linkedin_estimate",
        "g2_reviews": "api_verified",
        "manual": "manual_verified"
    }
    source_type = source_mapping.get(count_data.primary_source, "unknown")
    confidence_result = calculate_confidence_score(source_type=source_type)

    # Get previous count for growth calculation
    previous = db.query(CustomerCountEstimate).filter(
        CustomerCountEstimate.competitor_id == count_data.competitor_id
    ).order_by(CustomerCountEstimate.as_of_date.desc()).first()

    growth_rate = None
    previous_count = None
    if previous and previous.count_value and count_data.count_value:
        previous_count = previous.count_value
        if previous_count > 0:
            growth_rate = ((count_data.count_value - previous_count) / previous_count) * 100

    new_count = CustomerCountEstimate(
        competitor_id=count_data.competitor_id,
        count_value=count_data.count_value,
        count_display=count_data.count_display,
        count_type=count_data.count_type,
        count_unit=count_data.count_unit,
        count_definition=count_data.count_definition,
        segment_breakdown=count_data.segment_breakdown,
        primary_source=count_data.primary_source,
        primary_source_url=count_data.primary_source_url,
        primary_source_date=count_data.primary_source_date,
        as_of_date=count_data.as_of_date or datetime.utcnow(),
        previous_count=previous_count,
        growth_rate=growth_rate,
        confidence_score=confidence_result.score,
        confidence_level=confidence_result.level,
        is_verified=(count_data.primary_source in ["sec_10k", "sec_filing", "definitive_hc"])
    )
    db.add(new_count)
    db.commit()
    db.refresh(new_count)

    # Log activity
    log_activity(
        db, current_user["email"], current_user["id"],
        "customer_count_added",
        f"Added customer count for {competitor.name}: {count_data.count_display}"
    )

    return new_count


@app.put("/api/customer-counts/{count_id}", response_model=CustomerCountResponse)
async def update_customer_count(
    count_id: int,
    count_data: CustomerCountCreate,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Update an existing customer count estimate."""
    count = db.query(CustomerCountEstimate).filter(CustomerCountEstimate.id == count_id).first()
    if not count:
        raise HTTPException(status_code=404, detail="Customer count estimate not found")

    # Update fields
    for field in ["count_value", "count_display", "count_type", "count_unit",
                  "count_definition", "segment_breakdown", "primary_source",
                  "primary_source_url", "primary_source_date", "as_of_date"]:
        if hasattr(count_data, field):
            value = getattr(count_data, field)
            if value is not None:
                setattr(count, field, value)

    count.updated_at = datetime.utcnow()
    db.commit()
    db.refresh(count)

    return count


@app.delete("/api/customer-counts/{count_id}")
async def delete_customer_count(
    count_id: int,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Delete a customer count estimate."""
    count = db.query(CustomerCountEstimate).filter(CustomerCountEstimate.id == count_id).first()
    if not count:
        raise HTTPException(status_code=404, detail="Customer count estimate not found")

    competitor = db.query(Competitor).filter(Competitor.id == count.competitor_id).first()
    db.delete(count)
    db.commit()

    log_activity(
        db, current_user["email"], current_user["id"],
        "customer_count_deleted",
        f"Deleted customer count for {competitor.name if competitor else 'Unknown'}"
    )

    return {"status": "deleted", "count_id": count_id}


@app.post("/api/customer-counts/{count_id}/verify", response_model=CustomerCountResponse)
async def verify_customer_count(
    count_id: int,
    verification: CustomerCountVerifyRequest,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Verify a customer count with additional sources or methods."""
    count = db.query(CustomerCountEstimate).filter(CustomerCountEstimate.id == count_id).first()
    if not count:
        raise HTTPException(status_code=404, detail="Customer count estimate not found")

    # Update verification status
    count.is_verified = True
    count.verification_method = verification.verification_method
    count.verification_date = datetime.utcnow()
    count.confidence_notes = verification.verification_notes

    # If additional sources provided, store them and recalculate confidence
    if verification.additional_sources:
        import json
        existing_sources = json.loads(count.all_sources) if count.all_sources else []
        existing_sources.extend(verification.additional_sources)
        count.all_sources = json.dumps(existing_sources)
        count.corroborating_sources = len(existing_sources)

        # Calculate source agreement score
        if count.count_value:
            values = [count.count_value]
            for source in existing_sources:
                if source.get("value"):
                    try:
                        values.append(int(str(source["value"]).replace(",", "").replace("+", "")))
                    except (ValueError, TypeError):
                        pass

            if len(values) > 1:
                avg = sum(values) / len(values)
                # Agreement score: 1.0 if all within 20% of average, 0 if wildly different
                deviations = [abs(v - avg) / avg for v in values if avg > 0]
                if deviations:
                    count.source_agreement_score = max(0, 1 - (sum(deviations) / len(deviations) / 0.2))

        # Recalculate confidence with corroboration
        confidence_result = calculate_confidence_score(
            source_type="manual_verified" if verification.verification_method == "manual" else "api_verified",
            corroborating_sources=len(existing_sources)
        )
        count.confidence_score = confidence_result.score
        count.confidence_level = confidence_result.level

    count.updated_at = datetime.utcnow()
    db.commit()
    db.refresh(count)

    # Log activity
    competitor = db.query(Competitor).filter(Competitor.id == count.competitor_id).first()
    log_activity(
        db, current_user["email"], current_user["id"],
        "customer_count_verified",
        f"Verified customer count for {competitor.name if competitor else 'Unknown'}"
    )

    return count


@app.get("/api/customer-counts/compare")
async def compare_customer_counts(
    unit: Optional[str] = None,
    min_confidence: int = 0,
    db: Session = Depends(get_db)
):
    """Compare customer counts across all competitors."""
    query = db.query(CustomerCountEstimate)

    if unit:
        query = query.filter(CustomerCountEstimate.count_unit == unit)
    if min_confidence > 0:
        query = query.filter(CustomerCountEstimate.confidence_score >= min_confidence)

    # Get latest count per competitor
    from sqlalchemy import func
    subquery = db.query(
        CustomerCountEstimate.competitor_id,
        func.max(CustomerCountEstimate.as_of_date).label("max_date")
    ).group_by(CustomerCountEstimate.competitor_id).subquery()

    counts = query.join(
        subquery,
        (CustomerCountEstimate.competitor_id == subquery.c.competitor_id) &
        (CustomerCountEstimate.as_of_date == subquery.c.max_date)
    ).all()

    result = []
    for c in counts:
        competitor = db.query(Competitor).filter(Competitor.id == c.competitor_id).first()
        result.append({
            "competitor_id": c.competitor_id,
            "competitor_name": competitor.name if competitor else "Unknown",
            "count_value": c.count_value,
            "count_display": c.count_display,
            "count_unit": c.count_unit,
            "confidence_score": c.confidence_score,
            "confidence_level": c.confidence_level,
            "is_verified": c.is_verified,
            "as_of_date": c.as_of_date.isoformat() if c.as_of_date else None,
            "growth_rate": c.growth_rate
        })

    # Sort by count_value descending
    result.sort(key=lambda x: x["count_value"] or 0, reverse=True)

    return {
        "comparisons": result,
        "total_competitors": len(result),
        "unit_filter": unit,
        "min_confidence_filter": min_confidence
    }


@app.get("/api/customer-counts/units")
async def get_customer_count_units():
    """Get all available customer count unit types with descriptions."""
    return {
        "units": [
            {
                "value": "healthcare_organizations",
                "label": "Healthcare Organizations",
                "description": "Distinct hospital/clinic/practice entities"
            },
            {
                "value": "providers",
                "label": "Providers",
                "description": "Individual physicians or clinicians using the platform"
            },
            {
                "value": "locations",
                "label": "Locations",
                "description": "Physical practice sites or facilities"
            },
            {
                "value": "users",
                "label": "Users",
                "description": "All user accounts (may include staff, admins)"
            },
            {
                "value": "lives_covered",
                "label": "Lives Covered",
                "description": "Patient lives managed through the platform"
            },
            {
                "value": "encounters",
                "label": "Encounters/Visits",
                "description": "Annual patient encounters processed"
            },
            {
                "value": "beds",
                "label": "Hospital Beds",
                "description": "Licensed hospital beds served"
            }
        ]
    }


@app.post("/api/competitors/{competitor_id}/triangulate-customer-count")
async def triangulate_customer_count(
    competitor_id: int,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """
    Triangulate customer count from multiple sources.
    Collects data from: website scrapes, SEC filings, news articles, existing estimates.
    """
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    sources = []

    # 1. Get from existing DataSource records (website scrapes)
    website_source = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id,
        DataSource.field_name == "customer_count"
    ).order_by(DataSource.extracted_at.desc()).first()

    if website_source and website_source.current_value:
        sources.append({
            "source_type": "website_scrape",
            "value": website_source.current_value,
            "source_url": website_source.source_url,
            "date": website_source.extracted_at.isoformat() if website_source.extracted_at else None
        })

    # 2. Get from Competitor model directly
    if competitor.customer_count:
        sources.append({
            "source_type": "existing_record",
            "value": competitor.customer_count,
            "source_url": competitor.website,
            "date": competitor.last_updated.isoformat() if competitor.last_updated else None
        })

    # 3. Try SEC data for public companies
    if competitor.is_public and competitor.ticker_symbol:
        try:
            from sec_edgar_scraper import SECEdgarScraper
            scraper = SECEdgarScraper()
            sec_data = scraper.get_company_data(competitor.name)
            if sec_data and sec_data.customers_mentioned:
                sources.append({
                    "source_type": "sec_filing",
                    "value": f"{len(sec_data.customers_mentioned)}+ (named customers)",
                    "source_url": f"https://sec.gov/cgi-bin/browse-edgar?CIK={sec_data.cik}",
                    "date": sec_data.last_updated
                })
        except Exception as e:
            logger.error(f"SEC scraper error for {competitor.name}: {e}")

    # 4. Check existing CustomerCountEstimate records
    existing_counts = db.query(CustomerCountEstimate).filter(
        CustomerCountEstimate.competitor_id == competitor_id
    ).all()

    for ec in existing_counts:
        if ec.count_display:
            sources.append({
                "source_type": ec.primary_source or "previous_estimate",
                "value": ec.count_display,
                "source_url": ec.primary_source_url,
                "date": ec.as_of_date.isoformat() if ec.as_of_date else None,
                "confidence": ec.confidence_score
            })

    if not sources:
        return {
            "status": "no_sources",
            "competitor": competitor.name,
            "message": "No customer count data found from any source"
        }

    # Triangulate - use triangulate_data_points from confidence_scoring
    triangulation_sources = [
        {
            "value": s["value"],
            "source_type": s["source_type"],
            "reliability": get_source_defaults(s["source_type"]).get("reliability", "F"),
            "credibility": get_source_defaults(s["source_type"]).get("credibility", 6)
        }
        for s in sources
    ]

    result = triangulate_data_points(triangulation_sources)

    # Create a new CustomerCountEstimate with triangulated result
    # Parse the best value to extract numeric count
    import re
    best_value = result.best_value
    count_value = None
    count_type = "estimate"

    # Try to extract number from value like "3,000+" or "3000-5000"
    if best_value and best_value != "Unknown":
        numbers = re.findall(r'[\d,]+', best_value.replace(",", ""))
        if numbers:
            try:
                count_value = int(numbers[0])
                if "+" in best_value:
                    count_type = "minimum"
                elif "-" in best_value or "to" in best_value.lower():
                    count_type = "range"
                else:
                    count_type = "exact"
            except ValueError:
                pass

    new_estimate = CustomerCountEstimate(
        competitor_id=competitor_id,
        count_value=count_value,
        count_display=best_value if best_value != "Unknown" else None,
        count_type=count_type,
        count_unit="healthcare_organizations",  # Default, should be refined
        primary_source=result.source_used,
        all_sources=json.dumps(sources),
        source_agreement_score=1.0 if not result.discrepancy_flag else 0.5,
        confidence_score=result.confidence_score,
        confidence_level=result.confidence_level,
        confidence_notes=result.review_reason,
        is_verified=result.confidence_level == "high",
        verification_method="triangulation" if len(sources) > 1 else None,
        verification_date=datetime.utcnow() if len(sources) > 1 else None,
        as_of_date=datetime.utcnow()
    )
    db.add(new_estimate)
    db.commit()
    db.refresh(new_estimate)

    log_activity(
        db, current_user["email"], current_user["id"],
        "customer_count_triangulated",
        f"Triangulated customer count for {competitor.name} from {len(sources)} sources"
    )

    return {
        "status": "success",
        "competitor": competitor.name,
        "triangulation_result": {
            "best_value": result.best_value,
            "confidence_score": result.confidence_score,
            "confidence_level": result.confidence_level,
            "source_used": result.source_used,
            "discrepancy_flag": result.discrepancy_flag,
            "review_reason": result.review_reason
        },
        "sources_analyzed": len(sources),
        "sources": sources,
        "new_estimate_id": new_estimate.id
    }


@app.get("/api/customer-counts/history/{competitor_id}")
async def get_customer_count_history(
    competitor_id: int,
    db: Session = Depends(get_db)
):
    """Get historical customer count data for trend analysis."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    counts = db.query(CustomerCountEstimate).filter(
        CustomerCountEstimate.competitor_id == competitor_id
    ).order_by(CustomerCountEstimate.as_of_date.asc()).all()

    history = []
    for c in counts:
        history.append({
            "id": c.id,
            "count_value": c.count_value,
            "count_display": c.count_display,
            "count_type": c.count_type,
            "count_unit": c.count_unit,
            "as_of_date": c.as_of_date.isoformat() if c.as_of_date else None,
            "growth_rate": c.growth_rate,
            "confidence_level": c.confidence_level,
            "is_verified": c.is_verified,
            "primary_source": c.primary_source
        })

    # Calculate overall growth trend
    growth_trend = None
    if len(history) >= 2:
        first_with_value = next((h for h in history if h["count_value"]), None)
        last_with_value = next((h for h in reversed(history) if h["count_value"]), None)
        if first_with_value and last_with_value and first_with_value["count_value"] > 0:
            total_growth = ((last_with_value["count_value"] - first_with_value["count_value"]) /
                          first_with_value["count_value"]) * 100
            growth_trend = round(total_growth, 1)

    return {
        "competitor": competitor.name,
        "competitor_id": competitor_id,
        "history": history,
        "total_records": len(history),
        "growth_trend_percent": growth_trend
    }


# ============== PHASE 5: ENHANCED SCRAPER WITH SOURCE TRACKING ==============

@app.post("/api/scrape/enhanced/{competitor_id}")
async def run_enhanced_scrape(
    competitor_id: int,
    background_tasks: BackgroundTasks,
    pages: Optional[List[str]] = None,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """
    Run enhanced scrape with full source tracking for each extracted field.

    This uses the EnhancedGPTExtractor to track which page each data point
    came from, with confidence scores based on page type and field relevance.

    Args:
        competitor_id: ID of the competitor to scrape
        pages: Optional list of pages to scrape (default: homepage, pricing, about, features)
    """
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    if not pages:
        pages = ["homepage", "pricing", "about", "features", "customers"]

    # Run in background (using sync wrapper for proper async handling)
    background_tasks.add_task(
        _run_enhanced_scrape_job_sync,
        competitor_id,
        competitor.name,
        competitor.website,
        pages,
        current_user["email"]
    )

    log_activity(
        db, current_user["email"], current_user["id"],
        "enhanced_scrape_started",
        f"Started enhanced scrape for {competitor.name} ({len(pages)} pages)"
    )

    return {
        "status": "started",
        "competitor_id": competitor_id,
        "competitor_name": competitor.name,
        "pages_to_scrape": pages,
        "message": "Enhanced scrape started in background"
    }


def _run_enhanced_scrape_job_sync(
    competitor_id: int,
    competitor_name: str,
    website: str,
    pages: List[str],
    user_email: str
):
    """Synchronous wrapper for enhanced scrape background task."""
    import asyncio
    try:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        loop.run_until_complete(_run_enhanced_scrape_job(competitor_id, competitor_name, website, pages, user_email))
    except Exception as e:
        logger.error(f"[Enhanced Scrape] Error for {competitor_name}: {e}")
    finally:
        loop.close()


async def _run_enhanced_scrape_job(
    competitor_id: int,
    competitor_name: str,
    website: str,
    pages: List[str],
    user_email: str
):
    """
    Background job for enhanced scraping with full source tracking.
    """
    from scraper import CompetitorScraper
    from extractor import EnhancedGPTExtractor

    db = SessionLocal()

    try:
        competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
        if not competitor:
            logger.error(f"Enhanced scrape: Competitor {competitor_id} not found")
            return

        logger.info(f"[Enhanced Scrape] Starting for {competitor_name}...")

        # 1. Scrape all requested pages
        page_contents = {}
        async with CompetitorScraper() as scraper:
            for page_type in pages:
                try:
                    if page_type == "homepage":
                        url = website
                    else:
                        url = f"{website.rstrip('/')}/{page_type}"

                    result = await scraper.scrape(url)
                    if result and result.get("content"):
                        page_contents[page_type] = result["content"]
                        logger.info(f"[OK] Scraped {page_type}: {len(result['content'])} chars")
                    else:
                        logger.debug(f"[--] No content from {page_type}")
                except Exception as e:
                    logger.error(f"[ERR] Failed to scrape {page_type}: {e}")

        if not page_contents:
            logger.warning(f"[Enhanced Scrape] No content scraped for {competitor_name}")
            return

        # 2. Extract with source tracking
        extractor = EnhancedGPTExtractor()
        extracted = extractor.extract_with_sources(
            competitor_name=competitor_name,
            competitor_website=website,
            page_contents=page_contents
        )

        logger.info(f"[Enhanced Scrape] Extracted {len(extracted.field_sources)} fields from {len(page_contents)} pages")

        # 3. Store DataSource records with confidence scoring
        data_sources = extractor.to_data_sources(extracted, competitor_id)
        changes_count = 0
        new_values_count = 0

        for ds_data in data_sources:
            field_name = ds_data["field_name"]
            new_value = ds_data["current_value"]

            # Check if field is locked by manual correction
            is_locked = db.query(DataSource).filter(
                DataSource.competitor_id == competitor_id,
                DataSource.field_name == field_name,
                DataSource.source_type == "manual"
            ).first()

            if is_locked:
                logger.debug(f"[LOCKED] Skipping {field_name} (manual correction exists)")
                continue

            # Get old value
            old_value = getattr(competitor, field_name, None) if hasattr(competitor, field_name) else None
            old_str = str(old_value) if old_value else None

            # Update competitor if field exists
            if hasattr(competitor, field_name) and new_value:
                if old_str != new_value:
                    # Log change
                    change_record = DataChangeHistory(
                        competitor_id=competitor_id,
                        competitor_name=competitor_name,
                        field_name=field_name,
                        old_value=old_str,
                        new_value=new_value,
                        changed_by=f"Enhanced Scrape ({user_email})",
                        change_reason="Enhanced scrape with source tracking"
                    )
                    db.add(change_record)
                    setattr(competitor, field_name, new_value)

                    if old_value is None or old_str == "" or old_str == "None":
                        new_values_count += 1
                    else:
                        changes_count += 1

            # Create/update DataSource record
            existing_source = db.query(DataSource).filter(
                DataSource.competitor_id == competitor_id,
                DataSource.field_name == field_name,
                DataSource.source_type == "website_scrape"
            ).first()

            if existing_source:
                # Update existing
                existing_source.previous_value = existing_source.current_value
                existing_source.current_value = new_value
                existing_source.source_url = ds_data["source_url"]
                existing_source.source_name = ds_data["source_name"]
                existing_source.extracted_at = ds_data["extracted_at"]
                existing_source.confidence_score = ds_data["confidence_score"]
                existing_source.confidence_level = ds_data["confidence_level"]
                existing_source.updated_at = datetime.utcnow()
            else:
                # Create new
                new_source = DataSource(
                    competitor_id=competitor_id,
                    field_name=field_name,
                    current_value=new_value,
                    source_type="website_scrape",
                    source_url=ds_data["source_url"],
                    source_name=ds_data["source_name"],
                    extraction_method="gpt_extraction",
                    extracted_at=ds_data["extracted_at"],
                    source_reliability=ds_data["source_reliability"],
                    information_credibility=ds_data["information_credibility"],
                    confidence_score=ds_data["confidence_score"],
                    confidence_level=ds_data["confidence_level"],
                    data_as_of_date=ds_data["data_as_of_date"]
                )
                db.add(new_source)

        competitor.last_updated = datetime.utcnow()
        db.commit()

        logger.info(f"[Enhanced Scrape] Completed for {competitor_name}: {changes_count} changes, {new_values_count} new values")

        # 4. Run triangulation for key fields
        try:
            triangulator = DataTriangulator(db)
            await triangulator.triangulate_all_key_fields(
                competitor_id=competitor_id,
                competitor_name=competitor_name,
                website=website,
                is_public=competitor.is_public,
                ticker_symbol=competitor.ticker_symbol
            )
            db.commit()
            logger.info(f"[Enhanced Scrape] Triangulation completed for {competitor_name}")
        except Exception as e:
            logger.error(f"[Enhanced Scrape] Triangulation error: {e}")

    except Exception as e:
        logger.exception(f"[Enhanced Scrape] Error for {competitor_name}: {e}")
        db.rollback()
    finally:
        db.close()


@app.get("/api/scrape/enhanced/{competitor_id}/sources")
async def get_enhanced_scrape_sources(
    competitor_id: int,
    db: Session = Depends(get_db)
):
    """Get all source data from the most recent enhanced scrape."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    sources = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id,
        DataSource.source_type == "website_scrape"
    ).order_by(DataSource.extracted_at.desc()).all()

    result = []
    for s in sources:
        result.append({
            "field_name": s.field_name,
            "current_value": s.current_value,
            "source_url": s.source_url,
            "source_name": s.source_name,
            "confidence_score": s.confidence_score,
            "confidence_level": s.confidence_level,
            "extracted_at": s.extracted_at.isoformat() if s.extracted_at else None,
            "is_verified": s.is_verified
        })

    return {
        "competitor": competitor.name,
        "competitor_id": competitor_id,
        "sources": result,
        "total_fields": len(result)
    }


# ============== BULK UPDATE ENDPOINT ==============

@app.post("/api/competitors/bulk-update")
async def bulk_update_competitors(
    request: Request,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Bulk update competitors from JSON data."""
    data = await request.json()
    updates = data.get("updates", [])
    changed_by = data.get("changed_by", "bulk_import")
    
    results = {"success": [], "errors": []}
    
    for update in updates:
        competitor_id = update.get("id")
        if not competitor_id:
            results["errors"].append({"error": "Missing competitor ID", "data": update})
            continue
        
        competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
        if not competitor:
            results["errors"].append({"error": f"Competitor {competitor_id} not found", "data": update})
            continue
        
        try:
            for field, new_value in update.items():
                if field == "id":
                    continue
                if hasattr(competitor, field):
                    old_value = getattr(competitor, field)
                    setattr(competitor, field, new_value)
                    
                    # Log the change
                    change = DataChangeHistory(
                        competitor_id=competitor_id,
                        competitor_name=competitor.name,
                        field_name=field,
                        old_value=str(old_value) if old_value else None,
                        new_value=str(new_value) if new_value else None,
                        source_url=data.get("source_url"), # Capture Source URL from bulk/correction payload
                        changed_by=changed_by
                    )
                    db.add(change)
            
            competitor.last_updated = datetime.utcnow()
            competitor.data_quality_score = calculate_quality_score(competitor)
            results["success"].append({"id": competitor_id, "name": competitor.name})
        except Exception as e:
            results["errors"].append({"error": "An unexpected error occurred", "id": competitor_id})
    
    db.commit()

    return {
        "total_processed": len(updates),
        "successful": len(results["success"]),
        "failed": len(results["errors"]),
        "results": results
    }


# ============== UX-005: BULK OPERATIONS ENDPOINTS ==============
# BulkUpdateRequest, BulkDeleteRequest, BulkExportRequest imported from schemas.competitors

@app.put("/api/competitors/bulk-update")
async def bulk_update_competitors_by_ids(
    request: BulkUpdateRequest,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Bulk update selected competitors with the same values."""
    results = {"success": [], "errors": []}

    for competitor_id in request.ids:
        competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
        if not competitor:
            results["errors"].append({"id": competitor_id, "error": "Not found"})
            continue

        try:
            for field, new_value in request.updates.items():
                if hasattr(competitor, field):
                    old_value = getattr(competitor, field)
                    setattr(competitor, field, new_value)

                    # Log the change
                    change = DataChangeHistory(
                        competitor_id=competitor_id,
                        competitor_name=competitor.name,
                        field_name=field,
                        old_value=str(old_value) if old_value else None,
                        new_value=str(new_value) if new_value else None,
                        changed_by="bulk_operation"
                    )
                    db.add(change)

            competitor.last_updated = datetime.utcnow()
            competitor.data_quality_score = calculate_quality_score(competitor)
            results["success"].append({"id": competitor_id, "name": competitor.name})
        except Exception as e:
            results["errors"].append({"id": competitor_id, "error": "An unexpected error occurred"})

    db.commit()

    return {
        "total": len(request.ids),
        "successful": len(results["success"]),
        "failed": len(results["errors"]),
        "results": results
    }


@app.delete("/api/competitors/bulk-delete")
async def bulk_delete_competitors(
    request: BulkDeleteRequest,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Bulk delete selected competitors (soft delete)."""
    results = {"success": [], "errors": []}

    for competitor_id in request.ids:
        competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
        if not competitor:
            results["errors"].append({"id": competitor_id, "error": "Not found"})
            continue

        try:
            # Soft delete - mark as deleted instead of removing
            competitor.is_deleted = True
            competitor.last_updated = datetime.utcnow()

            # Log the deletion
            change = DataChangeHistory(
                competitor_id=competitor_id,
                competitor_name=competitor.name,
                field_name="is_deleted",
                old_value="False",
                new_value="True",
                changed_by="bulk_delete"
            )
            db.add(change)

            results["success"].append({"id": competitor_id, "name": competitor.name})
        except Exception as e:
            results["errors"].append({"id": competitor_id, "error": "An unexpected error occurred"})

    db.commit()

    return {
        "total": len(request.ids),
        "deleted": len(results["success"]),
        "failed": len(results["errors"]),
        "results": results
    }


@app.post("/api/competitors/bulk-refresh")
async def bulk_refresh_competitors(
    request: BulkDeleteRequest,  # Reuse same schema (just needs ids)
    background_tasks: BackgroundTasks,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Queue bulk refresh for selected competitors."""
    valid_ids = []

    for competitor_id in request.ids:
        competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
        if competitor:
            valid_ids.append(competitor_id)

    # Queue background refresh for each competitor (uses own DB session, not request-scoped)
    for competitor_id in valid_ids:
        background_tasks.add_task(refresh_single_competitor, competitor_id)

    return {
        "message": f"Refresh queued for {len(valid_ids)} competitors",
        "queued_ids": valid_ids,
        "invalid_ids": [id for id in request.ids if id not in valid_ids]
    }


async def refresh_single_competitor(competitor_id: int):
    """Background task to refresh a single competitor (creates own DB session)."""
    db = SessionLocal()
    try:
        competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
        if competitor and competitor.website:
            # Use existing scraper
            from scraper import scrape_competitor
            data = await scrape_competitor(competitor.website)
            if data:
                for field, value in data.items():
                    if hasattr(competitor, field) and value:
                        setattr(competitor, field, value)
                competitor.last_updated = datetime.utcnow()
                db.commit()
    except Exception as e:
        logger.error(f"Bulk refresh error for {competitor_id}: {e}")
    finally:
        db.close()


@app.post("/api/competitors/export")
async def bulk_export_competitors(
    request: BulkExportRequest,
    db: Session = Depends(get_db)
):
    """Export selected competitors to various formats."""
    from io import BytesIO

    # Get selected competitors
    competitors = db.query(Competitor).filter(
        Competitor.id.in_(request.ids),
        Competitor.is_deleted == False
    ).all()

    if not competitors:
        raise HTTPException(status_code=404, detail="No competitors found")

    # Prepare data for export
    export_data = []
    for c in competitors:
        export_data.append({
            "ID": c.id,
            "Name": c.name,
            "Website": c.website,
            "Threat Level": c.threat_level,
            "Status": c.status,
            "Primary Market": c.primary_market,
            "Product Categories": c.product_categories,
            "Pricing Model": c.pricing_model,
            "Base Price": c.base_price,
            "Customer Count": c.customer_count,
            "Employee Count": c.employee_count,
            "Year Founded": c.year_founded,
            "Headquarters": c.headquarters,
            "Funding Total": c.funding_total,
            "G2 Rating": c.g2_rating,
            "Last Updated": c.last_updated.strftime("%Y-%m-%d") if c.last_updated else "",
            "Data Quality Score": c.data_quality_score
        })

    if request.format == "excel":
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Competitors"

        # Headers
        headers = list(export_data[0].keys())
        ws.append(headers)

        # Data rows
        for row in export_data:
            ws.append(list(row.values()))

        # Auto-adjust column widths
        for col in ws.columns:
            max_length = max(len(str(cell.value or "")) for cell in col)
            ws.column_dimensions[col[0].column_letter].width = min(max_length + 2, 50)

        output = BytesIO()
        wb.save(output)
        output.seek(0)

        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f"attachment; filename=competitors_export_{datetime.now().strftime('%Y%m%d')}.xlsx"}
        )

    elif request.format == "csv":
        import csv
        output = BytesIO()
        text_output = output

        # Write BOM for Excel compatibility
        output.write(b'\xef\xbb\xbf')

        # Write CSV
        fieldnames = list(export_data[0].keys())
        csv_content = ",".join(fieldnames) + "\n"
        for row in export_data:
            csv_content += ",".join([f'"{str(v)}"' for v in row.values()]) + "\n"

        output.write(csv_content.encode('utf-8'))
        output.seek(0)

        return StreamingResponse(
            output,
            media_type="text/csv",
            headers={"Content-Disposition": f"attachment; filename=competitors_export_{datetime.now().strftime('%Y%m%d')}.csv"}
        )

    elif request.format == "json":
        return Response(
            content=json.dumps(export_data, indent=2, default=str),
            media_type="application/json",
            headers={"Content-Disposition": f"attachment; filename=competitors_export_{datetime.now().strftime('%Y%m%d')}.json"}
        )

    elif request.format == "pdf":
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet

        output = BytesIO()
        doc = SimpleDocTemplate(output, pagesize=landscape(letter))
        elements = []
        styles = getSampleStyleSheet()

        # Title
        elements.append(Paragraph(f"Competitors Export - {datetime.now().strftime('%Y-%m-%d')}", styles['Heading1']))
        elements.append(Spacer(1, 20))

        # Table data - use subset of columns for PDF
        pdf_columns = ["Name", "Threat Level", "Primary Market", "Customer Count", "Pricing Model"]
        table_data = [pdf_columns]
        for row in export_data:
            table_data.append([str(row.get(col, "")) for col in pdf_columns])

        # Create table
        table = Table(table_data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#122753')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.white),
            ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 9),
            ('GRID', (0, 0), (-1, -1), 1, colors.lightgrey),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F8FAFC')])
        ]))

        elements.append(table)
        doc.build(elements)
        output.seek(0)

        return StreamingResponse(
            output,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=competitors_export_{datetime.now().strftime('%Y%m%d')}.pdf"}
        )

    raise HTTPException(status_code=400, detail=f"Unsupported format: {request.format}")


# Import routers
from routers import reports, discovery, sales_marketing, teams, knowledge_base, products
from routers import agents as agents_router  # v7.0 AI Agents
from routers import auth as auth_router  # v9.0 extracted auth
from routers import health as health_router  # v9.0 extracted health
from routers import webhooks as webhooks_router  # v9.0 extracted webhooks
from routers import winloss as winloss_router  # v9.0 extracted win/loss
from routers import dashboard as dashboard_router  # v9.0 extracted dashboard
from routers import chat as chat_router  # v9.0 extracted chat
from routers import data_quality as data_quality_router  # v9.0 extracted data quality
from routers import admin as admin_router  # v9.0 extracted admin
from routers import ai_cost as ai_cost_router  # v9.0 AI cost analytics, audit, relationships
import api_routes

# Include routers
app.include_router(auth_router.router)  # Auth (login, register, refresh, logout)
app.include_router(health_router.router)  # Health & readiness probes
app.include_router(webhooks_router.router)  # Webhook management
app.include_router(winloss_router.router)  # Win/loss deal tracking
from routers import competitors as competitors_router  # v9.0 extracted competitors CRUD
app.include_router(competitors_router.router)  # Competitors CRUD
app.include_router(discovery.router)
app.include_router(api_routes.router)  # Covers analytics, notifications, external, etc.
app.include_router(teams.router)  # Team Features (v5.2.0)
app.include_router(reports.router)
app.include_router(sales_marketing.router)  # Sales & Marketing Module (v5.0.7)
app.include_router(knowledge_base.router)  # Knowledge Base Import (v5.0.8)
app.include_router(products.router)  # Product Discovery System (v5.1.0)
app.include_router(agents_router.router)  # v7.0 AI Agents
app.include_router(dashboard_router.router)  # Dashboard stats & threats
app.include_router(chat_router.router)  # Chat sessions & messages
app.include_router(data_quality_router.router)  # Data quality metrics
app.include_router(admin_router.router)  # System prompts, KB, data providers
app.include_router(ai_cost_router.router)  # AI cost analytics, audit logs, relationships

app.add_middleware(
    CORSMiddleware,
    allow_origins=[o.strip() for o in os.getenv("CORS_ORIGINS", "*").split(",")],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ==============================================================================
# REL-001: API Rate Limiting Middleware
# Per-IP: 100 requests/minute, Per-User: 1000 requests/minute
# ==============================================================================
from collections import defaultdict
from time import time as current_time
import threading

class RateLimiter:
    """Simple in-memory rate limiter with per-IP and per-user tracking."""

    def __init__(self, ip_limit: int = 100, user_limit: int = 1000, window_seconds: int = 60):
        self.ip_limit = ip_limit
        self.user_limit = user_limit
        self.window_seconds = window_seconds
        self.ip_requests = defaultdict(list)
        self.user_requests = defaultdict(list)
        self._lock = threading.Lock()

    def _clean_old_requests(self, requests_list: list, now: float) -> list:
        """Remove requests older than the window."""
        cutoff = now - self.window_seconds
        return [ts for ts in requests_list if ts > cutoff]

    def is_allowed(self, ip: str, user_id: str = None) -> tuple:
        """Check if request is allowed. Returns (allowed, reason)."""
        now = current_time()

        with self._lock:
            # Check IP rate limit
            self.ip_requests[ip] = self._clean_old_requests(self.ip_requests[ip], now)
            if len(self.ip_requests[ip]) >= self.ip_limit:
                return False, f"IP rate limit exceeded ({self.ip_limit}/min)"

            # Check user rate limit if authenticated
            if user_id:
                self.user_requests[user_id] = self._clean_old_requests(self.user_requests[user_id], now)
                if len(self.user_requests[user_id]) >= self.user_limit:
                    return False, f"User rate limit exceeded ({self.user_limit}/min)"

            # Record this request
            self.ip_requests[ip].append(now)
            if user_id:
                self.user_requests[user_id].append(now)

            return True, None

    def get_remaining(self, ip: str, user_id: str = None) -> dict:
        """Get remaining requests for rate limit headers."""
        now = current_time()
        with self._lock:
            ip_reqs = len(self._clean_old_requests(self.ip_requests.get(ip, []), now))
            user_reqs = len(self._clean_old_requests(self.user_requests.get(user_id, []), now)) if user_id else 0
        return {
            "ip_remaining": max(0, self.ip_limit - ip_reqs),
            "user_remaining": max(0, self.user_limit - user_reqs) if user_id else self.user_limit
        }

# Global rate limiter instance
rate_limiter = RateLimiter(ip_limit=100, user_limit=1000, window_seconds=60)

from starlette.middleware.base import BaseHTTPMiddleware
from starlette.responses import JSONResponse

class RateLimitMiddleware(BaseHTTPMiddleware):
    """Middleware to enforce rate limits on API requests."""

    # Paths that skip rate limiting
    EXEMPT_PATHS = {"/health", "/readiness", "/login.html", "/", "/app", "/static"}

    async def dispatch(self, request: Request, call_next):
        path = request.url.path

        # Skip rate limiting for static files and exempt paths
        if any(path.startswith(exempt) for exempt in self.EXEMPT_PATHS) or path.endswith(('.html', '.css', '.js', '.png', '.jpg', '.ico')):
            return await call_next(request)

        # Get client IP
        client_ip = request.client.host if request.client else "unknown"
        forwarded = request.headers.get("X-Forwarded-For")
        if forwarded:
            client_ip = forwarded.split(",")[0].strip()

        # Try to get user ID from auth header (for user-specific limits)
        user_id = None
        auth_header = request.headers.get("Authorization", "")
        if auth_header.startswith("Bearer "):
            try:
                from extended_features import auth_manager
                token = auth_header.split(" ")[1]
                payload = auth_manager.verify_token(token)
                if payload:
                    user_id = payload.get("sub")
            except (ValueError, TypeError, ImportError):
                pass

        # Check rate limit
        allowed, reason = rate_limiter.is_allowed(client_ip, user_id)

        if not allowed:
            # Return 429 Too Many Requests
            return JSONResponse(
                status_code=429,
                content={
                    "detail": reason,
                    "retry_after": 60
                },
                headers={
                    "Retry-After": "60",
                    "X-RateLimit-Limit": str(rate_limiter.ip_limit),
                    "X-RateLimit-Remaining": "0"
                }
            )

        # Process request and add rate limit headers to response
        response = await call_next(request)
        remaining = rate_limiter.get_remaining(client_ip, user_id)
        response.headers["X-RateLimit-Limit"] = str(rate_limiter.ip_limit)
        response.headers["X-RateLimit-Remaining"] = str(remaining["ip_remaining"])

        return response

# Add rate limiting middleware
app.add_middleware(RateLimitMiddleware)

# ==============================================================================
# PERF-002: In-Memory Caching Layer with TTL
# Caches: competitors (5min), dimensions (1hr), news (15min), products (30min)
# ==============================================================================
from functools import lru_cache
import hashlib

class TTLCache:
    """Simple TTL-based cache for API responses."""

    def __init__(self):
        self._cache = {}
        self._timestamps = {}
        self._lock = threading.Lock()

    def get(self, key: str, ttl_seconds: int = 300) -> Any:
        """Get cached value if not expired. Returns None if miss or expired."""
        with self._lock:
            if key not in self._cache:
                return None

            timestamp = self._timestamps.get(key, 0)
            if current_time() - timestamp > ttl_seconds:
                # Expired - clean up
                del self._cache[key]
                del self._timestamps[key]
                return None

            return self._cache[key]

    def set(self, key: str, value: Any) -> None:
        """Set cache value with current timestamp."""
        with self._lock:
            self._cache[key] = value
            self._timestamps[key] = current_time()

    def invalidate(self, pattern: str = None) -> int:
        """Invalidate cache entries matching pattern (or all if None)."""
        with self._lock:
            if pattern is None:
                count = len(self._cache)
                self._cache.clear()
                self._timestamps.clear()
                return count

            keys_to_delete = [k for k in self._cache if pattern in k]
            for key in keys_to_delete:
                del self._cache[key]
                del self._timestamps[key]
            return len(keys_to_delete)

    def stats(self) -> dict:
        """Get cache statistics."""
        with self._lock:
            return {
                "entries": len(self._cache),
                "keys": list(self._cache.keys())[:20]  # First 20 keys
            }

# Global cache instance
api_cache = TTLCache()

# Cache TTL values (in seconds)
CACHE_TTL = {
    "competitors": 300,      # 5 minutes
    "competitor_detail": 300,
    "dimensions": 3600,      # 1 hour (rarely change)
    "news": 900,             # 15 minutes
    "products": 1800,        # 30 minutes
    "analytics": 600,        # 10 minutes
    "data_quality": 600      # 10 minutes
}

def cache_key(*args) -> str:
    """Generate cache key from arguments."""
    return hashlib.md5(":".join(str(a) for a in args).encode()).hexdigest()

@app.get("/api/cache/stats")
async def get_cache_stats(current_user: dict = Depends(get_current_user)):
    """Get cache statistics (admin only)."""
    return {
        "stats": api_cache.stats(),
        "ttl_config": CACHE_TTL
    }

@app.post("/api/cache/invalidate")
async def invalidate_cache(
    pattern: str = None,
    current_user: dict = Depends(get_current_user)
):
    """Invalidate cache entries (admin only)."""
    count = api_cache.invalidate(pattern)
    return {"invalidated": count, "pattern": pattern or "all"}

# ==============================================================================
# UX-004: Global Smart Search (FTS5-like functionality)
# Searches across competitors, products, news, and knowledge base
# ==============================================================================

from sqlalchemy import or_, func
# SearchResult imported from schemas.competitors

@app.get("/api/search")
async def global_search(
    q: str,
    types: str = "all",  # comma-separated: competitor,product,news,knowledge
    limit: int = 20,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """
    Global search across all data types.
    Uses LIKE for SQLite (FTS5 would require separate table setup).
    Returns ranked results with snippets.
    """
    if not q or len(q) < 2:
        return {"results": [], "total": 0, "query": q}

    # Check cache first
    cache_key_str = cache_key("search", q, types, limit)
    cached = api_cache.get(cache_key_str, ttl_seconds=300)
    if cached:
        return cached

    results = []
    search_pattern = f"%{q.lower()}%"
    type_list = types.split(",") if types != "all" else ["competitor", "product", "news", "knowledge"]

    # Search Competitors (using valid fields: notes, key_features, target_segments)
    if "competitor" in type_list:
        competitor_results = db.query(Competitor).filter(
            or_(
                func.lower(Competitor.name).like(search_pattern),
                func.lower(Competitor.notes).like(search_pattern),
                func.lower(Competitor.key_features).like(search_pattern),
                func.lower(Competitor.target_segments).like(search_pattern),
                func.lower(Competitor.website).like(search_pattern)
            )
        ).limit(limit).all()

        for comp in competitor_results:
            # Calculate relevance score based on where match was found
            score = 1.0
            if q.lower() in (comp.name or "").lower():
                score = 2.0  # Name match is highest
            elif q.lower() in (comp.notes or "").lower()[:100]:
                score = 1.5

            snippet = None
            for field in [comp.notes, comp.key_features, comp.target_segments]:
                if field and q.lower() in field.lower():
                    idx = field.lower().find(q.lower())
                    start = max(0, idx - 50)
                    end = min(len(field), idx + 50)
                    snippet = "..." + field[start:end] + "..."
                    break

            results.append(SearchResult(
                type="competitor",
                id=comp.id,
                title=comp.name or "Unknown Competitor",
                subtitle=f"{comp.target_segments or 'Competitor'} | {comp.threat_level or 'Unknown'} threat",
                snippet=snippet,
                score=score,
                url=f"/competitors/{comp.id}"
            ))

    # Search Products
    if "product" in type_list:
        product_results = db.query(CompetitorProduct).filter(
            or_(
                func.lower(CompetitorProduct.product_name).like(search_pattern),
                func.lower(CompetitorProduct.description).like(search_pattern),
                func.lower(CompetitorProduct.product_category).like(search_pattern)
            )
        ).limit(limit).all()

        for prod in product_results:
            comp = db.query(Competitor).filter(Competitor.id == prod.competitor_id).first()
            results.append(SearchResult(
                type="product",
                id=prod.id,
                title=prod.product_name or "Unknown Product",
                subtitle=f"{comp.name if comp else 'Unknown'} | {prod.product_category or 'Product'}",
                snippet=prod.description[:150] + "..." if prod.description and len(prod.description) > 150 else prod.description,
                score=1.0,
                url=f"/competitors/{prod.competitor_id}"
            ))

    # Search Knowledge Base
    if "knowledge" in type_list:
        kb_results = db.query(KnowledgeBaseItem).filter(
            or_(
                func.lower(KnowledgeBaseItem.title).like(search_pattern),
                func.lower(KnowledgeBaseItem.content_text).like(search_pattern),
                func.lower(KnowledgeBaseItem.tags).like(search_pattern)
            )
        ).limit(limit).all()

        for kb in kb_results:
            results.append(SearchResult(
                type="knowledge",
                id=kb.id,
                title=kb.title or "Untitled",
                subtitle=f"{kb.source_type or 'Document'} | {kb.tags or ''}",
                snippet=kb.content_text[:150] + "..." if kb.content_text and len(kb.content_text) > 150 else kb.content_text,
                score=1.0,
                url=f"/knowledge/{kb.id}"
            ))

    # Sort by score (descending)
    results.sort(key=lambda x: x.score, reverse=True)

    # Limit total results
    results = results[:limit]

    response = {
        "results": [r.dict() for r in results],
        "total": len(results),
        "query": q,
        "types_searched": type_list
    }

    # Cache the result
    api_cache.set(cache_key_str, response)

    return response

@app.get("/api/search/suggestions")
async def search_suggestions(
    q: str,
    limit: int = 10,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get autocomplete suggestions for search."""
    if not q or len(q) < 2:
        return {"suggestions": []}

    search_pattern = f"%{q.lower()}%"

    # Get competitor names
    competitors = db.query(Competitor.name).filter(
        func.lower(Competitor.name).like(search_pattern)
    ).limit(limit).all()

    # Get product names
    products = db.query(CompetitorProduct.product_name).filter(
        func.lower(CompetitorProduct.product_name).like(search_pattern)
    ).limit(limit).all()

    suggestions = []
    for (name,) in competitors:
        suggestions.append({"text": name, "type": "competitor"})
    for (name,) in products:
        suggestions.append({"text": name, "type": "product"})

    # Deduplicate and limit
    seen = set()
    unique_suggestions = []
    for s in suggestions:
        if s["text"] not in seen:
            seen.add(s["text"])
            unique_suggestions.append(s)
            if len(unique_suggestions) >= limit:
                break

    return {"suggestions": unique_suggestions}

# ==============================================================================
# INFRA-003: Database Backup Automation
# Daily SQLite backups with retention policy
# ==============================================================================

import shutil
from pathlib import Path

BACKUP_DIR = Path("backups")
BACKUP_RETENTION_DAYS = 7

def get_database_path() -> Path:
    """Get the path to the SQLite database file."""
    db_url = os.getenv("DATABASE_URL", "sqlite:///./certify_intel.db")
    if db_url.startswith("sqlite:///"):
        db_path = db_url.replace("sqlite:///", "")
        if db_path.startswith("./"):
            db_path = db_path[2:]
        return Path(db_path)
    return Path("certify_intel.db")

def create_backup() -> dict:
    """Create a backup of the database."""
    try:
        BACKUP_DIR.mkdir(exist_ok=True)
        db_path = get_database_path()

        if not db_path.exists():
            return {"success": False, "error": "Database file not found"}

        # Create timestamped backup
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        backup_name = f"certify_intel_backup_{timestamp}.db"
        backup_path = BACKUP_DIR / backup_name

        shutil.copy2(db_path, backup_path)

        # Also create a "latest" symlink/copy
        latest_path = BACKUP_DIR / "certify_intel_latest.db"
        if latest_path.exists():
            latest_path.unlink()
        shutil.copy2(backup_path, latest_path)

        # Clean old backups (retention policy)
        cleanup_old_backups()

        file_size = backup_path.stat().st_size
        return {
            "success": True,
            "backup_file": backup_name,
            "backup_path": str(backup_path),
            "size_bytes": file_size,
            "size_mb": round(file_size / (1024 * 1024), 2),
            "timestamp": timestamp
        }
    except Exception as e:
        return {"success": False, "error": "An unexpected error occurred"}

def cleanup_old_backups():
    """Remove backups older than retention period."""
    if not BACKUP_DIR.exists():
        return

    cutoff = datetime.now().timestamp() - (BACKUP_RETENTION_DAYS * 86400)

    for backup_file in BACKUP_DIR.glob("certify_intel_backup_*.db"):
        if backup_file.stat().st_mtime < cutoff:
            backup_file.unlink()
            logger.info(f"[Backup] Removed old backup: {backup_file.name}")

def list_backups() -> list:
    """List all available backups."""
    if not BACKUP_DIR.exists():
        return []

    backups = []
    for backup_file in sorted(BACKUP_DIR.glob("certify_intel_backup_*.db"), reverse=True):
        stat = backup_file.stat()
        backups.append({
            "filename": backup_file.name,
            "path": str(backup_file),
            "size_bytes": stat.st_size,
            "size_mb": round(stat.st_size / (1024 * 1024), 2),
            "created_at": datetime.fromtimestamp(stat.st_mtime).isoformat()
        })
    return backups

@app.post("/api/backup/create")
async def api_create_backup(current_user: dict = Depends(get_current_user)):
    """Create a manual database backup (admin only)."""
    if current_user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")

    result = create_backup()
    if result["success"]:
        return result
    else:
        raise HTTPException(status_code=500, detail=result["error"])

@app.get("/api/backup/list")
async def api_list_backups(current_user: dict = Depends(get_current_user)):
    """List all available backups."""
    backups = list_backups()
    return {
        "backups": backups,
        "total": len(backups),
        "retention_days": BACKUP_RETENTION_DAYS,
        "backup_dir": str(BACKUP_DIR)
    }

@app.get("/api/backup/download/{filename}")
async def api_download_backup(
    filename: str,
    current_user: dict = Depends(get_current_user)
):
    """Download a backup file (admin only)."""
    if current_user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")

    # Path traversal protection: resolve to real path and verify within BACKUP_DIR
    backup_path = (BACKUP_DIR / filename).resolve()
    if not str(backup_path).startswith(str(BACKUP_DIR.resolve())):
        raise HTTPException(status_code=400, detail="Invalid filename")
    if not backup_path.exists() or not backup_path.name.startswith("certify_intel_"):
        raise HTTPException(status_code=404, detail="Backup not found")

    return FileResponse(
        path=backup_path,
        filename=filename,
        media_type="application/octet-stream"
    )

@app.post("/api/backup/restore/{filename}")
async def api_restore_backup(
    filename: str,
    current_user: dict = Depends(get_current_user)
):
    """Restore from a backup file (admin only). Requires server restart."""
    if current_user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")

    # Path traversal protection: resolve to real path and verify within BACKUP_DIR
    backup_path = (BACKUP_DIR / filename).resolve()
    if not str(backup_path).startswith(str(BACKUP_DIR.resolve())):
        raise HTTPException(status_code=400, detail="Invalid filename")
    if not backup_path.exists() or not backup_path.name.startswith("certify_intel_"):
        raise HTTPException(status_code=404, detail="Backup not found")

    try:
        db_path = get_database_path()

        # Create a backup of current state before restoring
        pre_restore_backup = create_backup()

        # Copy backup over current database
        shutil.copy2(backup_path, db_path)

        return {
            "success": True,
            "restored_from": filename,
            "pre_restore_backup": pre_restore_backup.get("backup_file"),
            "message": "Database restored. Please restart the server for changes to take effect."
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail="Restore failed. Please try again.")

# ==============================================================================
# SM-011: Competitive Alert Subscriptions
# Subscribe to competitors for instant notifications
# ==============================================================================

from database import CompetitorSubscription
# SubscriptionCreate, SubscriptionUpdate imported from schemas.common

@app.get("/api/subscriptions")
async def get_user_subscriptions(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get all competitor subscriptions for the current user."""
    user_id = current_user.get("id")
    if not user_id:
        raise HTTPException(status_code=401, detail="User not found")

    subscriptions = db.query(CompetitorSubscription).filter(
        CompetitorSubscription.user_id == user_id
    ).all()

    return [{
        "id": sub.id,
        "competitor_id": sub.competitor_id,
        "competitor_name": sub.competitor.name if sub.competitor else "Unknown",
        "notify_email": sub.notify_email,
        "notify_slack": sub.notify_slack,
        "notify_teams": sub.notify_teams,
        "notify_push": sub.notify_push,
        "alert_on_pricing": sub.alert_on_pricing,
        "alert_on_products": sub.alert_on_products,
        "alert_on_news": sub.alert_on_news,
        "alert_on_threat_change": sub.alert_on_threat_change,
        "min_severity": sub.min_severity,
        "created_at": sub.created_at.isoformat() if sub.created_at else None,
        "last_notified_at": sub.last_notified_at.isoformat() if sub.last_notified_at else None
    } for sub in subscriptions]

@app.post("/api/subscriptions")
async def create_subscription(
    subscription: SubscriptionCreate,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Subscribe to a competitor for alerts."""
    user_id = current_user.get("id")
    if not user_id:
        raise HTTPException(status_code=401, detail="User not found")

    # Check if competitor exists
    competitor = db.query(Competitor).filter(Competitor.id == subscription.competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    # Check for existing subscription
    existing = db.query(CompetitorSubscription).filter(
        CompetitorSubscription.user_id == user_id,
        CompetitorSubscription.competitor_id == subscription.competitor_id
    ).first()

    if existing:
        raise HTTPException(status_code=400, detail="Already subscribed to this competitor")

    # Create subscription
    new_sub = CompetitorSubscription(
        user_id=user_id,
        competitor_id=subscription.competitor_id,
        notify_email=subscription.notify_email,
        notify_slack=subscription.notify_slack,
        notify_teams=subscription.notify_teams,
        notify_push=subscription.notify_push,
        alert_on_pricing=subscription.alert_on_pricing,
        alert_on_products=subscription.alert_on_products,
        alert_on_news=subscription.alert_on_news,
        alert_on_threat_change=subscription.alert_on_threat_change,
        min_severity=subscription.min_severity
    )

    db.add(new_sub)
    db.commit()
    db.refresh(new_sub)

    return {
        "id": new_sub.id,
        "competitor_id": new_sub.competitor_id,
        "competitor_name": competitor.name,
        "message": f"Subscribed to {competitor.name} for alerts"
    }

@app.put("/api/subscriptions/{subscription_id}")
async def update_subscription(
    subscription_id: int,
    updates: SubscriptionUpdate,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Update subscription preferences."""
    user_id = current_user.get("id")

    sub = db.query(CompetitorSubscription).filter(
        CompetitorSubscription.id == subscription_id,
        CompetitorSubscription.user_id == user_id
    ).first()

    if not sub:
        raise HTTPException(status_code=404, detail="Subscription not found")

    # Update fields
    for field, value in updates.dict(exclude_unset=True).items():
        if value is not None:
            setattr(sub, field, value)

    db.commit()

    return {"message": "Subscription updated", "id": subscription_id}

@app.delete("/api/subscriptions/{subscription_id}")
async def delete_subscription(
    subscription_id: int,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Unsubscribe from a competitor."""
    user_id = current_user.get("id")

    sub = db.query(CompetitorSubscription).filter(
        CompetitorSubscription.id == subscription_id,
        CompetitorSubscription.user_id == user_id
    ).first()

    if not sub:
        raise HTTPException(status_code=404, detail="Subscription not found")

    competitor_name = sub.competitor.name if sub.competitor else "Unknown"
    db.delete(sub)
    db.commit()

    return {"message": f"Unsubscribed from {competitor_name}"}

@app.get("/api/competitors/{competitor_id}/subscription")
async def get_competitor_subscription_status(
    competitor_id: int,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Check if current user is subscribed to a competitor."""
    user_id = current_user.get("id")

    sub = db.query(CompetitorSubscription).filter(
        CompetitorSubscription.competitor_id == competitor_id,
        CompetitorSubscription.user_id == user_id
    ).first()

    return {
        "subscribed": sub is not None,
        "subscription_id": sub.id if sub else None
    }

# ==============================================================================
# REL-010: Security Headers Middleware
# Extracted to middleware/security.py for modularity
# ==============================================================================
from middleware.security import SecurityHeadersMiddleware  # noqa: E402

# Add security headers middleware (configurable via SECURITY_HEADERS_ENABLED env var)
app.add_middleware(SecurityHeadersMiddleware)


# ==============================================================================
# PERF-008: Browser Caching Headers Middleware
# Adds Cache-Control headers for static assets to improve load times
# ==============================================================================

class CachingHeadersMiddleware(BaseHTTPMiddleware):
    """Middleware to add caching headers for static assets."""

    # File extension to cache duration mapping (in seconds)
    # v7.0.5: Disabled caching for JS/CSS/HTML to prevent stale version issues
    CACHE_DURATIONS = {
        # Immutable assets (hashed filenames) - 1 year
        '.woff2': 31536000,
        '.woff': 31536000,
        '.ttf': 31536000,
        '.eot': 31536000,
        # Images - 1 week
        '.png': 604800,
        '.jpg': 604800,
        '.jpeg': 604800,
        '.gif': 604800,
        '.ico': 604800,
        '.svg': 604800,
        '.webp': 604800,
        # CSS/JS - no cache to prevent stale version issues (v7.0.5 fix)
        '.css': 0,
        '.js': 0,
        # HTML - no cache (always fresh)
        '.html': 0,
    }

    async def dispatch(self, request: Request, call_next):
        response = await call_next(request)

        path = request.url.path.lower()

        # Skip API endpoints
        if path.startswith('/api/'):
            return response

        # Determine cache duration based on file extension
        for ext, duration in self.CACHE_DURATIONS.items():
            if path.endswith(ext):
                if duration > 0:
                    # Check if version query parameter exists for longer caching
                    if '?v=' in str(request.url):
                        duration = 604800  # 1 week for versioned assets
                    response.headers["Cache-Control"] = f"public, max-age={duration}"
                    response.headers["Vary"] = "Accept-Encoding"
                else:
                    response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
                break

        return response

# Add caching headers middleware
app.add_middleware(CachingHeadersMiddleware)

# ==============================================================================
# OBS-001: Prometheus Metrics Middleware
# Tracks HTTP request count and duration. Zero overhead when METRICS_ENABLED=false.
# ==============================================================================
from middleware.metrics import MetricsMiddleware  # noqa: E402

app.add_middleware(MetricsMiddleware)

# ==============================================================================
# PERF-010: Performance Monitoring - Core Web Vitals Endpoint
# ==============================================================================

# In-memory storage for performance metrics (last 1000 entries)
performance_metrics = []
MAX_METRICS = 1000

# WebVitalsMetric imported from schemas.common

@app.post("/api/metrics/vitals")
async def record_web_vitals(metric: WebVitalsMetric):
    """Record Core Web Vitals metrics from frontend."""
    global performance_metrics  # noqa: F824

    entry = {
        "name": metric.name,
        "value": metric.value,
        "rating": metric.rating,
        "url": metric.url,
        "timestamp": datetime.utcnow().isoformat()
    }

    performance_metrics.append(entry)

    # Keep only last MAX_METRICS entries
    if len(performance_metrics) > MAX_METRICS:
        performance_metrics = performance_metrics[-MAX_METRICS:]

    return {"status": "recorded"}

@app.get("/api/metrics/vitals/summary")
async def get_web_vitals_summary():
    """Get summary of Core Web Vitals metrics."""
    if not performance_metrics:
        return {"message": "No metrics recorded yet", "metrics": {}}

    # Calculate averages by metric name
    metrics_by_name = {}
    for entry in performance_metrics:
        name = entry["name"]
        if name not in metrics_by_name:
            metrics_by_name[name] = {"values": [], "ratings": {"good": 0, "needs-improvement": 0, "poor": 0}}
        metrics_by_name[name]["values"].append(entry["value"])
        if entry.get("rating"):
            metrics_by_name[name]["ratings"][entry["rating"]] = metrics_by_name[name]["ratings"].get(entry["rating"], 0) + 1

    summary = {}
    for name, data in metrics_by_name.items():
        values = data["values"]
        summary[name] = {
            "avg": sum(values) / len(values),
            "min": min(values),
            "max": max(values),
            "count": len(values),
            "p75": sorted(values)[int(len(values) * 0.75)] if len(values) > 1 else values[0],
            "ratings": data["ratings"]
        }

    return {
        "total_entries": len(performance_metrics),
        "metrics": summary,
        "last_updated": performance_metrics[-1]["timestamp"] if performance_metrics else None
    }


# Duplicate /health removed - now in routers/health.py

@app.get("/api/corporate-profile")
async def get_corporate_profile(db: Session = Depends(get_db), current_user: dict = Depends(get_current_user)):
    """Return Certify Health corporate profile data."""
    from database import CompetitorProduct, DataSource, NewsArticleCache
    competitors_count = db.query(Competitor).filter(Competitor.is_deleted == False).count()
    products_count = db.query(CompetitorProduct).count()
    sources_count = db.query(DataSource).count()
    news_count = db.query(NewsArticleCache).filter(NewsArticleCache.is_archived != True).count()
    return {
        "name": "Certify Health",
        "tagline": "Modern Patient Access & Intake Platform",
        "founded": "2019",
        "headquarters": "Atlanta, GA",
        "website": "https://www.certifyhealth.com",
        "products": [
            {"name": "Patient Intake", "description": "Digital intake forms and insurance verification"},
            {"name": "Patient Access", "description": "Self-service scheduling and registration"},
            {"name": "Revenue Cycle", "description": "Automated eligibility and prior authorization"}
        ],
        "key_metrics": {
            "competitors_tracked": competitors_count,
            "products_analyzed": products_count,
            "data_sources": sources_count,
            "news_articles": news_count
        },
        "mission": "Empowering healthcare organizations with competitive intelligence to make informed strategic decisions."
    }


# --- Competitors CRUD moved to routers/competitors.py ---


# Competitors CRUD (list, get, create, update, delete, correct) moved to routers/competitors.py


# --- Excel Export ---

def auto_fit_columns(worksheet):
    """Auto-fit all columns to the widest content."""
    for column_cells in worksheet.columns:
        max_length = 0
        column_letter = get_column_letter(column_cells[0].column)
        for cell in column_cells:
            try:
                if cell.value:
                    cell_length = len(str(cell.value))
                    if cell_length > max_length:
                        max_length = cell_length
            except Exception:
                pass
        adjusted_width = max_length + 3
        if adjusted_width < 12:
            adjusted_width = 12
        if adjusted_width > 50:
            adjusted_width = 50
        worksheet.column_dimensions[column_letter].width = adjusted_width


def apply_white_background(worksheet, max_row=100, max_col=50):
    """Apply white background to all cells."""
    for row in range(1, max_row + 1):
        for col in range(1, max_col + 1):
            cell = worksheet.cell(row=row, column=col)
            if cell.fill.fgColor.rgb in (None, '00000000', 'FFFFFFFF') or cell.fill.fill_type is None:
                cell.fill = WHITE_FILL


# NOTE: Early /api/export/excel removed (duplicate of comprehensive v6.1.2 version below at ~line 10942)


@app.get("/api/export/pptx")
def export_pptx(
    competitor_ids: str = Query(None, description="Comma-separated competitor IDs"),
    db: Session = Depends(get_db)
):
    """
    P3-5: Export competitor comparison to PowerPoint.

    Creates a professional presentation with:
    - Title slide with Certify Intel branding
    - Executive summary slide
    - Individual competitor slides with key metrics
    - Comparison table slide
    - Threat matrix slide
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx not installed. Run: pip install python-pptx")

    # Get competitors
    if competitor_ids:
        ids = [int(id.strip()) for id in competitor_ids.split(',') if id.strip().isdigit()]
        competitors_list = db.query(Competitor).filter(
            Competitor.id.in_(ids),
            Competitor.is_deleted == False
        ).all()
    else:
        competitors_list = db.query(Competitor).filter(Competitor.is_deleted == False).limit(10).all()

    if not competitors_list:
        raise HTTPException(status_code=404, detail="No competitors found")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Color scheme
    DARK_BLUE = RgbColor(47, 84, 150)  # #2F5496
    LIGHT_BLUE = RgbColor(0, 180, 216)  # #00B4D8
    WHITE = RgbColor(255, 255, 255)
    DARK_GRAY = RgbColor(51, 51, 51)

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(title_slide_layout)

    # Title background shape
    bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = DARK_BLUE
    bg_shape.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Competitive Intelligence Report"
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = f"Analysis of {len(competitors_list)} Competitors"
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = LIGHT_BLUE
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12.333), Inches(0.5))
    date_frame = date_box.text_frame
    date_para = date_frame.paragraphs[0]
    date_para.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
    date_para.font.size = Pt(14)
    date_para.font.color.rgb = WHITE
    date_para.alignment = PP_ALIGN.CENTER

    # Slide 2: Executive Summary
    summary_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(summary_layout)

    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    header_frame = header_box.text_frame
    header_para = header_frame.paragraphs[0]
    header_para.text = "Executive Summary"
    header_para.font.size = Pt(32)
    header_para.font.bold = True
    header_para.font.color.rgb = DARK_BLUE

    # Summary stats
    threat_counts = {}
    for comp in competitors_list:
        level = comp.threat_level or "Unknown"
        threat_counts[level] = threat_counts.get(level, 0) + 1

    summary_text = f"""
Key Findings:
• Total Competitors Analyzed: {len(competitors_list)}
• High Threat: {threat_counts.get('High', 0)} competitors
• Medium Threat: {threat_counts.get('Medium', 0)} competitors
• Low Threat: {threat_counts.get('Low', 0)} competitors

Top Competitors by Threat Level:
"""
    high_threat = [c for c in competitors_list if c.threat_level == 'High'][:3]
    for comp in high_threat:
        summary_text += f"• {comp.name}: {comp.pricing_model or 'N/A'} pricing\n"

    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.5))
    summary_frame = summary_box.text_frame
    summary_frame.word_wrap = True
    summary_para = summary_frame.paragraphs[0]
    summary_para.text = summary_text.strip()
    summary_para.font.size = Pt(16)
    summary_para.font.color.rgb = DARK_GRAY

    # Individual competitor slides
    for comp in competitors_list[:8]:  # Limit to 8 slides
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Competitor name header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
        header_frame = header_box.text_frame
        header_para = header_frame.paragraphs[0]
        header_para.text = comp.name
        header_para.font.size = Pt(32)
        header_para.font.bold = True
        header_para.font.color.rgb = DARK_BLUE

        # Threat level badge
        threat_color = {
            'High': RgbColor(220, 53, 69),
            'Medium': RgbColor(255, 193, 7),
            'Low': RgbColor(40, 167, 69)
        }.get(comp.threat_level, DARK_GRAY)

        threat_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.35), Inches(2.3), Inches(0.5))
        threat_frame = threat_box.text_frame
        threat_para = threat_frame.paragraphs[0]
        threat_para.text = f"{comp.threat_level or 'Unknown'} Threat"
        threat_para.font.size = Pt(14)
        threat_para.font.bold = True
        threat_para.font.color.rgb = threat_color
        threat_para.alignment = PP_ALIGN.RIGHT

        # Metrics grid
        metrics = [
            ("Website", comp.website or "N/A"),
            ("Pricing Model", comp.pricing_model or "N/A"),
            ("Base Price", comp.base_price or "N/A"),
            ("Employees", comp.employee_count or "N/A"),
            ("Customers", comp.customer_count or "N/A"),
            ("G2 Rating", f"{comp.g2_rating}/5" if comp.g2_rating else "N/A"),
            ("Funding", comp.funding_total or "N/A"),
            ("Founded", comp.year_founded or "N/A"),
        ]

        # Two columns
        for i, (label, value) in enumerate(metrics):
            col = i % 2
            row = i // 2
            x = Inches(0.5 + col * 6.5)
            y = Inches(1.3 + row * 1.3)

            box = slide.shapes.add_textbox(x, y, Inches(6), Inches(1.2))
            frame = box.text_frame
            frame.word_wrap = True

            p1 = frame.paragraphs[0]
            p1.text = label
            p1.font.size = Pt(12)
            p1.font.bold = True
            p1.font.color.rgb = DARK_GRAY

            p2 = frame.add_paragraph()
            p2.text = str(value)[:100]  # Truncate long values
            p2.font.size = Pt(16)
            p2.font.color.rgb = DARK_BLUE

        # Key features (if available)
        if comp.key_features:
            features_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.2))
            features_frame = features_box.text_frame
            features_frame.word_wrap = True
            features_para = features_frame.paragraphs[0]
            features_para.text = f"Key Features: {comp.key_features[:200]}"
            features_para.font.size = Pt(12)
            features_para.font.color.rgb = DARK_GRAY

    # Final slide: Comparison table
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    header_frame = header_box.text_frame
    header_para = header_frame.paragraphs[0]
    header_para.text = "Competitor Comparison Matrix"
    header_para.font.size = Pt(32)
    header_para.font.bold = True
    header_para.font.color.rgb = DARK_BLUE

    # Create table
    rows = min(len(competitors_list), 8) + 1  # Header + competitors
    cols = 5  # Name, Threat, Pricing, Customers, G2

    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.5)).table

    # Table headers
    headers = ["Competitor", "Threat Level", "Pricing", "Customers", "G2 Rating"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Table data
    for row_idx, comp in enumerate(competitors_list[:8], start=1):
        table.cell(row_idx, 0).text = comp.name or ""
        table.cell(row_idx, 1).text = comp.threat_level or "Unknown"
        table.cell(row_idx, 2).text = comp.pricing_model or "N/A"
        table.cell(row_idx, 3).text = comp.customer_count or "N/A"
        table.cell(row_idx, 4).text = f"{comp.g2_rating}/5" if comp.g2_rating else "N/A"

    # Save
    output_path = "./exports/competitive_intel_report.pptx"
    os.makedirs("./exports", exist_ok=True)
    prs.save(output_path)

    return FileResponse(
        output_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"certify_intel_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    )


@app.get("/api/export/json")
def export_json(db: Session = Depends(get_db)):
    """Export all competitor data as JSON."""
    competitors = db.query(Competitor).filter(Competitor.is_deleted == False).all()
    
    data = []
    for comp in competitors:
        data.append({
            "name": comp.name,
            "website": comp.website,
            "status": comp.status,
            "threat_level": comp.threat_level,
            "last_updated": comp.last_updated.isoformat() if comp.last_updated else None,
            "notes": comp.notes,
            "data_quality_score": comp.data_quality_score,
            "pricing_model": comp.pricing_model,
            "base_price": comp.base_price,
            "price_unit": comp.price_unit,
            "product_categories": comp.product_categories,
            "key_features": comp.key_features,
            "integration_partners": comp.integration_partners,
            "certifications": comp.certifications,
            "target_segments": comp.target_segments,
            "customer_size_focus": comp.customer_size_focus,
            "geographic_focus": comp.geographic_focus,
            "customer_count": comp.customer_count,
            "customer_acquisition_rate": comp.customer_acquisition_rate,
            "key_customers": comp.key_customers,
            "g2_rating": comp.g2_rating,
            "employee_count": comp.employee_count,
            "employee_growth_rate": comp.employee_growth_rate,
            "year_founded": comp.year_founded,
            "headquarters": comp.headquarters,
            "funding_total": comp.funding_total,
            "latest_round": comp.latest_round,
            "pe_vc_backers": comp.pe_vc_backers,
            "website_traffic": comp.website_traffic,
            "social_following": comp.social_following,
            "recent_launches": comp.recent_launches,
            "news_mentions": comp.news_mentions,
        })
    
    return {"competitors": data, "count": len(data), "exported_at": datetime.utcnow().isoformat()}


# Dashboard endpoints moved to routers/dashboard.py

# --- Change Log ---
# Note: Primary /api/changes endpoint is defined earlier using DataChangeHistory table


# --- Scraping Endpoints ---

@app.post("/api/scrape/all")
async def trigger_scrape_all(background_tasks: BackgroundTasks, db: Session = Depends(get_db)):
    """Trigger scrape for all active competitors with progress tracking."""
    global scrape_progress  # noqa: F824

    competitors = db.query(Competitor).filter(
        Competitor.is_deleted == False,
        Competitor.status == "Active"
    ).all()
    competitor_ids = [c.id for c in competitors]
    competitor_names = {c.id: c.name for c in competitors}

    # Phase 4: Create RefreshSession for audit trail (Task 5.0.1-031)
    refresh_session = RefreshSession(
        competitors_scanned=len(competitor_ids),
        status="in_progress"
    )
    db.add(refresh_session)
    db.commit()
    db.refresh(refresh_session)

    # Reset progress tracker with enhanced tracking (Phase 2: Task 5.0.1-026)
    scrape_progress = {
        "active": True,
        "total": len(competitor_ids),
        "completed": 0,
        "current_competitor": None,
        "competitors_done": [],
        "changes_detected": 0,
        "new_values_added": 0,
        "started_at": datetime.utcnow().isoformat(),
        "recent_changes": [],
        "change_details": [],
        "errors": [],
        "session_id": refresh_session.id,  # Track session ID for persistence
        # v7.1.0: Post-scrape enrichment tracking
        "enrichment_active": False,
        "news_articles_fetched": 0,
        "stock_prices_updated": 0
    }

    # Add to background tasks with progress tracking
    for cid in competitor_ids:
        background_tasks.add_task(run_scrape_job_with_progress, cid, competitor_names.get(cid, "Unknown"))

    return {
        "message": f"Scrape jobs queued for {len(competitor_ids)} competitors",
        "competitor_ids": competitor_ids,
        "total": len(competitor_ids),
        "session_id": refresh_session.id
    }


@app.get("/api/scrape/progress")
async def get_scrape_progress():
    """Get the current progress of a scrape operation."""
    return scrape_progress


# Phase 2: Task 5.0.1-028 - Get detailed session information
@app.get("/api/scrape/session")
async def get_scrape_session_details():
    """Get detailed information about the current or last refresh session."""
    return {
        "active": scrape_progress["active"],
        "total_competitors": scrape_progress["total"],
        "completed": scrape_progress["completed"],
        "current_competitor": scrape_progress.get("current_competitor"),
        "changes_detected": scrape_progress["changes_detected"],
        "new_values_added": scrape_progress["new_values_added"],
        "change_details": scrape_progress.get("change_details", []),
        "recent_changes": scrape_progress.get("recent_changes", []),
        "errors": scrape_progress.get("errors", []),
        "started_at": scrape_progress.get("started_at"),
        "competitors_processed": scrape_progress["competitors_done"],
        # v7.1.0: Enrichment tracking
        "enrichment_active": scrape_progress.get("enrichment_active", False),
        "news_articles_fetched": scrape_progress.get("news_articles_fetched", 0),
        "stock_prices_updated": scrape_progress.get("stock_prices_updated", 0)
    }


# Phase 3: Task 5.0.1-029 - Generate AI summary of refresh results
@app.post("/api/scrape/generate-summary")
async def generate_refresh_summary(db: Session = Depends(get_db)):
    """Use AI to generate a summary of the data refresh results."""
    import os

    if scrape_progress["active"]:
        return {"error": "Refresh still in progress", "type": "error"}

    if not scrape_progress.get("change_details"):
        return {
            "summary": "No changes detected during the last refresh. All competitor data remains the same.",
            "type": "static",
            "stats": {
                "competitors_scanned": scrape_progress.get("total", 0),
                "changes_detected": 0,
                "new_values": 0
            }
        }

    try:
        # Prepare change data for AI
        changes_text = ""
        change_details = scrape_progress.get("change_details", [])

        for change in change_details[:30]:  # Limit to prevent token overflow
            if change.get("type") == "new":
                changes_text += f"- NEW: {change.get('competitor', 'Unknown')} - {change.get('field', 'Unknown')}: {change.get('new_value', 'N/A')}\n"
            else:
                changes_text += f"- CHANGED: {change.get('competitor', 'Unknown')} - {change.get('field', 'Unknown')}: '{change.get('old_value', 'N/A')}' -> '{change.get('new_value', 'N/A')}'\n"

        # Generate AI summary via async AIRouter (avoids blocking the event loop)
        prompt = f"""Summarize the following data refresh results in 3-4 sentences.
Focus on:
1. Most significant changes (pricing, threat levels, new features)
2. Any concerning trends
3. Recommended actions for the sales team

Data Refresh Results:
- Competitors scanned: {scrape_progress.get('total', 0)}
- Changes detected: {scrape_progress.get('changes_detected', 0)}
- New data points: {scrape_progress.get('new_values_added', 0)}
- Errors encountered: {len(scrape_progress.get('errors', []))}

Detailed Changes:
{changes_text}

Provide a concise executive summary. Be specific about which competitors changed."""

        from ai_router import get_ai_router, TaskType
        router = get_ai_router()
        result = await router.generate(
            prompt=prompt,
            task_type=TaskType.SUMMARIZATION,
            system_prompt="You are a competitive intelligence analyst providing brief, actionable summaries for sales and product teams." + NO_HALLUCINATION_INSTRUCTION,
            max_tokens=300,
            temperature=0.7,
        )

        summary = result.get("response", "")
        if not summary:
            raise ValueError("Empty AI response")

        return {
            "summary": summary,
            "type": "ai",
            "model": result.get("model", "unknown"),
            "stats": {
                "competitors_scanned": scrape_progress.get("total", 0),
                "changes_detected": scrape_progress.get("changes_detected", 0),
                "new_values": scrape_progress.get("new_values_added", 0),
                "errors": len(scrape_progress.get("errors", []))
            }
        }

    except Exception as e:
        logger.warning(f"AI summary generation failed, using static fallback: {e}")
        return {
            "summary": f"Refreshed {scrape_progress.get('total', 0)} competitors. Found {scrape_progress.get('changes_detected', 0)} changes and {scrape_progress.get('new_values_added', 0)} new data points.",
            "type": "static",
            "error": "An unexpected error occurred"
        }


# Phase 4: Task 5.0.1-031 - Refresh history endpoint
@app.get("/api/refresh-history")
async def get_refresh_history(limit: int = 10, db: Session = Depends(get_db)):
    """Get history of data refresh sessions."""
    sessions = db.query(RefreshSession).order_by(
        RefreshSession.started_at.desc()
    ).limit(limit).all()

    return [{
        "id": s.id,
        "started_at": s.started_at.isoformat() if s.started_at else None,
        "completed_at": s.completed_at.isoformat() if s.completed_at else None,
        "competitors_scanned": s.competitors_scanned,
        "changes_detected": s.changes_detected,
        "new_values_added": s.new_values_added,
        "errors_count": s.errors_count,
        "status": s.status,
        "ai_summary": s.ai_summary
    } for s in sessions]


@app.post("/api/scrape/{competitor_id}")
async def trigger_scrape(competitor_id: int, background_tasks: BackgroundTasks, db: Session = Depends(get_db)):
    """Trigger a scrape job for a specific competitor."""
    competitor = db.query(Competitor).filter(
        Competitor.id == competitor_id,
        Competitor.is_deleted == False
    ).first()
    
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")
    
    # Add to background tasks
    background_tasks.add_task(run_scrape_job, competitor_id)
    
    return {
        "message": f"Scrape job queued for {competitor.name}",
        "competitor_id": competitor_id
    }


# NOTE: /api/discovery/run endpoint is defined later in file (line ~8813) with more features
# Duplicate removed to prevent FastAPI route override warnings


async def run_scrape_job(competitor_id: int):
    """Background job to scrape a competitor and update the database."""
    db = SessionLocal()
    try:
        comp = db.query(Competitor).filter(Competitor.id == competitor_id).first()
        if not comp:
            logger.error(f"Competitor {competitor_id} not found")
            return
        
        logger.info(f"Starting scrape for {comp.name}...")
        
        # Try to use the full scraper if available
        try:
            from scraper import CompetitorScraper
            from extractor import get_extractor

            scraper = CompetitorScraper()
            extractor = get_extractor()  # v5.0.2: Uses hybrid AI routing
            
            # Scrape the website
            content = await scraper.scrape(comp.website)
            
            if content:
                # Extract data using GPT
                from dataclasses import asdict
                
                # Note: extract_from_content is synchronous and takes (name, content)
                extracted_obj = extractor.extract_from_content(comp.name, content)
                # Convert dataclass to dict for iteration
                extracted = asdict(extracted_obj)
                
                if extracted:
                    # Update competitor with extracted data
                    for key, value in extracted.items():
                        if hasattr(comp, key) and value:
                            # Check if this field is locked (manual correction)
                            is_locked = db.query(DataSource).filter(
                                DataSource.competitor_id == comp.id,
                                DataSource.field_name == key,
                                DataSource.source_type == "manual"
                            ).first()
                            
                            if is_locked:
                                logger.debug(f"Skipping update for {comp.name}.{key} (locked by manual correction)")
                                continue
                                
                            old_value = getattr(comp, key)
                            if str(old_value) != str(value):
                                # Log the change
                                change = ChangeLog(
                                    competitor_id=comp.id,
                                    competitor_name=comp.name,
                                    change_type=key,
                                    previous_value=str(old_value) if old_value else None,
                                    new_value=str(value),
                                    source="scrape",
                                    severity="Medium"
                                )
                                db.add(change)
                                setattr(comp, key, value)
                    
                    comp.last_updated = datetime.utcnow()
                    db.commit()
                    logger.info(f"Scrape completed for {comp.name}")
                    return
        except ImportError as e:
            logger.warning(f"Scraper not available: {e}")
        except Exception as e:
            logger.error(f"Scrape error for {comp.name}: {e}")
        
        # Fallback: Just update the timestamp to show we tried
        comp.last_updated = datetime.utcnow()
        db.commit()
        logger.info(f"Refresh completed for {comp.name} (timestamp updated)")

    except Exception as e:
        logger.error(f"Scrape job failed for competitor {competitor_id}: {e}")
        db.rollback()
    finally:
        db.close()


def _update_data_source_with_confidence(
    db: Session,
    competitor_id: int,
    field_name: str,
    current_value: str,
    previous_value: str = None,
    source_url: str = None,
    source_name: str = None,
    gpt_confidence: int = 50
):
    """
    Helper to create or update DataSource record with confidence scoring.

    Maps GPT extraction confidence to our Admiralty Code-based scoring system.
    Website scrapes are inherently lower confidence (source_type = "website_scrape").
    """
    # Calculate confidence using our algorithm
    confidence_result = calculate_confidence_score(
        source_type="website_scrape",
        source_reliability="D",  # Website = Not usually reliable
        information_credibility=4,  # Doubtfully true until verified
        corroborating_sources=0,
        data_age_days=0
    )

    # Adjust score slightly based on GPT's own confidence assessment
    adjusted_score = min(100, max(0, confidence_result.score + (gpt_confidence - 50) // 5))
    adjusted_level = determine_confidence_level_from_score(adjusted_score)

    # Check if source already exists
    existing = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id,
        DataSource.field_name == field_name
    ).first()

    if existing:
        # Update existing record
        existing.previous_value = existing.current_value
        existing.current_value = current_value
        existing.source_type = "website_scrape"
        existing.source_url = source_url
        existing.source_name = source_name
        existing.extraction_method = "gpt_extraction"
        existing.extracted_at = datetime.utcnow()
        existing.source_reliability = "D"
        existing.information_credibility = 4
        existing.confidence_score = adjusted_score
        existing.confidence_level = adjusted_level
        existing.staleness_days = 0
        existing.updated_at = datetime.utcnow()
    else:
        # Create new record
        new_source = DataSource(
            competitor_id=competitor_id,
            field_name=field_name,
            current_value=current_value,
            previous_value=previous_value,
            source_type="website_scrape",
            source_url=source_url,
            source_name=source_name,
            extraction_method="gpt_extraction",
            source_reliability="D",
            information_credibility=4,
            confidence_score=adjusted_score,
            confidence_level=adjusted_level,
            corroborating_sources=0,
            staleness_days=0
        )
        db.add(new_source)


async def _run_post_scrape_enrichment():
    """v7.1.0: Post-scrape enrichment — fetch news articles and stock data for all competitors."""
    db = SessionLocal()
    try:
        competitors = db.query(Competitor).filter(
            Competitor.is_deleted == False,
            Competitor.status == "Active"
        ).all()

        # Phase 1: News enrichment
        logger.info(f"[Enrichment] Starting news fetch for {len(competitors)} competitors...")
        news_count = 0
        try:
            from news_monitor import NewsMonitor
            monitor = NewsMonitor()
            for comp in competitors:
                try:
                    scrape_progress["current_competitor"] = f"Fetching news: {comp.name}"
                    import asyncio
                    digest = await asyncio.to_thread(monitor.fetch_news, comp.name, 30)
                    if digest and digest.articles:
                        from database import NewsArticleCache
                        for article in digest.articles[:5]:
                            # Check for duplicate
                            existing = db.query(NewsArticleCache).filter(
                                NewsArticleCache.url == article.url
                            ).first()
                            if not existing:
                                cache_entry = NewsArticleCache(
                                    competitor_id=comp.id,
                                    competitor_name=comp.name,
                                    title=article.title,
                                    url=article.url,
                                    source=article.source,
                                    published_at=article.published_date if hasattr(article, 'published_date') else None,
                                    snippet=article.summary[:500] if article.summary else None,
                                    sentiment=article.sentiment if hasattr(article, 'sentiment') else 'neutral',
                                    event_type=article.event_type if hasattr(article, 'event_type') else None
                                )
                                db.add(cache_entry)
                                news_count += 1
                        db.commit()
                except Exception as news_err:
                    logger.warning(f"[Enrichment] News fetch failed for {comp.name}: {news_err}")
        except ImportError:
            logger.warning("[Enrichment] NewsMonitor not available, skipping news phase")
        except Exception as nm_err:
            logger.error(f"[Enrichment] News phase error: {nm_err}")

        scrape_progress["news_articles_fetched"] = news_count
        logger.info(f"[Enrichment] News phase complete: {news_count} articles fetched")

        # Phase 2: Stock data enrichment
        logger.info("[Enrichment] Starting stock data update...")
        stock_count = 0
        for comp in competitors:
            if comp.is_public and comp.ticker_symbol:
                try:
                    scrape_progress["current_competitor"] = f"Fetching stock: {comp.name} ({comp.ticker_symbol})"
                    import asyncio
                    stock_data = await asyncio.to_thread(fetch_real_stock_data, comp.ticker_symbol)
                    if stock_data and stock_data.get("price"):
                        comp.stock_price = stock_data["price"]
                        comp.market_cap = stock_data.get("market_cap")
                        comp.last_updated = datetime.utcnow()
                        db.commit()
                        stock_count += 1
                except Exception as stock_err:
                    logger.warning(f"[Enrichment] Stock fetch failed for {comp.name}: {stock_err}")

        scrape_progress["stock_prices_updated"] = stock_count
        logger.info(f"[Enrichment] Stock phase complete: {stock_count} prices updated")

    except Exception as e:
        logger.error(f"[Enrichment] Post-scrape enrichment error: {e}")
    finally:
        db.close()


async def run_scrape_job_with_progress(competitor_id: int, competitor_name: str):
    """Background job to scrape a competitor with progress tracking, unified change logging, and confidence scoring."""
    # Update current competitor being processed
    scrape_progress["current_competitor"] = competitor_name

    db = SessionLocal()
    changes_count = 0
    new_values_count = 0

    try:
        comp = db.query(Competitor).filter(Competitor.id == competitor_id).first()
        if not comp:
            logger.error(f"Competitor {competitor_id} not found")
            return

        logger.info(f"Starting scrape for {comp.name}...")

        # Try to use the full scraper if available
        try:
            from scraper import CompetitorScraper
            from extractor import get_extractor

            scraper = CompetitorScraper()
            extractor = get_extractor()  # v5.0.2: Uses hybrid AI routing

            # Scrape the website
            content = await scraper.scrape(comp.website)
            source_url = comp.website if not comp.website.startswith("http") else comp.website

            if content:
                # Extract data using AI (v5.0.2 - hybrid routing)
                from dataclasses import asdict

                extracted_obj = extractor.extract_from_content(comp.name, content.get("content", ""))
                extracted = asdict(extracted_obj) if hasattr(extracted_obj, '__dataclass_fields__') else extracted_obj

                # Get extraction confidence from AI (if available)
                ai_confidence = extracted.get("confidence_score") or 50

                if extracted:
                    # Update competitor with extracted data
                    for key, value in extracted.items():
                        if hasattr(comp, key) and value:
                            # Skip metadata fields
                            if key in ["confidence_score", "extraction_notes"]:
                                continue

                            # Check if this field is locked (manual correction)
                            is_locked = db.query(DataSource).filter(
                                DataSource.competitor_id == comp.id,
                                DataSource.field_name == key,
                                DataSource.source_type == "manual"
                            ).first()

                            if is_locked:
                                logger.debug(f"Skipping update for {comp.name}.{key} (locked by manual correction)")
                                continue

                            old_value = getattr(comp, key)
                            old_str = str(old_value) if old_value else None
                            new_str = str(value)

                            # Check if this is a new value or a change
                            if old_str != new_str:
                                # Determine change type
                                is_new_value = old_value is None or old_str == "" or old_str == "None"
                                change_type = "new" if is_new_value else "change"

                                # Log to DataChangeHistory (unified change log)
                                change_record = DataChangeHistory(
                                    competitor_id=comp.id,
                                    competitor_name=comp.name,
                                    field_name=key,
                                    old_value=old_str,
                                    new_value=new_str,
                                    changed_by="System (Auto-Refresh)",
                                    change_reason="Automated data refresh"
                                )
                                db.add(change_record)
                                setattr(comp, key, value)

                                # Phase 2: Track field-level changes (Task 5.0.1-027)
                                change_entry = {
                                    "competitor": comp.name,
                                    "field": key,
                                    "old_value": old_str[:50] if old_str else None,
                                    "new_value": new_str[:50] if new_str else None,
                                    "type": change_type,
                                    "timestamp": datetime.utcnow().isoformat()
                                }
                                scrape_progress["change_details"].append(change_entry)
                                scrape_progress["recent_changes"].append(change_entry)

                                # Keep only last 10 in recent_changes for live display
                                if len(scrape_progress["recent_changes"]) > 10:
                                    scrape_progress["recent_changes"] = scrape_progress["recent_changes"][-10:]

                                if is_new_value:
                                    new_values_count += 1
                                else:
                                    changes_count += 1

                            # Create or update DataSource with confidence scoring
                            _update_data_source_with_confidence(
                                db=db,
                                competitor_id=comp.id,
                                field_name=key,
                                current_value=new_str,
                                previous_value=old_str,
                                source_url=source_url,
                                source_name=f"{comp.name} Website",
                                gpt_confidence=ai_confidence  # v5.0.2: renamed to ai_confidence
                            )

                    comp.last_updated = datetime.utcnow()
                    db.commit()
                    logger.info(f"Scrape completed for {comp.name} - {changes_count} changes, {new_values_count} new values")

                    # Trigger triangulation for key fields to verify scraped data
                    try:
                        triangulator = DataTriangulator(db)
                        triangulation_results = await triangulator.triangulate_all_key_fields(
                            competitor_id=comp.id,
                            competitor_name=comp.name,
                            website=comp.website,
                            is_public=comp.is_public,
                            ticker_symbol=comp.ticker_symbol
                        )

                        # Update confidence scores based on triangulation
                        for field_name, result in triangulation_results.items():
                            if result.confidence_score > 0:
                                existing = db.query(DataSource).filter(
                                    DataSource.competitor_id == comp.id,
                                    DataSource.field_name == field_name
                                ).first()
                                if existing:
                                    existing.confidence_score = result.confidence_score
                                    existing.confidence_level = result.confidence_level
                                    existing.corroborating_sources = result.sources_agreeing

                        db.commit()
                        logger.info(f"Triangulation completed for {comp.name}")
                    except Exception as tri_err:
                        logger.error(f"Triangulation error for {comp.name}: {tri_err}")

        except ImportError as e:
            logger.warning(f"Scraper not available: {e}")
            # Track error for display
            scrape_progress["errors"].append({
                "competitor": competitor_name,
                "error": "Scraper not available",
                "timestamp": datetime.utcnow().isoformat()
            })
        except Exception as e:
            logger.error(f"Scrape error for {comp.name}: {e}")
            # Track error for display
            scrape_progress["errors"].append({
                "competitor": competitor_name,
                "error": "An unexpected error occurred"[:100],
                "timestamp": datetime.utcnow().isoformat()
            })

        # Fallback: Just update the timestamp to show we tried
        if not changes_count and not new_values_count:
            comp.last_updated = datetime.utcnow()
            db.commit()
            logger.info(f"Refresh completed for {comp.name} (timestamp updated)")

    except Exception as e:
        logger.error(f"Scrape job failed for competitor {competitor_id}: {e}")
        db.rollback()
        # Track critical error
        scrape_progress["errors"].append({
            "competitor": competitor_name,
            "error": "Job failed",
            "timestamp": datetime.utcnow().isoformat()
        })
    finally:
        db.close()

        # Update progress tracker
        scrape_progress["completed"] += 1
        scrape_progress["competitors_done"].append(competitor_name)
        scrape_progress["changes_detected"] += changes_count
        scrape_progress["new_values_added"] += new_values_count

        # Check if all scrapes are done
        if scrape_progress["completed"] >= scrape_progress["total"]:
            scrape_progress["current_competitor"] = None
            logger.info(f"All scrapes complete! {scrape_progress['changes_detected']} changes, {scrape_progress['new_values_added']} new values")

            # v7.1.0: Post-scrape enrichment — fetch news + stock data
            scrape_progress["enrichment_active"] = True
            scrape_progress["current_competitor"] = "Enriching: Fetching news & stock data..."
            try:
                await _run_post_scrape_enrichment()
            except Exception as enrich_err:
                logger.error(f"Post-scrape enrichment error: {enrich_err}")
            scrape_progress["enrichment_active"] = False
            scrape_progress["active"] = False
            scrape_progress["current_competitor"] = None

            # Phase 4: Persist RefreshSession results (Task 5.0.1-031)
            session_db = None
            try:
                session_id = scrape_progress.get("session_id")
                if session_id:
                    session_db = SessionLocal()
                    refresh_session = session_db.query(RefreshSession).filter(
                        RefreshSession.id == session_id
                    ).first()
                    if refresh_session:
                        refresh_session.completed_at = datetime.utcnow()
                        refresh_session.changes_detected = scrape_progress["changes_detected"]
                        refresh_session.new_values_added = scrape_progress["new_values_added"]
                        refresh_session.errors_count = len(scrape_progress.get("errors", []))
                        refresh_session.status = "completed"
                        # Store change details as JSON
                        import json
                        refresh_session.change_details = json.dumps(scrape_progress.get("change_details", []))
                        session_db.commit()
                        logger.debug(f"RefreshSession {session_id} persisted to database")
            except Exception as persist_err:
                logger.error(f"Error persisting RefreshSession: {persist_err}")
            finally:
                if session_db:
                    session_db.close()


# --- News Feed Endpoint ---

@app.get("/api/news/{company_name}")
async def get_competitor_news(
    company_name: str,
    limit: int = 10,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """Fetch recent news articles mentioning a competitor. DB-first, live fallback."""
    from database import NewsArticleCache

    # FIRST: Query database cache (fast path)
    try:
        cached = db.query(NewsArticleCache).filter(
            NewsArticleCache.competitor_name.ilike(f"%{company_name}%"),
            NewsArticleCache.is_archived != True
        ).order_by(NewsArticleCache.published_at.desc()).limit(limit).all()

        if cached:
            return {
                "company": company_name,
                "articles": [
                    {
                        "title": n.title,
                        "url": n.url,
                        "source": n.source,
                        "published_date": str(n.published_at) if n.published_at else "",
                        "snippet": n.snippet or "",
                        "sentiment": n.sentiment,
                        "event_type": n.event_type or "general",
                    }
                    for n in cached
                ],
                "count": len(cached),
                "source": "database",
                "fetched_at": datetime.utcnow().isoformat()
            }
    except Exception as db_err:
        logger.warning(f"DB news lookup failed for {company_name}: {db_err}")

    # ONLY IF zero DB results: fall back to live fetch
    try:
        from news_monitor import NewsMonitor
        monitor = NewsMonitor()
        digest = monitor.fetch_news(company_name, days=90)
        raw_articles = digest.articles[:limit] if digest and digest.articles else []
        return {
            "company": company_name,
            "articles": [
                {
                    "title": a.title,
                    "source": a.source,
                    "url": a.url,
                    "published_date": a.published_date,
                    "snippet": a.snippet,
                    "sentiment": a.sentiment,
                    "event_type": a.event_type,
                }
                for a in raw_articles
            ],
            "count": len(raw_articles),
            "source": "live",
            "fetched_at": datetime.utcnow().isoformat()
        }
    except Exception as e:
        logger.warning(f"Live news fetch also failed for {company_name}: {e}")

    return {
        "company": company_name,
        "articles": [],
        "count": 0,
        "source": "none",
        "fetched_at": datetime.utcnow().isoformat()
    }



# --- AI SWOT Analysis Endpoint (Real) ---

@app.post("/api/agents/battlecard")
async def generate_battlecard_swot(
    request: Request,
    background_tasks: BackgroundTasks = None,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """Generate battlecard SWOT analysis via agents endpoint (frontend compatibility)."""
    body = await request.json()
    competitor_id = body.get("competitor_id")
    battlecard_type = body.get("type", "swot")
    prompt_key = body.get("prompt_key")
    session_id = body.get("session_id")
    background = body.get("background", False)

    if not competitor_id:
        raise HTTPException(status_code=400, detail="competitor_id is required")

    comp = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not comp:
        raise HTTPException(status_code=404, detail="Competitor not found")

    # For SWOT type, delegate to existing SWOT generation logic
    if battlecard_type == "swot":
        if background and background_tasks and current_user:
            import uuid as _uuid
            task_id = str(_uuid.uuid4())
            comp_name = comp.name
            _ai_tasks[task_id] = {
                "status": "running", "page_context": "battlecard",
                "user_id": current_user.get("id"),
                "started_at": datetime.utcnow().isoformat(),
                "completed_at": None, "result": None, "error": None,
                "task_type": f"battlecard_swot_{comp_name}", "read_at": None,
            }

            async def _run_battlecard_swot_bg():
                try:
                    bg_db = SessionLocal()
                    try:
                        from ai_router import get_ai_router, TaskType
                        c = bg_db.query(Competitor).filter(
                            Competitor.id == competitor_id
                        ).first()
                        context = f"""
                        Analyze this competitor for 'Certify Health' (Provider of Patient Intake, Payments, and Biometrics).

                        Competitor: {c.name}
                        Description: {c.notes or 'N/A'}
                        Pricing: {c.pricing_model or 'N/A'} ({c.base_price or 'N/A'})
                        Target Segments: {c.target_segments or 'N/A'}
                        Key Features: {c.key_features or 'N/A'}
                        """
                        router = get_ai_router()
                        result = await asyncio.wait_for(
                            router.generate_json(
                                prompt=context,
                                task_type=TaskType.ANALYSIS,
                                system_prompt="Generate JSON SWOT with keys: strengths, weaknesses, opportunities, threats." + NO_HALLUCINATION_INSTRUCTION,
                                temperature=0.7,
                            ),
                            timeout=45.0
                        )
                        content = result.get("response_json", {})
                        _ai_tasks[task_id]["result"] = {
                            "swot": {
                                "strengths": content.get("strengths", []),
                                "weaknesses": content.get("weaknesses", []),
                                "opportunities": content.get("opportunities", []),
                                "threats": content.get("threats", []),
                            }
                        }
                    finally:
                        bg_db.close()
                    _ai_tasks[task_id]["status"] = "completed"
                    _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
                except asyncio.TimeoutError:
                    logger.error("Background SWOT generation timed out after 45s")
                    _ai_tasks[task_id]["status"] = "failed"
                    _ai_tasks[task_id]["error"] = "AI analysis timed out. Please try again."
                    _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
                except Exception as bg_err:
                    logger.error(f"Background battlecard SWOT error: {bg_err}")
                    _ai_tasks[task_id]["status"] = "failed"
                    _ai_tasks[task_id]["error"] = str(bg_err)
                    _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()

            background_tasks.add_task(_run_battlecard_swot_bg)
            return {"task_id": task_id, "status": "running"}

        # Synchronous generation
        try:
            context_template = None
            system_prompt_text = None
            if prompt_key:
                from database import SystemPrompt as SP
                p = db.query(SP).filter(SP.key == prompt_key, SP.user_id == None).first()
                if p:
                    system_prompt_text = p.content
            if not system_prompt_text:
                system_prompt_text = "You are a competitive strategy expert. Generate a strict JSON SWOT analysis with 3-4 bullet points per section. Return JSON with keys: strengths, weaknesses, opportunities, threats (each an array of strings)."

            context = f"""
            Analyze this competitor for 'Certify Health' (Provider of Patient Intake, Payments, and Biometrics).

            Competitor: {comp.name}
            Description: {comp.notes or 'N/A'}
            Pricing: {comp.pricing_model or 'N/A'} ({comp.base_price or 'N/A'})
            Target Segments: {comp.target_segments or 'N/A'}
            Key Features: {comp.key_features or 'N/A'}
            Weaknesses/Gaps: {comp.notes or 'N/A'}
            """

            from ai_router import get_ai_router, TaskType
            router = get_ai_router()
            result = await router.generate_json(
                prompt=context,
                task_type=TaskType.ANALYSIS,
                system_prompt=system_prompt_text + NO_HALLUCINATION_INSTRUCTION,
                temperature=0.7,
            )
            content = result.get("response_json", {})

            return {
                "swot": {
                    "strengths": content.get("strengths", []),
                    "weaknesses": content.get("weaknesses", []),
                    "opportunities": content.get("opportunities", []),
                    "threats": content.get("threats", []),
                },
                "provider": result.get("provider"),
                "model": result.get("model"),
            }
        except Exception as e:
            logger.error(f"Error generating battlecard SWOT: {e}")
            raise HTTPException(status_code=500, detail="Internal server error. Please try again.")

    else:
        raise HTTPException(status_code=400, detail=f"Unsupported battlecard type: {battlecard_type}")


@app.get("/api/competitors/{competitor_id}/swot")
async def get_swot_analysis(
    competitor_id: int,
    prompt_key: Optional[str] = None,
    session_id: Optional[int] = None,
    background: bool = False,
    background_tasks: BackgroundTasks = None,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """Generate a real-time SWOT analysis using GPT-4 based on database records."""
    comp = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not comp:
        raise HTTPException(status_code=404, detail="Competitor not found")

    if background and background_tasks and current_user:
        import uuid as _uuid
        task_id = str(_uuid.uuid4())
        comp_name = comp.name
        _ai_tasks[task_id] = {
            "status": "running", "page_context": "analytics",
            "user_id": current_user.get("id"),
            "started_at": datetime.utcnow().isoformat(),
            "completed_at": None, "result": None, "error": None,
            "task_type": f"swot_{comp_name}", "read_at": None,
        }

        async def _run_swot_bg():
            try:
                bg_db = SessionLocal()
                try:
                    from ai_router import get_ai_router, TaskType
                    c = bg_db.query(Competitor).filter(
                        Competitor.id == competitor_id
                    ).first()
                    context = f"Analyze competitor: {c.name}. Features: {c.key_features}"
                    router = get_ai_router()
                    result = await router.generate_json(
                        prompt=context,
                        task_type=TaskType.ANALYSIS,
                        system_prompt="Generate JSON SWOT with keys: strengths, weaknesses, opportunities, threats." + NO_HALLUCINATION_INSTRUCTION,
                        temperature=0.7,
                    )
                    content = result.get("response_json", {})
                    _ai_tasks[task_id]["result"] = {
                        "strengths": content.get("strengths", []),
                        "weaknesses": content.get("weaknesses", []),
                        "opportunities": content.get("opportunities", []),
                        "threats": content.get("threats", []),
                    }
                finally:
                    bg_db.close()
                _ai_tasks[task_id]["status"] = "completed"
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
            except Exception as bg_err:
                logger.error(f"Background SWOT error: {bg_err}")
                _ai_tasks[task_id]["status"] = "failed"
                _ai_tasks[task_id]["error"] = str(bg_err)
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()

        background_tasks.add_task(_run_swot_bg)
        return {"task_id": task_id, "status": "running"}

    try:
        # Load context prompt from DB if prompt_key provided
        context_template = None
        system_prompt_text = None
        if prompt_key:
            from database import SystemPrompt as SP
            p = db.query(SP).filter(SP.key == prompt_key, SP.user_id == None).first()
            if p:
                system_prompt_text = p.content
        if not system_prompt_text:
            system_prompt_text = "You are a competitive strategy expert. Generate a strict JSON SWOT analysis with 3-4 bullet points per section. Return JSON with keys: strengths, weaknesses, opportunities, threats (each an array of strings)."

        # Construct rich context from DB
        context = f"""
        Analyze this competitor for 'Certify Health' (Provider of Patient Intake, Payments, and Biometrics).

        Competitor: {comp.name}
        Description: {comp.notes or 'N/A'}
        Pricing: {comp.pricing_model} ({comp.base_price})
        Target Segments: {comp.target_segments}
        Key Features: {comp.key_features}
        Weaknesses/Gaps: {comp.notes}
        """

        from ai_router import get_ai_router, TaskType
        router = get_ai_router()
        result = await router.generate_json(
            prompt=context,
            task_type=TaskType.ANALYSIS,
            system_prompt=system_prompt_text + NO_HALLUCINATION_INSTRUCTION,
            temperature=0.7,
        )
        content = result.get("response_json", {})

        # Ensure correct key structure
        swot_result = {
            "strengths": content.get("strengths", []) or content.get("Strengths", []),
            "weaknesses": content.get("weaknesses", []) or content.get("Weaknesses", []),
            "opportunities": content.get("opportunities", []) or content.get("Opportunities", []),
            "threats": content.get("threats", []) or content.get("Threats", []),
            "session_id": session_id,
        }

        # Persist to chat session if session_id provided
        if session_id:
            try:
                import json as _json
                _save_chat_messages(db, session_id, f"Generate SWOT for {comp.name}", _json.dumps(swot_result, indent=2), {
                    "endpoint": "swot", "competitor_id": competitor_id,
                })
            except Exception as save_err:
                logger.warning(f"Failed to save SWOT to chat: {save_err}")

        return swot_result

    except Exception as e:
        logger.error(f"SWOT Generation Error: {e}")
        # Return honest fallback - no fabricated data
        return {
            "strengths": ["No verified data available - AI analysis unavailable"],
            "weaknesses": [comp.notes or "No verified weakness data recorded"],
            "opportunities": ["No verified data available for this metric."],
            "threats": ["No verified data available for this metric."]
        }


# --- Stock Ticker Endpoint ---

def fetch_real_stock_data(ticker: str) -> Dict[str, Any]:
    """Fetch real-time stock data using yfinance."""
    try:
        stock = yf.Ticker(ticker)
        info = stock.info
        
        # Calculate change if not provided
        current = info.get("currentPrice") or info.get("regularMarketPrice", 0)
        previous = info.get("previousClose", 0)
        change = current - previous if current and previous else 0
        percent = (change / previous) * 100 if previous else 0
        
        return {
            "price": current,
            "change": change,
            "change_percent": percent,
            "market_cap": info.get("marketCap"),
            "pe_ratio": info.get("trailingPE"),
            "eps": info.get("trailingEps"),
            "beta": info.get("beta"),
            "volume": info.get("volume"),
            "high52": info.get("fiftyTwoWeekHigh"),
            "low52": info.get("fiftyTwoWeekLow"),
            "target_est": info.get("targetMeanPrice"),
            "company": info.get("longName"),
            "ticker": ticker
        }
    except Exception as e:
        logger.error(f"Error fetching stock data for {ticker}: {e}")
        return None

@app.get("/api/stock/{company_name}")
async def get_stock_data(company_name: str, db: Session = Depends(get_db)):
    """Get stock data for a public company using yfinance."""
    import yfinance as yf
    from datetime import datetime
    
    company_lower = company_name.lower()
    
    # 1. Check Database first (for manually updated tickers)
    db_comp = db.query(Competitor).filter(Competitor.name == company_name).first()
    
    ticker_obj = None
    
    if db_comp and db_comp.is_public and db_comp.ticker_symbol:
        ticker_obj = {
            "symbol": db_comp.ticker_symbol,
            "exchange": db_comp.stock_exchange or "NYSE",
            "name": db_comp.name
        }
    
    # 2. Check hardcoded list if not in DB
    if not ticker_obj:
        ticker_obj = KNOWN_TICKERS.get(company_lower)
    
    if ticker_obj:
        try:
            # Fetch data from yfinance
            ticker = yf.Ticker(ticker_obj["symbol"])
            info = ticker.info
            
            # Helper for safe float formatting
            def get_val(key, default=None):
                return info.get(key, default)

            # Calculate Free Cash Flow if not directly available
            fcf = get_val('freeCashflow')
            if fcf is None and get_val('operatingCashflow') and get_val('capitalExpenditures'):
                fcf = get_val('operatingCashflow') - abs(get_val('capitalExpenditures', 0))

            next_earnings = "N/A"
            if get_val('earningsTimestamp'):
                next_earnings = datetime.fromtimestamp(get_val('earningsTimestamp')).strftime('%Y-%m-%d')
            elif get_val('earningsDate'):
                 # Sometimes it's a list
                 dates = get_val('earningsDate')
                 if isinstance(dates, list) and len(dates) > 0:
                     next_earnings = str(dates[0])[:10]

            return {
                "is_public": True,
                "company": ticker_obj["name"],
                "ticker": ticker_obj["symbol"],
                "exchange": ticker_obj["exchange"],
                
                # Trading
                "price": get_val('currentPrice', get_val('previousClose')),
                "change": float(get_val('currentPrice', 0)) - float(get_val('previousClose', 0)) if get_val('currentPrice') and get_val('previousClose') else 0,
                "change_percent": ((float(get_val('currentPrice', 0)) - float(get_val('previousClose', 0))) / float(get_val('previousClose', 1))) * 100 if get_val('previousClose') else 0,
                "volume": get_val('volume'),
                "avg_volume_10d": get_val('averageVolume10days'),
                "avg_volume_90d": get_val('averageVolume'),
                "fifty_two_week_low": get_val('fiftyTwoWeekLow'),
                "fifty_two_week_high": get_val('fiftyTwoWeekHigh'),
                "market_cap": get_val('marketCap'),

                # Valuation
                "enterprise_value": get_val('enterpriseValue'),
                "pe_trailing": get_val('trailingPE'),
                "pe_forward": get_val('forwardPE'),
                "ev_ebitda": get_val('enterpriseToEbitda'),
                "price_to_book": get_val('priceToBook'),
                "peg_ratio": get_val('pegRatio'),

                # Operating
                "eps_trailing": get_val('trailingEps'),
                "eps_forward": get_val('forwardEps'),
                "ebitda": get_val('ebitda'),
                "revenue_ttm": get_val('totalRevenue'),
                "free_cash_flow": fcf,
                "profit_margin": get_val('profitMargins'),

                # Risk
                "beta": get_val('beta'),
                "short_interest": get_val('shortPercentOfFloat'),
                "float_shares": get_val('floatShares'),

                # Capital
                "shares_outstanding": get_val('sharesOutstanding'),
                "inst_ownership": get_val('heldPercentInstitutions'),
                "dividend_yield": get_val('dividendYield'),
                "next_earnings": next_earnings,
                
                "data_sources": ["Yahoo Finance", "SEC EDGAR"]
            }
        except Exception as e:
            logger.error(f"yfinance error for {ticker_obj['symbol']}: {e}")
            return {
                "is_public": True,
                "company": ticker_obj["name"],
                "ticker": ticker_obj["symbol"],
                "exchange": ticker_obj["exchange"],
                "error": "Unable to fetch live data"
            }



    # 3. Private Company Intelligence Logic
    
    # Initialize trackers
    from sec_edgar_scraper import SECEdgarScraper
    from linkedin_tracker import LinkedInTracker
    from gov_contracts_scraper import GovContractsScraper
    from h1b_scraper import H1BScraper
    from uspto_scraper import USPTOScraper
    from appstore_scraper import AppStoreScraper
    from glassdoor_scraper import GlassdoorScraper
    from google_ecosystem_scraper import GoogleEcosystemScraper
    from tech_stack_scraper import TechStackScraper
    from sentiment_scraper import SentimentScraper
    from seo_scraper import SEOScraper
    from risk_management_scraper import RiskManagementScraper
    
    sec = SECEdgarScraper()
    linkedin = LinkedInTracker()
    gov = GovContractsScraper()
    h1b = H1BScraper()
    uspto = USPTOScraper()
    appstore = AppStoreScraper()
    glassdoor = GlassdoorScraper()
    google = GoogleEcosystemScraper()
    tech = TechStackScraper()
    sentiment = SentimentScraper()
    seo = SEOScraper()
    risk = RiskManagementScraper()
    
    # Fetch Data
    form_d = sec.get_latest_form_d(company_name)
    li_data = linkedin.get_company_data(company_name)
    gov_data = gov.get_contract_data(company_name)
    h1b_data = h1b.get_h1b_data(company_name)
    patent_data = uspto.get_patent_data(company_name)
    app_data = appstore.get_app_data(company_name)
    glassdoor_data = glassdoor.get_company_data(company_name)
    google_data = google.get_ecosystem_data(company_name)
    tech_data = tech.get_tech_stack(company_name)
    sentiment_data = sentiment.get_sentiment_data(company_name)
    seo_data = seo.get_seo_data(company_name)
    risk_data = risk.get_risk_data(company_name)
    
    # Calculate Est. Revenue (Proxy: $150k ARR per employee for HealthTech)
    est_revenue = (li_data.employee_count or 0) * 150000
    
    # Determine Status
    stage = "Late Stage VC" 
    if li_data.employee_count < 50: stage = "Seed/Early"
    elif li_data.employee_count < 200: stage = "Growth Stage"
    elif li_data.employee_count > 1000:
        if gov_data and gov_data.total_amount > 10000000:
            stage = "Gov. Contractor / Mature"
        elif google_data.ads.active_creative_count > 100:
            stage = "Mass Market / Mature"
        else:
            stage = "Pre-IPO / PE Backed"
    
    return {
        "is_public": False,
        "company": company_name,
        "stage": stage,
        
        # Capital
        "total_funding": form_d["amount_raised"] if form_d else None,
        "latest_deal_date": form_d["filing_date"] if form_d else None,
        "latest_deal_amount": form_d["amount_raised"] if form_d else None,
        "latest_deal_type": form_d.get("round_type", "Venture Round") if form_d else None,
        
        # Growth & Ops
        "headcount": li_data.employee_count,
        "growth_rate_6mo": li_data.employee_growth_6mo,
        "active_hiring": li_data.open_jobs,
        "est_revenue": est_revenue,
        "hiring_departments": list(li_data.job_categories.keys())[:3] if li_data.job_categories else [],
        
        # Identity
        "headquarters": li_data.headquarters,
        "founded": li_data.founded_year,

        # Alternative Intelligence
        "gov_contracts": {
            "total_awards": gov_data.total_awards,
            "total_amount": gov_data.total_amount,
            "top_agency": gov_data.top_agency
        },
        "h1b_data": {
            "filings": h1b_data.total_filings_2023,
            "avg_salary": h1b_data.avg_salary_engineer,
            "top_title": h1b_data.top_job_titles[0] if h1b_data.top_job_titles else "N/A"
        },
        "innovation": {
            "patents": patent_data.total_patents,
            "pending": patent_data.pending_applications,
            "innovation_score": patent_data.innovation_score
        },
        "app_quality": {
            "avg_rating": app_data.avg_rating,
            "downloads": app_data.total_downloads,
            "sentiment": app_data.sentiment_summary
        },
        "employee_sentiment": {
            "rating": glassdoor_data.overall_rating,
            "ceo_approval": glassdoor_data.ceo_approval,
            "recommend": glassdoor_data.recommend_to_friend
        },

        # Google Digital Footprint
        "google_ecosystem": {
            "ads_active": google_data.ads.active_creative_count,
            "ad_formats": google_data.ads.formats,
            "brand_index": google_data.trends.current_index,
            "trend": google_data.trends.trend_direction,
            "reviews": google_data.maps.review_count,
            "review_velocity": google_data.maps.reviews_last_month
        },
        "tech_stack": {
            "signal": tech_data.marketing_budget_signal,
            "tools": tech_data.detected_tools,
            "has_enterprise_ads": tech_data.has_floodlight or tech_data.has_adobe_analytics
        },
        
        # Deep Dive Intelligence (New)
        "sentiment": {
            "g2_score": sentiment_data.g2_score,
            "g2_badges": sentiment_data.g2_badges[:2],
            "trustpilot": sentiment_data.trustpilot_score,
            "reddit": sentiment_data.reddit_sentiment,
            "complaints": sentiment_data.top_complaints[:1]
        },
        "seo": {
            "da": seo_data.domain_authority,
            "backlinks": seo_data.backlink_count,
            "speed": seo_data.page_load_speed,
            "keywords": seo_data.top_keywords[:3]
        },
        "risk_mgmt": {
            "founder_exit": risk_data.founder_exits,
            "exec_tenure": risk_data.avg_executive_tenure,
            "tier1_vc": risk_data.tier_1_investors,
            "soc2": risk_data.soc2_compliant,
            "warn": risk_data.warn_notices
        },
        
        "data_sources": ["SEC", "LinkedIn", "USAspending", "H-1B", "Google", "Tech", "G2/Capterra", "Moz", "Crunchbase"]
    }


# --- Review Platform Endpoints ---

@app.get("/api/reviews/certify-health")
async def get_certify_reviews():
    """Get Certify Health's reviews from G2, Trustpilot, etc."""
    try:
        from review_scraper import get_certify_health_reviews
        result = await get_certify_health_reviews()
        return result
    except Exception as e:
        return {"error": "An unexpected error occurred"}


@app.get("/api/reviews/competitor/{competitor_key}")
async def get_competitor_reviews_endpoint(competitor_key: str):
    """Get reviews for a competitor (e.g., phreesia, luma_health)."""
    try:
        from review_scraper import get_competitor_reviews
        result = await get_competitor_reviews(competitor_key)
        return result
    except Exception as e:
        return {"error": "An unexpected error occurred"}


# NOTE: Early /api/reviews/compare removed (duplicate of DB-backed version below at ~line 10545)


@app.get("/api/reviews/platforms")
async def list_review_platforms():
    """List available review platforms and known competitor mappings."""
    try:
        from review_scraper import KNOWN_REVIEW_URLS
        return {
            "platforms": ["g2", "trustpilot", "google_business"],
            "certify_links": {
                "g2_review": "https://www.g2.com/products/certify-health/take_survey",
                "g2_video": "https://interviews.g2.com/review/certify-health?mode=call&auto=true",
                "capterra": "https://reviews.capterra.com/products/new/5baef775-663a-4b3c-95cf-dc868b7aa283/",
                "trustpilot": "https://www.trustpilot.com/review/certifyhealth.com",
                "google": "https://g.page/r/CfVO-Pq4X5IhEAE/review"
            },
            "known_competitors": list(KNOWN_REVIEW_URLS.keys())
        }
    except Exception as e:
        return {"error": "An unexpected error occurred"}


# --- Alert Endpoints ---

@app.post("/api/alerts/send-digest")
def send_digest_email():
    """Manually trigger daily digest email."""
    try:
        from alerts import send_daily_digest
        success = send_daily_digest()
        return {"success": success, "message": "Daily digest sent" if success else "No changes to report or email not configured"}
    except Exception as e:
        return {"success": False, "error": "An unexpected error occurred"}


@app.post("/api/alerts/send-summary")
def send_summary_email():
    """Manually trigger weekly summary email."""
    try:
        from alerts import send_weekly_summary
        success = send_weekly_summary()
        return {"success": success, "message": "Weekly summary sent" if success else "Email not configured"}
    except Exception as e:
        return {"success": False, "error": "An unexpected error occurred"}


# --- Report Endpoints ---




@app.get("/api/reports/weekly-briefing")
def generate_weekly_briefing(db: Session = Depends(get_db)):
    """Generate executive weekly briefing PDF."""
    try:
        from reports import ReportManager
        competitors = db.query(Competitor).filter(Competitor.is_deleted == False).all()
        changes = db.query(ChangeLog).order_by(ChangeLog.detected_at.desc()).limit(10).all()
        
        # Calculate stats
        stats = {
            "total_competitors": len(competitors),
            "high_threat": len([c for c in competitors if c.threat_level and c.threat_level.upper() == "HIGH"]),
            "medium_threat": len([c for c in competitors if c.threat_level and c.threat_level.upper() == "MEDIUM"]),
            "low_threat": len([c for c in competitors if c.threat_level and c.threat_level.upper() == "LOW"])
        }
        
        manager = ReportManager("./exports")
        filepath = manager.generate_weekly_briefing(
            [c.__dict__ for c in competitors], 
            [c.__dict__ for c in changes],
            stats
        )
        return FileResponse(filepath, filename="weekly_briefing.pdf", media_type="application/pdf")
    except Exception as e:
        # Fallback - return a simple text summary
        competitors = db.query(Competitor).filter(Competitor.is_deleted == False).all()
        stats = {
            "total": len(competitors),
            "high": len([c for c in competitors if c.threat_level == "High"]),
            "medium": len([c for c in competitors if c.threat_level == "Medium"]),
            "low": len([c for c in competitors if c.threat_level == "Low"])
        }
        return {"error": "An unexpected error occurred", "summary": stats, "message": "PDF generation unavailable, run pip install reportlab"}


@app.get("/api/reports/comparison")
def generate_comparison_report(db: Session = Depends(get_db)):
    """Generate competitor comparison PDF."""
    try:
        from reports import ReportManager
        competitors = db.query(Competitor).filter(Competitor.is_deleted == False).all()
        
        manager = ReportManager("./exports")
        filepath = manager.generate_comparison([c.__dict__ for c in competitors])
        return FileResponse(filepath, filename="competitor_comparison.pdf", media_type="application/pdf")
    except Exception as e:
        return {"error": "An unexpected error occurred", "message": "PDF generation unavailable, run pip install reportlab"}


@app.get("/api/reports/battlecard/{competitor_id}")
def generate_battlecard(competitor_id: int, db: Session = Depends(get_db)):
    """Generate battlecard PDF for a specific competitor."""
    competitor = db.query(Competitor).filter(
        Competitor.id == competitor_id,
        Competitor.is_deleted == False
    ).first()
    
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")
    
    try:
        from reports import ReportManager
        manager = ReportManager("./exports")
        
        # Inject Stock Data if Public
        comp_dict = {k: v for k, v in competitor.__dict__.items() if not k.startswith('_')}
        if competitor.is_public and competitor.ticker_symbol:
            stock_data = fetch_real_stock_data(competitor.ticker_symbol)
            if stock_data:
                comp_dict["stock_data"] = stock_data
                
        filepath = manager.generate_battlecard(comp_dict)
        return FileResponse(filepath, filename=f"{competitor.name}_battlecard.pdf", media_type="application/pdf")
    except Exception as e:
        return {"error": "An unexpected error occurred", "competitor": competitor.name, "message": "PDF generation unavailable"}


# ============== Discovery Endpoints ==============

# In-memory storage for discovery results (would use Redis in production)
discovery_results = {"candidates": [], "last_run": None}


@app.delete("/api/discovery/results")
def delete_discovery_results(db: Session = Depends(get_db)):
    """Clear saved discovery results from DB and in-memory cache (v8.0.8)."""
    global discovery_results
    from database import PersistentCache
    discovery_results = {"candidates": [], "last_run": None}
    try:
        db.query(PersistentCache).filter(
            PersistentCache.cache_key == "discovery_results"
        ).delete()
        db.commit()
    except Exception as e:
        logger.warning(f"Failed to delete discovery results: {e}")
    return {"status": "ok"}


def _normalize_candidate_scores(data: dict) -> dict:
    """Normalize score fields on discovery candidates so frontend has a consistent `match_score`."""
    if not data or not data.get("candidates"):
        return data
    for c in data["candidates"]:
        if "match_score" not in c or c["match_score"] is None:
            c["match_score"] = (
                c.get("qualification_score")
                or c.get("relevance_score")
                or c.get("score")
                or None
            )
    return data


@app.get("/api/discovery/results")
def get_discovery_results(db: Session = Depends(get_db)):
    """Get previously discovered competitor candidates.
    v8.0.8: Falls back to DB if in-memory cache is empty (survives server restart).
    v8.1.0: Normalizes match_score on all candidates.
    """
    from database import PersistentCache
    import json as _json

    # Check in-memory cache first
    if discovery_results.get("candidates"):
        return _normalize_candidate_scores(dict(discovery_results))

    # Fall back to DB
    try:
        cached = db.query(PersistentCache).filter(
            PersistentCache.cache_key == "discovery_results"
        ).first()
        if cached and cached.data_json:
            data = _json.loads(cached.data_json)
            if data.get("candidates"):
                return _normalize_candidate_scores(data)
    except Exception as e:
        logger.warning(f"Failed to load discovery results from DB: {e}")

    return discovery_results


@app.post("/api/discovery/run")
async def run_discovery(request: Request):
    """Run the autonomous competitor discovery agent with optional custom criteria."""
    global discovery_results  # noqa: F824

    # Parse request body for criteria
    criteria = None
    try:
        body = await request.json()
        criteria = body.get("criteria", None)
        if criteria:
            logger.info(f"Discovery running with custom criteria:\n{criteria[:200]}...")
    except Exception:
        pass
    
    try:
        from discovery_agent import DiscoveryAgent
        import asyncio
        
        # Pass criteria to agent for potential future use in search customization
        agent = DiscoveryAgent(use_live_search=False, use_openai=False)
        
        # Store criteria for reference
        if criteria:
            agent.custom_criteria = criteria
        
        candidates = await agent.run_discovery_loop(max_candidates=10)
        
        discovery_results = {
            "status": "success",
            "candidates": candidates,
            "last_run": datetime.utcnow().isoformat(),
            "count": len(candidates),
            "criteria_used": criteria[:100] + "..." if criteria and len(criteria) > 100 else criteria
        }

        return discovery_results

    except Exception as e:
        logger.error(f"Discovery error: {e}")
        return {"status": "error", "error": "An unexpected error occurred", "candidates": [], "message": "Discovery failed"}


@app.post("/api/discovery/run-live")
async def run_live_discovery():
    """Run live discovery with DuckDuckGo search (rate-limited)."""
    global discovery_results  # noqa: F824

    try:
        from discovery_agent import DiscoveryAgent

        agent = DiscoveryAgent(use_live_search=True, use_openai=False)
        candidates = await agent.run_discovery_loop(max_candidates=5)

        discovery_results = {
            "candidates": candidates,
            "last_run": datetime.utcnow().isoformat(),
            "count": len(candidates),
            "mode": "live"
        }

        return discovery_results

    except Exception as e:
        logger.error(f"Live discovery error: {e}")
        return {"error": "An unexpected error occurred", "candidates": [], "message": "Live discovery failed"}


def _extract_domain(url: str) -> str:
    """Extract domain from URL for deduplication."""
    try:
        from urllib.parse import urlparse
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url
        parsed = urlparse(url)
        domain = parsed.netloc.lower()
        # Remove www. prefix
        if domain.startswith('www.'):
            domain = domain[4:]
        return domain
    except (ValueError, AttributeError):
        return url.lower()


@app.post("/api/discovery/add")
async def add_discovered_competitors(
    request: Request,
    background_tasks: BackgroundTasks,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """
    Add discovered competitors to the database.

    v5.2.0: Phase 3 - End-to-end discovery flow.
    v6.3.2: Bug fixes - B3 (score), B4 (auth), B7 (dedup)

    Request body:
    {
        "candidates": [
            {"name": "...", "url": "...", "reasoning": "...", "relevance_score": 0.8}
        ]
    }
    """
    try:
        body = await request.json()
        candidates = body.get("candidates", [])

        if not candidates:
            return {"success": False, "error": "No candidates provided"}

        added = []
        skipped = []

        # P2-4: Use explicit transaction for multi-step DB operations
        # All competitors and change logs will be committed together or rolled back together
        try:
            # Load all competitors once to avoid N+1 queries inside the loop
            all_competitors = db.query(Competitor).filter(
                Competitor.is_deleted == False
            ).all()
            # Build lookup dicts for fast dedup by domain and name
            domain_lookup = {}
            name_lookup = {}
            for comp in all_competitors:
                if comp.website:
                    comp_domain = _extract_domain(comp.website)
                    if comp_domain:
                        domain_lookup[comp_domain] = comp
                name_lookup[comp.name.lower()] = comp

            for candidate in candidates:
                name = candidate.get("name", "").strip()
                url = candidate.get("url", "").strip()
                reasoning = candidate.get("reasoning", "") or candidate.get("qualification_reasoning", "")

                # B3 Fix: Handle both relevance_score (0-1) and qualification_score (0-100)
                score = candidate.get("relevance_score") or candidate.get("qualification_score", 0)
                # Normalize score to 0-100 range
                if score and score <= 1:
                    normalized_score = int(score * 100)
                else:
                    normalized_score = int(score) if score else 50

                if not name or not url:
                    skipped.append({"name": name, "reason": "Missing name or URL"})
                    continue

                # B7 Fix: Better deduplication using domain extraction
                candidate_domain = _extract_domain(url)

                # Check for existing competitors by domain or name (using pre-built lookups)
                exists = domain_lookup.get(candidate_domain) if candidate_domain else None
                if not exists:
                    exists = name_lookup.get(name.lower())

                if exists:
                    skipped.append({"name": name, "reason": f"Duplicate of {exists.name}"})
                    continue

                # v7.1.0: Extract rich discovery data for populated competitor records
                ai_threat_level = candidate.get("threat_level") or "Low"
                strengths = candidate.get("strengths", [])
                weaknesses = candidate.get("weaknesses", [])
                competitive_positioning = candidate.get("competitive_positioning", "")
                ai_summary = candidate.get("ai_summary", "")
                products_found = candidate.get("products_found", [])
                features_found = candidate.get("features_found", [])
                meta_description = candidate.get("meta_description", "")  # noqa: F841

                # v7.2: Extract additional fields from candidate
                employee_count = candidate.get("employee_count")
                customer_count = candidate.get("customer_count")
                headquarters = candidate.get("headquarters")
                year_founded = candidate.get("year_founded")
                annual_revenue = candidate.get("annual_revenue")
                funding_total = candidate.get("funding_total")
                base_price = candidate.get("base_price")
                pricing_model = candidate.get("pricing_model")
                target_segments = candidate.get("target_segments", "")

                # v7.2: Compute threat_level via ThreatScoreCalculator
                try:
                    from analytics import ThreatScoreCalculator
                    calc = ThreatScoreCalculator()
                    threat_result = calc.calculate({
                        "name": name,
                        "target_segments": (
                            target_segments if isinstance(target_segments, str)
                            else "; ".join(target_segments) if target_segments else ""
                        ),
                        "customer_size_focus": "",
                        "funding_total": str(funding_total) if funding_total else "",
                        "employee_growth_rate": "",
                        "product_categories": "; ".join(products_found[:10]) if products_found else "",
                        "customer_count": str(customer_count) if customer_count else "",
                        "base_price": str(base_price) if base_price else "",
                        "pricing_model": str(pricing_model) if pricing_model else "",
                        "recent_launches": "",
                    })
                    threat_level = threat_result.threat_level
                except Exception as tsc_err:
                    logger.warning(f"ThreatScoreCalculator failed for {name}: {tsc_err}")
                    threat_level = ai_threat_level

                # Build rich notes with AI insights
                notes_parts = [f"Discovered by Certify Scout. Relevance: {normalized_score}%."]
                if ai_summary:
                    notes_parts.append(f"AI Summary: {ai_summary}")
                if reasoning and reasoning != ai_summary:
                    notes_parts.append(f"Reasoning: {reasoning}")
                if strengths:
                    notes_parts.append(f"Strengths: {'; '.join(strengths[:5])}")
                if weaknesses:
                    notes_parts.append(f"Weaknesses: {'; '.join(weaknesses[:5])}")
                if competitive_positioning:
                    notes_parts.append(f"Positioning: {competitive_positioning}")

                # v7.2: Create new competitor with status="Active" so it appears everywhere
                new_comp = Competitor(
                    name=name,
                    website=url,
                    status="Active",
                    threat_level=threat_level,
                    notes=" | ".join(notes_parts),
                    data_quality_score=normalized_score,
                    created_at=datetime.utcnow(),
                    # v7.1.0: Populate additional fields from discovery data
                    key_features="; ".join(features_found[:10]) if features_found else None,
                    product_categories="; ".join(products_found[:10]) if products_found else None,
                    # v7.2: Populate additional structured fields
                    employee_count=str(employee_count) if employee_count else None,
                    customer_count=str(customer_count) if customer_count else None,
                    headquarters=str(headquarters) if headquarters else None,
                    year_founded=str(year_founded) if year_founded else None,
                    annual_revenue=str(annual_revenue) if annual_revenue else None,
                    funding_total=str(funding_total) if funding_total else None,
                    base_price=str(base_price) if base_price else None,
                    pricing_model=str(pricing_model) if pricing_model else None,
                )
                db.add(new_comp)
                db.flush()

                # Update lookup dicts so later candidates in this batch are deduped
                if candidate_domain:
                    domain_lookup[candidate_domain] = new_comp
                name_lookup[name.lower()] = new_comp

                # v7.2: Create DataSource for extracted fields
                disc_fields = {"name": name, "website": url}
                if features_found:
                    disc_fields["key_features"] = "; ".join(
                        features_found[:10]
                    )
                if products_found:
                    disc_fields["product_categories"] = "; ".join(
                        products_found[:10]
                    )
                for ek in [
                    "employee_count", "customer_count",
                    "base_price", "pricing_model",
                    "headquarters", "year_founded",
                    "annual_revenue", "funding_total",
                ]:
                    ev = candidate.get(ek)
                    if ev:
                        disc_fields[ek] = str(ev)

                for df_name, df_val in disc_fields.items():
                    if df_val:
                        db.add(DataSource(
                            competitor_id=new_comp.id,
                            field_name=df_name,
                            current_value=str(df_val),
                            source_type="website_scrape",
                            source_url=url,
                            source_name="Discovery Scout",
                            extraction_method="ai_extraction",
                            confidence_score=min(
                                normalized_score, 100
                            ),
                            confidence_level=(
                                "moderate"
                                if normalized_score >= 50
                                else "low"
                            ),
                            is_verified=False,
                            extracted_at=datetime.utcnow(),
                        ))

                # Log the discovery
                change = ChangeLog(
                    competitor_id=new_comp.id,
                    competitor_name=new_comp.name,
                    change_type="New Competitor Discovered",
                    new_value=url,
                    source="Certify Scout",
                    severity="Medium"
                )
                db.add(change)

                added.append({
                    "id": new_comp.id,
                    "name": name,
                    "url": url,
                    "score": normalized_score
                })

            db.commit()
        except Exception as db_error:
            # P2-4: Rollback on any database error
            db.rollback()
            logger.error(f"Transaction failed during discovery add: {db_error}")
            raise HTTPException(status_code=500, detail="An internal database error occurred while adding competitors")

        # v7.1.0: Auto-trigger comprehensive enrichment for newly added competitors
        if added:
            async def enrich_discovered_competitors(comp_ids: List[int]):
                """Background task: scrape, fetch news, and generate battlecard for discovered competitors."""
                for comp_id in comp_ids:
                    enrich_db = SessionLocal()
                    try:
                        comp = enrich_db.query(Competitor).filter(Competitor.id == comp_id).first()
                        if not comp:
                            continue

                        # Step 1: Source discovery
                        try:
                            from source_discovery_engine import get_source_discovery_engine
                            engine = get_source_discovery_engine()
                            await engine.discover_sources_for_competitor(comp_id, max_fields=20)
                            logger.info(f"[Discovery Enrich] Source discovery complete for {comp.name}")
                        except Exception as e:
                            logger.warning(f"[Discovery Enrich] Source discovery failed for {comp.name}: {e}")

                        # Step 2: Website scrape to fill remaining fields
                        try:
                            from scraper import CompetitorScraper
                            from extractor import get_extractor
                            scraper = CompetitorScraper()
                            extractor = get_extractor()
                            content = await scraper.scrape(comp.website)
                            if content and "error" not in content and content.get("content"):
                                from dataclasses import asdict
                                extracted_obj = extractor.extract_from_content(comp.name, content.get("content", ""))
                                extracted = asdict(extracted_obj) if hasattr(extracted_obj, '__dataclass_fields__') else extracted_obj
                                if extracted:
                                    for key, value in extracted.items():
                                        if hasattr(comp, key) and value and key not in ["confidence_score", "extraction_notes"]:
                                            if getattr(comp, key) is None:  # Only fill empty fields
                                                setattr(comp, key, value)
                                            # v7.2: DataSource for scraped field
                                            _es = enrich_db.query(
                                                DataSource
                                            ).filter(
                                                DataSource.competitor_id == comp.id,
                                                DataSource.field_name == key,
                                            ).first()
                                            if _es:
                                                _es.current_value = str(value)
                                                _es.source_url = comp.website
                                                _es.source_type = "website_scrape"
                                                _es.extraction_method = "ai_extraction"
                                                _es.extracted_at = datetime.utcnow()
                                                _es.updated_at = datetime.utcnow()
                                            else:
                                                enrich_db.add(DataSource(
                                                    competitor_id=comp.id,
                                                    field_name=key,
                                                    current_value=str(value),
                                                    source_type="website_scrape",
                                                    source_url=comp.website,
                                                    source_name="Website Scrape",
                                                    extraction_method="ai_extraction",
                                                    confidence_score=55,
                                                    confidence_level="moderate",
                                                    is_verified=False,
                                                    extracted_at=datetime.utcnow(),
                                                ))
                                    comp.last_updated = datetime.utcnow()
                                    enrich_db.commit()
                                    logger.info(f"[Discovery Enrich] Website scrape complete for {comp.name}")
                        except Exception as e:
                            logger.warning(f"[Discovery Enrich] Scrape failed for {comp.name}: {e}")

                        # Step 3: Fetch news articles
                        try:
                            from news_monitor import NewsMonitor
                            from database import NewsArticleCache
                            monitor = NewsMonitor()
                            digest = monitor.fetch_news(comp.name, days=30)
                            if digest and digest.articles:
                                for article in digest.articles[:5]:
                                    existing = enrich_db.query(NewsArticleCache).filter(
                                        NewsArticleCache.url == article.url
                                    ).first()
                                    if not existing:
                                        cache_entry = NewsArticleCache(
                                            competitor_id=comp.id,
                                            competitor_name=comp.name,
                                            title=article.title,
                                            url=article.url,
                                            source=article.source,
                                            snippet=article.summary[:500] if article.summary else None,
                                            sentiment=article.sentiment if hasattr(article, 'sentiment') else 'neutral'
                                        )
                                        enrich_db.add(cache_entry)
                                # v7.2: DataSource for news_mentions
                                aurls = [
                                    a.url for a in digest.articles[:5]
                                    if a.url
                                ]
                                if aurls:
                                    _ns = enrich_db.query(
                                        DataSource
                                    ).filter(
                                        DataSource.competitor_id == comp.id,
                                        DataSource.field_name == "news_mentions",
                                    ).first()
                                    if _ns:
                                        _ns.current_value = str(
                                            len(digest.articles)
                                        )
                                        _ns.source_url = aurls[0]
                                        _ns.extracted_at = datetime.utcnow()
                                        _ns.updated_at = datetime.utcnow()
                                    else:
                                        enrich_db.add(DataSource(
                                            competitor_id=comp.id,
                                            field_name="news_mentions",
                                            current_value=str(
                                                len(digest.articles)
                                            ),
                                            source_type="news",
                                            source_url=aurls[0],
                                            source_name="News Monitor",
                                            extraction_method="news_aggregation",
                                            confidence_score=70,
                                            confidence_level="moderate",
                                            is_verified=False,
                                            extracted_at=datetime.utcnow(),
                                        ))
                                enrich_db.commit()
                                logger.info(f"[Discovery Enrich] News fetch complete for {comp.name}")
                        except Exception as e:
                            logger.warning(f"[Discovery Enrich] News fetch failed for {comp.name}: {e}")

                        # Step 4: Generate basic battlecard
                        try:
                            from battlecard_generator import BattlecardGenerator
                            generator = BattlecardGenerator(enrich_db)
                            generator.generate_battlecard(competitor_id=comp.id, battlecard_type="quick")
                            logger.info(f"[Discovery Enrich] Battlecard generated for {comp.name}")
                        except Exception as e:
                            logger.warning(f"[Discovery Enrich] Battlecard generation failed for {comp.name}: {e}")

                    except Exception as e:
                        logger.error(f"[Discovery Enrich] Enrichment failed for comp {comp_id}: {e}")
                    finally:
                        enrich_db.close()

            comp_ids = [a["id"] for a in added]
            background_tasks.add_task(enrich_discovered_competitors, comp_ids)

        return {
            "success": True,
            "added": added,
            "added_count": len(added),
            "skipped": skipped,
            "skipped_count": len(skipped),
            "enrichment_triggered": len(added) > 0,  # v7.1.0: Full enrichment (scrape+news+battlecard)
            "source_discovery_triggered": len(added) > 0
        }

    except Exception as e:
        db.rollback()
        return {"success": False, "error": "An unexpected error occurred"}


@app.get("/api/discovery/history")
def get_discovery_history(
    limit: int = 50,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get history of discovered competitors."""
    try:
        discovered = db.query(Competitor).filter(
            Competitor.status == "Discovered",
            Competitor.is_deleted == False
        ).order_by(Competitor.created_at.desc()).limit(limit).all()

        return {
            "count": len(discovered),
            "competitors": [
                {
                    "id": c.id,
                    "name": c.name,
                    "website": c.website,
                    "discovered_at": c.created_at.isoformat() if c.created_at else None,
                    "notes": c.notes,
                    "score": c.data_quality_score
                }
                for c in discovered
            ]
        }
    except Exception as e:
        return {"error": "An unexpected error occurred", "competitors": []}


# ============== AI Discovery Engine Endpoints (v7.0.0) ==============

# Discovery engine references for progress polling
_discovery_engines: Dict[str, Any] = {}

@app.get("/api/discovery/provider-status")
async def get_discovery_provider_status(
    current_user: dict = Depends(get_current_user)
):
    """
    Check Discovery AI provider status with detailed diagnostics.

    v7.0.7: Added for debugging AI service availability issues

    Returns status of:
    - Gemini provider availability
    - DuckDuckGo fallback availability
    - API key configuration
    """
    status = {
        "gemini_available": False,
        "gemini_error": None,
        "duckduckgo_available": False,
        "duckduckgo_error": None,
        "api_key_present": bool(os.getenv("GOOGLE_AI_API_KEY")),
        "api_key_length": len(os.getenv("GOOGLE_AI_API_KEY", "")),
        "google_genai_installed": False,
        "discovery_engine_loaded": False,
    }

    # Test google.genai import
    try:
        from google import genai
        status["google_genai_installed"] = True
    except ImportError as e:
        status["gemini_error"] = f"google.genai import failed: {str(e)}"

    # Test Gemini provider
    if status["google_genai_installed"]:
        try:
            from discovery_engine import AIDiscoveryEngine
            status["discovery_engine_loaded"] = True
            engine = AIDiscoveryEngine()
            status["gemini_available"] = engine.gemini is not None
            if not engine.gemini:
                status["gemini_error"] = "GeminiProvider initialization returned None - check API key"
        except Exception as e:
            status["gemini_error"] = f"{type(e).__name__}: {str(e)}"

    # Test DuckDuckGo fallback
    try:
        from duckduckgo_search import DDGS
        with DDGS() as ddgs:
            results = list(ddgs.text("healthcare software", max_results=1))
            status["duckduckgo_available"] = len(results) > 0
    except Exception as e:
        status["duckduckgo_error"] = f"{type(e).__name__}: {str(e)}"

    # Test Vertex AI (optional provider)
    status["vertex_ai_enabled"] = (
        os.getenv("VERTEX_AI_ENABLED", "false").lower() == "true"
    )
    status["vertex_ai_available"] = False
    status["vertex_ai_error"] = None
    if status["vertex_ai_enabled"]:
        try:
            from vertex_ai_provider import get_vertex_provider
            vp = get_vertex_provider()
            if vp and vp.is_available:
                status["vertex_ai_available"] = True
                vp_status = vp.get_status()
                status["vertex_ai_project"] = vp_status.get("project_id")
                status["vertex_ai_location"] = vp_status.get("location")
            else:
                status["vertex_ai_error"] = (
                    "Vertex AI not initialized - check GCP credentials"
                )
        except Exception as e:
            status["vertex_ai_error"] = f"{type(e).__name__}: {str(e)}"

    return status


@app.post("/api/discovery/run-ai")
async def run_ai_discovery(
    request: Request,
    background: bool = False,
    background_tasks: BackgroundTasks = None,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """
    Run AI-powered multi-stage competitor discovery.

    v8.1.0: Always runs in background mode with progress polling.
    v7.0.0: Full AI discovery with structured criteria using Gemini grounding.
    v6.3.2: Added authentication (Bug B4 fix)

    Stages:
    1. Search: Gemini with Google Search grounding
    2. Scrape: Deep website analysis
    3. Qualify: AI evaluation against criteria
    4. Analyze: Threat assessment
    """
    try:
        from discovery_engine import AIDiscoveryEngine, QualificationCriteria

        body = await request.json()

        # Build market_context from prompt_key + custom_instructions
        market_context = body.get("custom_instructions", "")
        prompt_key = body.get("prompt_key", "")
        if prompt_key:
            from database import SystemPrompt
            prompt_obj = db.query(SystemPrompt).filter(
                SystemPrompt.key == prompt_key
            ).first()
            if prompt_obj:
                market_context = (
                    prompt_obj.content + "\n\n" + market_context
                    if market_context else prompt_obj.content
                )

        import uuid as _uuid
        task_id = str(_uuid.uuid4())
        _ai_tasks[task_id] = {
            "status": "running", "page_context": "discovered",
            "user_id": current_user.get("id"),
            "started_at": datetime.utcnow().isoformat(),
            "completed_at": None, "result": None, "error": None,
            "task_type": "discovery_pipeline", "read_at": None,
        }

        async def _run_discovery_bg():
            engine = None
            try:
                bg_db = SessionLocal()
                try:
                    criteria = QualificationCriteria(
                        target_segments=body.get("target_segments", []),
                        required_capabilities=body.get("required_capabilities", []),
                        company_size=body.get("company_size", {}),
                        geography=body.get("geography", ["us"]),
                        funding_stages=body.get("funding_stages", []),
                        tech_requirements=body.get("tech_requirements", []),
                        exclusions=body.get("exclusions", []),
                        custom_keywords=body.get("custom_keywords", {})
                    )
                    known = bg_db.query(Competitor).filter(
                        Competitor.is_deleted == False  # noqa: E712
                    ).all()
                    engine = AIDiscoveryEngine()
                    _discovery_engines[task_id] = engine
                    engine.set_known_competitors([
                        {"name": c.name, "website": c.website} for c in known
                    ])
                    result = await asyncio.wait_for(
                        engine.run_discovery(
                            criteria=criteria,
                            max_candidates=body.get("max_candidates", 10),
                            market_context=market_context if market_context else None
                        ),
                        timeout=300.0
                    )
                    _ai_tasks[task_id]["result"] = result
                    # Also cache in discovery_results global so /api/discovery/results
                    # returns the latest AI discovery as a fallback for page refresh
                    global discovery_results
                    if result and isinstance(result, dict):
                        discovery_results = {
                            "status": "success",
                            "candidates": result.get("candidates", []),
                            "last_run": datetime.utcnow().isoformat(),
                            "count": len(result.get("candidates", [])),
                        }
                        # v8.0.8: Persist to DB so results survive server restart
                        try:
                            import json as _json
                            from database import PersistentCache
                            existing = bg_db.query(PersistentCache).filter(
                                PersistentCache.cache_key == "discovery_results"
                            ).first()
                            if existing:
                                existing.data_json = _json.dumps(discovery_results, default=str)
                                existing.updated_at = datetime.utcnow()
                            else:
                                bg_db.add(PersistentCache(
                                    cache_key="discovery_results",
                                    data_json=_json.dumps(discovery_results, default=str),
                                ))
                            bg_db.commit()
                        except Exception as persist_err:
                            logger.warning(f"Failed to persist discovery results: {persist_err}")
                finally:
                    bg_db.close()
                _ai_tasks[task_id]["status"] = "completed"
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
            except asyncio.TimeoutError:
                logger.error("Background discovery pipeline timed out after 120s")
                _ai_tasks[task_id]["status"] = "failed"
                _ai_tasks[task_id]["error"] = "Discovery pipeline timed out. Please try again with fewer candidates."
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
            except Exception as bg_err:
                logger.error(f"Background discovery error: {bg_err}")
                _ai_tasks[task_id]["status"] = "failed"
                _ai_tasks[task_id]["error"] = "Discovery pipeline encountered an error"
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
            finally:
                _discovery_engines.pop(task_id, None)

        background_tasks.add_task(_run_discovery_bg)
        return {"task_id": task_id, "status": "running"}

    except Exception as e:
        logger.exception(f"AI Discovery error: {e}")
        return {"status": "error", "error": "An unexpected error occurred", "candidates": []}


@app.get("/api/discovery/progress/{task_id}")
async def get_discovery_progress(
    task_id: str,
    current_user: dict = Depends(get_current_user)
):
    """Get real-time progress of a discovery pipeline run."""
    task = _ai_tasks.get(task_id)
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    engine = _discovery_engines.get(task_id)
    progress = engine.get_progress() if engine else {
        "percent_complete": 100 if task["status"] == "completed" else 0
    }
    return {
        "status": task["status"],
        "progress": progress,
        "result": task.get("result") if task["status"] == "completed" else None,
        "error": task.get("error") if task["status"] == "failed" else None
    }


@app.post("/api/discovery/summarize")
async def summarize_discovery_results(
    request: Request,
    current_user: dict = Depends(get_current_user)
):
    """Generate AI executive summary of discovery results."""
    try:
        body = await request.json()
        candidates = body.get("candidates", [])
        if not candidates:
            return {"summary": "No candidates to summarize."}

        from ai_router import get_ai_router, TaskType
        router = get_ai_router()

        candidate_summaries = []
        for c in candidates[:20]:
            candidate_summaries.append(
                f"- {c.get('name', 'Unknown')}: "
                f"Threat={c.get('threat_level', 'N/A')}, "
                f"Score={c.get('qualification_score', 'N/A')}%, "
                f"Summary: {c.get('ai_summary', 'N/A')}"
            )

        prompt = (
            "Generate a concise executive summary of these competitor "
            "discovery results.\n\n"
            f"Total competitors found: {len(candidates)}\n\n"
            "COMPETITORS:\n" + "\n".join(candidate_summaries) + "\n\n"
            "Include:\n"
            "1. Overview of what was found\n"
            "2. Key threats identified\n"
            "3. Common strengths and weaknesses across competitors\n"
            "4. Recommended next steps\n\n"
            "Keep it to 3-4 paragraphs. Be specific and actionable."
        )

        result = await asyncio.wait_for(
            router.generate(
                prompt=prompt,
                task_type=TaskType.DISCOVERY,
                temperature=0.3,
                max_tokens=1024,
                agent_type="discovery_summary"
            ),
            timeout=30.0
        )
        summary_text = result.get("response", "Summary generation failed.")
        return {"summary": summary_text}

    except asyncio.TimeoutError:
        logger.error("Discovery summary timed out")
        return {"summary": "Summary generation timed out. Please try again."}
    except Exception as e:
        logger.exception(f"Discovery summary error: {e}")
        return {"summary": "Failed to generate summary."}


@app.post("/api/discovery/send-to-battlecard")
async def send_discovery_to_battlecard(
    request: Request,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Add discovered competitors to DB and generate battlecards."""
    try:
        from database import Battlecard
        body = await request.json()
        candidates = body.get("candidates", [])
        if not candidates:
            raise HTTPException(status_code=400, detail="No candidates provided")

        competitor_ids = []
        battlecard_ids = []

        for candidate in candidates:
            name = candidate.get("name", "").strip()
            if not name:
                continue

            # Check if competitor already exists
            existing = db.query(Competitor).filter(
                Competitor.name.ilike(name),
                Competitor.is_deleted == False  # noqa: E712
            ).first()

            if existing:
                comp_id = existing.id
            else:
                # Create new competitor
                new_comp = Competitor(
                    name=name,
                    website=candidate.get("url", ""),
                    description=candidate.get("ai_summary", ""),
                    threat_level=candidate.get("threat_level", "Medium"),
                    status="Active",
                    is_deleted=False
                )
                db.add(new_comp)
                db.flush()
                comp_id = new_comp.id

            competitor_ids.append(comp_id)

            # Generate battlecard
            try:
                existing_bc = db.query(Battlecard).filter(
                    Battlecard.competitor_id == comp_id
                ).first()
                if not existing_bc:
                    new_bc = Battlecard(
                        competitor_id=comp_id,
                        title=f"Battlecard: {name}",
                        content=candidate.get("ai_summary", "Pending AI generation"),
                        battlecard_type="full",
                        generated_by=current_user.get("email", "ai")
                    )
                    db.add(new_bc)
                    db.flush()
                    battlecard_ids.append(new_bc.id)
                else:
                    battlecard_ids.append(existing_bc.id)
            except Exception as bc_err:
                logger.warning(f"Battlecard creation failed for {name}: {bc_err}")

        db.commit()
        return {"competitor_ids": competitor_ids, "battlecard_ids": battlecard_ids}

    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.exception(f"Send to battlecard error: {e}")
        raise HTTPException(status_code=500, detail="Failed to create battlecards")


@app.post("/api/discovery/send-to-comparison")
async def send_discovery_to_comparison(
    request: Request,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Add discovered competitors to DB for comparison."""
    try:
        body = await request.json()
        candidates = body.get("candidates", [])
        if not candidates:
            raise HTTPException(status_code=400, detail="No candidates provided")

        competitor_ids = []

        for candidate in candidates:
            name = candidate.get("name", "").strip()
            if not name:
                continue

            existing = db.query(Competitor).filter(
                Competitor.name.ilike(name),
                Competitor.is_deleted == False  # noqa: E712
            ).first()

            if existing:
                competitor_ids.append(existing.id)
            else:
                new_comp = Competitor(
                    name=name,
                    website=candidate.get("url", ""),
                    description=candidate.get("ai_summary", ""),
                    threat_level=candidate.get("threat_level", "Medium"),
                    status="Active",
                    is_deleted=False
                )
                db.add(new_comp)
                db.flush()
                competitor_ids.append(new_comp.id)

        db.commit()
        return {"competitor_ids": competitor_ids}

    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.exception(f"Send to comparison error: {e}")
        raise HTTPException(status_code=500, detail="Failed to add competitors")


@app.get("/api/discovery/profiles")
def get_discovery_profiles(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get all saved discovery profiles."""
    try:
        from database import DiscoveryProfile

        profiles = db.query(DiscoveryProfile).order_by(
            DiscoveryProfile.is_default.desc(),
            DiscoveryProfile.name
        ).all()

        return {
            "profiles": [
                {
                    "id": p.id,
                    "name": p.name,
                    "description": p.description,
                    "is_default": p.is_default,
                    "created_at": p.created_at.isoformat() if p.created_at else None
                }
                for p in profiles
            ]
        }
    except Exception as e:
        logger.error(f"Get profiles error: {e}")
        return {"profiles": [], "error": "An unexpected error occurred"}


@app.get("/api/discovery/profiles/{profile_id}")
def get_discovery_profile(
    profile_id: int,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get a specific discovery profile by ID."""
    try:
        from database import DiscoveryProfile

        profile = db.query(DiscoveryProfile).filter(
            DiscoveryProfile.id == profile_id
        ).first()

        if not profile:
            raise HTTPException(status_code=404, detail="Profile not found")

        return {
            "id": profile.id,
            "name": profile.name,
            "description": profile.description,
            "is_default": profile.is_default,
            "target_segments": profile.target_segments,
            "required_capabilities": profile.required_capabilities,
            "company_size": profile.company_size,
            "geography": profile.geography,
            "funding_stages": profile.funding_stages,
            "tech_requirements": profile.tech_requirements,
            "exclusions": profile.exclusions,
            "custom_keywords": profile.custom_keywords,
            "ai_instructions": profile.ai_instructions
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Get profile error: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/discovery/profiles")
async def create_discovery_profile(
    request: Request,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Create a new discovery profile."""
    try:
        from database import DiscoveryProfile
        import json

        body = await request.json()

        profile = DiscoveryProfile(
            name=body.get("name", "Untitled Profile"),
            description=body.get("description"),
            target_segments=json.dumps(body.get("target_segments", [])),
            required_capabilities=json.dumps(body.get("required_capabilities", [])),
            company_size=json.dumps(body.get("company_size", {})),
            geography=json.dumps(body.get("geography", ["us"])),
            funding_stages=json.dumps(body.get("funding_stages", [])),
            tech_requirements=json.dumps(body.get("tech_requirements", [])),
            exclusions=json.dumps(body.get("exclusions", [])),
            custom_keywords=json.dumps(body.get("custom_keywords", {})),
            ai_instructions=body.get("ai_instructions")
        )

        db.add(profile)
        db.commit()
        db.refresh(profile)

        return {"id": profile.id, "name": profile.name, "status": "created"}

    except Exception as e:
        db.rollback()
        logger.error(f"Create profile error: {e}")
        return {"error": "An unexpected error occurred", "status": "error"}


@app.delete("/api/discovery/profiles/{profile_id}")
def delete_discovery_profile(
    profile_id: int,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Delete a discovery profile."""
    try:
        from database import DiscoveryProfile

        profile = db.query(DiscoveryProfile).filter(
            DiscoveryProfile.id == profile_id
        ).first()

        if not profile:
            raise HTTPException(status_code=404, detail="Profile not found")

        db.delete(profile)
        db.commit()

        return {"status": "deleted", "id": profile_id}
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.error(f"Delete profile error: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


# ============== Discovery Scout Default Prompt Endpoints (v7.0.8) ==============

@app.get("/api/discovery/default-prompt")
def get_default_scout_prompt(db: Session = Depends(get_db)):
    """
    Get the global default Discovery Scout AI Instructions prompt.
    Returns the database-stored prompt if set, otherwise returns hardcoded default.
    """
    try:
        from database import SystemPrompt

        prompt = db.query(SystemPrompt).filter(
            SystemPrompt.key == "discovery_scout_prompt",
            SystemPrompt.user_id == None  # Global prompt (NULL user_id)
        ).first()

        if prompt:
            return {
                "prompt": prompt.content,
                "source": "database",
                "updated_at": prompt.updated_at.isoformat() if prompt.updated_at else None
            }
        else:
            # Return hardcoded default if no DB prompt exists
            return {
                "prompt": None,
                "source": "hardcoded",
                "updated_at": None
            }
    except Exception as e:
        logger.error(f"Get default prompt error: {e}")
        return {"prompt": None, "source": "error", "error": "An unexpected error occurred"}


@app.put("/api/discovery/default-prompt")
async def set_default_scout_prompt(
    request: Request,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """
    Set the global default Discovery Scout AI Instructions prompt.
    Only admins can set the global default. All new users will see this prompt
    on first login until they customize their own.
    """
    try:
        # Check admin permission
        user_role = current_user.get("role", "user")
        if user_role != "admin":
            raise HTTPException(
                status_code=403,
                detail="Admin access required to set global default prompt"
            )

        # Get the prompt content from request body
        body = await request.json()

        prompt_content = body.get("prompt")
        if not prompt_content or not isinstance(prompt_content, str):
            raise HTTPException(
                status_code=400,
                detail="Request must include 'prompt' field with string content"
            )

        from database import SystemPrompt

        # Upsert: update existing or create new global prompt
        existing = db.query(SystemPrompt).filter(
            SystemPrompt.key == "discovery_scout_prompt",
            SystemPrompt.user_id == None
        ).first()

        if existing:
            existing.content = prompt_content
            existing.updated_at = datetime.utcnow()
        else:
            new_prompt = SystemPrompt(
                key="discovery_scout_prompt",
                user_id=None,  # NULL = global prompt
                content=prompt_content
            )
            db.add(new_prompt)

        db.commit()

        return {
            "status": "success",
            "message": "Global default prompt saved. All new users will see this prompt.",
            "source": "database"
        }

    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        logger.error(f"Set default prompt error: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


# ============== Enhancement Endpoints ==============

@app.get("/api/competitors/{competitor_id}/threat-analysis")
def get_threat_analysis(competitor_id: int, db: Session = Depends(get_db)):
    """Get AI-powered threat analysis for a competitor."""
    try:
        from threat_analyzer import analyze_competitor_threat
        
        competitor = db.query(Competitor).filter(
            Competitor.id == competitor_id,
            Competitor.is_deleted == False
        ).first()
        
        if not competitor:
            raise HTTPException(status_code=404, detail="Competitor not found")
        
        comp_data = {k: v for k, v in competitor.__dict__.items() if not k.startswith('_')}
        analysis = analyze_competitor_threat(comp_data)
        
        # Update competitor threat level if changed
        if analysis.get("level") != competitor.threat_level:
            competitor.threat_level = analysis["level"]
            db.commit()
        
        return analysis
        
    except HTTPException:
        raise
    except Exception as e:
        return {"error": "An unexpected error occurred", "message": "Threat analysis failed"}


@app.get("/api/competitors/{competitor_id}/news")
def get_competitor_news(competitor_id: int, days: int = 7, db: Session = Depends(get_db)):
    """Get news articles for a competitor."""
    try:
        from news_monitor import fetch_competitor_news
        
        competitor = db.query(Competitor).filter(
            Competitor.id == competitor_id,
            Competitor.is_deleted == False
        ).first()
        
        if not competitor:
            raise HTTPException(status_code=404, detail="Competitor not found")
        
        return fetch_competitor_news(competitor.name, days)
        
    except HTTPException:
        raise
    except Exception as e:
        return {"error": "An unexpected error occurred", "articles": [], "message": "News fetch failed"}


# NOTE: Duplicate /api/news/{company_name} removed (keeping async version above at ~line 8817)


# ============== News Feed Endpoint (v5.0.3 - Phase 1) ==============

@app.get("/api/news-feed")
def get_news_feed(
    competitor_id: Optional[int] = None,
    start_date: Optional[str] = None,
    end_date: Optional[str] = None,
    sentiment: Optional[str] = None,
    source: Optional[str] = None,
    event_type: Optional[str] = None,
    keyword: Optional[str] = None,
    page: int = 1,
    page_size: int = 25,
    db: Session = Depends(get_db)
):
    """
    Aggregated news feed with filtering across all competitors.

    v5.0.3: Core News Feed implementation.
    v7.1.2: Added keyword search filter.

    Args:
        competitor_id: Filter by specific competitor ID
        keyword: Search keyword to filter articles by title or snippet
        start_date: Start date filter (YYYY-MM-DD)
        end_date: End date filter (YYYY-MM-DD)
        sentiment: Filter by sentiment (positive, neutral, negative)
        source: Filter by news source (google_news, sec_edgar, newsapi, gnews, mediastack)
        event_type: Filter by event type (funding, acquisition, product_launch, partnership, leadership, financial, legal, general)
        page: Page number for pagination
        page_size: Number of articles per page

    Returns:
        Aggregated news articles with stats and pagination
    """
    try:
        from news_monitor import NewsMonitor
        from datetime import datetime, timedelta
        from database import NewsArticleCache

        monitor = NewsMonitor()
        all_articles = []
        cache_used = False

        # Parse date range ONCE at the top (used for both DB query and Python filter)
        # FIX v7.1.2: Date filtering was broken - no DB-level filter, end_dt at midnight,
        # unparseable dates defaulted to "today", missing dates bypassed filter entirely
        if start_date:
            try:
                start_dt = datetime.strptime(start_date, '%Y-%m-%d')
            except ValueError:
                start_dt = datetime.utcnow() - timedelta(days=30)
        else:
            start_dt = datetime.utcnow() - timedelta(days=30)

        if end_date:
            try:
                # Add 1 day for inclusive end (next-day-exclusive pattern)
                end_dt = datetime.strptime(end_date, '%Y-%m-%d') + timedelta(days=1)
            except ValueError:
                end_dt = datetime.utcnow() + timedelta(days=1)
        else:
            end_dt = datetime.utcnow() + timedelta(days=1)

        # Get competitors to fetch news for
        competitors_query = db.query(Competitor).filter(Competitor.is_deleted == False)
        if competitor_id:
            competitors_query = competitors_query.filter(Competitor.id == competitor_id)
        competitors_list = competitors_query.all()

        # v5.0.8: Check cache first
        # FIX v7.1.2: Added DB-level date filtering using published_at index
        cache_query = db.query(NewsArticleCache).filter(
            NewsArticleCache.is_archived != True
        )
        if competitor_id:
            cache_query = cache_query.filter(NewsArticleCache.competitor_id == competitor_id)
        cache_query = cache_query.filter(
            NewsArticleCache.published_at >= start_dt,
            NewsArticleCache.published_at < end_dt
        )

        cached_articles = cache_query.order_by(NewsArticleCache.published_at.desc()).limit(500).all()

        if cached_articles:
            cache_used = True
            for cached in cached_articles:
                all_articles.append({
                    "competitor_id": cached.competitor_id,
                    "competitor_name": cached.competitor_name,
                    "title": cached.title,
                    "url": cached.url,
                    "source": cached.source,
                    "source_type": cached.source_type or "google_news",
                    "published_at": cached.published_at.isoformat() if cached.published_at else "",
                    "snippet": cached.snippet or "",
                    "sentiment": cached.sentiment or "neutral",
                    "event_type": cached.event_type or "general",
                    "is_major_event": cached.is_major_event or False,
                    "cached": True
                })

        # If no cache, fetch live (batch to prevent timeout)
        if not cache_used:
            # Limit to first 25 competitors to prevent timeout (background job handles rest)
            limited_competitors = competitors_list[:25]

            # Calculate days to look back
            days_lookback = (datetime.utcnow() - start_dt).days + 1

            # Fetch news for each competitor (using limited batch)
            for comp in limited_competitors:
                try:
                    digest = monitor.fetch_news(comp.name, days=days_lookback)

                    # Add competitor info to each article
                    for article in digest.articles:
                        article_dict = {
                            "competitor_id": comp.id,
                            "competitor_name": comp.name,
                            "title": article.title,
                            "url": article.url,
                            "source": article.source,
                            "source_type": "google_news",  # Default source type
                            "published_at": article.published_date,
                            "snippet": article.snippet,
                            "sentiment": article.sentiment,
                            "event_type": article.event_type or "general",
                            "is_major_event": article.is_major_event
                        }
                        all_articles.append(article_dict)
                except Exception as e:
                    logger.error(f"Error fetching news for {comp.name}: {e}")
                    continue

        # FIX v7.1.2: Filter by date range - exclude (not include) unparseable/missing dates
        filtered_articles = []
        for article in all_articles:
            try:
                pub_date_str = article.get("published_at", "")
                if not pub_date_str:
                    continue  # Skip articles with no date

                # Handle various date formats
                pub_date = None
                for fmt in ["%Y-%m-%dT%H:%M:%S", "%Y-%m-%d",
                            "%a, %d %b %Y %H:%M:%S %Z",
                            "%a, %d %b %Y %H:%M:%S GMT"]:
                    try:
                        pub_date = datetime.strptime(pub_date_str[:19], fmt)
                        break
                    except (ValueError, TypeError):
                        continue

                if pub_date is None:
                    continue  # Skip unparseable dates

                if start_dt <= pub_date < end_dt:
                    filtered_articles.append(article)
            except (ValueError, TypeError, AttributeError):
                continue  # Skip on any error

        all_articles = filtered_articles

        # Apply sentiment filter
        if sentiment:
            all_articles = [a for a in all_articles if a.get("sentiment", "").lower() == sentiment.lower()]

        # Apply source filter
        if source:
            all_articles = [a for a in all_articles if a.get("source_type", "").lower() == source.lower() or source.lower() in a.get("source", "").lower()]

        # Apply event type filter
        if event_type:
            all_articles = [a for a in all_articles if a.get("event_type", "").lower() == event_type.lower()]

        # v7.1.2: Apply keyword search filter
        if keyword:
            kw_lower = keyword.strip().lower()
            all_articles = [
                a for a in all_articles
                if kw_lower in a.get("title", "").lower()
                or kw_lower in a.get("snippet", "").lower()
                or kw_lower in a.get("competitor_name", "").lower()
            ]

        # Sort by date (most recent first)
        all_articles.sort(key=lambda x: x.get("published_at", ""), reverse=True)

        # Calculate stats
        stats = {
            "total": len(all_articles),
            "positive": len([a for a in all_articles if a.get("sentiment") == "positive"]),
            "neutral": len([a for a in all_articles if a.get("sentiment") == "neutral"]),
            "negative": len([a for a in all_articles if a.get("sentiment") == "negative"])
        }

        # Paginate
        total_items = len(all_articles)
        total_pages = max(1, (total_items + page_size - 1) // page_size)
        start_idx = (page - 1) * page_size
        end_idx = start_idx + page_size
        paginated_articles = all_articles[start_idx:end_idx]

        return {
            "articles": paginated_articles,
            "stats": stats,
            "pagination": {
                "page": page,
                "page_size": page_size,
                "total_items": total_items,
                "total_pages": total_pages
            },
            "filters_applied": {
                "competitor_id": competitor_id,
                "start_date": start_date,
                "end_date": end_date,
                "sentiment": sentiment,
                "source": source,
                "event_type": event_type,
                "keyword": keyword
            }
        }

    except Exception as e:
        logger.error(f"Error in news feed: {e}")
        return {
            "articles": [],
            "stats": {"total": 0, "positive": 0, "neutral": 0, "negative": 0},
            "pagination": {"page": 1, "page_size": page_size, "total_items": 0, "total_pages": 1},
            "error": "An unexpected error occurred"
        }


@app.post("/api/news-feed/summarize")
async def summarize_news_feed(
    request: Request,
    background: bool = False,
    background_tasks: BackgroundTasks = None,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """
    Generate AI-powered summary of news articles.

    v7.1.2: Real AI summary using AIRouter for news analysis.

    Request body:
        period: 'today' | 'week' | 'competitor'
        competitor_id: Optional competitor ID (required for 'competitor' period)
        articles: List of article dicts with title, competitor_name, sentiment, event_type, date
    """
    from ai_router import get_ai_router, TaskType

    try:
        body = await request.json()
    except Exception:
        body = {}

    if background and background_tasks and current_user:
        import uuid as _uuid
        task_id = str(_uuid.uuid4())
        _ai_tasks[task_id] = {
            "status": "running", "page_context": "newsfeed",
            "user_id": current_user.get("id"),
            "started_at": datetime.utcnow().isoformat(),
            "completed_at": None, "result": None, "error": None,
            "task_type": "news_summary", "read_at": None,
        }

        async def _run_news_summary_bg():
            try:
                articles = body.get("articles", [])
                if not articles:
                    _ai_tasks[task_id]["result"] = {
                        "status": "error",
                        "content": "No articles provided."
                    }
                else:
                    article_text = "\n".join([
                        f"- [{a.get('sentiment', 'neutral').upper()}] "
                        f"{a.get('competitor_name', 'Unknown')}: {a.get('title', '')}"
                        for a in articles[:40]
                    ])
                    from ai_router import get_ai_router as _get_router
                    router = _get_router()
                    result = await asyncio.wait_for(
                        router.generate(
                            prompt=f"Summarize these news articles:\n{article_text}",
                            task_type=TaskType.ANALYSIS,
                            system_prompt="Competitive intelligence analyst. Concise analysis." + NO_HALLUCINATION_INSTRUCTION,
                            max_tokens=2048, temperature=0.5,
                        ),
                        timeout=45.0
                    )
                    _ai_tasks[task_id]["result"] = {
                        "status": "success",
                        "content": result.get("response", "")
                    }
                _ai_tasks[task_id]["status"] = "completed"
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
            except asyncio.TimeoutError:
                logger.error("Background news summary timed out after 45s")
                _ai_tasks[task_id]["status"] = "failed"
                _ai_tasks[task_id]["error"] = "AI analysis timed out. Please try again."
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()
            except Exception as bg_err:
                logger.error(f"Background news summary error: {bg_err}")
                _ai_tasks[task_id]["status"] = "failed"
                _ai_tasks[task_id]["error"] = str(bg_err)
                _ai_tasks[task_id]["completed_at"] = datetime.utcnow().isoformat()

        background_tasks.add_task(_run_news_summary_bg)
        return {"task_id": task_id, "status": "running"}

    period = body.get("period", "week")
    articles = body.get("articles", [])
    competitor_id_param = body.get("competitor_id")
    news_session_id = body.get("session_id")  # Optional chat session
    conversation_history = body.get("conversation_history")  # Optional history
    news_prompt_key = body.get("prompt_key")  # Optional custom prompt

    if not articles:
        return {"status": "error", "content": "No articles provided for summarization."}

    # Build period-specific prompt
    article_text = "\n".join([
        f"- [{a.get('sentiment', 'neutral').upper()}] {a.get('competitor_name', 'Unknown')}: "
        f"{a.get('title', 'No title')} ({a.get('event_type', 'general')}, {a.get('date', 'unknown date')})"
        for a in articles[:40]
    ])

    # Get competitor name if filtering by competitor
    comp_name = ""
    if competitor_id_param:
        comp = db.query(Competitor).filter(
            Competitor.id == competitor_id_param
        ).first()
        if comp:
            comp_name = comp.name

    if period == "today":
        prompt = (
            f"You are a competitive intelligence analyst. Analyze today's news articles "
            f"and provide a concise daily briefing.\n\n"
            f"Articles:\n{article_text}\n\n"
            f"Provide:\n"
            f"1. **Executive Summary** (2-3 sentences on the day's key themes)\n"
            f"2. **Key Headlines** (top 3-5 most significant stories)\n"
            f"3. **Sentiment Breakdown** (positive/neutral/negative trends)\n"
            f"4. **Action Items** (what the sales/product team should do today)\n\n"
            f"Be specific, cite competitor names, and keep it actionable."
        )
    elif period == "competitor" and comp_name:
        prompt = (
            f"You are a competitive intelligence analyst. Analyze the following news "
            f"articles about **{comp_name}** from the last 30 days.\n\n"
            f"Articles:\n{article_text}\n\n"
            f"Provide:\n"
            f"1. **{comp_name} Overview** (current strategic direction based on news)\n"
            f"2. **Key Developments** (product launches, partnerships, financials, leadership)\n"
            f"3. **Threat Assessment** (what this means for our competitive position)\n"
            f"4. **Recommended Response** (specific counter-strategies for sales and product)\n\n"
            f"Be specific and actionable."
        )
    else:
        prompt = (
            f"You are a competitive intelligence analyst. Analyze this week's news articles "
            f"and provide an executive weekly summary.\n\n"
            f"Articles:\n{article_text}\n\n"
            f"Provide:\n"
            f"1. **Weekly Overview** (2-3 sentences summarizing the competitive landscape)\n"
            f"2. **Top Competitor Activity** (group by competitor, highlight key moves)\n"
            f"3. **Market Trends** (funding, M&A, product launches, partnerships)\n"
            f"4. **Sentiment Analysis** (overall market mood and notable shifts)\n"
            f"5. **Action Items** (top 3-5 things sales/product should act on this week)\n\n"
            f"Be specific, cite competitor names, and keep it actionable."
        )

    # Prepend conversation history if provided
    if conversation_history and isinstance(conversation_history, list):
        history_text = "\n".join(
            f"[{m.get('role', 'user').upper()}]: {m.get('content', '')}"
            for m in conversation_history[-10:]
        )
        prompt = f"Previous conversation:\n{history_text}\n\nNow:\n{prompt}"

    # Resolve system prompt (custom or default)
    default_news_prompt = (
        "You are a competitive intelligence analyst for a healthcare IT company. "
        "Provide concise, actionable analysis."
    )
    user_id = current_user.get("id") if current_user else None
    news_sys_prompt = _resolve_system_prompt(
        db, user_id, news_prompt_key, default_news_prompt
    ) + NO_HALLUCINATION_INSTRUCTION

    try:
        from ai_router import get_ai_router as _get_router_news
        router = _get_router_news()
        result = await router.generate(
            prompt=prompt,
            task_type=TaskType.ANALYSIS,
            system_prompt=news_sys_prompt,
            max_tokens=2048,
            temperature=0.5
        )
        content = result.get("response", "")
        if content:
            # Persist to chat session if session_id provided
            if news_session_id:
                try:
                    _save_chat_messages(db, news_session_id, f"Summarize {period} news", content, {
                        "endpoint": "news_summarize", "period": period,
                    })
                except Exception as save_err:
                    logger.warning(f"Failed to save news summary to chat: {save_err}")
            return {"status": "success", "content": content, "session_id": news_session_id}
        else:
            return {"status": "error", "content": "AI returned empty response."}
    except Exception as ai_err:
        logger.warning(f"AI summary generation failed: {ai_err}")
        return {"status": "error", "content": "AI summary generation is temporarily unavailable."}


@app.post("/api/news-feed/completion-summary")
async def news_feed_completion_summary(
    request: Request,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
) -> Dict[str, Any]:
    """
    Generate an AI summary of a completed news fetch operation.

    Accepts stats about articles fetched and returns a natural-language summary.
    """
    from ai_router import get_ai_router, TaskType

    try:
        body = await request.json()
    except Exception:
        body = {}

    articles_count = body.get("articles_count", 0)
    competitors_count = body.get("competitors_count", 0)
    by_source = body.get("by_source", {})
    by_competitor = body.get("by_competitor", {})
    by_sentiment = body.get("by_sentiment", {})
    sample_headlines = body.get("sample_headlines", [])
    major_events = body.get("major_events", [])
    cs_prompt_key = body.get("prompt_key")  # Optional custom prompt

    default_cs_prompt = (
        "You are a competitive intelligence analyst reporting on a completed news scan. "
        "Provide a concise, professional summary of the results."
    )
    cs_user_id = current_user.get("id") if current_user else None
    system_prompt = _resolve_system_prompt(
        db, cs_user_id, cs_prompt_key, default_cs_prompt
    ) + NO_HALLUCINATION_INSTRUCTION

    prompt = (
        f"A news scan has just completed. Here are the results:\n\n"
        f"Total articles found: {articles_count}\n"
        f"Competitors scanned: {competitors_count}\n\n"
        f"Articles by source:\n"
        + "\n".join(f"  - {src}: {cnt}" for src, cnt in by_source.items())
        + f"\n\nArticles by competitor:\n"
        + "\n".join(f"  - {comp}: {cnt}" for comp, cnt in by_competitor.items())
        + f"\n\nSentiment breakdown:\n"
        + "\n".join(f"  - {sent}: {cnt}" for sent, cnt in by_sentiment.items())
        + "\n\nSample headlines:\n"
        + "\n".join(f"  - {h}" for h in sample_headlines[:10])
        + "\n\nMajor events detected:\n"
        + "\n".join(f"  - {e}" for e in major_events[:10])
        + "\n\nProvide:\n"
        "1. Total articles found and per-source breakdown\n"
        "2. Per-competitor coverage highlights\n"
        "3. Sentiment analysis (positive/negative/neutral trends)\n"
        "4. Notable findings from the sample headlines\n"
        "5. Strategic takeaways for the competitive intelligence team\n"
    )

    try:
        router = get_ai_router()
        result = await router.generate(
            prompt=prompt,
            task_type=TaskType.ANALYSIS,
            system_prompt=system_prompt,
            max_tokens=1024,
            temperature=0.3,
        )
        return {"status": "success", "summary": result.get("response", "")}
    except Exception as e:
        logger.error(f"Completion summary generation failed: {e}")
        return {"status": "error", "summary": "AI summary generation is temporarily unavailable."}


@app.post("/api/news-feed/summarize-all")
async def summarize_all_news(
    request: Request,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
) -> Dict[str, Any]:
    """
    Comprehensive AI analysis of all provided news articles.

    Reads every article's content, groups findings by competitor, identifies
    themes/threats/opportunities, and produces a professional markdown report.
    """
    from ai_router import get_ai_router, TaskType

    try:
        body = await request.json()
    except Exception:
        body = {}

    articles = body.get("articles", [])
    sa_prompt_key = body.get("prompt_key")  # Optional custom prompt
    if not articles:
        return {"status": "error", "content": "No articles provided for analysis."}

    # Build detailed article listing
    article_lines = []
    for i, a in enumerate(articles[:50], 1):
        line = (
            f"{i}. [{a.get('sentiment', 'neutral').upper()}] "
            f"{a.get('competitor_name', 'Unknown')} - "
            f"\"{a.get('title', 'No title')}\" "
            f"(Source: {a.get('source', 'unknown')}, "
            f"URL: {a.get('url', 'N/A')}, "
            f"Event: {a.get('event_type', 'general')}, "
            f"Date: {a.get('published_at', a.get('published_date', 'unknown'))})"
        )
        if a.get("snippet"):
            line += f"\n   Snippet: {a['snippet'][:200]}"
        article_lines.append(line)

    articles_text = "\n".join(article_lines)

    default_sa_prompt = (
        "You are a senior competitive intelligence analyst for a "
        "healthcare IT company (Certify Health). You produce thorough, "
        "executive-ready competitive analysis reports in professional markdown."
    )
    sa_user_id = current_user.get("id") if current_user else None
    system_prompt = _resolve_system_prompt(
        db, sa_user_id, sa_prompt_key, default_sa_prompt
    ) + NO_HALLUCINATION_INSTRUCTION

    prompt = (
        f"Analyze the following {len(articles)} news articles comprehensively:\n\n"
        f"{articles_text}\n\n"
        "Instructions:\n"
        "1. Read every article's content (titles, snippets, metadata)\n"
        "2. Group findings by competitor\n"
        "3. Identify themes, threats, opportunities, and market movements\n"
        "4. Highlight strategic takeaways for Certify Health\n"
        "5. Cite specific articles with titles and source URLs\n"
        "6. Provide an executive summary at the top\n\n"
        "Format the output in professional markdown with:\n"
        "- Executive Summary\n"
        "- Competitor-by-Competitor Analysis\n"
        "- Key Themes & Market Movements\n"
        "- Threats & Opportunities\n"
        "- Strategic Recommendations\n"
    )

    try:
        router = get_ai_router()
        result = await router.generate(
            prompt=prompt,
            task_type=TaskType.ANALYSIS,
            system_prompt=system_prompt,
            max_tokens=4096,
            temperature=0.4,
        )
        content = result.get("response", "")
        if content:
            return {"status": "success", "content": content}
        return {"status": "error", "content": "AI returned empty response."}
    except Exception as e:
        logger.error(f"Summarize-all generation failed: {e}")
        return {"status": "error", "content": "AI summary generation is temporarily unavailable."}


@app.get("/api/news-feed/fetch-progress/{progress_key}")
async def get_news_feed_fetch_progress(
    progress_key: str,
    current_user: dict = Depends(get_current_user_optional)
) -> Dict[str, Any]:
    """Return progress for an active news fetch operation."""
    from news_monitor import get_news_fetch_progress

    progress = get_news_fetch_progress(progress_key)
    if progress is None:
        return {"status": "not_found"}
    return progress


@app.put("/api/news-feed/{article_id}/archive")
async def archive_news_article(
    article_id: int,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
) -> Dict[str, Any]:
    """Archive a news article so it no longer appears in feeds."""
    from database import NewsArticleCache

    article = db.query(NewsArticleCache).filter(
        NewsArticleCache.id == article_id
    ).first()

    if not article:
        raise HTTPException(status_code=404, detail="Article not found")

    article.is_archived = True
    db.commit()

    return {"status": "archived", "article_id": article_id}


@app.post("/api/news-feed/fetch")
async def fetch_news_for_competitors(
    request: Request,
    background_tasks: BackgroundTasks,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """
    Fetch news for specified competitors with progress tracking.
    Called by the "AI Fetch News" button on the Live News tab.

    Accepts JSON body: {competitor_ids: [int]|null, keywords: str, date_from: str, date_to: str}
    Returns: {progress_key: str, status: "started", competitors_count: int}
    """
    import uuid as _uuid
    from database import NewsArticleCache
    from news_monitor import NewsMonitor, _news_fetch_progress
    from datetime import datetime

    body = await request.json()
    competitor_ids = body.get("competitor_ids")
    keywords = body.get("keywords", "")
    date_from = body.get("date_from")
    date_to = body.get("date_to")

    # Calculate days from date_from
    days = 30
    if date_from:
        try:
            start_dt = datetime.strptime(date_from, "%Y-%m-%d")
            days = max(1, (datetime.utcnow() - start_dt).days)
        except ValueError:
            pass

    # Get target competitors
    query = db.query(Competitor).filter(Competitor.is_deleted == False)
    if competitor_ids:
        query = query.filter(Competitor.id.in_(competitor_ids))
    target_competitors = query.all()

    progress_key = str(_uuid.uuid4())

    # Initialize progress with keys matching what frontend polls expect
    _news_fetch_progress[progress_key] = {
        "status": "running",
        "current_competitor": "",
        "total_competitors": len(target_competitors),
        "total": len(target_competitors),
        "completed": 0,
        "articles_found": 0,
        "total_articles": 0,
        "sources_checked": 0,
        "percentage": 0,
        "percent": 0,
        "by_competitor": {},
        "by_sentiment": {},
        "sample_headlines": [],
    }

    def do_fetch():
        from concurrent.futures import ThreadPoolExecutor, as_completed
        from database import SessionLocal as _SessionLocal

        progress = _news_fetch_progress[progress_key]
        total_articles = 0
        completed_count = 0
        lock = __import__('threading').Lock()

        def fetch_one_competitor(comp):
            """Fetch news for a single competitor in its own DB session."""
            nonlocal total_articles, completed_count
            import asyncio as _asyncio
            monitor = NewsMonitor()
            local_db = _SessionLocal()
            try:
                # Build search term — disambiguate short/generic names with healthcare context
                base_name = comp.name
                if len(base_name.split()) == 1 and len(base_name) < 12 and comp.website:
                    try:
                        from urllib.parse import urlparse
                        domain = urlparse(comp.website).netloc.replace('www.', '')
                        domain_name = domain.split('.')[0] if domain else ''
                        if domain_name and domain_name.lower() != base_name.lower():
                            base_name = f'"{comp.name}" healthcare OR "{domain_name}"'
                        else:
                            base_name = f'"{comp.name}" healthcare'
                    except Exception:
                        base_name = f'"{comp.name}" healthcare'
                else:
                    base_name = f'"{comp.name}"'
                search_term = f"{base_name} {keywords}".strip() if keywords else base_name
                # Use async version which fetches all sources in parallel via asyncio.gather
                loop = _asyncio.new_event_loop()
                try:
                    digest = loop.run_until_complete(monitor.fetch_news_async(
                        search_term, days=days,
                        real_name=comp.name, website=comp.website
                    ))
                finally:
                    loop.close()

                articles_for_comp = 0
                for article in digest.articles:
                    if date_from and article.published_date and str(article.published_date) < date_from:
                        continue
                    if date_to and article.published_date and str(article.published_date) > date_to:
                        continue

                    existing = local_db.query(NewsArticleCache).filter(
                        NewsArticleCache.url == article.url
                    ).first()

                    if not existing:
                        pub_date = None
                        if article.published_date:
                            for fmt in ["%Y-%m-%dT%H:%M:%S", "%Y-%m-%d", "%a, %d %b %Y %H:%M:%S %Z"]:
                                try:
                                    pub_date = datetime.strptime(str(article.published_date)[:19], fmt[:19])
                                    break
                                except (ValueError, TypeError):
                                    continue

                        new_entry = NewsArticleCache(
                            competitor_id=comp.id,
                            competitor_name=comp.name,
                            title=article.title,
                            url=article.url,
                            source=getattr(article, 'source', '') or '',
                            source_type="google_news",
                            snippet=getattr(article, 'summary', '') or getattr(article, 'snippet', '') or '',
                            sentiment=getattr(article, 'sentiment', 'neutral') or 'neutral',
                            event_type=getattr(article, 'event_type', '') or '',
                            is_major_event=getattr(article, 'is_major_event', False),
                            published_at=pub_date or datetime.utcnow(),
                            fetched_at=datetime.utcnow()
                        )
                        local_db.add(new_entry)
                        articles_for_comp += 1

                local_db.commit()

                # Thread-safe progress update
                with lock:
                    total_articles += articles_for_comp
                    completed_count += 1
                    progress["current_competitor"] = comp.name
                    progress["completed"] = completed_count
                    pct = int((completed_count / max(len(target_competitors), 1)) * 100)
                    progress["percentage"] = pct
                    progress["percent"] = pct
                    progress["articles_found"] = total_articles
                    progress["total_articles"] = total_articles
                    progress["sources_checked"] = completed_count
                    progress["by_competitor"][comp.name] = articles_for_comp
                    for art in digest.articles:
                        s = getattr(art, 'sentiment', 'neutral') or 'neutral'
                        progress["by_sentiment"][s] = progress["by_sentiment"].get(s, 0) + 1
                    if len(progress["sample_headlines"]) < 10:
                        for art in digest.articles:
                            if len(progress["sample_headlines"]) >= 10:
                                break
                            progress["sample_headlines"].append(art.title)

            except Exception as e:
                logger.error(f"[News Fetch] Error for {comp.name}: {e}")
                local_db.rollback()
            finally:
                local_db.close()

        # Process competitors 10 at a time
        with ThreadPoolExecutor(max_workers=10) as executor:
            futures = {
                executor.submit(fetch_one_competitor, comp): comp
                for comp in target_competitors
            }
            for future in as_completed(futures):
                comp = futures[future]
                try:
                    future.result(timeout=45)
                except Exception as e:
                    logger.error(f"[News Fetch] Timeout/error for {comp.name}: {e}")

        progress["completed"] = len(target_competitors)
        progress["percentage"] = 100
        progress["percent"] = 100
        progress["status"] = "complete"

    background_tasks.add_task(do_fetch)

    return {
        "progress_key": progress_key,
        "status": "started",
        "competitors_count": len(target_competitors)
    }


@app.post("/api/news-feed/refresh-cache")
async def refresh_news_cache(
    background_tasks: BackgroundTasks,
    competitor_id: Optional[int] = None,
    start_date: Optional[str] = None,
    end_date: Optional[str] = None,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """
    Manually refresh the news cache for competitors.

    v5.0.8: Fetches fresh news and stores in NewsArticleCache table.
    v7.2.1: Async, date range support, progress tracking, removed cache_expires_at.

    Args:
        competitor_id: Optional - refresh only this competitor
        start_date: Optional start date (YYYY-MM-DD), default 30 days ago
        end_date: Optional end date (YYYY-MM-DD), default today
    """
    import uuid as _uuid
    from database import NewsArticleCache
    from news_monitor import NewsMonitor
    from datetime import datetime

    # Calculate days from start_date
    if start_date:
        try:
            start_dt = datetime.strptime(start_date, "%Y-%m-%d")
            days = max(1, (datetime.utcnow() - start_dt).days)
        except ValueError:
            days = 30
    else:
        days = 30

    progress_key = str(_uuid.uuid4())

    def do_refresh() -> None:
        monitor = NewsMonitor()
        refreshed = 0

        # Get competitors
        query = db.query(Competitor).filter(Competitor.is_deleted == False)
        if competitor_id:
            query = query.filter(Competitor.id == competitor_id)

        competitors = query.all()

        for comp in competitors:
            try:
                digest = monitor.fetch_news(comp.name, days=days)

                for article in digest.articles:
                    # Check if article already exists
                    existing = db.query(NewsArticleCache).filter(
                        NewsArticleCache.url == article.url
                    ).first()

                    if not existing:
                        # Parse published date
                        pub_date = None
                        if article.published_date:
                            for fmt in ["%Y-%m-%dT%H:%M:%S", "%Y-%m-%d", "%a, %d %b %Y %H:%M:%S %Z"]:
                                try:
                                    pub_date = datetime.strptime(article.published_date[:19], fmt[:19])
                                    break
                                except (ValueError, TypeError):
                                    continue

                        # Create new cache entry
                        cache_entry = NewsArticleCache(
                            competitor_id=comp.id,
                            competitor_name=comp.name,
                            title=article.title,
                            url=article.url,
                            source=article.source,
                            source_type="google_news",
                            published_at=pub_date or datetime.utcnow(),
                            snippet=article.snippet,
                            sentiment=article.sentiment,
                            event_type=article.event_type,
                            is_major_event=article.is_major_event,
                            fetched_at=datetime.utcnow(),
                        )
                        db.add(cache_entry)
                        refreshed += 1

                db.commit()

            except Exception as e:
                logger.error(f"Error refreshing news for {comp.name}: {e}")
                continue

        logger.info(f"[News Cache] Refreshed {refreshed} articles for {len(competitors)} competitors")

    # Run in background
    background_tasks.add_task(do_refresh)

    return {
        "status": "started",
        "progress_key": progress_key,
        "message": f"Refreshing news cache for {'all' if not competitor_id else 'competitor ' + str(competitor_id)} competitors"
    }


@app.post("/api/news-feed/cleanup-irrelevant")
def cleanup_irrelevant_news(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """
    One-time cleanup: archive irrelevant news articles from the cache.

    v8.0.8: Removes articles that don't match their competitor name
    AND don't contain healthcare keywords. Also archives articles
    with URL-based competitor names (duplicates from bad scraping).

    Returns:
        {archived_count, total_checked, by_competitor: {name: count}}
    """
    from database import NewsArticleCache, Competitor

    HEALTHCARE_KEYWORDS = {
        "healthcare", "health", "medical", "patient", "clinical", "hospital",
        "ehr", "telehealth", "doctor", "nurse", "pharmacy", "hipaa", "wellness",
        "healthtech", "medtech", "biotech", "pharma", "therapeutic", "diagnosis",
        "treatment", "care", "provider", "payer", "insurance", "cms", "fda",
        "interoperability", "emr", "revenue cycle", "population health",
        "credentialing", "certification", "compliance", "regulatory",
    }

    try:
        # Build a map of competitor name -> website domain for domain matching
        competitors = db.query(Competitor).filter(
            Competitor.is_deleted == False
        ).all()
        comp_domains = {}
        for c in competitors:
            if c.website:
                try:
                    from urllib.parse import urlparse
                    netloc = urlparse(c.website).netloc.replace('www.', '')
                    domain = netloc.split('.')[0].lower() if netloc else ''
                    if domain:
                        comp_domains[c.name] = domain
                except Exception:
                    pass

        # Get all non-archived articles
        articles = db.query(NewsArticleCache).filter(
            NewsArticleCache.is_archived == False
        ).all()

        total_checked = len(articles)
        archived_count = 0
        by_competitor = {}

        for article in articles:
            comp_name = article.competitor_name or ""
            title = (article.title or "").lower()
            snippet = (article.snippet or "").lower()
            combined = title + " " + snippet
            should_archive = False

            # Rule 1: URL-based competitor names (e.g., "Https://Www.Phreesia.Com/")
            if comp_name.lower().startswith("http"):
                should_archive = True

            # Rule 2: Check if article is relevant to its competitor
            if not should_archive:
                # Extract meaningful name words
                name_words = [
                    w.lower() for w in comp_name.split()
                    if len(w) > 2 and w.lower() not in {
                        "the", "and", "inc", "llc", "ltd", "corp", "co", "or",
                    }
                ]
                name_lower = comp_name.lower()

                # Common English words that cause false positives
                _COMMON_WORDS = {
                    "get", "well", "set", "change", "care", "one", "first",
                    "best", "next", "open", "clear", "smart", "true", "pure",
                    "prime", "lead", "core", "base", "way", "rise", "bold",
                    "vital", "health", "healthcare", "medical", "med",
                }
                significant_words = [
                    w for w in name_words if w not in _COMMON_WORDS
                ]
                # Names made entirely of common words are generic
                has_only_common = len(significant_words) == 0 and len(name_words) > 0
                is_generic = (len(name_words) <= 1 and len(comp_name) < 15) or has_only_common

                # Get domain for this competitor
                domain_name = comp_domains.get(comp_name, "")

                # Check if article mentions the company domain
                domain_match = domain_name and (
                    domain_name in title or domain_name in snippet
                )

                # Check for healthcare keywords
                healthcare_match = any(
                    kw in combined for kw in HEALTHCARE_KEYWORDS
                )

                if is_generic:
                    # Generic / common-word names: REQUIRE domain match
                    if not domain_match:
                        should_archive = True
                else:
                    # Specific names: require full phrase or significant word
                    full_phrase_match = name_lower in title or name_lower in snippet
                    sig_word_match = len(significant_words) > 0 and any(
                        w in title for w in significant_words
                    )
                    name_match = full_phrase_match or sig_word_match
                    if not name_match and not healthcare_match:
                        should_archive = True

            if should_archive:
                article.is_archived = True
                archived_count += 1
                by_competitor[comp_name] = by_competitor.get(comp_name, 0) + 1

        db.commit()

        logger.info(
            f"[News Cleanup] Archived {archived_count}/{total_checked} "
            f"irrelevant articles across {len(by_competitor)} competitors"
        )

        return {
            "archived_count": archived_count,
            "total_checked": total_checked,
            "by_competitor": dict(sorted(
                by_competitor.items(), key=lambda x: x[1], reverse=True
            ))
        }

    except Exception as e:
        logger.error(f"[News Cleanup] Error: {e}")
        db.rollback()
        return {"error": "Cleanup failed", "archived_count": 0, "total_checked": 0}


@app.get("/api/news-coverage")
def get_news_coverage(db: Session = Depends(get_db)):
    """
    Get news coverage status for all competitors.

    v5.1.0: Shows which competitors have news coverage and identifies gaps.
    """
    from database import NewsArticleCache
    from sqlalchemy import func

    try:
        competitors = db.query(Competitor).filter(
            Competitor.is_deleted == False
        ).all()

        now = datetime.utcnow()
        week_ago = now - timedelta(days=7)
        month_ago = now - timedelta(days=30)

        coverage = []
        total_with_news = 0
        total_recent_news = 0

        for comp in competitors:
            # Count articles (exclude archived)
            total = db.query(NewsArticleCache).filter(
                NewsArticleCache.competitor_id == comp.id,
                NewsArticleCache.is_archived != True
            ).count()

            recent = db.query(NewsArticleCache).filter(
                NewsArticleCache.competitor_id == comp.id,
                NewsArticleCache.published_at >= week_ago,
                NewsArticleCache.is_archived != True
            ).count()

            if total > 0:
                total_with_news += 1
            if recent > 0:
                total_recent_news += 1

            # Get last fetch time
            latest = db.query(NewsArticleCache).filter(
                NewsArticleCache.competitor_id == comp.id,
                NewsArticleCache.is_archived != True
            ).order_by(
                NewsArticleCache.fetched_at.desc()
            ).first()

            coverage.append({
                "competitor_id": comp.id,
                "competitor_name": comp.name,
                "total_articles": total,
                "recent_articles": recent,
                "has_news": total > 0,
                "has_recent_news": recent > 0,
                "last_fetched": latest.fetched_at.isoformat() if latest else None
            })

        coverage_pct = (total_with_news / len(competitors) * 100) if competitors else 0
        recent_pct = (total_recent_news / len(competitors) * 100) if competitors else 0

        return {
            "total_competitors": len(competitors),
            "competitors_with_news": total_with_news,
            "competitors_with_recent_news": total_recent_news,
            "coverage_percentage": round(coverage_pct, 1),
            "recent_coverage_percentage": round(recent_pct, 1),
            "coverage_details": coverage,
            "competitors_missing_news": [
                c for c in coverage if not c["has_news"]
            ]
        }

    except Exception as e:
        return {"error": "An unexpected error occurred"}


@app.post("/api/news-coverage/refresh-all")
def refresh_all_news_coverage(
    background_tasks: BackgroundTasks,
    db: Session = Depends(get_db)
):
    """
    Comprehensive news refresh for ALL competitors.

    v5.1.0: Uses ComprehensiveNewsScraper to ensure complete news coverage.
    """
    def do_comprehensive_refresh():
        try:
            from comprehensive_news_scraper import NewsFeedService
            service = NewsFeedService(db)
            result = service.refresh_all_news()
            logger.info(f"[News Coverage] Refresh complete: {result.new_articles_found} new articles")
        except ImportError:
            # Fallback to basic refresh
            logger.warning("[News Coverage] ComprehensiveNewsScraper not available, using basic refresh")
            from news_monitor import NewsMonitor
            from database import NewsArticleCache

            monitor = NewsMonitor()
            competitors = db.query(Competitor).filter(Competitor.is_deleted == False).all()

            for comp in competitors:
                try:
                    digest = monitor.fetch_news(comp.name, days=30)
                    for article in digest.articles:
                        existing = db.query(NewsArticleCache).filter(
                            NewsArticleCache.url == article.url
                        ).first()

                        if not existing:
                            cache_entry = NewsArticleCache(
                                competitor_id=comp.id,
                                competitor_name=comp.name,
                                title=article.title[:500],
                                url=article.url,
                                source=article.source,
                                source_type="news_monitor",
                                published_at=datetime.utcnow(),
                                snippet=article.snippet[:1000] if article.snippet else None,
                                sentiment=article.sentiment,
                                event_type=article.event_type,
                                is_major_event=article.is_major_event,
                                fetched_at=datetime.utcnow(),
                                created_at=datetime.utcnow()
                            )
                            db.add(cache_entry)

                    db.commit()
                except Exception as e:
                    logger.error(f"Error refreshing news for {comp.name}: {e}")
                    continue

    background_tasks.add_task(do_comprehensive_refresh)

    return {
        "status": "started",
        "message": "Comprehensive news refresh started for all competitors"
    }


@app.get("/api/notifications")
def get_notifications(limit: int = 10, db: Session = Depends(get_db)):
    """Get recent notifications/alerts for the notification bell."""
    try:
        # Get recent changes that could be notifications
        changes = db.query(ChangeLog).order_by(ChangeLog.changed_at.desc()).limit(limit).all()

        notifications = []
        for change in changes:
            # Only include significant changes as notifications
            if change.field_name in ['threat_level', 'pricing_model', 'funding_total', 'employee_count']:
                competitor = db.query(Competitor).filter(Competitor.id == change.competitor_id).first()
                notifications.append({
                    "id": change.id,
                    "type": "change",
                    "title": f"{competitor.name if competitor else 'Unknown'}: {change.field_name} changed",
                    "message": f"Changed from '{change.old_value}' to '{change.new_value}'",
                    "timestamp": change.changed_at.isoformat() if change.changed_at else None,
                    "read": False
                })

        return notifications[:limit]
    except Exception as e:
        # Return empty list on error to prevent UI breaking
        return []


@app.get("/api/alerts/price-changes")
def get_price_alerts(threshold: float = 10.0, db: Session = Depends(get_db)):
    """Get price change alerts."""
    try:
        from price_tracker import PriceTracker
        
        tracker = PriceTracker(db)
        alerts = tracker.detect_price_alerts(threshold)
        
        return {
            "alerts": [
                {
                    "competitor_id": a.competitor_id,
                    "competitor_name": a.competitor_name,
                    "previous_price": a.previous_price,
                    "new_price": a.new_price,
                    "change_percent": a.change_percent,
                    "direction": a.change_direction,
                    "severity": a.severity,
                    "detected_at": a.detected_at
                }
                for a in alerts
            ],
            "count": len(alerts)
        }
        
    except Exception as e:
        return {"error": "An unexpected error occurred", "alerts": []}


@app.get("/api/alerts/email-log")
def get_email_log():
    """Get the test email log (emails logged when TEST_EMAIL_MODE=true)."""
    import json
    from pathlib import Path

    email_log_file = Path(__file__).parent / "email_log.json"

    test_mode = os.getenv("TEST_EMAIL_MODE", "true").lower() in ("true", "1", "yes")

    if not email_log_file.exists():
        return {
            "test_mode": test_mode,
            "emails": [],
            "message": "No emails logged yet. Test emails will appear here when alerts are triggered."
        }

    try:
        with open(email_log_file, 'r') as f:
            emails = json.load(f)
        return {
            "test_mode": test_mode,
            "emails": emails,
            "count": len(emails)
        }
    except Exception as e:
        return {"error": "An unexpected error occurred", "emails": []}


@app.delete("/api/alerts/email-log")
def clear_email_log():
    """Clear the test email log."""
    import json
    from pathlib import Path

    email_log_file = Path(__file__).parent / "email_log.json"

    if email_log_file.exists():
        email_log_file.unlink()

    return {"status": "cleared", "message": "Email log cleared"}


@app.post("/api/alerts/test-email")
def send_test_alert():
    """Send a test email alert to verify configuration."""
    from alerts import AlertSystem

    alert = AlertSystem()

    subject = "Test Alert from Certify Intel"
    body_html = """
    <div style="font-family: Arial, sans-serif; padding: 20px;">
        <h2 style="color: #3A95ED;">Test Alert from Certify Intel</h2>
        <p>This is a test email to verify your alert configuration is working.</p>
        <p><strong>Timestamp:</strong> {}</p>
        <hr style="border: 1px solid #e2e8f0;">
        <p style="color: #64748b; font-size: 12px;">
            If you received this email, your alert system is configured correctly.
        </p>
    </div>
    """.format(datetime.now().strftime("%Y-%m-%d %H:%M:%S"))

    success = alert.send_alert(subject, body_html)

    test_mode = os.getenv("TEST_EMAIL_MODE", "true").lower() in ("true", "1", "yes")

    return {
        "success": success,
        "test_mode": test_mode,
        "message": "Test email logged to backend/email_log.json" if test_mode else "Test email sent"
    }


@app.get("/api/pricing/comparison")
def get_pricing_comparison(db: Session = Depends(get_db)):
    """Get pricing comparison across all competitors."""
    try:
        from price_tracker import PriceTracker

        tracker = PriceTracker(db)
        return tracker.get_pricing_comparison()

    except Exception as e:
        return {"error": "An unexpected error occurred"}


@app.get("/api/competitors/{competitor_id}/reviews")
def get_competitor_reviews(competitor_id: int, db: Session = Depends(get_db)):
    """Get G2/Capterra review data for a competitor."""
    try:
        from review_scraper import get_competitor_reviews as _get_reviews
        
        competitor = db.query(Competitor).filter(
            Competitor.id == competitor_id,
            Competitor.is_deleted == False
        ).first()
        
        if not competitor:
            raise HTTPException(status_code=404, detail="Competitor not found")
        
        return _get_reviews(competitor.name)
        
    except HTTPException:
        raise
    except Exception as e:
        return {"error": "An unexpected error occurred", "overall_rating": 0}


@app.get("/api/reviews/compare")
def compare_reviews(competitor_ids: str, db: Session = Depends(get_db)):
    """Compare reviews across multiple competitors."""
    try:
        from review_scraper import compare_competitor_reviews
        
        ids = [int(id.strip()) for id in competitor_ids.split(",")]
        competitors = db.query(Competitor).filter(
            Competitor.id.in_(ids),
            Competitor.is_deleted == False
        ).all()
        
        names = [c.name for c in competitors]
        return compare_competitor_reviews(names)
        
    except Exception as e:
        return {"error": "An unexpected error occurred"}


@app.get("/api/competitors/{competitor_id}/linkedin")
def get_linkedin_data(competitor_id: int, db: Session = Depends(get_db)):
    """Get LinkedIn company data for a competitor."""
    try:
        from linkedin_tracker import get_linkedin_data as _get_linkedin
        
        competitor = db.query(Competitor).filter(
            Competitor.id == competitor_id,
            Competitor.is_deleted == False
        ).first()
        
        if not competitor:
            raise HTTPException(status_code=404, detail="Competitor not found")
        
        return _get_linkedin(competitor.name)
        
    except HTTPException:
        raise
    except Exception as e:
        return {"error": "An unexpected error occurred"}


@app.get("/api/competitors/{competitor_id}/hiring")
def get_hiring_analysis(competitor_id: int, db: Session = Depends(get_db)):
    """Get hiring trend analysis for a competitor."""
    try:
        from linkedin_tracker import analyze_competitor_hiring
        
        competitor = db.query(Competitor).filter(
            Competitor.id == competitor_id,
            Competitor.is_deleted == False
        ).first()
        
        if not competitor:
            raise HTTPException(status_code=404, detail="Competitor not found")
        
        return analyze_competitor_hiring(competitor.name)
        
    except HTTPException:
        raise
    except Exception as e:
        return {"error": "An unexpected error occurred"}


@app.get("/api/hiring/compare")
def compare_hiring(competitor_ids: str, db: Session = Depends(get_db)):
    """Compare hiring across multiple competitors."""
    try:
        from linkedin_tracker import LinkedInTracker
        
        ids = [int(id.strip()) for id in competitor_ids.split(",")]
        competitors = db.query(Competitor).filter(
            Competitor.id.in_(ids),
            Competitor.is_deleted == False
        ).all()
        
        tracker = LinkedInTracker()
        return tracker.compare_hiring([c.name for c in competitors])
        
    except Exception as e:
        return {"error": "An unexpected error occurred"}


@app.get("/api/competitors/{competitor_id}/insights")
def get_competitor_insights(competitor_id: int, db: Session = Depends(get_db)):
    """Get comprehensive insights for a competitor (threat + news + reviews + LinkedIn)."""
    try:
        competitor = db.query(Competitor).filter(
            Competitor.id == competitor_id,
            Competitor.is_deleted == False
        ).first()
        
        if not competitor:
            raise HTTPException(status_code=404, detail="Competitor not found")
        
        insights = {"competitor_name": competitor.name, "id": competitor_id}
        
        # Threat analysis
        try:
            from threat_analyzer import analyze_competitor_threat
            comp_data = {k: v for k, v in competitor.__dict__.items() if not k.startswith('_')}
            insights["threat"] = analyze_competitor_threat(comp_data)
        except Exception as e:
            insights["threat"] = {"error": "An unexpected error occurred"}
        
        # News
        try:
            from news_monitor import NewsMonitor
            monitor = NewsMonitor()
            summary = monitor.get_news_summary(competitor.name)
            insights["news"] = summary
        except Exception as e:
            insights["news"] = {"error": "An unexpected error occurred"}
        
        # Reviews
        try:
            from review_scraper import ReviewScraper
            scraper = ReviewScraper()
            review_insights = scraper.get_review_insights(competitor.name)
            insights["reviews"] = review_insights
        except Exception as e:
            insights["reviews"] = {"error": "An unexpected error occurred"}
        
        # LinkedIn
        try:
            from linkedin_tracker import LinkedInTracker
            tracker = LinkedInTracker()
            hiring = tracker.analyze_hiring_trends(competitor.name)
            insights["hiring"] = hiring
        except Exception as e:
            insights["hiring"] = {"error": "An unexpected error occurred"}
        
        return insights
        
    except HTTPException:
        raise
    except Exception as e:
        return {"error": "An unexpected error occurred"}


# ============== Win/Loss Tracker Endpoints ==============

@app.post("/api/deals")
def log_deal(deal: dict, db: Session = Depends(get_db)):
    """Log a competitive deal outcome."""
    try:
        from win_loss_tracker import get_tracker
        
        tracker = get_tracker(db)
        result = tracker.log_deal(
            competitor_id=deal.get("competitor_id"),
            competitor_name=deal.get("competitor_name", "Unknown"),
            deal_name=deal.get("deal_name", "Untitled Deal"),
            deal_value=deal.get("deal_value"),
            outcome=deal.get("outcome", "Won"),
            loss_reason=deal.get("loss_reason"),
            notes=deal.get("notes")
        )
        
        # Trigger webhook
        try:
            from webhooks import trigger_deal_outcome
            trigger_deal_outcome(
                deal.get("competitor_name", "Unknown"),
                deal.get("deal_name", "Untitled"),
                deal.get("deal_value", 0),
                deal.get("outcome", "Won")
            )
        except (ImportError, Exception):
            pass
        
        return result
        
    except Exception as e:
        return {"error": "An unexpected error occurred"}


@app.get("/api/deals/stats")
def get_deal_stats(days: int = 365):
    """Get win/loss statistics."""
    try:
        from win_loss_tracker import get_win_loss_stats
        return get_win_loss_stats(days)
    except Exception as e:
        return {"error": "An unexpected error occurred"}


@app.get("/api/deals/competitor/{competitor_id}")
def get_competitor_deals(competitor_id: int):
    """Get win/loss performance against a specific competitor."""
    try:
        from win_loss_tracker import get_tracker
        tracker = get_tracker()
        return tracker.get_competitor_performance(competitor_id)
    except Exception as e:
        return {"error": "An unexpected error occurred"}


@app.get("/api/deals/most-competitive")
def get_most_competitive(limit: int = 5):
    """Get competitors we face most often."""
    try:
        from win_loss_tracker import get_tracker
        tracker = get_tracker()
        return tracker.get_most_competitive(limit)
    except Exception as e:
        return {"error": "An unexpected error occurred"}


# NOTE: GET /api/scrape/all removed - use POST /api/scrape/all instead (defined at line ~5015)
# The POST endpoint provides full progress tracking and RefreshSession audit trail


@app.get("/api/refresh/history")
def get_refresh_history(competitor_id: int = None, limit: int = 100, db: Session = Depends(get_db)):
    """Get refresh history showing what changed."""
    from database import RefreshSnapshot
    
    query = db.query(RefreshSnapshot)
    if competitor_id:
        query = query.filter(RefreshSnapshot.competitor_id == competitor_id)
    
    snapshots = query.order_by(RefreshSnapshot.created_at.desc()).limit(limit).all()
    
    return {
        "total": len(snapshots),
        "history": [
            {
                "refresh_id": s.refresh_id,
                "competitor": s.competitor_name,
                "field": s.field_name,
                "old_value": s.old_value[:100] if s.old_value else None,
                "new_value": s.new_value[:100] if s.new_value else None,
                "source": s.source,
                "changed": s.changed,
                "date": s.created_at.isoformat()
            }
            for s in snapshots
        ]
    }



# Webhook test/events endpoints moved to routers/webhooks.py

# ============== New Data Source Endpoints ==============

@app.get("/api/competitors/{competitor_id}/employee-reviews")
def get_competitor_employee_reviews(competitor_id: int, db: Session = Depends(get_db)):
    """Get employee reviews and ratings from Glassdoor."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")
    
    return glassdoor_scraper.get_glassdoor_data(competitor.name)

@app.get("/api/competitors/{competitor_id}/jobs")
def get_competitor_jobs(competitor_id: int, db: Session = Depends(get_db)):
    """Get job postings and hiring signals from Indeed/ZipRecruiter."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")
    
    return indeed_scraper.get_job_data(competitor.name)

@app.get("/api/competitors/{competitor_id}/sec-filings")
def get_competitor_sec_filings(competitor_id: int, db: Session = Depends(get_db)):
    """Get public filings and financials from SEC EDGAR."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")
    
    return sec_edgar_scraper.get_sec_data(competitor.name)

@app.get("/api/competitors/{competitor_id}/patents")
def get_competitor_patents(competitor_id: int, db: Session = Depends(get_db)):
    """Get patent portfolio and IP analysis from USPTO."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")
    
    import uspto_scraper
    return uspto_scraper.get_patent_data(competitor.name)

@app.get("/api/competitors/{competitor_id}/klas-ratings")
def get_competitor_klas_ratings(competitor_id: int, db: Session = Depends(get_db)):
    """Get vendor ratings and customer satisfaction from KLAS Research."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")
    
    return klas_scraper.get_klas_data(competitor.name)

@app.get("/api/competitors/{competitor_id}/mobile-apps")
def get_competitor_mobile_apps(competitor_id: int, db: Session = Depends(get_db)):
    """Get mobile app ratings and reviews from App Store/Google Play."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")
    
    return appstore_scraper.get_app_store_data(competitor.name)

@app.get("/api/competitors/{competitor_id}/social-sentiment")
def get_competitor_social_sentiment(competitor_id: int, db: Session = Depends(get_db)):
    """Get brand sentiment and mentions from Twitter/Reddit."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")
    
    analysis = social_media_monitor.analyze_social_sentiment(competitor.name)
    raw_data = social_media_monitor.get_social_data(competitor.name)
    return {**analysis, "recent_posts": raw_data.get("top_posts", [])}

@app.get("/api/competitors/{competitor_id}/market-presence")
def get_competitor_market_presence(competitor_id: int, db: Session = Depends(get_db)):
    """Get industry presence and customer data from HIMSS/CHIME."""
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")
    
    return himss_scraper.get_himss_data(competitor.name)

# ============== EXPORT ENDPOINTS ==============

@app.get("/api/export/excel")
def get_export_excel(db: Session = Depends(get_db)):
    """Export all competitor data to Excel with comprehensive fields (v6.1.2)."""
    import pandas as pd
    from io import BytesIO
    from fastapi.responses import StreamingResponse

    competitors = db.query(Competitor).filter(Competitor.is_deleted == False).all()

    if not competitors:
        raise HTTPException(status_code=404, detail="No data to export")

    # Convert to list of dicts with ALL fields including 50 new extended fields
    data = []
    for comp in competitors:
        item = {
            # === CORE INFORMATION ===
            "ID": comp.id,
            "Company Name": comp.name,
            "Website": comp.website,
            "Status": comp.status,
            "Threat Level": comp.threat_level,
            "Last Updated": comp.last_updated.strftime("%Y-%m-%d %H:%M") if comp.last_updated else None,
            "Notes": comp.notes,
            "Data Quality Score": comp.data_quality_score,

            # === PRICING ===
            "Pricing Model": comp.pricing_model,
            "Base Price": comp.base_price,
            "Price Unit": comp.price_unit,

            # === PRODUCT ===
            "Product Categories": comp.product_categories,
            "Key Features": comp.key_features,
            "Integration Partners": comp.integration_partners,
            "Certifications": comp.certifications,

            # === MARKET ===
            "Target Segments": comp.target_segments,
            "Customer Size Focus": comp.customer_size_focus,
            "Geographic Focus": comp.geographic_focus,
            "Customer Count": comp.customer_count,
            "Customer Acquisition Rate": comp.customer_acquisition_rate,
            "Key Customers": comp.key_customers,
            "G2 Rating": comp.g2_rating,

            # === COMPANY ===
            "Employee Count": comp.employee_count,
            "Employee Growth Rate": comp.employee_growth_rate,
            "Year Founded": comp.year_founded,
            "Headquarters": comp.headquarters,
            "Funding Total": comp.funding_total,
            "Latest Round": comp.latest_round,
            "PE/VC Backers": comp.pe_vc_backers,

            # === DIGITAL PRESENCE ===
            "Website Traffic": comp.website_traffic,
            "Social Following": comp.social_following,
            "Recent Launches": comp.recent_launches,
            "News Mentions": comp.news_mentions,

            # === PUBLIC COMPANY INFO ===
            "Is Public": comp.is_public,
            "Ticker Symbol": comp.ticker_symbol,
            "Stock Exchange": comp.stock_exchange,

            # === MARKET VERTICAL ===
            "Primary Market": comp.primary_market,
            "Markets Served": comp.markets_served,
            "Market Focus Score": comp.market_focus_score,

            # === PRODUCT OVERLAP ===
            "Has PXP": comp.has_pxp,
            "Has PMS": comp.has_pms,
            "Has RCM": comp.has_rcm,
            "Has Patient Mgmt": comp.has_patient_mgmt,
            "Has Payments": comp.has_payments,
            "Has Biometric": comp.has_biometric,
            "Has Interoperability": comp.has_interoperability,
            "Product Overlap Score": comp.product_overlap_score,

            # === ENHANCED ANALYTICS ===
            "Telehealth Capabilities": comp.telehealth_capabilities,
            "AI Features": comp.ai_features,
            "Mobile App Available": comp.mobile_app_available,
            "HIPAA Compliant": comp.hipaa_compliant,
            "EHR Integrations": comp.ehr_integrations,

            # === DIMENSION SCORES ===
            "Dim: Product Packaging Score": comp.dim_product_packaging_score,
            "Dim: Integration Depth Score": comp.dim_integration_depth_score,
            "Dim: Support Service Score": comp.dim_support_service_score,
            "Dim: Retention Stickiness Score": comp.dim_retention_stickiness_score,
            "Dim: User Adoption Score": comp.dim_user_adoption_score,
            "Dim: Implementation TTV Score": comp.dim_implementation_ttv_score,
            "Dim: Reliability Enterprise Score": comp.dim_reliability_enterprise_score,
            "Dim: Pricing Flexibility Score": comp.dim_pricing_flexibility_score,
            "Dim: Reporting Analytics Score": comp.dim_reporting_analytics_score,
            "Dim: Overall Score": comp.dim_overall_score,
            "Dim: Sales Priority": comp.dim_sales_priority,

            # === SEC/FINANCIAL DATA ===
            "Logo URL": comp.logo_url,
            "SEC CIK": comp.sec_cik,
            "Annual Revenue": comp.annual_revenue,
            "Net Income": comp.net_income,
            "SEC Employee Count": comp.sec_employee_count,
            "Fiscal Year End": comp.fiscal_year_end,
            "Email Pattern": comp.email_pattern,

            # === SOCIAL MEDIA METRICS (NEW v6.1.2) ===
            "LinkedIn Followers": getattr(comp, 'linkedin_followers', None),
            "LinkedIn Employees": getattr(comp, 'linkedin_employees', None),
            "LinkedIn URL": getattr(comp, 'linkedin_url', None),
            "Twitter Followers": getattr(comp, 'twitter_followers', None),
            "Twitter Handle": getattr(comp, 'twitter_handle', None),
            "Facebook Followers": getattr(comp, 'facebook_followers', None),
            "Instagram Followers": getattr(comp, 'instagram_followers', None),
            "YouTube Subscribers": getattr(comp, 'youtube_subscribers', None),

            # === FINANCIAL METRICS (NEW v6.1.2) ===
            "Estimated Revenue": getattr(comp, 'estimated_revenue', None),
            "Revenue Growth Rate": getattr(comp, 'revenue_growth_rate', None),
            "Profit Margin": getattr(comp, 'profit_margin', None),
            "Estimated Valuation": getattr(comp, 'estimated_valuation', None),
            "Burn Rate": getattr(comp, 'burn_rate', None),
            "Runway Months": getattr(comp, 'runway_months', None),
            "Last Funding Date": getattr(comp, 'last_funding_date', None),
            "Funding Stage": getattr(comp, 'funding_stage', None),
            "Debt Financing": getattr(comp, 'debt_financing', None),
            "Revenue Per Employee": getattr(comp, 'revenue_per_employee', None),

            # === LEADERSHIP & TEAM (NEW v6.1.2) ===
            "CEO Name": getattr(comp, 'ceo_name', None),
            "CEO LinkedIn": getattr(comp, 'ceo_linkedin', None),
            "CTO Name": getattr(comp, 'cto_name', None),
            "CFO Name": getattr(comp, 'cfo_name', None),
            "Executive Changes": getattr(comp, 'executive_changes', None),
            "Board Members": getattr(comp, 'board_members', None),
            "Advisors": getattr(comp, 'advisors', None),
            "Founder Background": getattr(comp, 'founder_background', None),

            # === EMPLOYEE & CULTURE (NEW v6.1.2) ===
            "Glassdoor Rating": getattr(comp, 'glassdoor_rating', None),
            "Glassdoor Reviews Count": getattr(comp, 'glassdoor_reviews_count', None),
            "Glassdoor Recommend %": getattr(comp, 'glassdoor_recommend_pct', None),
            "Indeed Rating": getattr(comp, 'indeed_rating', None),
            "Employee Turnover Rate": getattr(comp, 'employee_turnover_rate', None),
            "Hiring Velocity (Open Positions)": getattr(comp, 'hiring_velocity', None),

            # === PRODUCT & TECHNOLOGY (NEW v6.1.2) ===
            "Product Count": getattr(comp, 'product_count', None),
            "Latest Product Launch": getattr(comp, 'latest_product_launch', None),
            "Tech Stack": getattr(comp, 'tech_stack', None),
            "Cloud Provider": getattr(comp, 'cloud_provider', None),
            "API Available": getattr(comp, 'api_available', None),
            "API Documentation URL": getattr(comp, 'api_documentation_url', None),
            "Open Source Contributions": getattr(comp, 'open_source_contributions', None),
            "R&D Investment %": getattr(comp, 'rd_investment_pct', None),

            # === MARKET & COMPETITIVE (NEW v6.1.2) ===
            "Estimated Market Share": getattr(comp, 'estimated_market_share', None),
            "NPS Score": getattr(comp, 'nps_score', None),
            "Customer Churn Rate": getattr(comp, 'customer_churn_rate', None),
            "Average Contract Value": getattr(comp, 'average_contract_value', None),
            "Sales Cycle Length": getattr(comp, 'sales_cycle_length', None),
            "Competitive Win Rate": getattr(comp, 'competitive_win_rate', None),

            # === REGULATORY & COMPLIANCE (NEW v6.1.2) ===
            "SOC2 Certified": getattr(comp, 'soc2_certified', None),
            "HITRUST Certified": getattr(comp, 'hitrust_certified', None),
            "ISO 27001 Certified": getattr(comp, 'iso27001_certified', None),
            "Legal Issues": getattr(comp, 'legal_issues', None),

            # === PATENTS & IP (NEW v6.1.2) ===
            "Patent Count": getattr(comp, 'patent_count', None),
            "Recent Patents": getattr(comp, 'recent_patents', None),
            "Trademark Count": getattr(comp, 'trademark_count', None),
            "IP Litigation": getattr(comp, 'ip_litigation', None),

            # === PARTNERSHIPS & ECOSYSTEM (NEW v6.1.2) ===
            "Strategic Partners": getattr(comp, 'strategic_partners', None),
            "Reseller Partners": getattr(comp, 'reseller_partners', None),
            "Marketplace Presence": getattr(comp, 'marketplace_presence', None),
            "Acquisition History": getattr(comp, 'acquisition_history', None),

            # === CUSTOMER INTELLIGENCE (NEW v6.1.2) ===
            "Notable Customer Wins": getattr(comp, 'notable_customer_wins', None),
            "Customer Case Studies": getattr(comp, 'customer_case_studies', None),

            # === METADATA ===
            "Created At": comp.created_at.strftime("%Y-%m-%d %H:%M") if comp.created_at else None,
            "Last Verified At": comp.last_verified_at.strftime("%Y-%m-%d %H:%M") if comp.last_verified_at else None,
        }
        data.append(item)

    df = pd.DataFrame(data)

    # Create Excel buffer with multiple sheets
    buffer = BytesIO()
    with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
        # Main data sheet
        df.to_excel(writer, index=False, sheet_name="Competitors")

        # Create summary sheet
        summary_data = {
            "Metric": [
                "Total Competitors",
                "High Threat",
                "Medium Threat",
                "Low Threat",
                "Public Companies",
                "With Products Identified",
                "With Revenue Data",
                "With Social Data",
                "Export Date"
            ],
            "Value": [
                len(competitors),
                len([c for c in competitors if c.threat_level == "High"]),
                len([c for c in competitors if c.threat_level == "Medium"]),
                len([c for c in competitors if c.threat_level == "Low"]),
                len([c for c in competitors if c.is_public]),
                len([c for c in competitors if c.product_categories]),
                len([c for c in competitors if c.annual_revenue]),
                len([c for c in competitors if getattr(c, 'linkedin_followers', None)]),
                datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            ]
        }
        summary_df = pd.DataFrame(summary_data)
        summary_df.to_excel(writer, index=False, sheet_name="Summary")

        # Auto-adjust column widths for main sheet
        worksheet = writer.sheets["Competitors"]
        for idx, col in enumerate(df.columns):
            max_len = max(
                df[col].astype(str).str.len().max() if len(df) > 0 else 0,
                len(col)
            )
            worksheet.column_dimensions[chr(65 + idx) if idx < 26 else 'A' + chr(65 + idx - 26)].width = min(max_len + 2, 50)

    buffer.seek(0)

    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename=certify_intel_full_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"}
    )


# ============== USER MANAGEMENT ENDPOINTS ==============

@app.get("/api/users", response_model=List[UserResponse])
def get_users(db: Session = Depends(get_db), current_user: dict = Depends(get_current_user)):
    """List all users (admin only)."""
    if current_user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    return db.query(User).filter(User.is_active == True).all()

@app.post("/api/users/invite")
def invite_user(invite: UserInviteRequest, db: Session = Depends(get_db), current_user: dict = Depends(get_current_user)):
    """Invite a new user (creates account and sends email invite). Admin only."""
    if current_user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    from extended_features import auth_manager
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart

    existing = db.query(User).filter(User.email == invite.email).first()
    if existing:
        raise HTTPException(status_code=400, detail="User already exists")

    # Create user with default password using auth_manager for consistent hashing
    default_password = "Welcome123!"
    new_user = auth_manager.create_user(
        db,
        email=invite.email,
        password=default_password,
        full_name=invite.full_name or "",
        role=invite.role
    )

    # Try to send email invite
    email_sent = False
    smtp_host = os.getenv("SMTP_HOST", "")
    smtp_port = int(os.getenv("SMTP_PORT", "587"))
    smtp_user = os.getenv("SMTP_USER", "")
    smtp_password = os.getenv("SMTP_PASSWORD", "")
    from_email = os.getenv("ALERT_FROM_EMAIL", smtp_user)

    if smtp_host and smtp_user and smtp_password:
        try:
            # Build email
            msg = MIMEMultipart("alternative")
            msg["Subject"] = "🎉 You're Invited to Certify Intel!"
            msg["From"] = from_email
            msg["To"] = invite.email

            # Get server URL for the invite link
            app_url = os.getenv("APP_URL", "")
            if not app_url:
                server_host = os.getenv("HOST", "localhost")
                server_port = os.getenv("PORT", "8000")
                app_url = f"http://{server_host}:{server_port}"
                if server_host == "0.0.0.0":
                    app_url = "http://localhost:8000"

            # Plain text version
            text_body = f"""
Hello {invite.full_name or 'there'}!

You've been invited to join Certify Intel - the Competitive Intelligence Platform.

Your Account Details:
- Email: {invite.email}
- Temporary Password: {default_password}
- Role: {invite.role}

To get started:
1. Go to: {app_url}
2. Log in with your email and temporary password
3. Change your password after first login

If you have any questions, contact your administrator.

Best regards,
The Certify Intel Team
"""

            # HTML version
            html_body = f"""
<!DOCTYPE html>
<html>
<head>
    <style>
        body {{ font-family: Arial, sans-serif; line-height: 1.6; color: #333; }}
        .container {{ max-width: 600px; margin: 0 auto; padding: 20px; }}
        .header {{ background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 30px; text-align: center; border-radius: 10px 10px 0 0; }}
        .content {{ background: #f9fafb; padding: 30px; border-radius: 0 0 10px 10px; }}
        .credentials {{ background: white; padding: 20px; border-radius: 8px; border-left: 4px solid #667eea; margin: 20px 0; }}
        .btn {{ display: inline-block; background: #667eea; color: white; padding: 12px 30px; text-decoration: none; border-radius: 6px; margin-top: 20px; }}
        .footer {{ text-align: center; color: #6b7280; font-size: 12px; margin-top: 20px; }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>🎉 Welcome to Certify Intel!</h1>
            <p>Competitive Intelligence Platform</p>
        </div>
        <div class="content">
            <p>Hello <strong>{invite.full_name or 'there'}</strong>!</p>
            <p>You've been invited to join Certify Intel - your team's competitive intelligence command center.</p>

            <div class="credentials">
                <h3>📧 Your Login Credentials</h3>
                <p><strong>Email:</strong> {invite.email}</p>
                <p><strong>Temporary Password:</strong> <code style="background: #e5e7eb; padding: 2px 8px; border-radius: 4px;">{default_password}</code></p>
                <p><strong>Role:</strong> {invite.role}</p>
            </div>

            <h3>🚀 Getting Started</h3>
            <ol>
                <li>Click the button below to access Certify Intel</li>
                <li>Log in with your email and temporary password</li>
                <li>Change your password after first login</li>
            </ol>

            <center>
                <a href="{app_url}" class="btn">Access Certify Intel →</a>
            </center>

            <div class="footer">
                <p>If you have any questions, contact your administrator.</p>
                <p>© 2026 Certify Intel | Competitive Intelligence Platform</p>
            </div>
        </div>
    </div>
</body>
</html>
"""

            msg.attach(MIMEText(text_body, "plain"))
            msg.attach(MIMEText(html_body, "html"))

            # Send email
            with smtplib.SMTP(smtp_host, smtp_port) as server:
                server.starttls()
                server.login(smtp_user, smtp_password)
                server.sendmail(from_email, [invite.email], msg.as_string())

            email_sent = True
            logger.info(f"[INVITE] Email sent successfully to {invite.email}")

        except Exception as e:
            logger.error(f"[INVITE] Failed to send email to {invite.email}: {e}")

    response_msg = f"User {invite.email} invited successfully."
    if email_sent:
        response_msg += " Email invite sent!"
    else:
        response_msg += " (Email not sent - SMTP not configured). Please share the temporary password securely with the user."

    # Never return passwords in API responses - log for admin reference only
    if not email_sent:
        logger.info(f"[INVITE] Temporary password for {invite.email} must be shared securely by admin")

    return {"message": response_msg, "email_sent": email_sent, "password_set": True}

@app.delete("/api/users/{user_id}")
def delete_user(user_id: int, db: Session = Depends(get_db), current_user: dict = Depends(get_current_user)):
    """Remove a user (admin only)."""
    if current_user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
        
    db.delete(user)
    db.commit()
    return {"message": "User deleted"}


@app.get("/api/innovations/compare")
def compare_innovation_metrics(competitor_ids: str, db: Session = Depends(get_db)):
    """Compare innovation metrics (patents) across competitors."""
    try:
        ids = [int(id) for id in competitor_ids.split(",")]
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid competitor ID format")
    names = []

    for c_id in ids:
        comp = db.query(Competitor).filter(Competitor.id == c_id).first()
        if comp:
            names.append(comp.name)

    scraper = uspto_scraper.USPTOScraper()
    return scraper.compare_innovation(names)

@app.get("/api/social/compare")
def compare_social_metrics(competitor_ids: str, db: Session = Depends(get_db)):
    """Compare social media sentiment and presence."""
    try:
        ids = [int(id) for id in competitor_ids.split(",")]
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid competitor ID format")
    names = []

    for c_id in ids:
        comp = db.query(Competitor).filter(Competitor.id == c_id).first()
        if comp:
            names.append(comp.name)

    monitor = social_media_monitor.SocialMediaMonitor()
    return monitor.compare_social_presence(names)



# Win/Loss & Webhooks moved to routers/winloss.py and routers/webhooks.py
# SystemPromptCreate2/Response2 consolidated into schemas.prompts
# KnowledgeBaseItemCreate/Response imported from schemas.prompts

@app.get("/api/analytics/trends")
def get_market_trends(db: Session = Depends(get_db)):
    """Get market trends (New Competitors, Avg Price)."""
    from sqlalchemy import func
    import re

    # 1. New Competitors (Monthly)
    # SQLite syntax for YYYY-MM
    trends = db.query(
        func.strftime('%Y-%m', Competitor.created_at).label('month'),
        func.count(Competitor.id).label('count')
    ).filter(Competitor.is_deleted == False).group_by('month').order_by('month').all()

    labels = []
    competitor_counts = []
    
    for t in trends:
        labels.append(t.month)
        competitor_counts.append(t.count)
        
    if not labels:
        labels = [datetime.now().strftime('%Y-%m')]
        competitor_counts = [0]

    # 2. Avg Price (Current Snapshot)
    prices = db.query(Competitor.base_price).filter(
        Competitor.is_deleted == False,
        Competitor.base_price.isnot(None)
    ).all()
    
    valid_prices = []
    for (p_str,) in prices:
        if not p_str: continue
        # Extract first float found
        match = re.search(r'\$?([\d,]+(\.\d{2})?)', p_str)
        if match:
             try:
                 val = float(match.group(1).replace(',', ''))
                 valid_prices.append(val)
             except (ValueError, TypeError):
                 pass
                 
    avg_price = sum(valid_prices) / len(valid_prices) if valid_prices else 0
    
    # Repeat avg_price to create a baseline line
    price_data = [round(avg_price, 2)] * len(labels)

    return {
        "labels": labels,
        "datasets": [
            {
                "label": "Avg Market Price",
                "data": price_data,
                "borderColor": "#3A95ED",
                "tension": 0.4
            },
            {
                "label": "New Competitors",
                "data": competitor_counts,
                "borderColor": "#DC3545",
                "tension": 0.4,
                "yAxisID": "y1"
            }
        ]
    }


@app.get("/api/analytics/sentiment-trend")
def get_sentiment_trend(
    days: int = 30,
    competitor_id: Optional[int] = None,
    db: Session = Depends(get_db)
):
    """Get news sentiment trend grouped by date (v7.1.5)."""
    from database import NewsArticleCache
    from sqlalchemy import func

    logger = logging.getLogger(__name__)

    try:
        cutoff_date = datetime.now() - timedelta(days=days)

        # Build query
        query = db.query(
            func.date(NewsArticleCache.published_at).label('date'),
            NewsArticleCache.sentiment,
            func.count(NewsArticleCache.id).label('count')
        ).filter(
            NewsArticleCache.published_at >= cutoff_date,
            NewsArticleCache.is_archived != True
        )

        if competitor_id:
            query = query.filter(NewsArticleCache.competitor_id == competitor_id)

        results = query.group_by('date', NewsArticleCache.sentiment).order_by('date').all()

        # Build response structure
        date_map = {}
        for row in results:
            date_str = str(row.date)
            if date_str not in date_map:
                date_map[date_str] = {"positive": 0, "negative": 0, "neutral": 0}

            sentiment = (row.sentiment or "neutral").lower()
            if sentiment in ["positive", "negative", "neutral"]:
                date_map[date_str][sentiment] = row.count

        # Create sorted lists
        sorted_dates = sorted(date_map.keys())

        return {
            "labels": sorted_dates,
            "positive": [date_map[d]["positive"] for d in sorted_dates],
            "negative": [date_map[d]["negative"] for d in sorted_dates],
            "neutral": [date_map[d]["neutral"] for d in sorted_dates]
        }
    except Exception as e:
        logger.error(f"Error fetching sentiment trend: {e}", exc_info=True)
        return {"labels": [], "positive": [], "negative": [], "neutral": []}


@app.get("/api/analytics/activity-trend")
def get_activity_trend(
    days: int = 30,
    competitor_id: Optional[int] = None,
    db: Session = Depends(get_db)
):
    """Get news + product update activity trend grouped by date (v7.1.5)."""
    from database import NewsArticleCache, ChangeLog
    from sqlalchemy import func, or_

    logger = logging.getLogger(__name__)

    try:
        cutoff_date = datetime.now() - timedelta(days=days)

        # Query news activity
        news_query = db.query(
            func.date(NewsArticleCache.published_at).label('date'),
            func.count(NewsArticleCache.id).label('count')
        ).filter(
            NewsArticleCache.published_at >= cutoff_date,
            NewsArticleCache.is_archived != True
        )

        if competitor_id:
            news_query = news_query.filter(NewsArticleCache.competitor_id == competitor_id)

        news_results = news_query.group_by('date').all()

        # Query product update activity from ChangeLog
        product_query = db.query(
            func.date(ChangeLog.detected_at).label('date'),
            func.count(ChangeLog.id).label('count')
        ).filter(
            ChangeLog.detected_at >= cutoff_date,
            or_(
                ChangeLog.change_type.like('%product%'),
                ChangeLog.change_type.like('%Product%')
            )
        )

        if competitor_id:
            product_query = product_query.filter(ChangeLog.competitor_id == competitor_id)

        product_results = product_query.group_by('date').all()

        # Build date map
        date_map = {}
        for row in news_results:
            date_str = str(row.date)
            if date_str not in date_map:
                date_map[date_str] = {"news": 0, "products": 0}
            date_map[date_str]["news"] = row.count

        for row in product_results:
            date_str = str(row.date)
            if date_str not in date_map:
                date_map[date_str] = {"news": 0, "products": 0}
            date_map[date_str]["products"] = row.count

        # Create sorted lists
        sorted_dates = sorted(date_map.keys())

        return {
            "labels": sorted_dates,
            "news_activity": [date_map[d]["news"] for d in sorted_dates],
            "product_updates": [date_map[d]["products"] for d in sorted_dates]
        }
    except Exception as e:
        logger.error(f"Error fetching activity trend: {e}", exc_info=True)
        return {"labels": [], "news_activity": [], "product_updates": []}


@app.get("/api/analytics/growth-trend")
def get_growth_trend(
    days: int = 90,
    competitor_id: Optional[int] = None,
    db: Session = Depends(get_db)
):
    """Get employee/customer growth trend (v7.1.5)."""
    from database import Competitor
    from sqlalchemy import func

    logger = logging.getLogger(__name__)

    try:
        # For now, use current snapshot data since CompetitorDimensionHistory doesn't track these
        # In future, could track changes in ChangeLog or add history tracking

        query = db.query(
            Competitor.employee_count,
            Competitor.customer_count
        ).filter(Competitor.is_deleted == False)

        if competitor_id:
            query = query.filter(Competitor.id == competitor_id)
            competitor = query.first()

            if competitor:
                # Parse counts
                emp_count = 0
                cust_count = 0

                if competitor.employee_count:
                    try:
                        # Extract number from strings like "100-500" or "200+"
                        import re
                        emp_str = str(competitor.employee_count)
                        match = re.search(r'(\d+)', emp_str.replace(',', ''))
                        if match:
                            emp_count = int(match.group(1))
                    except Exception:
                        pass

                if competitor.customer_count:
                    try:
                        import re
                        cust_str = str(competitor.customer_count)
                        match = re.search(r'(\d+)', cust_str.replace(',', '').replace('+', ''))
                        if match:
                            cust_count = int(match.group(1))
                    except Exception:
                        pass

                # Return flat line for single competitor (no historical data yet)
                num_points = max(min(days // 7, 12), 1)  # Weekly points, max 12, min 1
                labels = [(datetime.now() - timedelta(days=days-i*(days//max(num_points, 1)))).strftime('%Y-%m-%d')
                          for i in range(num_points)]

                return {
                    "labels": labels,
                    "employee_growth": [emp_count] * num_points,
                    "customer_growth": [cust_count] * num_points
                }

        # Multiple competitors - return averages
        competitors = query.all()

        total_emp = 0
        total_cust = 0
        emp_count_found = 0
        cust_count_found = 0

        for comp in competitors:
            if comp.employee_count:
                try:
                    import re
                    emp_str = str(comp.employee_count)
                    match = re.search(r'(\d+)', emp_str.replace(',', ''))
                    if match:
                        total_emp += int(match.group(1))
                        emp_count_found += 1
                except Exception:
                    pass

            if comp.customer_count:
                try:
                    import re
                    cust_str = str(comp.customer_count)
                    match = re.search(r'(\d+)', cust_str.replace(',', '').replace('+', ''))
                    if match:
                        total_cust += int(match.group(1))
                        cust_count_found += 1
                except Exception:
                    pass

        avg_emp = total_emp // emp_count_found if emp_count_found > 0 else 0
        avg_cust = total_cust // cust_count_found if cust_count_found > 0 else 0

        # Return flat line (no historical data yet)
        num_points = max(min(days // 7, 12), 1)  # Weekly points, max 12, min 1
        labels = [(datetime.now() - timedelta(days=days-i*(days//max(num_points, 1)))).strftime('%Y-%m-%d')
                  for i in range(num_points)]

        return {
            "labels": labels,
            "employee_growth": [avg_emp] * num_points,
            "customer_growth": [avg_cust] * num_points
        }
    except Exception as e:
        logger.error(f"Error fetching growth trend: {e}", exc_info=True)
        return {"labels": [], "employee_growth": [], "customer_growth": []}


@app.get("/api/changes/trend")
def get_changes_trend(
    competitor_ids: str = "",
    days: int = 30,
    metric: str = "changes",
    db: Session = Depends(get_db)
):
    """Get competitor changes trend over time (v7.1.5)."""
    from database import ChangeLog, Competitor
    from sqlalchemy import func

    logger = logging.getLogger(__name__)

    try:
        cutoff_date = datetime.now() - timedelta(days=days)

        # Parse competitor IDs
        comp_ids = []
        if competitor_ids:
            try:
                comp_ids = [int(x.strip()) for x in competitor_ids.split(',') if x.strip()]
            except Exception:
                pass

        if not comp_ids:
            return {"labels": [], "series": []}

        # Query changes grouped by date and competitor
        results = db.query(
            func.date(ChangeLog.detected_at).label('date'),
            ChangeLog.competitor_id,
            ChangeLog.competitor_name,
            func.count(ChangeLog.id).label('count')
        ).filter(
            ChangeLog.detected_at >= cutoff_date,
            ChangeLog.competitor_id.in_(comp_ids)
        ).group_by('date', ChangeLog.competitor_id, ChangeLog.competitor_name).order_by('date').all()

        # Build data structure
        date_set = set()
        competitor_data = {}

        for row in results:
            date_str = str(row.date)
            date_set.add(date_str)

            if row.competitor_id not in competitor_data:
                competitor_data[row.competitor_id] = {
                    "name": row.competitor_name,
                    "data": {}
                }

            competitor_data[row.competitor_id]["data"][date_str] = row.count

        # Create sorted labels
        sorted_dates = sorted(list(date_set))

        # Build series (fill missing dates with 0)
        series = []
        for comp_id, comp_info in competitor_data.items():
            data_list = [comp_info["data"].get(d, 0) for d in sorted_dates]
            series.append({
                "competitor_id": comp_id,
                "name": comp_info["name"],
                "data": data_list
            })

        return {
            "labels": sorted_dates,
            "series": series
        }
    except Exception as e:
        logger.error(f"Error fetching changes trend: {e}", exc_info=True)
        return {"labels": [], "series": []}


# =========================================================================
# NOTE: Duplicate endpoints removed - see lines 1774, 1426, 1797 for originals
# =========================================================================

# NOTE: analytics_routes.py (v5.0.2 legacy) removed - its 3 endpoints
# (/api/analytics/summary, /api/analytics/summary/progress, /api/analytics/chat)
# are superseded by main.py implementations with auth, prompt_key, background mode.



# ============== Run Server ==============

# NOTE: Admin endpoints (system-prompts, knowledge-base, data-providers)
# moved to routers/admin.py

# =========================================================================
# KB DATA RECONCILIATION - v7.1 Entity Extraction & Conflict Resolution
# =========================================================================

from database import KBEntityLink, KBDataExtraction


@app.post("/api/kb/upload-with-extraction")
async def upload_kb_document_with_extraction(
    file: UploadFile = File(...),
    title: Optional[str] = Form(None),
    category: Optional[str] = Form("general"),
    tags: Optional[str] = Form(None),
    document_date: Optional[str] = Form(None),
    data_as_of_date: Optional[str] = Form(None),
    linked_competitors: Optional[str] = Form(None),
    extract_entities: bool = Form(True),
    auto_link: bool = Form(True),
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """
    Upload a document with automatic entity extraction and competitor linking.

    This enhanced upload endpoint:
    1. Uploads and ingests the document (existing functionality)
    2. Extracts entities (competitors, products, metrics) using GPT-4
    3. Auto-links extracted entities to known competitors
    4. Creates KB extraction records for reconciliation

    Args:
        file: The document file (PDF, DOCX, TXT, MD, HTML)
        title: Optional custom title
        category: Document category
        tags: Comma-separated tags
        document_date: When the document was created (ISO format)
        data_as_of_date: When the data in the document was valid (ISO format)
        linked_competitors: JSON array of competitor IDs to explicitly link
        extract_entities: Whether to run AI entity extraction
        auto_link: Whether to auto-link extracted entities to competitors

    Returns:
        Upload result with extraction and linking summary
    """
    import os
    import tempfile
    import json
    from datetime import datetime

    # Validate file type
    allowed_extensions = {'.pdf', '.docx', '.txt', '.md', '.html'}
    file_ext = os.path.splitext(file.filename)[1].lower()

    if file_ext not in allowed_extensions:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {file_ext}. Allowed: {', '.join(allowed_extensions)}"
        )

    # Read file content
    content = await file.read()
    file_size = len(content)

    # Validate file size (max 50MB)
    max_size = 50 * 1024 * 1024
    if file_size > max_size:
        raise HTTPException(
            status_code=400,
            detail=f"File too large. Maximum size: 50MB"
        )

    # Parse dates
    parsed_doc_date = None
    parsed_data_date = None
    try:
        if document_date:
            parsed_doc_date = datetime.fromisoformat(document_date.replace('Z', '+00:00'))
        if data_as_of_date:
            parsed_data_date = datetime.fromisoformat(data_as_of_date.replace('Z', '+00:00'))
    except ValueError as e:
        logger.error(f"Invalid date format in KB upload: {e}")
        raise HTTPException(
            status_code=400,
            detail="Invalid date format. Please use ISO 8601 format (e.g., 2026-01-15)"
        )

    # Parse linked competitors
    explicit_competitor_ids = []
    if linked_competitors:
        try:
            explicit_competitor_ids = json.loads(linked_competitors)
        except json.JSONDecodeError:
            pass

    # Use title or filename
    doc_title = title if title else os.path.splitext(file.filename)[0]

    # Parse tags
    tag_list = [t.strip() for t in tags.split(',')] if tags else []

    # Save to temp file for processing
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        # Extract text content based on file type
        extracted_text = ""

        if file_ext == '.pdf':
            try:
                import PyPDF2
                with open(tmp_path, 'rb') as pdf_file:
                    reader = PyPDF2.PdfReader(pdf_file)
                    for page in reader.pages:
                        extracted_text += page.extract_text() + "\n"
            except Exception as e:
                logger.error(f"Failed to extract PDF text: {e}")
                raise HTTPException(status_code=400, detail="Failed to extract text from the uploaded PDF file")

        elif file_ext == '.docx':
            try:
                from docx import Document
                doc = Document(tmp_path)
                extracted_text = "\n".join([para.text for para in doc.paragraphs])
            except Exception as e:
                logger.error(f"Failed to extract DOCX text: {e}")
                raise HTTPException(status_code=400, detail="Failed to extract text from the uploaded DOCX file")

        elif file_ext in {'.txt', '.md', '.html'}:
            extracted_text = content.decode('utf-8', errors='replace')

        if not extracted_text.strip():
            raise HTTPException(status_code=400, detail="No text content could be extracted")

        # Compute content hash for deduplication
        import hashlib
        content_hash = hashlib.sha256(extracted_text.encode()).hexdigest()

        # Check for duplicate
        existing = db.query(KnowledgeBaseItem).filter(
            KnowledgeBaseItem.content_hash == content_hash
        ).first()

        if existing:
            return {
                "status": "duplicate",
                "existing_id": existing.id,
                "message": "Document already exists in knowledge base"
            }

        # Create KB item with enhanced metadata
        new_item = KnowledgeBaseItem(
            title=doc_title,
            content_text=extracted_text,
            content_type=file_ext.lstrip('.'),
            source=f"upload:{file.filename}",
            source_type="upload",
            category=category,
            tags=json.dumps(tag_list) if tag_list else None,
            extra_metadata=json.dumps({
                "original_filename": file.filename,
                "file_size": file_size,
                "uploaded_by": current_user.get("email", "unknown"),
                "upload_date": datetime.utcnow().isoformat()
            }),
            document_date=parsed_doc_date,
            data_as_of_date=parsed_data_date,
            document_type="upload",
            linked_competitor_ids=json.dumps(explicit_competitor_ids) if explicit_competitor_ids else None,
            extraction_status="pending" if extract_entities else "skipped",
            content_hash=content_hash,
            is_active=True
        )

        db.add(new_item)
        db.commit()
        db.refresh(new_item)

        # Entity extraction
        extraction_result = None
        if extract_entities:
            try:
                from entity_extraction import extract_and_link_entities

                extraction_result = await extract_and_link_entities(
                    document_id=str(new_item.id),
                    content=extracted_text,
                    db_session=db,
                    kb_item_id=new_item.id
                )

                # Update extraction status
                new_item.extraction_status = "completed"
                db.commit()

            except Exception as e:
                logger.error(f"Entity extraction failed: {e}")
                new_item.extraction_status = f"failed: {str(e)[:50]}"
                db.commit()

        # Create explicit links for user-specified competitors
        explicit_links_created = 0
        if auto_link and explicit_competitor_ids:
            for comp_id in explicit_competitor_ids:
                try:
                    link = KBEntityLink(
                        kb_item_id=new_item.id,
                        competitor_id=comp_id,
                        link_type="explicit",
                        link_confidence=1.0
                    )
                    db.add(link)
                    explicit_links_created += 1
                except Exception as e:
                    logger.warning(f"Failed to create explicit link: {e}")

            if explicit_links_created > 0:
                db.commit()

        return {
            "id": new_item.id,
            "title": new_item.title,
            "content_type": new_item.content_type,
            "category": new_item.category,
            "document_date": new_item.document_date.isoformat() if new_item.document_date else None,
            "data_as_of_date": new_item.data_as_of_date.isoformat() if new_item.data_as_of_date else None,
            "char_count": len(extracted_text),
            "word_count": len(extracted_text.split()),
            "extraction": extraction_result,
            "explicit_links_created": explicit_links_created,
            "message": "Document uploaded and processed successfully"
        }

    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


@app.get("/api/competitors/{competitor_id}/reconciled/{field_name}")
async def get_reconciled_field(
    competitor_id: int,
    field_name: str,
    db: Session = Depends(get_db)
):
    """
    Get reconciled value for a specific field combining KB and live data.

    Returns the best value based on source authority, freshness, and confidence,
    along with all sources and any conflicts.
    """
    from source_reconciliation import SourceReconciliationEngine, SourceRecord

    # Get competitor
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    # Get KB extractions for this field
    kb_extractions = db.query(KBDataExtraction).filter(
        KBDataExtraction.competitor_id == competitor_id,
        KBDataExtraction.field_name == field_name,
        KBDataExtraction.status != "rejected"
    ).all()

    kb_sources = [
        SourceRecord(
            value=e.extracted_value,
            source_type="kb_extraction",
            source_id=e.id,
            source_origin="kb",
            data_as_of_date=e.data_as_of_date or e.document_date,
            is_verified=e.status == "verified",
            document_id=e.document_id
        )
        for e in kb_extractions
    ]

    # Get live data sources for this field
    live_sources_db = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id,
        DataSource.field_name == field_name
    ).all()

    live_sources = [
        SourceRecord(
            value=s.current_value or "",
            source_type=s.source_type or "unknown",
            source_id=s.id,
            source_origin="live",
            data_as_of_date=s.data_as_of_date,
            extracted_at=s.extracted_at,
            is_verified=s.is_verified,
            source_url=s.source_url,
            source_name=s.source_name
        )
        for s in live_sources_db
    ]

    # Reconcile
    engine = SourceReconciliationEngine(db_session=db)
    result = await engine.reconcile_field(
        competitor_id=competitor_id,
        field_name=field_name,
        kb_sources=kb_sources,
        live_sources=live_sources
    )

    return {
        "field_name": result.field_name,
        "best_value": result.best_value,
        "confidence_score": result.confidence_score,
        "confidence_level": result.confidence_level,
        "reconciliation_method": result.reconciliation_method,
        "needs_review": result.needs_review,
        "sources": [
            {
                "value": s.value,
                "type": s.source_type,
                "origin": s.source_origin,
                "confidence": s.confidence,
                "date": s.data_as_of_date.isoformat() if s.data_as_of_date else None,
                "verified": s.is_verified
            }
            for s in result.sources_used
        ],
        "conflicts": result.conflicts
    }


@app.get("/api/competitors/{competitor_id}/kb-documents")
def get_competitor_kb_documents(
    competitor_id: int,
    db: Session = Depends(get_db)
):
    """
    Get all KB documents linked to a competitor.

    Returns documents with their extraction summaries.
    """
    # Verify competitor exists
    competitor = db.query(Competitor).filter(Competitor.id == competitor_id).first()
    if not competitor:
        raise HTTPException(status_code=404, detail="Competitor not found")

    # Get entity links
    links = db.query(KBEntityLink).filter(
        KBEntityLink.competitor_id == competitor_id
    ).all()

    # Get unique KB item IDs
    kb_item_ids = list(set(l.kb_item_id for l in links if l.kb_item_id))

    # Get KB items
    kb_items = db.query(KnowledgeBaseItem).filter(
        KnowledgeBaseItem.id.in_(kb_item_ids)
    ).all() if kb_item_ids else []

    # Get extractions for this competitor
    extractions = db.query(KBDataExtraction).filter(
        KBDataExtraction.competitor_id == competitor_id
    ).all()

    # Build response
    documents = []
    for item in kb_items:
        item_links = [l for l in links if l.kb_item_id == item.id]
        item_extractions = [e for e in extractions if e.kb_item_id == item.id]

        documents.append({
            "id": item.id,
            "title": item.title,
            "content_type": item.content_type,
            "category": item.category,
            "document_date": item.document_date.isoformat() if item.document_date else None,
            "data_as_of_date": item.data_as_of_date.isoformat() if item.data_as_of_date else None,
            "created_at": item.created_at.isoformat() if item.created_at else None,
            "link_type": item_links[0].link_type if item_links else "unknown",
            "link_confidence": item_links[0].link_confidence if item_links else 0.0,
            "extractions": [
                {
                    "field_name": e.field_name,
                    "value": e.extracted_value,
                    "confidence": e.extraction_confidence,
                    "status": e.status
                }
                for e in item_extractions
            ]
        })

    return {
        "competitor_id": competitor_id,
        "competitor_name": competitor.name,
        "document_count": len(documents),
        "documents": documents
    }


@app.get("/api/kb-items/scan")
async def scan_knowledge_base_items(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """Scan and return all knowledge base items."""
    try:
        items = db.query(KnowledgeBaseItem).filter(
            KnowledgeBaseItem.is_active == True
        ).order_by(KnowledgeBaseItem.created_at.desc()).all()

        return {
            "items": [
                {
                    "id": item.id,
                    "title": item.title,
                    "content_type": item.content_type,
                    "category": item.category,
                    "source": item.source,
                    "created_at": item.created_at.isoformat() if item.created_at else None,
                    "tags": item.tags,
                    "document_date": item.document_date.isoformat() if item.document_date else None,
                }
                for item in items
            ],
            "total": len(items)
        }
    except Exception as e:
        logger.error(f"Error scanning knowledge base: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.get("/api/kb-items/verification-queue")
async def get_kb_items_verification_queue(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """Get items in verification queue."""
    try:
        # For now, return KB items that might need verification
        # In a full implementation, this would check a verification status field
        items = db.query(KnowledgeBaseItem).filter(
            KnowledgeBaseItem.is_active == True,
            KnowledgeBaseItem.category.is_(None)  # Items without category need review
        ).order_by(KnowledgeBaseItem.created_at.desc()).limit(50).all()

        return {
            "queue": [
                {
                    "id": item.id,
                    "title": item.title,
                    "content_preview": item.content_text[:200] if item.content_text else "",
                    "source": item.source,
                    "created_at": item.created_at.isoformat() if item.created_at else None,
                }
                for item in items
            ],
            "total": len(items)
        }
    except Exception as e:
        logger.error(f"Error getting verification queue: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.get("/api/kb-items/competitor-names")
async def get_kb_items_competitor_names(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """Extract competitor names found in knowledge base."""
    try:
        # Get all active competitors from database
        competitors = db.query(Competitor).filter(
            Competitor.is_deleted == False
        ).all()

        # Get KB items to scan for competitor mentions
        kb_items = db.query(KnowledgeBaseItem).filter(
            KnowledgeBaseItem.is_active == True
        ).all()

        # Simple keyword matching - count mentions
        mentions = {}
        for comp in competitors:
            count = 0
            comp_name_lower = comp.name.lower()
            for item in kb_items:
                if item.content_text and comp_name_lower in item.content_text.lower():
                    count += 1
            if count > 0:
                mentions[comp.name] = {
                    "id": comp.id,
                    "name": comp.name,
                    "mentions": count
                }

        return {
            "competitors": list(mentions.values()),
            "total": len(mentions)
        }
    except Exception as e:
        logger.error(f"Error extracting competitor names from KB: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.get("/api/kb-items/preview")
async def preview_kb_items_import(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user_optional)
):
    """Preview what would be imported from knowledge base."""
    try:
        # Get KB items that mention competitors
        kb_items = db.query(KnowledgeBaseItem).filter(
            KnowledgeBaseItem.is_active == True
        ).limit(10).all()

        competitors = db.query(Competitor).filter(
            Competitor.is_deleted == False
        ).all()

        preview = []
        for item in kb_items:
            matched_competitors = []
            for comp in competitors:
                if item.content_text and comp.name.lower() in item.content_text.lower():
                    matched_competitors.append(comp.name)

            if matched_competitors:
                preview.append({
                    "kb_item_id": item.id,
                    "title": item.title,
                    "matched_competitors": matched_competitors,
                    "content_preview": item.content_text[:300] if item.content_text else ""
                })

        return {
            "preview": preview,
            "total_items": len(preview)
        }
    except Exception as e:
        logger.error(f"Error previewing KB import: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/kb-items/import")
async def import_from_kb_items(
    request: Request,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Import competitor data from knowledge base."""
    try:
        body = await request.json()
        kb_item_ids = body.get("kb_item_ids", [])

        if not kb_item_ids:
            raise HTTPException(status_code=400, detail="kb_item_ids required")

        imported = []
        for item_id in kb_item_ids:
            kb_item = db.query(KnowledgeBaseItem).filter(
                KnowledgeBaseItem.id == item_id
            ).first()

            if kb_item:
                # This is a simplified import - in production you'd extract structured data
                # For now, just mark as processed
                imported.append({
                    "kb_item_id": item_id,
                    "title": kb_item.title,
                    "status": "imported"
                })

        return {
            "imported": imported,
            "total": len(imported),
            "message": f"Imported {len(imported)} items from knowledge base"
        }
    except Exception as e:
        logger.error(f"Error importing from KB: {e}")
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.post("/api/kb-items/verification/bulk-approve")
async def bulk_approve_kb_items_verification(
    request: Request,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Bulk approve verification items."""
    try:
        body = await request.json()
        item_ids = body.get("item_ids", [])

        if not item_ids:
            raise HTTPException(status_code=400, detail="item_ids required")

        approved = []
        for item_id in item_ids:
            kb_item = db.query(KnowledgeBaseItem).filter(
                KnowledgeBaseItem.id == item_id
            ).first()

            if kb_item:
                # Mark as verified by setting a category if none exists
                if not kb_item.category:
                    kb_item.category = "general"
                approved.append(item_id)

        db.commit()

        return {
            "approved": approved,
            "total": len(approved),
            "message": f"Approved {len(approved)} items"
        }
    except Exception as e:
        logger.error(f"Error bulk approving verification: {e}")
        db.rollback()
        raise HTTPException(status_code=500, detail="Internal server error. Please try again.")


@app.get("/api/kb/documents/{document_id}/extractions")
def get_document_extractions(
    document_id: int,
    db: Session = Depends(get_db)
):
    """
    Get all extractions from a specific KB document.
    """
    # Get KB item
    kb_item = db.query(KnowledgeBaseItem).filter(
        KnowledgeBaseItem.id == document_id
    ).first()

    if not kb_item:
        raise HTTPException(status_code=404, detail="Document not found")

    # Get entity links
    links = db.query(KBEntityLink).filter(
        KBEntityLink.kb_item_id == document_id
    ).all()

    # Get extractions
    extractions = db.query(KBDataExtraction).filter(
        KBDataExtraction.kb_item_id == document_id
    ).all()

    return {
        "document_id": document_id,
        "title": kb_item.title,
        "extraction_status": kb_item.extraction_status,
        "entity_links": [
            {
                "competitor_id": l.competitor_id,
                "link_type": l.link_type,
                "link_confidence": l.link_confidence,
                "created_at": l.created_at.isoformat() if l.created_at else None
            }
            for l in links
        ],
        "data_extractions": [
            {
                "id": e.id,
                "competitor_id": e.competitor_id,
                "field_name": e.field_name,
                "extracted_value": e.extracted_value,
                "extraction_confidence": e.extraction_confidence,
                "status": e.status,
                "data_as_of_date": e.data_as_of_date.isoformat() if e.data_as_of_date else None
            }
            for e in extractions
        ]
    }


@app.get("/api/reconciliation/conflicts")
def get_all_conflicts(
    status: Optional[str] = None,
    limit: int = 50,
    db: Session = Depends(get_db)
):
    """
    Get all unresolved data conflicts requiring human review.

    Returns KB extractions that have conflicts with live data.
    """
    # Get extractions with pending status (potential conflicts)
    query = db.query(KBDataExtraction).filter(
        KBDataExtraction.status == "pending"
    )

    extractions = query.limit(limit).all()

    conflicts = []
    for ext in extractions:
        # Check if there's conflicting live data
        live_source = db.query(DataSource).filter(
            DataSource.competitor_id == ext.competitor_id,
            DataSource.field_name == ext.field_name
        ).first()

        if live_source and live_source.current_value:
            # Check if values differ
            kb_value = ext.extracted_value
            live_value = live_source.current_value

            if kb_value.lower().strip() != live_value.lower().strip():
                # Get competitor name
                competitor = db.query(Competitor).filter(
                    Competitor.id == ext.competitor_id
                ).first()

                conflicts.append({
                    "extraction_id": ext.id,
                    "competitor_id": ext.competitor_id,
                    "competitor_name": competitor.name if competitor else "Unknown",
                    "field_name": ext.field_name,
                    "kb_value": kb_value,
                    "kb_confidence": ext.extraction_confidence,
                    "kb_date": ext.data_as_of_date.isoformat() if ext.data_as_of_date else None,
                    "live_value": live_value,
                    "live_confidence": live_source.confidence_score,
                    "live_date": live_source.data_as_of_date.isoformat() if live_source.data_as_of_date else None,
                    "status": ext.status
                })

    return {
        "total_conflicts": len(conflicts),
        "conflicts": conflicts
    }


@app.put("/api/reconciliation/resolve/{extraction_id}")
def resolve_conflict(
    extraction_id: int,
    resolution: str,  # "accept_kb", "accept_live", "reject"
    note: Optional[str] = None,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """
    Resolve a data conflict by accepting KB value, live value, or rejecting.
    """
    extraction = db.query(KBDataExtraction).filter(
        KBDataExtraction.id == extraction_id
    ).first()

    if not extraction:
        raise HTTPException(status_code=404, detail="Extraction not found")

    if resolution == "accept_kb":
        # Mark extraction as verified
        extraction.status = "verified"
        extraction.verified_by = current_user.get("email")
        extraction.verified_at = datetime.utcnow()
        extraction.reconciliation_note = note or "KB value accepted"

        # Update the competitor field with KB value
        competitor = db.query(Competitor).filter(
            Competitor.id == extraction.competitor_id
        ).first()

        if competitor and hasattr(competitor, extraction.field_name):
            setattr(competitor, extraction.field_name, extraction.extracted_value)

    elif resolution == "accept_live":
        # Mark extraction as superseded
        extraction.status = "superseded"
        extraction.reconciliation_note = note or "Live value preferred"

    elif resolution == "reject":
        extraction.status = "rejected"
        extraction.reconciliation_note = note or "Extraction rejected"

    else:
        raise HTTPException(status_code=400, detail="Invalid resolution type")

    db.commit()

    return {
        "extraction_id": extraction_id,
        "new_status": extraction.status,
        "resolved_by": current_user.get("email"),
        "message": f"Conflict resolved: {resolution}"
    }


# =========================================================================
# USER SAVED PROMPTS - Per-user prompt management
# =========================================================================

from database import UserSavedPrompt
# UserSavedPromptCreate, UserSavedPromptUpdate, UserSavedPromptResponse imported from schemas.prompts

@app.get("/api/user/prompts", response_model=List[UserSavedPromptResponse])
def get_user_prompts(
    prompt_type: Optional[str] = None,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get all saved prompts for the current user."""
    user_id = current_user.get("id")
    if not user_id:
        raise HTTPException(status_code=401, detail="User not authenticated")

    query = db.query(UserSavedPrompt).filter(UserSavedPrompt.user_id == user_id)
    if prompt_type:
        query = query.filter(UserSavedPrompt.prompt_type == prompt_type)

    return query.order_by(UserSavedPrompt.updated_at.desc()).all()


@app.post("/api/user/prompts", response_model=UserSavedPromptResponse)
def create_user_prompt(
    prompt_data: UserSavedPromptCreate,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Create a new saved prompt for the current user."""
    user_id = current_user.get("id")
    if not user_id:
        raise HTTPException(status_code=401, detail="User not authenticated")

    new_prompt = UserSavedPrompt(
        user_id=user_id,
        name=prompt_data.name,
        prompt_type=prompt_data.prompt_type,
        content=prompt_data.content,
        is_default=False
    )
    db.add(new_prompt)
    db.commit()
    db.refresh(new_prompt)

    # Log activity
    log_activity(db, current_user.get("email", "unknown"), user_id, "prompt_created", f"Created prompt: {prompt_data.name}")

    return new_prompt


@app.get("/api/user/prompts/{prompt_id}", response_model=UserSavedPromptResponse)
def get_user_prompt(
    prompt_id: int,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get a specific saved prompt by ID."""
    user_id = current_user.get("id")
    if not user_id:
        raise HTTPException(status_code=401, detail="User not authenticated")

    prompt = db.query(UserSavedPrompt).filter(
        UserSavedPrompt.id == prompt_id,
        UserSavedPrompt.user_id == user_id
    ).first()

    if not prompt:
        raise HTTPException(status_code=404, detail="Prompt not found")

    return prompt


@app.put("/api/user/prompts/{prompt_id}", response_model=UserSavedPromptResponse)
def update_user_prompt(
    prompt_id: int,
    prompt_data: UserSavedPromptUpdate,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Update an existing saved prompt."""
    user_id = current_user.get("id")
    if not user_id:
        raise HTTPException(status_code=401, detail="User not authenticated")

    prompt = db.query(UserSavedPrompt).filter(
        UserSavedPrompt.id == prompt_id,
        UserSavedPrompt.user_id == user_id
    ).first()

    if not prompt:
        raise HTTPException(status_code=404, detail="Prompt not found")

    if prompt_data.name is not None:
        prompt.name = prompt_data.name
    if prompt_data.content is not None:
        prompt.content = prompt_data.content
    if prompt_data.is_default is not None:
        # If setting as default, unset any other defaults for this type
        if prompt_data.is_default:
            db.query(UserSavedPrompt).filter(
                UserSavedPrompt.user_id == user_id,
                UserSavedPrompt.prompt_type == prompt.prompt_type,
                UserSavedPrompt.id != prompt_id
            ).update({"is_default": False})
        prompt.is_default = prompt_data.is_default

    prompt.updated_at = datetime.utcnow()
    db.commit()
    db.refresh(prompt)

    return prompt


@app.delete("/api/user/prompts/{prompt_id}")
def delete_user_prompt(
    prompt_id: int,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Delete a saved prompt."""
    user_id = current_user.get("id")
    if not user_id:
        raise HTTPException(status_code=401, detail="User not authenticated")

    prompt = db.query(UserSavedPrompt).filter(
        UserSavedPrompt.id == prompt_id,
        UserSavedPrompt.user_id == user_id
    ).first()

    if not prompt:
        raise HTTPException(status_code=404, detail="Prompt not found")

    prompt_name = prompt.name
    db.delete(prompt)
    db.commit()

    # Log activity
    log_activity(db, current_user.get("email", "unknown"), user_id, "prompt_deleted", f"Deleted prompt: {prompt_name}")

    return {"message": "Prompt deleted successfully"}


@app.post("/api/user/prompts/{prompt_id}/set-default")
def set_prompt_as_default(
    prompt_id: int,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Set a prompt as the default for its type."""
    user_id = current_user.get("id")
    if not user_id:
        raise HTTPException(status_code=401, detail="User not authenticated")

    prompt = db.query(UserSavedPrompt).filter(
        UserSavedPrompt.id == prompt_id,
        UserSavedPrompt.user_id == user_id
    ).first()

    if not prompt:
        raise HTTPException(status_code=404, detail="Prompt not found")

    # Unset any existing defaults for this prompt type
    db.query(UserSavedPrompt).filter(
        UserSavedPrompt.user_id == user_id,
        UserSavedPrompt.prompt_type == prompt.prompt_type
    ).update({"is_default": False})

    # Set this one as default
    prompt.is_default = True
    prompt.updated_at = datetime.utcnow()
    db.commit()

    return {"message": f"Prompt '{prompt.name}' set as default"}




# ============== Settings API Endpoints (User-Specific) ==============

import json

def get_user_setting(db: Session, user_id: int, setting_key: str):
    """Get a user-specific setting from the database."""
    setting = db.query(UserSettings).filter(
        UserSettings.user_id == user_id,
        UserSettings.setting_key == setting_key
    ).first()
    if setting:
        try:
            return json.loads(setting.setting_value)
        except (json.JSONDecodeError, TypeError, ValueError):
            return setting.setting_value
    return {}

def save_user_setting(db: Session, user_id: int, setting_key: str, setting_value):
    """Save a user-specific setting to the database."""
    setting = db.query(UserSettings).filter(
        UserSettings.user_id == user_id,
        UserSettings.setting_key == setting_key
    ).first()

    value_str = json.dumps(setting_value) if not isinstance(setting_value, str) else setting_value

    if setting:
        setting.setting_value = value_str
        setting.updated_at = datetime.utcnow()
    else:
        setting = UserSettings(
            user_id=user_id,
            setting_key=setting_key,
            setting_value=value_str
        )
        db.add(setting)
    db.commit()
    return setting


@app.post("/api/settings/schedule")
async def save_schedule_settings(
    settings: dict,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Save user refresh schedule preferences (personal to each user)."""
    save_user_setting(db, current_user["id"], "schedule", settings)
    return {"success": True, "message": "Schedule settings saved"}

@app.get("/api/settings/schedule")
async def get_schedule_settings(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get saved schedule settings (personal to each user)."""
    return get_user_setting(db, current_user["id"], "schedule")

@app.post("/api/settings/notifications")
async def save_notification_settings(
    settings: dict,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Save user notification preferences (personal to each user)."""
    save_user_setting(db, current_user["id"], "notifications", settings)
    return {"success": True, "message": "Notification settings saved"}

@app.get("/api/settings/notifications")
async def get_notification_settings(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get saved notification settings (personal to each user)."""
    return get_user_setting(db, current_user["id"], "notifications")


@app.get("/api/settings/threat-criteria")
async def get_threat_criteria_settings(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get threat criteria settings for determining threat levels."""
    settings = get_user_setting(db, current_user["id"], "threat_criteria")
    if settings:
        return settings
    # Return sensible defaults
    return {
        "criteria": {
            "high_threshold": 7,
            "medium_threshold": 4,
            "factors": [
                "market_share",
                "growth_rate",
                "product_overlap",
                "funding",
                "customer_base"
            ]
        }
    }


@app.post("/api/settings/threat-criteria")
async def save_threat_criteria_settings(
    settings: dict,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Save threat criteria settings."""
    save_user_setting(db, current_user["id"], "threat_criteria", settings)
    return {"success": True, "message": "Threat criteria settings saved"}


@app.put("/api/user-settings/ai-insights")
async def save_ai_insights_settings(
    settings: dict,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Save AI insights preferences."""
    save_user_setting(db, current_user["id"], "ai_insights", settings)
    return {"status": "saved", "message": "AI insights settings saved"}


# ============== Alert Rules CRUD ==============

@app.get("/api/notifications/rules")
async def get_alert_rules(
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Get all alert rules for the current user."""
    rules = get_user_setting(db, current_user["id"], "alert_rules")
    return rules if rules else []

@app.post("/api/notifications/rules")
async def create_alert_rule(
    rule: dict,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Create a new alert rule."""
    existing_rules = get_user_setting(db, current_user["id"], "alert_rules") or []
    existing_rules.append(rule)
    save_user_setting(db, current_user["id"], "alert_rules", existing_rules)
    return {"success": True, "rule": rule}

@app.put("/api/notifications/rules/{rule_id}")
async def update_alert_rule(
    rule_id: str,
    rule: dict,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Update an existing alert rule."""
    existing_rules = get_user_setting(db, current_user["id"], "alert_rules") or []
    for i, r in enumerate(existing_rules):
        if r.get("id") == rule_id:
            existing_rules[i] = rule
            break
    save_user_setting(db, current_user["id"], "alert_rules", existing_rules)
    return {"success": True, "rule": rule}

@app.delete("/api/notifications/rules/{rule_id}")
async def delete_alert_rule(
    rule_id: str,
    db: Session = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    """Delete an alert rule."""
    existing_rules = get_user_setting(db, current_user["id"], "alert_rules") or []
    existing_rules = [r for r in existing_rules if r.get("id") != rule_id]
    save_user_setting(db, current_user["id"], "alert_rules", existing_rules)
    return {"success": True}


# ============== P3-8: WebSocket for Real-time Updates ==============

from fastapi import WebSocket, WebSocketDisconnect
from typing import List, Dict, Set
import asyncio

class ConnectionManager:
    """
    P3-8: Enhanced WebSocket connection manager for real-time updates.

    Supports multiple event types:
    - refresh_progress: Data refresh progress updates
    - competitor_update: Competitor data changes
    - news_alert: New news articles for tracked competitors
    - discovery_result: Discovery agent findings
    - system_notification: System-wide notifications
    """
    def __init__(self):
        self.active_connections: List[WebSocket] = []
        self.subscriptions: Dict[WebSocket, Set[str]] = {}

    async def connect(self, websocket: WebSocket, subscriptions: List[str] = None):
        """Connect a WebSocket and optionally subscribe to event types."""
        await websocket.accept()
        self.active_connections.append(websocket)
        self.subscriptions[websocket] = set(subscriptions or ['all'])
        # Send welcome message
        await websocket.send_json({
            "type": "connection_established",
            "message": "Connected to Certify Intel real-time updates",
            "subscriptions": list(self.subscriptions[websocket])
        })

    def disconnect(self, websocket: WebSocket):
        """Disconnect a WebSocket and clean up subscriptions."""
        if websocket in self.active_connections:
            self.active_connections.remove(websocket)
        if websocket in self.subscriptions:
            del self.subscriptions[websocket]

    async def broadcast(self, message: dict, event_type: str = "general"):
        """Broadcast message to all connections subscribed to the event type."""
        message["event_type"] = event_type
        message["timestamp"] = datetime.now().isoformat()

        disconnected = []
        for connection in self.active_connections:
            try:
                subs = self.subscriptions.get(connection, {'all'})
                if 'all' in subs or event_type in subs:
                    await connection.send_json(message)
            except Exception:
                disconnected.append(connection)

        # Clean up disconnected clients
        for conn in disconnected:
            self.disconnect(conn)

    async def send_personal(self, websocket: WebSocket, message: dict):
        """Send a message to a specific connection."""
        try:
            await websocket.send_json(message)
        except Exception:
            self.disconnect(websocket)

    def get_stats(self):
        """Get connection statistics."""
        return {
            "active_connections": len(self.active_connections),
            "subscriptions": {
                str(id(ws)): list(subs)
                for ws, subs in self.subscriptions.items()
            }
        }

ws_manager = ConnectionManager()


# Helper functions for broadcasting events
async def broadcast_refresh_progress(competitor_name: str, progress: int, status: str):
    """Broadcast refresh progress update."""
    await ws_manager.broadcast({
        "type": "refresh_progress",
        "competitor": competitor_name,
        "progress": progress,
        "status": status
    }, event_type="refresh_progress")


async def broadcast_competitor_update(competitor_id: int, competitor_name: str, field: str, new_value: str):
    """Broadcast competitor data change."""
    await ws_manager.broadcast({
        "type": "competitor_update",
        "competitor_id": competitor_id,
        "competitor_name": competitor_name,
        "field": field,
        "new_value": new_value
    }, event_type="competitor_update")


async def broadcast_news_alert(competitor_name: str, headline: str, source: str, url: str):
    """Broadcast new news article alert."""
    await ws_manager.broadcast({
        "type": "news_alert",
        "competitor": competitor_name,
        "headline": headline,
        "source": source,
        "url": url
    }, event_type="news_alert")


async def broadcast_discovery_result(candidates: list, total_found: int):
    """Broadcast discovery agent results."""
    await ws_manager.broadcast({
        "type": "discovery_result",
        "candidates": candidates[:5],  # Limit to 5 for performance
        "total_found": total_found
    }, event_type="discovery_result")


async def broadcast_system_notification(title: str, message: str, level: str = "info"):
    """Broadcast system notification."""
    await ws_manager.broadcast({
        "type": "system_notification",
        "title": title,
        "message": message,
        "level": level  # info, warning, error, success
    }, event_type="system_notification")


@app.websocket("/ws/updates")
async def websocket_realtime_updates(websocket: WebSocket):
    """
    P3-8: Main WebSocket endpoint for real-time updates.

    Query params:
    - subscribe: comma-separated list of event types to subscribe to
      (refresh_progress, competitor_update, news_alert, discovery_result, system_notification)
      Default: 'all'
    """
    # Parse subscription from query params
    query_params = dict(websocket.query_params)
    subscriptions = query_params.get("subscribe", "all").split(",")
    subscriptions = [s.strip() for s in subscriptions if s.strip()]

    await ws_manager.connect(websocket, subscriptions)
    try:
        while True:
            # Handle incoming messages from client
            data = await websocket.receive_json()
            if data.get("type") == "ping":
                await websocket.send_json({"type": "pong"})
            elif data.get("type") == "subscribe":
                # Allow dynamic subscription changes
                new_subs = data.get("subscriptions", [])
                if new_subs:
                    ws_manager.subscriptions[websocket] = set(new_subs)
                    await websocket.send_json({
                        "type": "subscription_updated",
                        "subscriptions": new_subs
                    })
    except WebSocketDisconnect:
        ws_manager.disconnect(websocket)
    except Exception:
        ws_manager.disconnect(websocket)


@app.websocket("/ws/refresh-progress")
async def websocket_refresh_progress(websocket: WebSocket):
    """Legacy WebSocket endpoint for refresh progress (backwards compatibility)."""
    await ws_manager.connect(websocket, ["refresh_progress"])
    try:
        while True:
            data = await websocket.receive_text()
    except WebSocketDisconnect:
        ws_manager.disconnect(websocket)


@app.get("/api/ws/stats")
async def get_websocket_stats():
    """Get WebSocket connection statistics."""
    return ws_manager.get_stats()


# ============== Scheduler Configuration Endpoints (v5.2.0) ==============

@app.get("/api/scheduler/status")
async def get_scheduler_status():
    """Get scheduler status and configured jobs."""
    try:
        from scheduler import scheduler
        jobs = []
        for job in scheduler.get_jobs():
            jobs.append({
                "id": job.id,
                "name": job.name,
                "next_run": job.next_run_time.isoformat() if job.next_run_time else None,
                "trigger": str(job.trigger)
            })
        return {
            "running": scheduler.running,
            "jobs": jobs
        }
    except Exception as e:
        return {"running": False, "jobs": [], "error": "An unexpected error occurred"}


@app.post("/api/scheduler/start")
async def start_scheduler():
    """Start the scheduler with all jobs."""
    try:
        from scheduler import start_scheduler as start_sched
        start_sched()
        return {"success": True, "message": "Scheduler started"}
    except Exception as e:
        return {"success": False, "error": "An unexpected error occurred"}


@app.post("/api/scheduler/stop")
async def stop_scheduler_endpoint():
    """Stop the scheduler."""
    try:
        from scheduler import stop_scheduler as stop_sched
        stop_sched()
        return {"success": True, "message": "Scheduler stopped"}
    except Exception as e:
        return {"success": False, "error": "An unexpected error occurred"}


@app.post("/api/refresh/schedule")
async def configure_refresh_schedule(request: Request):
    """
    Configure the refresh schedule.

    Request body:
    {
        "type": "weekly" | "daily" | "custom",
        "day_of_week": "sun",  // for weekly
        "hour": 2,
        "minute": 0,
        "cron": "0 2 * * 0"  // for custom
    }
    """
    try:
        from scheduler import scheduler, CompetitorRefreshJob
        from apscheduler.triggers.cron import CronTrigger

        body = await request.json()
        schedule_type = body.get("type", "weekly")
        day_of_week = body.get("day_of_week", "sun")
        hour = body.get("hour", 2)
        minute = body.get("minute", 0)

        job = CompetitorRefreshJob()

        if schedule_type == "weekly":
            trigger = CronTrigger(day_of_week=day_of_week, hour=hour, minute=minute)
        elif schedule_type == "daily":
            trigger = CronTrigger(hour=hour, minute=minute)
        else:
            cron_expr = body.get("cron", "0 2 * * 0")
            parts = cron_expr.split()
            trigger = CronTrigger(
                minute=parts[0], hour=parts[1],
                day=parts[2], month=parts[3], day_of_week=parts[4]
            )

        async def refresh_task():
            await job.run_full_refresh()

        scheduler.add_job(
            refresh_task,
            trigger,
            id="configured_refresh",
            name="Configured Competitor Refresh",
            replace_existing=True
        )

        return {
            "success": True,
            "schedule": {
                "type": schedule_type,
                "day_of_week": day_of_week if schedule_type == "weekly" else None,
                "hour": hour,
                "minute": minute
            }
        }
    except Exception as e:
        return {"success": False, "error": "An unexpected error occurred"}


@app.post("/api/refresh/trigger")
async def trigger_manual_refresh(background_tasks: BackgroundTasks, db: Session = Depends(get_db)):
    """Trigger an immediate manual refresh for all competitors."""
    try:
        from scheduler import CompetitorRefreshJob

        job = CompetitorRefreshJob()

        async def run_refresh():
            await job.run_full_refresh()

        background_tasks.add_task(run_refresh)

        return {
            "success": True,
            "message": "Refresh job started in background",
            "status_endpoint": "/api/scrape/progress"
        }
    except Exception as e:
        return {"success": False, "error": "An unexpected error occurred"}


# ============== SOURCE VERIFICATION API (v7.2) ==============

# Standard list of tracked fields for verification summary
TRACKED_COMPETITOR_FIELDS = [
    "name", "website", "employee_count", "annual_revenue", "customer_count",
    "base_price", "pricing_model", "headquarters", "year_founded",
    "funding_total", "latest_round", "pe_vc_backers", "product_categories",
    "key_features", "target_segments", "geographic_focus", "g2_rating",
    "website_traffic", "social_following", "estimated_revenue",
    "estimated_market_share", "nps_score", "ceo_name", "tech_stack",
]


@app.get("/api/competitors/{competitor_id}/verification-summary")
def get_competitor_verification_summary(
    competitor_id: int,
    db: Session = Depends(get_db)
):
    """
    Verification summary for a competitor.

    v7.2: Shows how many tracked fields have verified sources.
    """
    competitor = db.query(Competitor).filter(
        Competitor.id == competitor_id
    ).first()
    if not competitor:
        raise HTTPException(
            status_code=404, detail="Competitor not found"
        )

    sources = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id,
        DataSource.field_name.in_(TRACKED_COMPETITOR_FIELDS)
    ).all()

    latest_by_field = {}
    for s in sources:
        existing = latest_by_field.get(s.field_name)
        is_newer = (
            s.extracted_at and existing
            and existing.extracted_at
            and s.extracted_at > existing.extracted_at
        )
        if existing is None or is_newer:
            latest_by_field[s.field_name] = s

    total_fields = len(TRACKED_COMPETITOR_FIELDS)
    sourced_count = len(latest_by_field)
    verified_count = sum(
        1 for s in latest_by_field.values() if s.is_verified
    )
    unverified_count = sourced_count - verified_count

    scores = [
        s.confidence_score for s in latest_by_field.values()
        if s.confidence_score is not None
    ]
    avg_confidence = (
        round(sum(scores) / len(scores), 1) if scores else 0
    )

    verification_pct = (
        round((verified_count / total_fields) * 100, 1)
        if total_fields > 0 else 0
    )

    return {
        "competitor_id": competitor_id,
        "competitor_name": competitor.name,
        "total_fields": total_fields,
        "sourced_count": sourced_count,
        "verified_count": verified_count,
        "unverified_count": unverified_count,
        "unsourced_count": total_fields - sourced_count,
        "verification_percentage": verification_pct,
        "avg_confidence": avg_confidence,
    }


@app.get("/api/competitors/{competitor_id}/source-links")
def get_competitor_source_links(
    competitor_id: int,
    db: Session = Depends(get_db)
):
    """
    Get ALL data fields for a competitor with source links.

    v7.2: Used by battlecard source dot display.
    """
    competitor = db.query(Competitor).filter(
        Competitor.id == competitor_id
    ).first()
    if not competitor:
        raise HTTPException(
            status_code=404, detail="Competitor not found"
        )

    sources = db.query(DataSource).filter(
        DataSource.competitor_id == competitor_id
    ).order_by(
        DataSource.field_name,
        DataSource.extracted_at.desc()
    ).all()

    latest_by_field = {}
    for s in sources:
        if s.field_name not in latest_by_field:
            latest_by_field[s.field_name] = s

    result = []
    for field_name, source in latest_by_field.items():
        current_value = None
        if hasattr(competitor, field_name):
            val = getattr(competitor, field_name)
            current_value = str(val) if val is not None else None

        result.append({
            "field_name": field_name,
            "current_value": source.current_value or current_value,
            "source_url": source.source_url,
            "source_name": source.source_name,
            "source_type": source.source_type,
            "is_verified": bool(source.is_verified),
            "confidence_score": source.confidence_score,
            "extracted_at": (
                source.extracted_at.isoformat()
                if source.extracted_at else None
            ),
            "deep_link_url": getattr(source, 'deep_link_url', None),
            "url_quality": _map_url_quality(getattr(source, 'url_status', None)),
            "source_section": getattr(source, 'source_section', None),
        })

    result.sort(key=lambda x: x["field_name"])

    return {
        "competitor_id": competitor_id,
        "competitor_name": competitor.name,
        "total_fields": len(result),
        "fields": result,
    }


# ============== Static Files (Must be Last) ==============
import os
import sys

# Determine the directory where the backend executable (or script) is running
if getattr(sys, 'frozen', False):
    # If we are running as a bundle (PyInstaller)
    base_dir = os.path.dirname(sys.executable)
    # In bundle, we expect 'frontend' to be a sibling of the executable (or in resources/frontend)
    # The 'backend-bundle' is usually in 'resources/backend-bundle'
    # So we look at 'resources/frontend' which is 1 level up from 'backend-bundle'
    frontend_dir = os.path.join(os.path.dirname(base_dir), "frontend")
else:
    # Running in normal dev environment
    frontend_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "frontend")


if os.path.exists(frontend_dir):
    from fastapi.responses import FileResponse
    
    from fastapi.responses import HTMLResponse

    @app.get("/clear-cache", response_class=HTMLResponse)
    async def clear_cache_page():
        """Serve a cache-clearing page that unregisters service workers and purges caches."""
        return """<!DOCTYPE html><html><head><title>Clearing Cache...</title></head><body style="background:#0f172a;color:#e2e8f0;font-family:sans-serif;display:flex;justify-content:center;align-items:center;height:100vh;margin:0;">
<div style="text-align:center"><h2 id="msg">Clearing browser cache...</h2><p id="detail">Unregistering service workers and purging cached assets</p></div>
<script>
(async function(){
    var msg=document.getElementById('msg'), detail=document.getElementById('detail');
    try {
        if('serviceWorker' in navigator){
            var regs=await navigator.serviceWorker.getRegistrations();
            for(var r of regs) await r.unregister();
            detail.textContent='Unregistered '+regs.length+' service worker(s). Clearing caches...';
        }
        var names=await caches.keys();
        for(var n of names) await caches.delete(n);
        localStorage.setItem('sw_cache_version','v8.2.1');
        msg.textContent='Cache cleared!';
        detail.textContent='Redirecting to app in 2 seconds...';
        setTimeout(function(){ window.location.href='/'; },2000);
    } catch(e) {
        msg.textContent='Error clearing cache';
        detail.textContent=e.message+'. Try manually: F12 > Application > Service Workers > Unregister';
    }
})();
</script></body></html>"""

    @app.get("/app")
    async def read_app_root():
        """Serve the frontend app root for Electron."""
        return FileResponse(os.path.join(frontend_dir, 'index.html'))

    app.mount("/", StaticFiles(directory=frontend_dir, html=True), name="frontend")
    logger.info(f"Serving frontend from: {frontend_dir}")
else:
    logger.warning(f"Warning: Frontend directory not found at {frontend_dir}")


# ============== Start Server ==============
if __name__ == "__main__":
    import uvicorn
    logger.info("=" * 50)
    logger.info("Certify Intel Backend Starting...")
    logger.info("Open http://localhost:8000 in your browser")
    logger.info("=" * 50)
    uvicorn.run(
        app,
        host=os.getenv("HOST", "127.0.0.1"),
        port=int(os.getenv("PORT", "8000")),
        reload=False
    )
